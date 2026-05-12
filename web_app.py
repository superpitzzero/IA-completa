"""
Web UI privada para Nexo, la IA local.

Ejecutar:
    python web_app.py --host 127.0.0.1 --port 7860
"""

from __future__ import annotations

import argparse
try:
    from stats_addon import init_stats_addon, track_message  # type: ignore
except Exception:
    init_stats_addon = None
    def track_message(**kwargs): pass
import base64
import hashlib
import io
import ipaddress
import json
import mimetypes
import os
import re
import secrets
import shutil
import socket
import tarfile
import threading
import time
import traceback
import unicodedata
import zipfile
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Any, Dict, Generator, Iterable, List, Optional
from urllib.parse import urlparse

import requests
from flask import (
    Flask,
    Response,
    jsonify,
    redirect,
    render_template_string,
    request,
    session,
    stream_with_context,
    url_for,
)
from werkzeug.security import check_password_hash, generate_password_hash
from werkzeug.utils import secure_filename

from orchestrator import (
    MODELS,
    OLLAMA_HOST,
    OLLAMA_KEEP_ALIVE,
    PROMPT_ARQUITECTO,
    PROMPT_PROGRAMADOR,
    PROMPT_CONVERSACIONAL,
    is_model_installed,
    is_ollama_running,
    ollama_options,
    start_ollama,
)

try:
    from dotenv import load_dotenv

    load_dotenv()
except Exception:
    pass


# Sesión HTTP reutilizable (keep-alive + pool) para Ollama, OpenAI y web fetch.
_HTTP: Optional[requests.Session] = None


def http_session() -> requests.Session:
    global _HTTP
    if _HTTP is not None:
        return _HTTP
    s = requests.Session()
    try:
        adapter = requests.adapters.HTTPAdapter(pool_connections=64, pool_maxsize=64, max_retries=0)
        s.mount("http://", adapter)
        s.mount("https://", adapter)
    except Exception:
        pass
    _HTTP = s
    return s


APP_TITLE = "Nexo"
DATA_DIR = Path(os.getenv("NEXO_WEB_DATA") or os.getenv("IA_COMBINADA_WEB_DATA", "web_data"))
USERS_FILE = DATA_DIR / "users.json"
CHATS_FILE = DATA_DIR / "chats.json"
MEMORY_FILE = DATA_DIR / "memory.json"
SECRET_FILE = DATA_DIR / "secret_key.txt"
SETTINGS_FILE = DATA_DIR / "settings.json"
UPLOAD_DIR = DATA_DIR / "uploads"
UPLOAD_TTL_DAYS = 7
MAX_TOTAL_UPLOAD_BYTES = 128 * 1024 * 1024
MAX_ATTACHMENTS_PER_MESSAGE = 8
MAX_IMAGE_BYTES = 12 * 1024 * 1024
MAX_DOCUMENT_BYTES = 20 * 1024 * 1024
MAX_VIDEO_BYTES = 100 * 1024 * 1024
MAX_AUDIO_BYTES = 50 * 1024 * 1024
MAX_ARCHIVE_BYTES = 80 * 1024 * 1024
MAX_UNKNOWN_BYTES = 128 * 1024 * 1024
MAX_DOCUMENT_CHARS = 24_000
MAX_ATTACHMENT_CONTEXT_CHARS = 30_000
MAX_WEB_CONTEXT_CHARS = 12_000
# Límites reducidos para modo Rápido (acceso básico a archivos, sin calidad completa)
MAX_DOCUMENT_CHARS_RAPIDO = 3_000   # Texto truncado: lectura rápida, no análisis profundo
MAX_ATTACHMENT_CONTEXT_CHARS_RAPIDO = 6_000
MAX_ARCHIVE_MEMBERS = 30
MAX_ARCHIVE_EXTRACT_BYTES = 40 * 1024 * 1024
VIDEO_FRAME_COUNT = 24
VISION_IMAGE_BATCH = 6
OPENAI_API_BASE = "https://api.openai.com/v1"

IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif"}
TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".jsonl", ".xml", ".html", ".htm", ".css", ".js", ".ts",
    ".tsx", ".jsx", ".py", ".java", ".c", ".cpp", ".h", ".hpp", ".cs", ".go", ".rs", ".php",
    ".rb", ".swift", ".kt", ".kts", ".sql", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".env", ".log", ".bat", ".ps1", ".sh", ".dockerfile", ".vue", ".svelte",
}
PDF_EXTENSIONS = {".pdf"}
OFFICE_EXTENSIONS = {".docx", ".doc", ".rtf", ".odt", ".pptx", ".ppt", ".xlsx", ".xls", ".ods"}
DOCUMENT_EXTENSIONS = TEXT_EXTENSIONS | PDF_EXTENSIONS | OFFICE_EXTENSIONS
VIDEO_EXTENSIONS = {".mp4", ".webm", ".mov", ".avi", ".mkv", ".m4v"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".aac", ".ogg", ".oga", ".flac", ".wma", ".webm"}
ARCHIVE_EXTENSIONS = {".zip", ".tar", ".gz", ".tgz", ".bz2", ".tbz2", ".xz", ".txz", ".7z"}

UPLOAD_LIMITS = {
    "image": MAX_IMAGE_BYTES,
    "document": MAX_DOCUMENT_BYTES,
    "video": MAX_VIDEO_BYTES,
    "audio": MAX_AUDIO_BYTES,
    "archive": MAX_ARCHIVE_BYTES,
    "unknown": MAX_UNKNOWN_BYTES,
}

DEFAULT_SETTINGS = {
    "ai_provider": "openai",
    "openai_api_key": "",
    "openai_model": "gpt-5.5",
    "openai_fast_model": "gpt-4.1-mini",
    "openai_transcribe_model": "gpt-4o-transcribe",
    "fallback_to_ollama": True,
    # Identidad (opcional). Si esta vacio, el asistente NO debe afirmar un nombre fijo.
    "identity_name": "",
    "identity_guard_enabled": False,
    # Router automatico (clasificador barato) para escoger modo por mensaje.
    "auto_router_enabled": True,
    "router_provider": "auto",  # auto | ollama | openai
    "router_model_ollama": "qwen2.5:1.5b",
    "router_model_openai": "gpt-4.1-mini",
    "donate_url": "",
}

PLAN_FREE = "gratis"
PLAN_BETA = "beta_tester"
PLAN_DEVELOPER = "developer"
NEW_USER_DEFAULT_PLAN = PLAN_FREE
LEGACY_USER_DEFAULT_PLAN = PLAN_FREE

# Promoción de lanzamiento: primeros N registros en X días → Developer gratis Y días
PROMO_SLOTS = 30       # cupos máximos
PROMO_WINDOW_DAYS = 3  # duración de la ventana de promoción
PROMO_PLAN_DAYS = 7    # días de Developer gratis

PLAN_DEFINITIONS: Dict[str, Dict[str, Any]] = {
    PLAN_FREE: {
        "key": PLAN_FREE,
        "label": "Plan Gratis",
        "price_eur": 0,
        "priority": 10,
        "allowed_modes": ["auto", "rapido"],
        "includes_api_key": False,
        "features": ["Modo rapido", "Prioridad estandar"],
    },
    PLAN_BETA: {
        "key": PLAN_BETA,
        "label": "Plan BETA Tester",
        "price_eur": 3,  # Descuento de lanzamiento (antes 5)
        "priority": 50,
        "allowed_modes": ["auto", "rapido", "combinado"],
        "includes_api_key": False,
        "features": ["Modo combinado", "Mas prioridad entre usuarios"],
    },
    PLAN_DEVELOPER: {
        "key": PLAN_DEVELOPER,
        "label": "Plan Developer",
        "price_eur": 7,  # Descuento de lanzamiento (antes 15)
        "priority": 100,
        "allowed_modes": ["auto", "rapido", "combinado", "codigo"],
        "includes_api_key": True,
        "features": ["Modo codigo", "Prioridad maxima", "API Key"],
    },
}

PLAN_ALIASES = {
    "free": PLAN_FREE,
    "gratis": PLAN_FREE,
    "plan gratis": PLAN_FREE,
    "beta": PLAN_BETA,
    "beta_tester": PLAN_BETA,
    "beta tester": PLAN_BETA,
    "plan beta tester": PLAN_BETA,
    "developer": PLAN_DEVELOPER,
    "dev": PLAN_DEVELOPER,
    "plan developer": PLAN_DEVELOPER,
}

# Usuarios con acceso al panel de administración
ADMIN_USERS = {"aerys"}

def is_admin_user(username: str) -> bool:
    return (username or "").strip().lower() in ADMIN_USERS


def identity_guard_from_settings(settings: Optional[Dict[str, Any]] = None) -> str:
    settings = settings or load_ai_settings()
    enabled = bool(settings.get("identity_guard_enabled", False))
    name = str(settings.get("identity_name") or "").strip()
    if not enabled or not name:
        return ""
    return (
        f"Identidad fija: te llamas {name}. Eres una IA local creada y personalizada en este proyecto. "
        "No digas que perteneces a ningun proveedor externo ni uses otra identidad. "
        f"Si te preguntan como te llamas, responde claramente: 'Soy {name}'. "
        f"Tu nombre visible y permanente es {name}."
    )

PLACEHOLDER_PASSWORDS = {
    "",
    "TU_CONTRASENA",
    "TU_CONTRASE\u00d1A",
    "TU_CONTRASE\u00c3\u2018A",
    "CAMBIA_ESTA_CONTRASENA",
}
USERNAME_RE = re.compile(r"^[A-Za-z0-9_.-]{3,32}$")
LOGIN_ATTEMPTS: Dict[str, Dict[str, float]] = {}
LOGIN_LOCK = threading.Lock()
DATA_LOCK = threading.RLock()
MEMORY_UPDATE_LOCK = threading.Lock()

# ═══ NEXO MEJORAS: Rate Limiter por usuario ═══
_RATE_LIMIT_LOCK = threading.Lock()
_RATE_LIMIT_DATA: Dict[str, List[float]] = {}
RATE_LIMIT_MAX    = int(os.getenv('NEXO_RATE_LIMIT_MAX',    '5'))
RATE_LIMIT_WINDOW = int(os.getenv('NEXO_RATE_LIMIT_WINDOW', '60'))

def check_rate_limit(user_id: str) -> Optional[int]:
    """Devuelve segundos de espera si supera el límite; None si OK."""
    now = time.time()
    with _RATE_LIMIT_LOCK:
        ts = _RATE_LIMIT_DATA.get(user_id, [])
        ts = [t for t in ts if now - t < RATE_LIMIT_WINDOW]
        if len(ts) >= RATE_LIMIT_MAX:
            wait = int(RATE_LIMIT_WINDOW - (now - ts[0])) + 1
            _RATE_LIMIT_DATA[user_id] = ts
            return wait
        ts.append(now)
        _RATE_LIMIT_DATA[user_id] = ts
        return None

# ═══ NEXO MEJORAS: Personality (thread-local) ═══
_personality_local = threading.local()

PERSONALITY_PROMPTS: Dict[str, str] = {
    'normal':      '',
    'programador': 'Tono técnico: responde directo, usa ejemplos de código cuando sea útil, evita explicaciones de relleno.',
    'creativo':    'Tono creativo: usa metáforas, analogías imaginativas y entusiasmo. Haz las respuestas amenas y originales.',
    'conciso':     'Tono ultra-conciso: responde siempre en máximo 3 líneas. Nada de introducciones, solo el núcleo de la respuesta.',
}

def set_personality(p: str) -> None:
    _personality_local.value = p if p in PERSONALITY_PROMPTS else 'normal'

def get_personality() -> str:
    return getattr(_personality_local, 'value', 'normal')


try:
    MAX_CONCURRENT_AI_REQUESTS = max(1, int(os.getenv("NEXO_MAX_CONCURRENT_AI_REQUESTS", "1")))
except ValueError:
    MAX_CONCURRENT_AI_REQUESTS = 1

AI_PRIORITY_CONDITION = threading.Condition()
AI_PRIORITY_QUEUE: List[Dict[str, Any]] = []
AI_ACTIVE_REQUESTS = 0
AI_TICKET_COUNTER = 0

# Cache corto para evitar lecturas de disco por request.
_AI_SETTINGS_CACHE: Dict[str, Any] = {"ts": 0.0, "mtime": 0.0, "value": None}
_AI_SETTINGS_TTL_S = 1.0

# Cache corto del router (evita repetir llamadas baratas si el usuario reintenta el mismo mensaje).
_ROUTER_CACHE: Dict[str, Dict[str, Any]] = {}
_ROUTER_CACHE_TTL_S = 90.0
_ROUTER_CACHE_MAX = 256


LOGIN_HTML = r"""
<!doctype html>
<html lang="es">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{{ page_title }} - Nexo</title>
  <style>
    :root {
      color-scheme: dark;
      --bg: #101010;
      --panel: #171717;
      --panel-2: #202020;
      --text: #f3f3f3;
      --muted: #a8a8a8;
      --line: #303030;
      --accent: #19c37d;
      --danger: #ff6b6b;
    }
    * { box-sizing: border-box; }
    body {
      margin: 0;
      min-height: 100vh;
      display: grid;
      place-items: center;
      background: var(--bg);
      color: var(--text);
      font-family: Inter, ui-sans-serif, system-ui, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
    }
    .login {
      width: min(420px, calc(100vw - 32px));
      padding: 28px;
      border: 1px solid var(--line);
      background: var(--panel);
      border-radius: 8px;
      box-shadow: 0 24px 70px rgba(0, 0, 0, .32);
    }
    h1 { margin: 0 0 8px; font-size: 26px; letter-spacing: 0; }
    p { margin: 0 0 24px; color: var(--muted); line-height: 1.5; }
    label { display: block; margin: 14px 0 8px; color: #dedede; font-size: 14px; }
    input {
      width: 100%;
      height: 44px;
      border: 1px solid var(--line);
      border-radius: 6px;
      background: #0f0f0f;
      color: var(--text);
      padding: 0 12px;
      font-size: 16px;
      outline: none;
    }
    input:focus { border-color: var(--accent); }
    button {
      width: 100%;
      height: 44px;
      margin-top: 22px;
      border: 0;
      border-radius: 6px;
      background: var(--accent);
      color: #06140e;
      font-weight: 700;
      cursor: pointer;
    }
    .error {
      margin: 0 0 16px;
      padding: 11px 12px;
      border-radius: 6px;
      font-size: 14px;
      line-height: 1.45;
    }
    .error { background: rgba(255, 107, 107, .12); color: #ffd0d0; border: 1px solid rgba(255, 107, 107, .3); }
    .switch {
      margin: 18px 0 0;
      color: var(--muted);
      text-align: center;
      font-size: 14px;
    }
    .switch a { color: var(--accent); text-decoration: none; font-weight: 700; }
    .switch a:hover { text-decoration: underline; }
  </style>
</head>
<body>
  <main class="login">
    <h1>Nexo</h1>
    <p>{{ subtitle }}</p>
    {% if error %}
      <div class="error">{{ error }}</div>
    {% endif %}
    <form method="post" action="{{ action }}">
      <label for="username">Usuario</label>
      <input id="username" name="username" autocomplete="username" required autofocus>
      <label for="password">Contrasena</label>
      <input id="password" name="password" type="password" autocomplete="{{ password_autocomplete }}" required>
      {% if register %}
        <label for="confirm_password">Repite contrasena</label>
        <input id="confirm_password" name="confirm_password" type="password" autocomplete="new-password" required>
      {% endif %}
      <button type="submit">{{ button_text }}</button>
    </form>
    <p class="switch">{{ switch_text }} <a href="{{ switch_href }}">{{ switch_label }}</a></p>
  </main>
</body>
</html>
"""


DONATE_HTML = r"""
<!doctype html>
<html lang="es">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>Donate - Nexo</title>
  <style>
    :root {
      color-scheme: dark;
      --bg: #161616;
      --panel: #202020;
      --panel-2: #292929;
      --text: #f4f4f4;
      --muted: #b8b8b8;
      --line: #3a3a3a;
      --accent: #19c37d;
      --accent-2: #2dd4bf;
      --danger: #ff9a9a;
    }
    * { box-sizing: border-box; }
    body {
      margin: 0;
      min-height: 100vh;
      background: var(--bg);
      color: var(--text);
      font-family: Inter, ui-sans-serif, system-ui, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
    }
    .page {
      width: min(980px, calc(100vw - 32px));
      margin: 0 auto;
      padding: 42px 0;
    }
    header {
      display: grid;
      gap: 14px;
      padding: 0 0 24px;
      border-bottom: 1px solid var(--line);
    }
    .nav {
      display: flex;
      justify-content: space-between;
      align-items: center;
      gap: 12px;
    }
    .brand {
      display: flex;
      align-items: center;
      gap: 10px;
      color: var(--text);
      text-decoration: none;
      font-weight: 800;
    }
    .mark {
      display: grid;
      place-items: center;
      width: 32px;
      height: 32px;
      border-radius: 6px;
      background: #f2f2f2;
      color: #111;
      font-size: 13px;
    }
    h1 {
      margin: 18px 0 0;
      max-width: 760px;
      font-size: clamp(34px, 6vw, 58px);
      line-height: 1.03;
      letter-spacing: 0;
    }
    .lead {
      margin: 0;
      max-width: 760px;
      color: var(--muted);
      font-size: 18px;
      line-height: 1.6;
    }
    .actions {
      display: flex;
      flex-wrap: wrap;
      gap: 10px;
      margin-top: 8px;
    }
    .btn {
      display: inline-flex;
      align-items: center;
      justify-content: center;
      min-height: 44px;
      border: 1px solid var(--line);
      border-radius: 6px;
      padding: 0 14px;
      background: var(--panel-2);
      color: var(--text);
      text-decoration: none;
      font-weight: 700;
    }
    .btn.primary { background: var(--accent); border-color: var(--accent); color: #06140e; }
    .btn.disabled { color: var(--muted); cursor: not-allowed; opacity: .75; }
    .grid {
      display: grid;
      grid-template-columns: minmax(0, 1.2fr) minmax(260px, .8fr);
      gap: 18px;
      margin-top: 22px;
    }
    .panel {
      border: 1px solid var(--line);
      border-radius: 8px;
      background: var(--panel);
      padding: 18px;
    }
    h2 { margin: 0 0 12px; font-size: 18px; letter-spacing: 0; }
    p { color: var(--muted); line-height: 1.58; }
    ul {
      margin: 0;
      padding: 0;
      list-style: none;
      display: grid;
      gap: 9px;
    }
    li {
      display: flex;
      justify-content: space-between;
      gap: 12px;
      border-bottom: 1px solid rgba(255, 255, 255, .08);
      padding-bottom: 9px;
      color: var(--muted);
    }
    li:last-child { border-bottom: 0; padding-bottom: 0; }
    li strong { color: var(--text); font-weight: 700; }
    .note {
      margin-top: 14px;
      padding: 12px;
      border-radius: 6px;
      border: 1px solid rgba(45, 212, 191, .35);
      background: rgba(45, 212, 191, .08);
      color: #d9fffa;
      line-height: 1.5;
    }
    .warning {
      border-color: rgba(255, 154, 154, .35);
      background: rgba(255, 154, 154, .08);
      color: #ffe1e1;
    }
    @media (max-width: 760px) {
      .page { padding: 24px 0; }
      .grid { grid-template-columns: 1fr; }
      .nav { align-items: flex-start; }
      li { display: grid; }
    }
  </style>
</head>
<body>
  <main class="page">
    <header>
      <div class="nav">
        <a class="brand" href="/"><span class="mark" style="background:transparent;display:inline-flex;align-items:center;justify-content:center;"><img src="data:image/png;base64,/9j/4AAQSkZJRgABAQAAAQABAAD/4gHYSUNDX1BST0ZJTEUAAQEAAAHIAAAAAAQwAABtbnRyUkdCIFhZWiAH4AABAAEAAAAAAABhY3NwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQAA9tYAAQAAAADTLQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAlkZXNjAAAA8AAAACRyWFlaAAABFAAAABRnWFlaAAABKAAAABRiWFlaAAABPAAAABR3dHB0AAABUAAAABRyVFJDAAABZAAAAChnVFJDAAABZAAAAChiVFJDAAABZAAAAChjcHJ0AAABjAAAADxtbHVjAAAAAAAAAAEAAAAMZW5VUwAAAAgAAAAcAHMAUgBHAEJYWVogAAAAAAAAb6IAADj1AAADkFhZWiAAAAAAAABimQAAt4UAABjaWFlaIAAAAAAAACSgAAAPhAAAts9YWVogAAAAAAAA9tYAAQAAAADTLXBhcmEAAAAAAAQAAAACZmYAAPKnAAANWQAAE9AAAApbAAAAAAAAAABtbHVjAAAAAAAAAAEAAAAMZW5VUwAAACAAAAAcAEcAbwBvAGcAbABlACAASQBuAGMALgAgADIAMAAxADb/2wBDAAUDBAQEAwUEBAQFBQUGBwwIBwcHBw8LCwkMEQ8SEhEPERETFhwXExQaFRERGCEYGh0dHx8fExciJCIeJBweHx7/2wBDAQUFBQcGBw4ICA4eFBEUHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh7/wAARCAMABYADASIAAhEBAxEB/8QAHQABAAEFAQEBAAAAAAAAAAAAAAEDBAUGCAcCCf/EAF0QAAIBAwEEBgQICQgDDQgDAQABAgMEEQUGEiExB0FRYXGBCBMikRQyQlJygqGxIzNDYpKUssHRFRYXJFOi0uFjs8IlNDU2RUZUVmR0g5PwGCZEVXN1hKNlw/Hi/8QAGwEBAAIDAQEAAAAAAAAAAAAAAAIFAQMEBgf/xAAtEQEAAgIBAwQABgIDAQEAAAAAAQIDEQQSITEFE0FRFBUiM1JhIyQyNEJxgf/aAAwDAQACEQMRAD8A7LAAAYAAEYJADAwAAwMIAAMAAAAAGAAAwAAwMAAMAAAMAAMDAAAYAADAADAwAAwMAAMDAADAAAAAAMAABgABggkAMDAAAYAAAABgYAAAAAMAAQTgAAAAGBgAAAAAAAYGAAAAAYGEAAGAAAwAAwAAGEMAAMDAAAYAADAADAwAAwMAAAAAAADAwAAwAAGBhAAMAABgYAAYGEAAAADAAAAAAAAAwAAGO4ABgYAABIABgYAAAABgYAAgnAAAAAMDAADAwAAAADAwAAwAAAwAAwMIABgjBIADAADAwAAwMAAAAAGAAGBgAAMAAMDAADAx3AAMDAAEYJwAAwMAAMIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAYAAAAAAAAAAAAAAAYAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAD5AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAwAACAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAABDJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAABxAAcQAAAAAAAAAAAADkAAAAAAAAAAAAAAAAAAHEAAAAAAAAAAAAAI4k8QBHEkAAAAAAAAAAAAAAAAABxAAAAAAAAAAAAAAAAAAAAAAAAAAAAABxAAAAAAAAAAAAAAAA4gAQSAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACOIAEkEgCCQAAAAAAAAAAAAcQAAHEAAAAAAAAAAAAAAAAAAAAAAAAAAGAAAAAAAAAAAAAAAAAAwAAAAAAAAAAAAAAAAMAAMAAAAAAAAAAAAAHAAABgAAAAAAAAAAAAAAADAAAAAAAAAAAAAMAABgAABgAAAAAAAAAAAAGCCQAwQkSAAAAAAAAAABDYE4HMpV69GhTdStVhSgucpySSMdLaHSE2o3iq466UJTXvimNMdUMtgGOtNb0u6moUb6i5vlCUt2XufEyCawCJiUgAMgAAAAAAABGCQAwAAAAAAAAAABHWSAAAAgkAAAAAAAAAAAAAAADAAAAAAAAAAAAAAAA4AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAGAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAHxUqQpwc6kowjFZcpPCRiqu0FnKThZRq381wfqFmC8ZvEftM6Rm0Qy+UU7m4oW9J1a9WnSgucpyUUveYGrealXi5Va9DT6S5+r9ueO+csRXuZZ0FZVavrLS1uNUrr8q/wiT+nP2V5GelH3N+GWq6/QmsWFvXvX86Ed2n+lLC92SxuLzVK0XOveULCiuaorLXjOfD7C5p6bqdzxr1qNnH5tJetn+lJYXuZdW+hadSmqlSlK5qrlUuJOo14Z4LySG4hHV7NepU7S4qqdtaXGqVl+Vlma/Tn7K8jJQs9anFYpWVCPVGVWUmvckjYlFJJJYSJE2IxR8tYurS+cHG80yldQ/0U1P8AuzS+xlrRqUaNRUrK+ubCr1UJ5S/Qn/sm4FK5tre6pOlcUadWD5xnFNfaIse1rxLC0dW1Kg925tad3FfLt3uT/Ql/Ev7XWtOuKipOt6ms/wAlWi6cvJPn5FCehUqfGxuq9t+Y36yn+jLl5NFvc299Gk6d3YUb6j1+pay/GE/3NjtJHVHln8onJqdq7elUVLT76vY1f+j1M4/8ufV9HBf09V1Cg926so3MV8u2eJecJfubHTKUZI+WdBj7HWLC7mqVOuo1v7KonCa+q+JfriRTiYnwkABkAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAgASAPMAAAAAAAAAAAAIJAAgkAAAAAAAAAAAABAEggnzAAAAAAAAAAAAAAABAEgEAOJJBIAAAABkABlFveXlrZ0fXXVenRpr5U5YyDa4bIzgwtbWa9fhp1nJxfKtc5pw8o/GfuXiY2+nGc1DVdRq15S5W1LMFLwhHMpebZmKoTePhmrzW9Ptqro+udeuvyVCO/PzS5eeDH19R1a5T9RTpWFP51T8LU9y9le9nxa2d9UpqlZWFHT7fqlWST8qcf3teBe0dAtpYd9VrXz+bUaVP9BYXvyZ7QjPVZgJfBrqvicrnWLiL5fjIxfgsU4+ZlaOnatcJKcqFjS5JJetqL7or7TP06VOlCNOlCNOEeCjFYS8j7HWRi+2Lt9CsYSVS4jO8qLlO4lv48F8VeSMpGKjFJLCXUggR22RWI8JBGUTkMgIcornJLzPl1afXUh7wPvqB8KpB8px96JynxT+0D6I5kNdhOe8Cjd2tvdU/V3NCnWh82cU19pjKuh+rWbC8r26/s5v1lP3S4ryaM0DMTMIzWJarf2t5ubl/ptO8pr5dFb+O/cfFeTZRsa9SEnHTdUqRcedvXzUUfGMsTj7zb8Itb7TrO9ildW9Orjk2vaj4PmvIz1Nc4teGNo63cUeGoWE1Fc6tv+Ej5x+MvczKWN/Z30HO0uadZLnuy4rxXNeZjKui3NDjYX0mlypXK315SXtLzyY69oqFRVNS02pRnHlc0cyS+vH2l5pDUSRa1fLbQa1Z3d9TpqpaXtK/odUaz4+VSP70y/o67bJqF9TqWM28L1y9h+E17P2pmNSnF4llgfKkpJNPOeKJy+8wmkEJ5JAABAAAAAAAAAACAJBBIAAAAAAAAAAAAAABBIAEEgAAAAAAAAAAAAIAkAAAAAAAAAAAAAAAAAjzAkAAAAAAAAAAAAAA4gAAAAAAAAAAAAAAAIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAIwBJBIAAAAAAAAAAAAAAAAAgkYAAAAAAAAAAAAAAAIJAEEgAAAAAAAAACCQAIJAEDBIADAAAAAAG8Ftf31rY0HXu68KUM4Tk+LfYlzb7kDelyWWo6nZ2CXwisozl8SnFb05+EVxZirnUNQveFunYW7/KVIp1pLui+EF3vL7kY6yiqlSS0m2ld1JPFS6qTe43+dUeXLwjnyJRX7apyfEMhcalqF2mqMY2FL59RKdXHh8WPnnwLC09VWuHUsaFXUrlcHcSlvJf+JLgvCPuMpabPwniep1neS5+q3d2ivq/K+s35GbpwhThGEIqMYrCSWEh1RHhiKWnyw1HSbuu96+vPVRfOla5j75v2n5bpk7Kws7KDja29Olnm4rjLxfN+ZccC3rXttSn6uVRSqf2cFvS9yMbmWyKxC4wgWnrb2t+KoRoRfyqzy/0V/EfAvWf75uK1btjndj7l+/JjTO1Wtd21KW7UrQUuqOePuKfwuU/xNrXqd7juL+9grUaFGgt2jRp01+bHBUlKMYtyaSXW2BaKWoTxinb0V+dJyf2Y+8K3vJ/jL1rup00vvyfUtQs08KvGb7IZk/sPn4ZOX4qzuZ97io/e0GOx8BT+PdXU/wDxWvuwP5Ns38anKf0pyf3seuv5fFtKUPp1v4JhR1GXy7WHhGUv3oHZK06xX/wlHzifSsLL/otD9BHx6m/fO8pLwo/5j1F5133/AOpBl9uwsnztaP6CPn+TbHOVbU0+5YIVC8/6cn40UT6q/i+FzRl40X/iAj+T7ZfE9bD6NaS/ePgU48ad7cx8ZKX3pjOoxfxbWa8ZR/iPX3kX7djvf/Tqp/fgMdj1V9F+zdU5rsnS4+9P9w9bew+Pawmu2nU/c0vvHw6EfxtG4pfSpNr3rKKtG8tazxTuKcn2KSz7gPj4dSjwrQq0e+cGl71w+0r0qtOrHepTjOPbF5R9PBQq2dtUlvujFT+dH2Ze9cQyuBhFp6i6pL8Bdb6+bWW99q4/eQ7upSf9ZtqkF8+n7cfs4/YDajd6NY16jrRhK3rPj62hLck/HHCXmmWNex1O3hKOKWoUWsNcKdTHh8WX90zdG4o1o71KpGce2LyVFxMxMwjNIlp9s40azhp11VsKy4u1qR9n/wAt9XfFoylHXZUPZ1S2dFL/AOIpZnS8+uPmsd5lL6xtL2l6u6oQqxXFZXGL7U+afejE3GmXtqnKzq/Cqa/JV5Ymvoz6/CXvJbiUNWr4Zq3rUa9KNWjVhUpyWYyjJNNdzKvmahbKEbqTs51dNvfjToyjje73DlJfnR95lLfW5UWqWqUlbZeI3EXmjJ+PyX3S97MTVKuSJ8s0SRGSkk0+ZJFsAAAAAAAcADI8yQBBIAAAAAAAAAAAAAABBKAAAAAAAAAAAAAAAIJAAgkAAAAAAAAAAAOAAgkAQSAAAAAAAAAAAAAAAMAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQ5Jcy21G/trC3de5qbkc7qSWZSk+Silxb7ka5qV1WvaM6upT+BWEeLoOeJSX+kkv2F5t8jMRMoWvEMjea26k5UNLhCvNPEq8vxNN+K+M+5ebRjI4V88KtqWpYw3wzTT/u04/a+8uLGyvL+EFGM9OsUsRSju1prsS/Jx/vfRNgsrS2srdULWjGlTTziPW+tt9b72Z8IxE27yxdvoruMVNVnGt1q2p5VFePXPz4dxmI04QgowioxSwklhJH1yLWpeJ1HStoOvVXBqL9mP0pdXhz7iO5lsiIhc5S58C1nfKb3LSnK4l2x4QXjLl7siNpKs1K9qet6/Vx4U15dfn7i6jGMYpJKKS4JLkDvKzVtc13m7uHGP9lRbivOXN/YXNvQo0IblGlCmuyKwUKl/Tc3Ttoyuai4NU+S8ZckfPqLyuv6xX9TB/k6HPzk+PuwGFevc0KCXrasIt8k3xfgubKLurirwtrSbXzqr3I+7i/sK1va0LfLpUoxk+cucn4t8WVwz3WSoXlX8bdqmuuNGGPtef3Ex021T3pwdaXbVk5/eXgBp8whCEcQiorsSwTgkBkI5kgAkAAGAAAAADBSrW9CssVaNOf0oplUAWnwCnHjQq1qD/Mnle55R841ClylRuYrqa3Je/ivuL0BjSzV9CHC4pVLd9s17P6S4FzCUJxUoSUovk08o+mk+ZbTsaO+50t6hN85Unu58Vyfmgd017KhVn6zdcKnVUg92XvXPzKf9dt/m3UF4RqL9z+wlyvbf49NXMO2Hsz9z4Mq211QuMqnP2l8aEliS8U+IEW93RrS3FJxqLi6c1iS8ivz4Mp3FvSrpKpBSxyfJrwfUUVG5t/it3FNdTwprz5PzB3h931la3tL1d1RjUiuKzzi+1Pmn3oxFzY3tnFum5X1tjDhLHrYrz4TXjh+Jmbe4pV09yXtL40WsSj4rqKw3piaxLVbCVS2j6zR60ZUU8TtKragn2Rzxpvu5dyM3puqW97KVLEqNzBZnQqLE49/eu9ZRGpaXRu5+vhKVvcpYVanzfdJcpLuf2GBvYSjVp0NUp+oqqX4C6pNxi5fmy5xl+a+feS7S1zM0bcga/Z6xWs5Ro6s06TeIXiWI+E18l9/J93Iz8WmsriYmNNlbRZJGOJIMJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAGQ3jmgJ5GL1bVo2s1bW1P4ReSWY0k8KK+dN/Jj9r6ky31LValWtUstMcXUg92tcNZhRfYl8qfdyXX2PH2lObqzstMhv1t7Nxc1XvKEu2b+VPHKK5deFglENdrfEPie/G7jUruV9qdVP1cILG6uvdXKEO2T597wjK6bou7VheajKFe5jxpxS/BUfop83+c+PhyL3TNPoWNJqDlUqz41a0+M6j7W/uS4LqLzkJt9FcffciWEUbq5pW8E6kuMniMUsyk+xLrKNe5q1KkqFnFTnF4nUl8Sn49r7l54KlrZ06MnVlJ1a0liVWfPwXYu5EU/8A4oqlc3nGu5UKL/JQl7TX50ly8F7y8o06dGmqdKEYQXBRisJFK6u6NthSblOXxKcFmUvBfvKCoXF3xu5eqpdVGnLn9KXX4Lh4gfdS9UqrpWtN3FRcHuvEY/Sl1eHF9x8qyqV3vX1b1i/sYezTXj1y8/cXdKlClTjTpwjCMeSisJH2DX2+acIU4KEIKMVwSSwkfQAZAAAAAAAAAAAAAAAAAAAAAAAAAAAKFza0LjHrKabXxZJ4lHwa4orgCyxd2vJu6p9jaVRefJ/YVra5o3Capy9qPxoNYlHxXUV2s8y3ubWlXxKScZx+LOLxKPgwxrRc21Ku1J5jUj8WpB4kvP8AdyKLuK9pwu479LqrQXL6S6vFcPAOtcWj/rKdWl/bQjxX0o/vX2F3CcKsFOElOMlwaeUwJhOM4qUWmmsprkz4uaFK4oypVqcKlOaxKMllNeBbytqlvJ1LJpLnKi3iEvD5r+z7ytbXNOvvJb0Kkfj05LEo+P8AEHnywl1p9xYKXqYzu7FrEqT9upTXdn48e58fHkUNPr1tPpxq6fJ3unvj6iLzKmv9G31L5j8scjaOZidR0uSqzvNOcadeXGpTk8U63j2S/OXnklE78tc013hkLG7t722jcW1WNSnLk19zXU+4rmp0aklc1LmxzbXkWvhNvV4KfdNLr7JrPmjPaVqdC/hJRUqVenwq0Z/Hpvv7U+prgzEwlW8SvgAYTAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADAAAAAAAAAAAAAAAAAEAcABIAAAAAAAAAAAACCQBHADgSBBIAAAAAAAAHUADBAADgABIAAAAAAAAAAAACBwAAAkCCQACAAAAAAAAIJKVetSo0p1atSNOEE5SlJ4UUutgfcmksvka1qOpVdTk7bT6sqVmm1UuoPEqvbGm+ztn7u1Rf3MtVhKVSbt9KS3nGb3ZV186fzafdzfXhcH9Weny1RKVWEqOncFGnjdlcLvXyYd3N9eFwcojXdqtabdoUdOtZ3sI0LD+rWFP2XWp8HPtjT/AHz92XxWyWVrQtLaFvb0406cFiMYr/1nxKsKcIQUIRUYxWEksJLuPm5q07elKrUmoQjzbMTO0q1ir7nKMIucpKMUstt8EWWat/8AEcqNr1zXCdTw7I9/N93MilRqXklWuouFFPNOg+vvn393V48rytVp0acqlWahCKy23wRhLymlTp0aUadOKhCKwklwRaVLmrcSdKx3Wk8Sry4wj3L5z+z7iFGtf8ailRtXyhynU8exd3Pt7C9hCNOChCKjFLCilhIHlQtLSlbtzWalWXx6k3mUvPs7lwLhcsEgMo8xwHAcAIZPVkDgA4doJCQEYQJAAciG12nzKpCKy5JGPLEzEeX2DDahtRs/YZV5rFjRkvkyrR3vdnJgb7pO2St1+DvKty+yjQm/taSJxS0/DRfl4aebQ3cg8wuumTRqcsUdK1Or3uMI/fIsZ9N1ipYWhXX1q0US9qzT+Zcf+T10Hk1HpnoVGlHZ65l9CvF/uL+h0rQq4xsrrrT66dHfMe3ZOOdht4l6UEaTZ9IdtXaVTZ3aSj3y02bX2ZM9YbRadd4Sjd0ZPhu1rSrTf2xIzEw3UzUv4ZjBJ805RnHejxTPow2gAAABgAQSBAwOAAjCLSdpKlUlVspKnJ8ZU38Sfl1PvX2l5hEoC2tbqFZulOLpVorMqcua712rvRN1awr4mpOnVh8SpHmv4ruJubanXS3sqUXmE48JRfcz4j8Pgt1/B6uPlNuLflhhhFC4nGqre7ioVX8WS+LU8Ox933l2WVeneV6bp1KVq4v8+Xv5cGVLGF1Tpbl1UhOSfsyjza7+8C11nTKd6o1ac3Ruqa/B1orLXc18qL60/sfE15+td5ClXzZanRTdKpDjGpHr3c/Hh2xfFd3Bm5tcC0v9PtL+l6q7owrQTUlGS5NdaMxbXlC2PfeFDQNTeo29RzpblWjLcnKPGnN9sJda+7kZM+KdOFKEYQioxisKKWEkfZhOPHcIJI4BkAAE4RBIAAAAAAAAAEEgCAOBIEDgCQI4EgAAAAAAAAACAAAJYAAAAAAAAABgAACAABIAYAAAAAAAAAAEEgAAAAAAAAAAAAAAAAAAgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAZAAAFpquoW+nWvr67fF7sIRWZVJPlGK62GJnSpfXdvZW07m5qRp0o85P7Eu1vsNYv7mV6/hmpr4PZ05J0raXNvPBzS5yzjEFnHe+Xxc1atStC/1PLqb27bWtP2txvqivlVH1y5LjjCy3ltH0mfrYX+oqLuI5dKknmFBPs7Z9svJYXOWtNM2m86h86fptW8nC61Gm4Uk96jay7eqVTtl2R5LvfLOYCPivWp0aUqtWajCKy2yO9t0REQXFanQpOrUliK7uL7kutlrb29StVV1drDXGlS5qn3vtl93UTbUqlesru5i44/E0n8hdr/Of2cu0r3dxC2pb88tt7sYrnN9SXeAuq9O2pesqPhnCS4uT7EutltQoVbipG4vEluvNOjzUO99svu6u0+rW3qOr8Ku8Os/ixXGNJdi7+1/uLzCxwQEjgFwGAyAgngAA6gAAYXIARkPgjxfpU6V9a0HW6ug2OkOwqRWY3d1ifrI9U6cU8Nd7b70SpSbzqHPyORTj06rPYby8tbOhK4uq9KhSisynUmoxXi2aHrXS3sxaTlR0+pV1SrH/o8fYz9N4Xuyc/6zreq63W9dquo3N7POV62eYx8IrgvJFrb1HSqxqccLmu1HdXh6j9TzXK9fvPbFGnrer9Kuu3rcLCnb6fB8sL1k/e+H2Gm6vqurao277VL25z8mdZ7v6K4fYY2LSSlFpp9aPr1hKuOsfCkzc/kZZ/VZTjTdNtxe75I+ncSjBrebwJ1E+GCjOO9F8OD4GyIhz9cz5e1dG+y+xmv7NWupT0uFW4x6u4jUrTnu1I8JcG8cefg0bza7K7N2sUqGh6fTxyxbxz78HinQlr70nan+S7iri01LEI55RrJey/Nez44Og4LgivzdVbae39K9nNhiemNwpUbS1orFKhTgvzYpFXcj81H0iOs07lbRSseIFGPYhhdg8xlDaWoSBnvHmGQAZwAAIYEgAAAAAAAAZCADHHII8wJaRGExwznIz1ASQwaJ0sdIVnsbYRo0acLvVq8W6FtvYUVy359kc+b5LralWs2nUNWbNXFWbWbhqWp6fptrK61C8oWtCHxqlaahFebNI1Hpl2DtJuFPU6t41z+DW85r34SfvObNqNb1jaa+d9rV9VvJ5e5CTxTp90Icl95iFCMeUV5nfj4Uf+peezeu23/jh0/Q6cdhqlTdq1tQtov5dS0lj+7k3LZvavZ7aOm6mi6taXqj8aNOot6PjHmvNHFzUWsOK9x92nrLW5heWletbXFN5hVozcJxfc1xJ24Ea/TLXj9dyRP647O6U8rgSeGdEPS9VrXNvoG11aPrarVO21DG6pyfBQqLkpPqkuD68dfuSaa4Ffkx2xzqXouNyacinVVIANboAAAAAAAAAAAAAAAAAAAAAAAAEAAAAAAAAAAAHAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAcQAAAAAAAAAAAAAAAGAAABAMx2sapTsKcIQputdVsqjRTw5Y5tvqiut9XjhBiZ0+9W1Gjp9KLmpVKtR7tKjD49SXYv3vkus16tKr8KhdXcXc6hVzGhQpPhBdcYZ5L503/BHzmtG7U5YvtWuY8Ir2Ywjnq+ZTXbzb7WbBpGmQslKtVn6+8qpetrNY+rFfJiupeby+JLw1d7yp6RpfqKnwu7lGreSjjMfiUl82Hd2vm+vqSyoWOoPkRltiIjw+ZSUIuUmkkstvkWNCLvq0bqomqEHmjBr4z+e/3e/wAFb+v13RXG2pSxVa5VJL5Pguv3dpe1J06NGVSclGEVlt8kgeXxc3EKFJ1JvhySXFt9SS62UbOhUnV+GXSXrmsQhnKpLs8e1/uPi0pyua6vK8XFL8TTkvir5z/Of2LzL3wAldhJBIZD5bwTy7S21G+tdNsa99fXFK2tqEHUq1aklGFOKWXJt8EkuseWJnXeVdz4ZawattZ0jbEbKzdPX9pNNsay4+olWUqv6Ecy+w5i6bPSG1TaK9uNE2Huqum6LBunO/p+xcXfa4PnTp9jXtPu5HhSowqVJVJOcqk3mU3Jtyfa3zZ2YuJNo3Kr5HqMUnVXcFb0jOjGE3GlqGoXEV8unp9XD96Rt+h9I+x+rqn6jWaNGpUScadynRlx5L2sH5/6fbesuaNFSl+EnGCWX1tHrNV09582uXHijong112lW5PWsmO3jbtOnOFSClCSlFrKafM+snJOxe3Ov7J3Mf5PuHWss+3Z1pN05L835j71w7UzpPYLazTNrtHjf6fNxnF7tehP49GXZJfc+TOLNx7Y1rwvUsfJ7eJbFlmp9J2xNjtloMraoo0r6hmpZ3GONOfY+2L5NfvSNuRBpraazuHdlxVy1mtocZX1vcWF7Wsrui6NzQqOnVg+cZJ8V/mUfWRXA9S9JjZ6VpqNptPa08U7nFtdYX5RL2JPxScfKJ47Cq285LjFk667eC5nE9nLNfhnNOu8p0H1cYfwLxZfX5Gv0ZyjOMoNxknlMz1vUjVoxqLhnn3PsMWq4bx0+H11n1lLhwPiTSXI+XLtwRa1WlUlCtCpSk4ThJShJc4yXFNeDOmtgNfhtHsxa6j7KrY9XcQXyKkeEl+9dzRzBGaXZk9A6DtpI6XtPLSbiri21LChl8I1kuH6SyvFI08inVXa89D5fs5uifEvf0QE01lcgyumdPbw1fbrbrRNjatjDWfhMfh0pRoSp0t6LlFJuLfU8PguvD7DBvpf2WXyNQz/AN2f8S/6ZNkYbabB3+j7kHdKPrrKUvkV4cY+T4xfdJnGekbS6lp0/gl9TqVqdKThOlUeKlNrg0n3PhhnZgw1yV38qP1Hl8nBf9Hh14ul7Zd4e5qH6s/4n1/S3sx1w1D9Wf8AE590a8stUoestK280vapy4Tj4r96Mh6rHW/JmyeNWPKkn17k1nT3T+lvZjONzUP1Zh9Ley/XDUOH/Zn/ABPCHBdUpe8pyjjhl+8fhqM/n/Ie9Ppc2WSzuah+rP8AiR/S9sv8zUf1Z/xPA3z6/efMufN+8fhqH5/yXvr6X9luuGo/qz/iQ+mHZVfI1H9Wf8TwCXn5s+GvHl2mfwtD8/5DoF9MOyuPxeo/qz/iQ+mLZX+z1L9Wf8Tn3Da6yfNiOLRn8+5LoD+mPZT+y1L9Vf8AEPpk2U/s9T/VX/E59kl1cT4lHL68eJL8JRn895DoR9M2yf8AZan+qv8AifL6Z9k/7HVP1X/M56lD86TPmS75e8z+Eoz+e8iXQz6adklzpap+q/5ny+mzZFfktU/VH/E54nDhjL95SlDvb8zP4OjP55yHRL6cNkM/itW/VP8AMj+nHZBPjS1b9U/zOc5QS7fefO7h837x+DozHrXIdG1unLZL1M/U0NTnU3W4RdvuqT6lnPDxPAta1W91rVrnVdSqupd3U3OpLPCPZFfmpcEu4xjiu8mL6sYwbcOCuOdw5uXz8vJrqz7qcfApSgn1cSq2kuSPnK5tHTEK6NwtpLHBIKo4lWe6lmTSXazHXddy4U8xXW+tmYjTZHdWuLqMU485fcdO+jvt5ParZyrpmo1vWappbjTnJ/GrUmvYm+18Gn3rPWcoT8De/Rx1WppnS/YUVJqlqFKpa1Fng/Z34/bD7Tl5WOL02t/S8s4csRHiXY6JPmPI+imewAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAIAkAAAAAAAAAAAAAAAAEEgAABHmSgAAAAAAAAAAIAEggkCASAAAAAAAAAABAEggIA+QD4k5ADrA7wLXVr2lp9jVuqqclBcIx5zk3hRXe20vM1NTuvhClKEa+rXrwo59iCXVnqpwzx7W+1mU2hqO41e2tOcLem7ma/Pb3Yf7b8kfWylsqjr6pNZlWk6VF/NpQbSx4vMvNdhOO0babfqtpkNH02FhSk3J1rirh160l7U3+6K6lyRkACDbEaCzv61RyjaW8sVqq4y/s4dcv4d/mV7uvC3oTq1M7sV1c2+pLvZR0+jOEZVq+PX1nvT/ADeyK7l/HtArW9KnQoxpU1uwisItGlf3Hba0ZeVSa/cvv8D7vqk6lWNlQlu1Ki3pyXyIdb8XyXv6i6o04UqUadOKjCCxFLqQH0lgYWSQGUE9QAEcjlj0yekGvcahT6PtMryja0VCtqrg/wAZJ4lTovuSxNrrzHsOo7urChbzrVHuwpxcpPsSWWfm7tFrdbaLaPU9duW3V1G7qXLy+SlJuK8EsLyOvh44tfc/Ct9SzTTHqPlZSpLfb3ePaRutdxXS4lRQTLfTzfX9rjZam6mu2keqNRzfkmz0Wl7Sw+t9ZpexVu3qs6vVTot+baRukGlgRHZx57bsrwguvBs2wm09xsltBR1Kk827ap3dNflKXX5rmv8AM1ZT44RFWSkkurK6zXkrFo004cl8WSLVl2pZ3FK7taVzQmqlKrBThJcpJrKaKuFzNE6CNRnqHRzYQqSbnaOdq2+yEsR/u7pvZRXr020+hcfJ7mOL/bWek/Qo7RbDarpaipValCUqPdUj7UP7yRx7b7soqWMZ4+HcdzzWYNdxxPtBbKx2j1azXBW9/XpLuSqSS+w7uFbzVQevYo7XUU8cmX+n3Lp1N2T9mfDwfaYxPrTPpTec5LCY28xMbbJLPXwKcpLPEtLK6dWhiTzOCw89a6mTUm314NGu6HTL7nU5ooO4q06katGpKnUg1KE0+MZJ5TXemfM5Z4lKos8ORLW0qfpnbq7oy2lp7UbI2WqJx9fKPq7iK+TVjwkvfxXc0bPnPWc5ej5tB/JG0c9FuZqFpqTSp55Rrrl+kuHikdHLGCpz4+i73vpvJ9/DEocU44ZyR6TexMNA29/ly0pKNjrWarwuEbiP4xfWTUvFyOuTTOmPZKO2Owd7pcIr4bTXwiyk/k1oJuK8Hxi+6TJcfJ7d4T53H97FMfLiq2qytq0KtGcqdSHGM4PDXmbdo21sWlR1WK7PhFOPFfSj+9e406pJb8oypyptSw4SWGmuDT78opzk+SLyYraHiMmDqnUvV1WoVKUa1CrCrSkvZlB5T8ylOfY+Z5fZalfabWdS0rOKfxoS4wl4r9/M9j6Ldmbzb/RKuoaff2VrUoVfU17eopOUHhNPh1NcvPsOXJHR3lCnBvedUYni+SPnD7Gz0mHQ1tDvZ/lTTMfRn/Arrod13HHU9Nf1Zmn3qOj8o5U+KvL91tdwcVFcD0+XQ5rv/wA0039GZ8S6Gtef/Kum+6Zn36fZ+T8r+LzBvCyvApTXHKZ6j/QxtB/800z3VP4HzLoW2gfLVdL/AEahmM+P7Z/KOV/F5c8rGT5lJcng9Rl0KbQvlqulv6tQpy6Edo3y1bS/0ahn8Rj+2fynlfxeYuSa4NHxLnzPT30IbS9WraV+jU/gR/QftJ/830r9Gp/AzHJx/ZHpPK/i8weOR8tcFwyepf0I7Sder6V+jU/gfcehLaF8Hqul/o1CX4nH9s/lXJ/i8mlFPOCnKGM4Z69/Qfr7/wCVtM/QqEPoN19/8r6Z5wqGPxOP7Tj0rlfxePNNLkfDzza9x65U6CdpG/Z1jSvONQovoG2nef8AdvSln8yoZ/E4/s/K+T81eU7/AHcCnVuKdNe005Y4RzxMr0k7P3GxmsU9HuNQtby7lR9bVdCMkqSbxFPPW8N//wCmpOTk8ybcu1nRS0WjcOa2CaW6bLmvWnV+M+9JckUZNdqKTn2+8b6Zs2x0a8JkuODYuh+hOp0ubNerftK/jJ47FCTf2Guppy5nqvov6HPUukiWquD9RpVtKbljh6yotyK/R335GjkTFccy7eDWb56xDqqK4I+iOGSSie2AAAAAAAgCQQABJHAkAOsAAAAAAAAAACABIIJAAAAAAAAAAAAAQBIIAEkEgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAMgAAABHkSAAAAAAAAAAAAAAAB5AAAAAAAAAAAAAI8iQAAAAAAAAADBRu68Le2q16nCFODnLwSyBq19cS9Zq9/Hi4ydOn37kd1L9NyNm0y2jZ6db2seVKlGHjhYNTsqcqtppVpUj+Eua8KlVe+tL7VjzN0RmWunnaSM46ieBa6jXnSoJUknWqy3KSfzn1+CWX5GGxQX9d1BvGbe2lw7J1P/APlfa+4uruvC2t5VZpvHKK5yfUl3tk2lCFtbQow4qK5vm31t97Zaw/rmoOpzoWzxHslU635Lh4t9gYVtPoSpU5VKuHXqveqvv6ku5ci6ADIAAAAAxO2Uak9lNWhSz6yVlWUcdvq5YPzYtOFvTx8yP3H6dV4xnTcJLMXwa7UfnRt9s9V2U201jZ+tBr4Fdzp08r41Nvepy84OLO7hWjcwp/VazqJYyl2YLiCy8lFLHeVocOrzLPagmG0bEwSo3dZrm4wz73+82DHBcDEbJrc0ne+fVk/dw/cZinlpfuJfCuyT+qSKecYYmt5JYxxRWhFbyWCqoxjHfnhRjxbfYa5a4ncxEPf/AEbISWxN05Z3ZahUa/Rh+89SNP6HtKnpOwGm0asNytWg7ionzTm95J+CaXkbgUmWd3l9C4VJpgrEolyZxl0izj/P/aPdxj+U6/7bOy68406M5yaSim2+44P1XVJanrWoajnhd3lauvCU5NfYzq4MfqmVV67O8cQrupFM+XW4lop8Ez6UuPAs5eX6dL23uZUa0ZpvC5rtRl3OMoqSbaayn2mvpvuz1GS0uq2nRm+fGP8AAhaqNl5jxCcVxwGuGeRTly5kYavKvRryo1IVaM5U6tOanTmucZJ5TXg0jrHo92hpbT7KWeqxaVWcNyvBfIqx4SXv5dzRyJOe7jjwPS/R02s/k7aqts/c1VG21L2qOXwVeK/2or3xRz8rF1V3Hwu/ReV7OXot4l0h5CSzFoJ8ESVT2flx76TOx72b29qana0nHT9Z3riDS4Qrr8ZHzypfWfYeXx4xznKO1+m/Y5bZbBXlhRhvX9v/AFmyfX62KeI/WWY+ZxlCgowSaafY1xXcy44mXrrr6eT9Uwe1k38StnTb6keh+j1tVPZTpFoU7mpu6bqija3WX7MJN/g5vwk8eEmaO4rd5EwSw12rDwdWTFF66lwYc84rxaHW3pC7T7YbFbJUNptmI2ta0tayjqVKrQc5KnLCjUjxXKWE1+dnhg8g07p7211Cj661vNIqR618DalHxW9wPaOiTXLPpE6KVZ6uo3VVUZ6dqdOXy3u7rb+lFqXn3HEe2uj6hsHt3qmztSrUhc6dcOFOquDq0nxpz71KLT8clfxcdJmaWju9Bzb5ZpXJit2l71Lpw2+i/wDfGlfqb/xnzPpx2+T4XGlfqT/xHjGj7UUbhqlqUY0JvlWivwb8V8n7vA2aNLejGUeKkspp8Gjp9jHHwpL8zk0nvaW/Ppz6QFyr6T+pS/xlN9O3SF/b6T+pS/xmiK3y848iHbp8MD2Mf01/mPI/k3l9PHSH/b6R+pS/xHxLp56RE+FfSP1J/wCI0Z2i57uT5drFcGn4j2MX0lHqOf8Ak3eXT30jJfj9I/Un/iPh9PvSRxxcaPn/ALi/8Ro7tFjlgpytI8jPsYvpn8xzfyb0+nzpIb/31o6//Bf+Mj+n3pIT/wB9aP8AqL/xGhytFjOSjUtmuozGDF9Ec/PP/p6Iun/pF67rSP1F/wCM+ZdP/SNnhd6T+oP/ABHnPwdyluxjl9hUp2lGDUquJ4+SuC/zHsYvpP8AH5/5PSrXpw6TblOcL/SYU1xlN2HBf3jo/ZjUtUsuju21na+4pK9haO6vZQperjTWHLG71NRwn3pnM/QpoENqdvLGwnTUrK0/rd0kvZ3INbsX9KTivDJ656UG1EbDZq32YtquLjUnv10nxjQg03+lLC8FI4s1KzkilYWvF5GSMFsuSf8A4542s1a52k2jv9evMqte1nU3PmQ5Qh5RSXkYrlwZWXBJPmfL4vJZ1rERqHnbZZtaZn5UJPtT5lObSWeJVrOMX7Tx48jKbK7IbS7W3Co6BpFxdxziVdx3KEPpVHw8ll9xC9op5luxUtl7VhhrWNxdXtGytKFS4ua9RU6NKmsyqTfKKR2h0K7FR2I2NpWVdxlqVzL199NcV6xr4qfzYrCXm+swPQx0P6fsVu6rqdSnqOvSjj1yj+Dtk+caSf2yfF9y4Hq3BIq+TyPc/THh6X07g+z+u3lIAORbAAAADrAEEgCCQAAAAAAAAAAAAAACB5EgAAAAAAAAAAAAAAEEgB5AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACORh9r6mNFqUIvErqcLdeE5JP7MmZZr200lV1TT7bi9z1lxLyjuR+2f2GYjujfw+NKirjaNSxmFpbNrHLeqSwvsg/ebIuRg9k6akr67/tblwi/zYJR+9S95nBbyxTwhllbL4TfTunxp0s0qXj8qXv4eT7SpqdadKhu0uFapJU6f0n1+Sy/IrW9KFChClD4sIpIwkpX9aVKkoUseuqy3Kfi+vwSy/IqWlCFvbwpU87sVzfNvrb7y2tf6zf1blv8AB0s0qXe/lP38PJl8CAABkAAAAAQznT0vujqvqVlT270W3lVubKl6rU6VNZlOgnmNVLrcMvP5r/NOi0j5qQjUi4yWU1hp8ieO80tuGnPijLSay/NqHf18fFFTw595656TmwuzmyW1tlU0Cq6E9ThUr1dOjH2KCTS34P5MZN4UeXB4wuB5XRt5TaiubwkXOO/XXcPKcjH7NprLcNDoqlpdvHk9xN+fEyMMJIp0qKpQUG+SS9xVptLBu8KW87mVSm3k2/ov0K12h2wstPvKlONtHNerCT41lDD3F25eM9yZqCnHzK9rqdzpl3QvrGrKjc29RVKU18mS+9dq61lGrJEzWYhPi2rTNW1o7OzYRUIKMUkkuCPo1ro72mobV7MW2qU92FZrcuKSf4qqvjR8OtdzRsmOBRWiYtqX0TFet6RavhgOkSlqlzsPrVtoai9Sq2VWFspPGZuLS49vZ34OFaEnCXqnTlTcPZcZLDi08NNdTP0InBSjg5k9Jjo4emajLbTR6DVndTS1GnFcKVV8FV8JPg/zsPrZ2cPLFZ1Kp9W41sleuPh45GXDgVItrmijCMoxWXl9pVjPHUy41vu8rb6XMH3o+1VcZb0Xh88lrvr/ADJ9Yt3kJiENM7RufXUlNcJfK7mRObZh7O7VKth8IS4P+JkKj78GqY1KE10VJ56ylSuKtpd0by2qSpV6ElUpTXOM4tNP3ktrsx3nxKMZdSaMTG4SrM1mJh2T0cbS0NrNkbHWqOFKtTxWgn+LqR4Tj5NPywbIc2+jHtOtN2iutmbmoo22oP1ttl8I14rjFfSivfDvOkimzU6LzD3fAz+9hiUNJo5M9IzZNbN7bzv7alu2Gr71xTwuEKufwsfe1L6z7DrQ0jps2V/nbsHeWdGmpX9t/WbLt9bFP2frLMfMlx8nt3iUfUOPGfDMfMONJJcUuRSct1PGCZ1U5bu649z4Pw8T4msrJd9W/Dx3RqdS9G9HTbF7NbfU7G5q7mn6xu2tXL9mNXP4KXvbh9ZdhsPpubDu50uw6QrCjmvY7tpqW6udCUvYm/ozePCfceH3G9Gm5QcoyjxTi+Kfau87F6NtXsuk7omVvq8IXE7i3np+qUvz93dk+7eTUl9JHBnicd4vC84GSMlJxS4Qoe1HEsJ5wZvQ9WvdMe5Taq22cyoT5fVfyWNqtmb7ZPavU9nL7edawuHSc5LHrIc4TXdKLi/MoUaa4LHHtO2urxtw566may9D0a9tNUpp28t2qlmVGfCcf4rvRkfgnW4nndo3SqQqU3KE48Yyi8NPuZuWi7SRlGNHU8Z5KvFftJfejVasx4VmSkxO4ZCVulyWCnKgl1YMolCrBTpyjKDWVJPKZDt21zZDq01RMsS6Cl1EStufAybt2fM1CjSnXqzhSo0/j1JyxCPi+3u5szFmYiZ8MZ8EzyiW2oUqFovw7frH+Sj8bz+avEo6ptLGf4LSVKEeTuJrE39BfJXe+PgYZTk17WW3xbfFtm+sT8tsRMeV1VrOfBKMIfMjy/zKU293OCnGRsPR3s5U2u2ysNCpuXq7ipmvJPjCjHjUfu4LvaMXmKxtuxVte8Vj5e/+jJs2tE2EqbQ3kVTudWfr96XDct45VNcep8ZfWPCOkraqe1u3Wo6upN20p+qtE+qjDhH38ZeMjoP0idep7LdHMdF03doXGpL4FQhDh6uil7bXco4j9ZHKkIJOPDG6cnGibWm8rT1C3t1jDC4k+I3steJTnJYGn0LrUtVtNMsI793d14UKMcc5yeF5dp3TfphT0xzeYiHvHoxbD2epU7zanWLKjc0Yy+DWMK1NSjlfHqJPhnOIp90joalQpUoRhSpxhCKwoxWEvIxexmh2uzey+n6HacaVpRjT3uub+VJ97eX5mYKLNknJaZe34nGrhxxGu4kGC21G9trC1lc3VVU6a4Lrcn1JLm2+xGp1rkFK1r07m2pXFJt06sFOLa6msoqgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADIAAcCAJAAAAAAAAAAAAAAAAA6ggAAAAAAAAAGRwAAgASGQSAAAAAAAAgABHXkCQQ/EASyMgeYDK7TVtQrReuahdTxuWtKFL3J1JftR9xtD4cTSY/1nTG+OdSuu35NSp/gRKrXf6bPs3byttDtKc17bpqc/pS9p/a2ZHPEiKwuHI+LmrChb1K8/iwi5PyI+U4jstI/1jVJz5wto7se+cuL9yx72VdSrzoWrdLjVm1CmvznwXu5+ROnUZUbSCqfjZZnU+k+L+8pca+q8eMLaGfry/gv2hoXNpRhb21OhDO7BYy+b7yq2OCAZBkjrHmBIIyEAJI8AgJKderTo0p1ak4whBOUpSeEkubZ99Z5B6VW2C2b6Op6Xb1HC+1uTtIbrw40sZqy/R9nxmiVKza0Q1ZskY6TaXOPS7tRLbHpB1LXIycrZyVGyT6reHCLx1bzzP6xgNJTqanbQecSqIs4vjnGOGEZbZqKer0pL5ClL7C9x16YiIeL5OSbzNpbZNttvtZTb4CU97/I+G+XE2yqYhKbfWfXFrDkj4z7ipHjHiyGk5bl0RbX1Nk9p4/CKmNLvHGld55QfyanlnD7n3I6jozVWEakZKUWspp8zjGlTi3xXimj3zoG2v8Ahtj/ADZ1Ctm6tIZtZyfGpRXye9x5eGOxldy8P/qHpPRef39m0vVi11bT7PVdMuNNv6ELi1uaUqVWnJcJRaw0XSBXx2enmImNS4k6VdlLvYfauro9xv1LWS9bZ3Evy1HPDP50eT9/Jo1b1meLOy+mjYa2262TqWS3KWpWzdawryXxKmPiv82S4PyfNI4quqd3Y31ewv6E7e6t6jpV6U/jQmnhxLriZ/cpqfLyPqPD9nJuPErp1CPWSzj9xQUkuOT6Ul2nVpWaVlMyVjX9ZS3Jv2oLh3ow+8u1lSlWdKanF8U8+Jia7YtXbOcU+LPmUspnzGpCpSU4yynx/wAj4nKKWcmvTXpd2F7c2N3RvLSo6dxQqRq0pr5M4vKfvOydgdorfanZOw1q3cV6+n+Fgn+LqLhOPlJM4mq1fYfF4wevei9tf/J+v3Gyt5Wxb6ivX2m8+Ea6XtR+tFZ8Y95ycvF1V3C59H5XtZOifEumG8FOolJPjkmMk45WcMlpFVD1vmHHvpC7IrZjpArXdtRcbDVt66o4XCNTP4WC82pfW7jz2KTOxunTZBbW7CXVC3pb+oWf9as8Li5xTzD60cx8WjjyjHKznPHh3IueLl66PI+pcecOXfxKnOnlcEj0/wBG/at7N7c09Nuqu5p+sbtvPL4Rrfk5eeXH6y7Dzfda4pce0+op78GnKEovKlF8U+prvN+XH11mHHx804skWh7d6ZexsZ2Vjt1ZUfbobtnqDiudOT/BzfhJuOfz12HNVOSeDuTYfUbDpQ6JZWuqxjUldW07HUYLnGoo4k12N5U12ZRxVtJod9sztHf6BqKfwrT7h0KjxjfS4xmu6UWpLuZz8PJ2nHPmFv6hii0Rlr4l9UOovqMd7+JY20stF9QlFdZ1zCjtDMaVdXVk06M96m+dOXGL/gbXp19QvEoxzCr/AGcufl2mn22HFF5W9iyuKi4OFOUk/CJptjizTaImWT2j2i0zRnKjOSvL5cFbUpr2H/pJ8o/RWZeHM8+1XVtQ1e4Va+qpxi/wdGC3adP6Mf3vi+tssKUIxhFY6itHHA2UxRV1VrFY1CrTk4/KzllzTqrGHz7SzTWcH3l9uCfgmsSyFOcW8HSfombLK20a+2uuaW7UvZO2tG+qjB+1JfSmv7iOaNDsL3WdodO0bT03c31eNvT/ADXJ8ZPuSy33JnZPSFq9p0cdEkqOmYpToW8LDT11+sa3YvvaScn4M4eXeZ1SPlZ+mYq0mc1vh4X0/bS/zj6QbiNCanZ6d/VKDT4Np5qS85cPCKPOqifYfcqm/wAc5eObecspuT5Y49R1YqRWsQqeTmtlyTeVvVUkuvB7B6KWyf8AKW1V1tPdUs2+lx9Vb5XB15ri/qwf988mjGc5RhCEqk5NRjFLjKT5LxbO1OijZanshsLp2juMfhEYetupL5VafGb8m8LuSObmZOmuln6RgnLk6p8Q2iHsxxxPpPiRjHWM4KiIesfWcmE1mGNfsa0/bi6NSME1ncmmnld7jn3GZ3jDbUycLa2uVw9Rcwk33Se4/sl9hKEb+FbZduOmytm+NtWnR8k8x/utGVMHoNVx1a8ovgqsIV148YS/Zj7zOCfJSdwkAGEgAAAAAAIygJHIhtDqAAe8cAJAAAAAAAAAyMgABwAAAAAAAAAAAAAAAAAAgkAAAAAAAAAAAAGQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQSBBIAAAAAAAAAAAAAABBPkAAAAAAAAAAAAAAACASAAAAAAAAAAAAAAAQSAA8gAIJAGP2hryttEvKsH7apSjD6TWI/a0YXTqEVq+m2ceMLanOp+jFQX7T9xkNqJb8LO1/trmLl9GCc/viijs5H1usX9x1U4U6K8eM5ftRJR2hqt3s2BFjqX4atbWi5VJ78/ox4/fuovnyLG1/C6nc1uqmlRj98vvS8iMNkruclCnKcniMU233FvpUGrX1s01UrydWXdnkvJYXkfOqtyoQtVzrzVP6vOX2Jl5HguWAykgZJAgkeQAAAAAHwAiTwjiL0idq3tf0j3tS3nv6fpjdjaYeVLdf4Sa8Z5WeyMTqDp52vexvRvqOo0Kijf3C+CWPb66plJr6K3pfVOIbdYpqKy0uHE7+Fj3PVKk9W5GoikPuMWuRl9msq7qVOHs08Z8WY5x4cjMbOUsQrVH1yUSziO7zmW36WZjLi/E+4rOFjifMUuwqw9lLCJSr+z6jT4rK6itFY5xwek7F7I/zq6J72dpSj/KllqFWpbPlv+xDNPPZJfakzzd+shNwqQcZJ4lGSw01wafY0aaZYvMw6M/GvirW8+JVIcOovdKvrjTL6hfWVV0bm3mqlOa6pLt7U+TXWiyjJdhG828JcOslaImNS5sd7VtFodXdH+1FrtXs7R1KjuwrL2LijnjSqLmvDrXc0bDzOUujna652Q1+N0nOdhWxC8orrh1SS+dH7eKOodLvra/sqN5a1oV6FaCnTqReVKL5MpuRinHb+nu/TefXk0iJ8rtxT5nhfpI9E9XaKnPazZq33tYoU0rq2guN5Tjycf8ASRXL5y4dSPdOYccriasWSaTuHbyMFc1emX52RbjJxeVx6+DXiVYvjg6O9Ifoid58I2v2VtHK8w56hY0o8a666kF8/tXyufPnzZv8E3jjx4F9gzVy128jyuJbBfUq2/hciHV48iipPHMb/Wb3LpktOusS9TLhGT9nuZd1JPPDjxMBvtdz5mWtLj19Deb9tcJfxNVoRtXXdVk+1Hzb3V1YXlvfWNR0bq3rQq0aifGM4vKfvR81JFGcnngY1Expik9Ntw7m6PNpLXazZDT9ctd1K4pL1kE/xdRcJw8pJo2A5d9FHbH+TtpbnY+8q7tvqKdez3nwVeK9uK+lBZ+ozqFcuZR58ft3mHtuFyPfxRKJxyuByL087JPZjbyvVt6W5p+p711btLhGWfwkF4SeUuyS7Drx8jQ+nLZL+dmwt1St6W9qFn/WrPC4ucU8w+tHMfFrsJcbJ7dmv1Hje9inXmHICiml1hLhyKmU4Jx+K1lZPiXDBeRMT3eMmJidS9T9GvataJtvLSLmpuWes4pLL4RuIp7j+ssx8d0yPpj7FxU7DbuypcPZstR3V1N/gqj824N98ew8XjVrUpqpQnKlWhJThOLw4STTUl3p8Tr3ZS+0/pW6I/V6lCEo39tK0v6cfydZLEmuxp4lHxRX56ziyRkh6DgZIz4Jw28uK6O7vezyLyhnPBMnWdJvdD2h1DRtQji6s7idGpw4Np/GXc1hruZNDHYd0TFo3Cmy1mtpiWStZPcXAudVq+r0G9knj8BJe/h+8tbf4iwj52nnu7N3HbJwgvOSMR5c3mzVIyS8D7TWMot4N4w+PeVqb7UbXZMKvPqRE5uEN7CeOoLHgVaNpXvalO0tacqtxXnGlRpxXGc5NKMV4tpEbdo2VrNpiHufofbLO+12/wBsr2lmjYp2dk2uDqyWak14Raj9ZlH0otq46ztlS2ftKilaaPH8Lh8JXE1l/oxwvFyPaLG0sOiboYjB7s3pdk5VGvy9xLi/0qkseZx7WrXN1dV7y7qurc3FWVWtN85Tk8yfvbK7DHu5JvK05lvw+GMUeZSpuCHrMvLKc5Y4pFvXqqPPhwznwO/elNFdvV/Ru2b/AJx9IsbqvTcrHR1G6qZXB1eVKPvzL6h11wPMvRw2Tlsv0eW87qluajqcvht1lcY7y9iH1Y44drZ6aUvIydd5ew9PwRhwxH2e8pXVaFvQnWmm4wWXhZb7ku0q8S0ufw97St18Wn+GqeXxV7+P1TQ7pXFCpGtRhWjxjOKkvBlprts7vSLq3ivanSko/Sxw+0q6d7Eatv8A2VVpeD9pfY8eRcy5Bie8NV0+5T1HS73lG4hKk/rxU19sftNrRpFTNrptVR56fdOX1YVN5f3GbrB5imuRKzXj+n0ACLaAAACOJPkBAJHHsAgkAAEAAAAAAAAABAJAEEoAAAAAAAAAAAR5ASQBx7ABIAAAAAAAAAAAAAABBIAAAAAAAAAAAAAAwACAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADA1/WKnrdfo0uqhaym/Gckl9kZFxsnBPT6tzj/fFxUqeW9ux+yKMTfXO5d6zfc/VNU4+FOGcfpSZsWi2/wTSLS1xxpUYRfikskp8NVY3ba5rSUKcpyeIxWX5FvpUHCypynnfqZqS8ZPP7xqz/AKlKkudWUaX6TSf2ZLpcFhLBFtWjaq6slzVCln60v8l9pdos9L9v4TX4/hK0kvCPsr7i9DEAADIAAAAAEN4XIeBrvSTtNbbIbFantDcpNWlBypwb/GVHwhDzk0vMzEbnSF7RSs2lzT6V21v8t7c0dnrWe/Z6LHFTHKVxNJy/RjhdzcjyChFQilEm4ubi91C5vL2o6tzXrSrVqjfGc5Pek/e2VYQT59ZeYqRSsQ8dyc05Mk2lUhxXIz2jU3CwUl8ubf7jCxjCK3pPGFk3PUtFudDr0tMvFu16dvRqTjjDg6kFPdfet7DNsTG9OLNE9O/haRUmuwr04PC4nxHCZUU1FYZOY7OKZdB+jMktjL+Oc/7o1P2IGuekDsU7O6ltXptL+rVpJX8IL8XN8FV8Hyffh9bM36MdXe2S1Fdmoz/YieqXltRvbSrbXNKFWjVi4VITWVKLWGmikvknFlmXs8fGryuFFJccL7j7ikbL0lbG3OyG0Do01KpplxmdnUfUuunJ/Oj9qw+01qHFcOss6Xi8bh5DPgthvNLJdPea5cOJ6H0R7dVNmryOl6pVctIrS9mT4u2k/lfQfWurn2nn64M+4tZy1xI5McXjUnG5N+Pki9XYtGpTrUo1KcozhJJxlF5TXaj6PAOibpFloVWnous1W9Lb3aNaXH4M31P8z9nw5e+0akKtONSElKMllNPKaKjLinHOnveFzacqm6+UyWUc8+kN0LyvXcbW7G2n9e41L7T6S/3x21Ka/tO2PyvHn0P1hrJjFltjncN2fj1zV6bPzkbcZOL5ru6yVJZ5nT3pDdDK1mNxtZsjapaqk53tlTWFeLrnBdVX9rx58sTqbkt2W9GSeGmsNPrTXNMvMGeuWv8AbyvK4lsFtfC4c1w4n3a3Pqayln2Xwku4s3NvkQ5S7Ebphya22KTTw08p8u8ozaRZ6XXcoujLnHjHw7C7fPiiMNU10q6Zd3Om6ta6rZVPVXdpUjVoT+bOLyn4dXg2d0bAbR2u1eyWn67aYUbqknOGfxdRcJwfepJryOD3NLqR7b6Jm2isNoLvY69rKNvqGbmy3nwjWS9uC+lFKX1X2nFzcXVXqj4W/pHJ9vJ0T4l1CfMlmIi01nPAniVG3qfMOQunvZZ7Mbe1nb09zT9SzdW2Fwi2/wAJBeEnnwkjQop+XgdddO2yMtqNh6/wWjv6jYP4TapLjJpe1BfSjlY7cHJKXLDfHivAueJk66aeP9U4/s5d/EvlRS5I9X9GbauGhbVS2fuam7Z6thQzyjcJez+kvZ8VE8ql2ZKbqVrdxrW85U61OUZU5xfGMk8prvTWTflpGSkw5OLyJw5YtD2b0utlI215Zba2dLEa+LO/cV8r8lN+WY5+ieF2rUlzOwtIuLDpV6IpUL3dXw+1lb3SivxNePBtdjUkpLyOQa2n3mj6td6TqFP1d3Z15UK8erfi8ZXc+a7mjm41+3RPwsPUscTPu18SvrbKS+wtts57mhUofPrx+xNl7bRct3gYrb6TVrY0sc5zl7kl+864jupqd7w1ylwgnnzKiZRg8R5cj6jLhjHE3adi430+PYey+iXsote26q6/d0lKy0SKlDPFSuZ5UP0Y70u57p4o5xWetrjhLi/A7c6IdAtejPogpS1VxoV40J6jqc31VJR3pL6qSj9U4eZk6a9MfKw9Pwxa/XPiGg+lhtWq13Y7IWtTMKKV3eJPnJ8KUH9sv0TwRpLgi/2m1m71/aTUNavU1WvazrNZ+Kn8WP1YpLyMa231k8GPopEOLmZpzZZtKKjfNI2fod2Vnth0j6ZptSG/ZUJfC71YyvVQae6/pS3Y+bNYqNbq49Z1N6LOyUdI2LqbQXFLF3rElODfONvHKprz9qXmiHJyRSjb6dgnNlj6h6/ShuxwlgqYzzCXYSUsf29hEajQyz0178atZ8JzqyUl2YeEvci8ZZWf4O+uqPU5Rqr6yw/ti/eZJSn6vVO6tS+2L/hL7C75otNR9h29f+zrRT8Jez+9F3ngYGt31CD1PUrWS9m4pwqY+lFwf7KMrs7cO40KyqzeZujFT+klh/amWOvJ09Xs6y/KUqlJ+KxJfdIqbJS3bW7tv7C7ml4SxNftE/hrr2szS5AhEkW0AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQAAAAAAAAAAAAAAAAAAAAAAACAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAIk8LLJLHX6/wXRL65T407ecl4qLDE+GrU18I0ukmuN9dKT71Uq7z/um7I1Owobt/o1phYpZm19Cnu/fJG2ErIY1pdPfvrSl1KUqj8lj75Ir3VRUbarVfKEHL3LJQh7erzf9nQS/Sbf+yNY42Eqa/KyjT/SkkRTVNNpulYUIP4ygs+L4suAuC4AMgAAAAAAQ+IE5OY/TG2rdxf6bsba1E4UGr2+SfOTyqUH4Lel5xOjtc1K00bRrvVb+qqVrZ0Z1q038mMU2/uOCdp9ZudpNpNQ1++T9ffXDrOLedxPhGPhGKivI7OHi677n4VXquf28fTHyxkIbrazkr0pYaeRuJnzOL3Gk9144PsLWezy8d+z0PoO2ajtb0gWFnVhv2Vni8u8rKcINbsX9KW6sdmTZenBtdJ+rpdXqV/8Aqgek+i3sotE2C/ly4p4vdZarcecaEeFNeazL6x5Z0318dK+tx680f9TA5MOTrzz/AEsOZgjFw4+5alhrtPiTeeGRGeeecn0odvWd9vDz0R3e/wDouxb2T1N//wAhL9iJ7ElwPIfReTWympr/APkJfsRPXzz/ACP3Je/9NjXHqxG1uz9htJolfS7+nmnUWYTS9qnNcpRfU1/lyZy9tVod/s3rFbS9Rp7tSHGE0vZqw6px7n9jyjrg1TpI2Nstr9GdtWao3dLMrW4S405dj7YvrX70jZx83tzqfDn9U9Ojk06q/wDKHLblxwmTFt9Z96xpmoaJq1bTNSoOhcUXiUW8prqkn1xfUynT48+KLWJiY3DxN8dqT0yrRWVwZ6Z0Ubf1NBlDSNZrSnpbeKNWXF23c+2H3eHLzSnhLgVoSTazjl1mrJji8als4vKvxrxarr6jVp1qUatKcZwkk4yTymu1H2l1ngPRd0hT2eq09J1apKekyeKdR5btm/vh3dXge929alXowq0pxnCaUoyi8pp8mmVWTFNJ7vecLm05VN18vtngfpGdCsNo6dfavZG1jT12K37q1hiMb5Lm11Kr2P5XJ9TPfSMIjjyWx23Dfmw1y16bPzWrwr0LidC4oVKNWnJwqU6kXGUZJ4cWnxTXJoRb6zrj0i+hmG1lGrtNszRhT2gpQ/DUItRjfQXU+pVEuCl18n1NciV1Wt61S3uKdSlWpScalOccShJPDjJPimutF3x80ZY/t5jlcS2C39KsKsqUlOLaaeUZONwqtJVIcE+a7GYF1OHNsrWNx6uq4Sb3Jvjx5PtOiYcdqbZWU8rhyFhfXOmahb6nYVnRu7SrGvRqL5M4vK8v3ZKTafcfMkpRafJ8DEx1RpGm6zuHfnRttLa7X7FaZtDaYjG7oKU6af4uouE4PwkmvI2M5c9ELbGOn6zdbF3lXdoX+bmw3nwVaK/CQX0opSX0Zdp1EuKKDPj9u8w9jw80ZcUSNJrHUcjdOuy/82Nu7hUKe5Yahm6tsL2Ytv8ACQXhJ5x2SR10eddPeyD2o2GrztaW/qWn5urTC4yaXtQ+tHK8cE+Nl9u7R6nxvew9vMOTcYR8y48PtKdCTnTTy2nxy+wq7nHjzLqJ28XManUvVvRn2r/kra6rs7dT3bTVVmll8I3EVw/SimvGMS/9KjY1WmsWe2dnTxSvN21vt1cFUS/Bzfik4t/mx7Tx+zrVbS8o3NrVlSuKM41aVRc4Ti04teaR13p1fTOlLoqca27GGoWzp1lHi6FePPHfGaTXgu04M8TiyRePC+4M15OC2GfMeHI1DMWkYLb6rvXFnTa4Royfvl/kbNqNjc6XrF1pd5FQurWrKlWS5KUXh47uteJpm3NXe1pQXyKEF78v953Y53O4Utcc1y9M/DD+saPqM+JRWF4nxOWN2O8lx5tm2Z07Yrvw9L9HLZeW2HSvp9KtT39P0vF/dtrg9x/g4P6U8cOtRke7eljtYrTQrLZK1q4r6hNVrpJ8VQg+Cf0pY/RZceipslT2S6LFruoRVG91pfD7iU+Dp0En6qL7lD2vGbOfukfaartftxqOvTcnRrVdy1i/k0I8ILuyvaffJlZH+bLv4hY5rfh+P0x5lim23lrmfLXPmfLqcePIb65o74UTL7F7PV9q9rNM2ftm/wCuVlGrOPyKS41JeUU/PB3bp9rQsrKjZ21ONKhQpxp04RXCMYrCS8jwL0Q9lnG01HbG6pYlXbs7LK+RF5qSXjJJfUZ0GuRUczL1319PV+lcf2sXVPmUgA5FqMs6v4PVaElyqU5QfisNfvLws9TzF21X5leOfB+z+8MS+9Tg6mn14x+NuNx8VxX2lajNVKMKkeUoprzJaTi01lMttJ/4PpQfFwzD9FtfuMfB8rLaaGLa2r9VG6ptvuk9x/tFDQZKnrl5Rzwq0KdVLvTlF/7Je7Swc9Cvd3jKNGU4+MfaX2oxVlUxtDY1o/FrUqtPxzuzX7LJx4a7drNnAQItoAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACCQAAAAAAAAAAAAAAAAAI8gAJAAAAAAAAAAAAAQABIAAAAAAAAAAAAACABJBI8gIMTte86DWpLnWlTpY+lOK+5mXMNtS80bGn8+8p5+rmX+yI8o28LbS4qrtGp9VK1k/Oc1/gNhfIwGzizq1/Pnu06MPsk/3mwdRm3linhZ2XtXt7PsnGHuin+8ah7Ve0p9tbL8ot/uQ0virmfzrif2cP3C446nax+bCpL7l+8wkvFyAAZCMkgAQCQIbIfVzPopXNSnQozrVZxhThFylJvCSXNthiZ1G3g3pf7X/ANnLPZC0qf1jU5euuknxVvB8E/pTwvCMjmem00uDTM50q7S1Ns9v9V13elK2q1fU2mfk0IcIeGeMvGTMBDgXXGp7dHkedm97JK+inyWOBsGwGzNTa7a3TtBo53LurmvNfk6EeNR+5NLvaNbpzxhHSvokbNqjpOobW3FPE7ubtbVtcqUH7bXjPh9QcjL0U7IcDj+7miJ8Pb7W0o2dlRtbanGlRowVOnCPBRilhJeRyb04Rb6Wdb4/Ko/6mB7r0b7YS2r6TNuaFCtKemaRK1sbVJ+zKcfXeumvr5jnrUEeKdONFf0q62038aj/qYHJwYmMk7WXrUR7ETHhpcVw5lWEuGGUcNdTZMXwLiXkoju6I9F552T1L/7hL9iJ68eP+i487Jany/4Ql+xA9gPPcj9yXv/AE3/AK9UPwJZINEu5pvSVsLYbXadl7tvqNFP4NcpcV+bLti+zq5o5x1bTLzR9RradqNtO2uaLxOEvskn1xfU0dgGpdJGxVjtbpmJbtDUKKfwa5S4xfzZdsX1rzR1YORNJ1PhS+p+l15FZvT/AJOY+GFjCJUlg+9esb/RdUraZqdvK2uaL9qL5NdTi+uL6mWsZNotImLRuHi8mO1LdNo7rlVG+Rv3RX0g3GztxDTNVnKrpM5ey+Lds31rth2rq5rsPOoy7HgqxeY4zjvIXxxeNS2cfk341+qrsK0uaF1b069vVhVpVIqUJwllST5NMq5Zzn0W7f1tmLiOm6jUnW0ecuC5ytm3xku2PavNdj6Gs7m3u7ancW9WFWlUipQnB5Uk+TTKrLinHL3fB51OVTceVXOVyPDfSK6E7fbGhW2l2ZpQtto6cc1KSajTv4r5MnyVTHBS6+T4Ya9yzjkQ1nnghjyWpO4deXFXLXps/M67t7qxu6tnfW1a1uqM3TrUasd2dOa5xknyZ8N5SyjtX0gehqz27s5a3okKNrtLQhiM37MLyK5U6j+d1Rn1cnw5cX6naXmm6hX0/ULaraXdvN061CrHdnTkucWi74/IjLH9vN8rizht/S+sq/rKOJcZx4P+J9zmv8zD0biVGsprivlJdhfSqp4cXwfFPtOiXBNe68stXvdJ1Oy1PTqvqr2zrwr0J9SnF5We7qfczv7YDaa02t2R03aCyaVO9oKbh1058pwffGSa8j88ZZb4HS3oSa1e1P5w7O1ZudpQVK8op/InNyjNLue5F+Oe04Odi3TqWnpmWa36Pt0v1HzUW9Fxzg+0g0U0RMPRTG405G6a9kYbLbdXXwemoWF/m7tkuUMvFSHlLiu6SNHlFZPbfSxlu6poi/7NcftUzxOWHgv+Lbqxxt4f1HHGPPMQpyxFZPafRO1i5/lfXdF3m7T1NO6jH5s8uLfmt39E8Wq8ermeseidHG2uuL/sEH/fMcuInHKfptpjPGmE9IRQo9MV/wCrgo+ttaFSWOuW61n3RXuPEdpqiqa/d8M7sow90Ue3+khGX9MNZYxmwof7R4DqVZVNWvZ81K4n9+CXGn9MJ56/57SnhnijauibYue3XSHpOguMnaTq+vvmvk28OM/De4R8ZI1Pe/BuXLgdeehzsnHTtirja+7oqN3rEtyg2uMbam2lj6Ut596UTPLy+3j/ALdXBwzky9/EMp6Tm0a2c6Pf5DsZKlc6rm2hGHDcoRX4RrsWMR+scpU/ipLC4e49E6c9qVtX0hX9xTmp2Vm3Z2mHlOEG96S+lLL8Ejz6quPBmvjY+mm5aOfmjJlmI8QSfDLZV0azutW1i00mwpqpd3tWNCjH8+Twm+5c33FtVfVk9q9EXZKWo7U3u191Szb6dH4PatrhKvNe3JfRg8fX7iWbJ0UmWviYPdyRDpHYrQbTZjZbTtBsl+AsqEaUX1ya+NJ97eW/EzIXIgpJnc7eyrWKxqEggGEklnrH/B1d8cxjv+5p/uLwo3sXO0rR+dTkvsAqQaayustdOe7K6p/Mry+3Ev3lWxk52VCfzqcX9hTteGoXce1wl/dx+4MK9emqtCdOS4Ti4vzRqGnTaoaJWb9qFanCXnCUH9rNyfI0p/g7FdXwe+/ZuP4Eq+Gu/mG7LkCFyJItoAAABAEgEAGASAAAAAAAAAAGQABAEggkCCQAAAAAAAAAABAD3kkEgB5AAAAAAAAAAAAAAZAE+QAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAHEAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADiAAAAAGD2lebzS4f6ecseFKf8TOGD2h/4T0zudZ/3BXyjfwnZlL4TqU0vy8Y+6nH+Jm+ow2zHLUH/ANsl+zEzJmfJTws9I/3lvfOqVH/fYlx1iC7LeX2yX8CNE46ZRfPO8/7zJ/5Y8Lf/AGjDK8AAZAAAAAA8m9KLa3+b3R3U0y2q7t/rUnaUsPjGnjNWX6Ps+MkesSeEcUekDtb/ADs6Tb2dGpv2GmN2Frh8Huv8JPznlZ7Io6ONj67uD1DP7WKfuXn9KnBU4YWFjkfe5hZfA+1jOT7jGPWs54F3p5Dqnaroel3mt67ZaNYx3ri9rwoUuHJyfxvBLi+5M6y6WtftOiDoOuKumNU61raxsNMi+cq8luxl3tcZv6LPPvRO2QV1rl9tdcU26Vina2mVwdWSzOS8ItL6zPP/AE4NsKmt7cWuyVrU3rHQ6e/XSfCV1USfH6MMLxnIrsk+7l6fiHoOHSMOHrt5l6F6DrktK1aU5OUqlvbTlJ85NyrZbMN04TT6VNb48pUf9TAzXoWR3NM1Jdlnaf8A9rNW6cZNdLWuLLxv0ev/AEMDZgiPflyc6/VxIj+2sPi+BDT5op0pZSy8PJcQ4rkWUvOQ6B9Fn/ilqf8A9wl+xA9iPIfRcS/mjqX/ANwl+xA9ePO8j9yXvvTv+vUKNtdW1xOtChXp1ZUanq6qjJNwnhPdfY8NPHeiszjHpJ6SNoOjz0m9przR6+/b1alurmyqSfqrheop8GuqXZJcV3rgRx45yTqHRmzRiiJl2cHxNW6Ntt9F252bo6zpFfKl7NejJr1lvUxxhNdTXua4o2lPJrtExOpTpeLxuGo9JGxGn7YaX6mslRvaSbtrqMcypvsfbF9a/ecza3ouo6BqtXTNUoSoXNJ8fmzj1Ti+uL7ffxOxzWtvtj9N2t0p213H1VzTy7a5ivapS/fF9a6/HDXTg5E4+0+FT6l6XXkR11/5OVVh8sH3vYWC/wBp9E1LZzVaum6nQ9VVg8qS+LUj1Ti+tP7OT4mL3uziWtbRaNw8blxWx26bK6e9JcFlcjeOjHb662Xu42V26lfSZy9qC4yoN85Q7u1ea489BjL/AP0rRmRyY4vGpS4/Ivx7xakuwNNvbbUbKleWVxTr0K0VKnUg8qSfWi54nNHRpt3d7J3it6ynX0ipLNWiuMqbfy4fvXX4nRulX9nqVjSvrG4p17etFSp1IPKkiozYZxy9zwOfTlU7eV15njnpE9DNn0h6fLVtH9VZbTW9PFGs/ZhdRXKlVx9kua8OB7JhZDXAhS80ncO3JjrkjVn5h6tYajomr3Gj6vZ1rO/tpblehVWJQl+9daa4NFTT6jeaMuK5xf7juHp86HNK6R9MV5bersdorWDVreYwqi/squOcex84viutPibXNK1TZ3WLjSNZsaljqFrPdq0ai4p9TXU4vmmuDRc8bkRljv5ed5nFnDP9PvEcP7zoL0JGv52bSpf9Bt/25nPUavrIqXLqfczoD0IZb22O0qT/AOT7f/WTJcz9mWn0/wDfh1hkELmChiXq3O/pZv8A3X0Ps+DV/wBqmeKvPBo9q9LJZ1jRP+7V/wBumeKtPCLziftw8R6pP+xZ9Y4cfuPXfRRgv5562/8AsFP9s8hbeD1/0UX/AO+Ws/8AcIf6wcqP8Us+l/8AZqw3pF0s9MNWTxhWFB/bM5onUcqtSo+c5t/adM+krUVLpQvqvzNKpv7JnMzi3TT7jPF/bh0Z+3IsyexWi321m1mlbM2Ckq+o3MaLmvydPnOf1YqT8jt/pg2gs+jnoj+A6Tu29WVCGmaZTi+MW47qkvoxTlntS7Tx/wBCDYrer6tt7e0niOdP09yXPGHVmvPdjnukYf0ntq3tB0kz0mhU3rHRIu3ik+Dryw6kvL2Y/VZzZJ97Lr4h37/D4Jn5l5zSqLc3UsJPEc9hDeXjJRUmlhPgPWYWXxRYR4eftEzO1aNGtcXNGhb03UrVZqFOmuc5N4il4s7o6LtlaGxuxGm6DRUXUoU964ml+MrS4zl+k3juwc1ei/sutoekNarXp79losFcNtcHXllU15e1Lxijrxciq5uXdumHpfR+P0065+UgA4V0AAARJZTXaSALTR3nTLbtVNL3Cnw1aqvnUYP7ZEaN/wAG0U+pNfayY4/lZ99BftMMLpml3/s22rxXyK9WX3SN16jTtRSzrcfz6j//AFRJUQyNwi8xTXWSUrR71rSl2wT+wqkU48AADIAAHEAAAAAAAAAAAAA4gABxAAAAAAAAAAAAAAAAAADiAAAAAAAAAAA4gAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAMFtFw1TTH31l/cM6YTaRL4bpkv9LUj76Uv4Ga+UL+H3sv8AF1Bf9sn+zEzD5GF2Zf4XUo9l0pe+nAzRifLNPCy0L/gqh4P72T/yy/8Au/8AtDRuGnwj82c17psmXDWI8OdB/ZJfxDK7AAZAAAAIbwBovTptf/M3o71DUKFVQ1Cuvgtin/bTTSf1VmX1Thqk5QkuMpdrlxz2vzPZvSx2qet7f0Nn7ao5Weiw3amOUriok5fox3V4uR496vL4Fvw8XTTf28v6nyPcydPxC4hPrwXdpCrWq0aNtB1K1WooU4Li5ybxGK8W0WD9lZbxhdh6v6L2y/8ALu3q1m5hvWWiwVfiuEq8sqmvJKUvFI35cnRWZcXHwe7kiHQNmtO6JOhiVa53ZR0mxdWs1w9fcS4tLvlUlheKOBteurvV7y91O/qOte3lWdxcT+dUm3KT97OkPTW223paZsHaVcqKWoagk+xtUYPz3pY/Nic3Uo+t3YL5bS97Ofh4+03n5WnqGbWqV8Q6y9EO1dvR1ik0vZtrSOPKoaN06QlHpZ1zs36P+pgel+i0krnaCK6oWy/1p5/07Jf0r63nnmj/AKmAwT/sTDk5Vf8ARrZotMuKUmljqLdPqx9pKqOKLKVBru6O9FuWdkdS/wDuMv2IHr+Txz0V23sfqTf/AMxl+xA9jXE85yf3Je89Nn/Xq+Jt44H5+elKp/097TSTw1O3aa/+hTP0GwcC+lJTS6ddpnj5Vv8A6imdPAjeRr9S7Y9sb0R7ba3sTrlLWtIqKSeIXlrN4p3EOe6+x9al1Puyjuno72z0TbfZ6lq+jV96Pxa1CfCpQn1wmup/Y1xR+c+nV40K633+Dlwkv3m87C7a61sLtFT1nRK6Uvi16Em/VXFP5k0vsfNff28riReOqvlWcPnzhv028P0DRDeDUOjHb/Rtvdn4anpknTrQxG6tKjXrLaePiy7U+alya80ttznkyktWazqXo8eSt46qtb2/2S03a3SHZ3sXCtDMre4ivboy7V2rtXJ+5rmDafQtU2a1mppeqUXCpHjCcfiVY/Pg+tfauTOw8ZRgNt9ktL2s0iVhqFNqS9qjXhwqUZfOi/vXJnTg5E451PhU+o+mRyY6q+XJybzho+4yeTK7W7N6lstrFTTdShjnKlVivYrR+dH966jD55NcUWtbReNw8dlxWxWmtlRyzwyzb+jXbu92S1FUarnX0utLNa3Ty4P58O/tXX4mkuSXVy7wpZfUZvii8alnBmvgvFqS7I0XVLLWNNoahp9xC4tq0d6nUg+DX7vDqL45V6ONudQ2N1F4UrjTK0s3Ftn+/Dsl9j9zXTWgaxp+uaXR1HTbmFxbVlmM4/amuprrXUU+bBOOXt+B6hXlU/tkOo836ceirSOkjQ/adOy1y1i3Y36jlx6/Vzx8am31dXNd/o/BkOKxjJpraazuHfekXrqX5q7V6XrGye0lxoG0FlKyvaLxKD4xmnynB/Ki+pr7+B7l6DUnLbTabj/yfQ/1kz3bpr6KtC6S9n1aXq+C6nbpysL+EczoS7H86D64+7DWTyP0Rtl9Y2P6T9rtA1+0+DX1vY0OK4wqw9ZPFSD64vqfinhplhfkxlxTE+VVThzizxMeHTyPoArYXDnn0r451jROv+rV/wBumeKyXJI9r9K541jReX+9a/7dM8WeHjgXnD/bh4f1X/sypy4Ya5nq/oo1H/PbWef/AAfD/WHk9RZXxcHqvooJ/wA9ta4f8nw/1hLlftyj6dOs8MF6Vl16rpA1R8v9yKSXDtcl+88D0izvdT1Wy0zTqfrr28rwtqEMZ3pze6vvR7P6YlZ0+kC/SXGWn2sffOX8CPQp2Qlre3d3tdeUt6y0OHqrdyXCV1UjjP1YN+c4mqmT28O1pXDOXkTv7dG6rPTuiHoWhb2W4/5Ls429snw9fcS4JtfnTbk/M4vulUqXM69arKrWqSc6s5PLnNvMpPvbye4elrtgtQ2ls9k7Srm303Fa6UXwdea9lP6MHn6/ceHylnr4GOJj1Xqn5Q9Uz9WSKV8QpTe7zLW4qqFOUs8EssvJYfcbl0EbHLbHpHsLS5pqpYWf9cvVjhKEGt2D+lLCx2ZOjJborMuLj09y8Vh076O2yb2V6NbGF1T3NQ1D+u3eVxUppbsX9GO6vFM9IKdKO7HGCoUV7dUzL2WKnRSKgAItgAAABD4JvsAtNG/4Npef3smP/C0v/oL9pjR1jTLfvpp+8U8PVqv5tGC97kGF2afqL9vXH+dP/UxNwNM1CWaGtz+dVqpfoRiZqhkbbYrFnQX+jj9xWPmkt2lGPYkj6MJx4AAGQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAPIeQAAAAAAAAAAAAQSABAAkAAAAAAAAAAAQMgSCCfIAAAAAAAAAAAABAEgAB1GF2oXDT6nzL2C/SjKP7zNGI2uWNFnV/salOr+jOLf2GY8o28KOzz3dT1GHzvVVPfHH+yZ18jX9Kap7Q1afL1lrn9GbX+0bB1C3linhaaXwp1ofNr1Ptk3+8iv7OqWz+dTqR/Zf7ibLhdXkP9Ipe+K/gxfYjc2dTsq7r84tfwMJLsBckAyAEASa30lbT2ux2xOq7R3WJRs6DlCDf4yo+EIfWk0vM2M5c9NXa1XFxp2xNrUzCg4318ov5TeKUH/el+ibcNOu8Q5uXljFimXh13e3V/e1r69qeturmrKtXm/lVJPMn72IotqbfYXNPi1wwy+jtGnjskzM7TPDg2/k8zsPoh2ftejnoije6w1bVPUz1LU5y5we7vbr+jFKOO1M5/wCgLZaO1XSPaULikp2Vg1eXOVwlGDW5F9uZ44diZ6l6Ze1y0/ZOx2OtKuLjV5+tuknxjbUmm0/pT3V3pSK/kz13jHC49PxRjxTms5a251292s2x1XaS+TjWv7h1VB/k4YShD6sVFeRY6TT9bqtpT+dcQXj7SJqwjKbcVxk8viXWzNPf2hsuxVN5+CTZ3RXpjSvy5Orcy6t9Fif9e2iWfkW331TQenieeljW02vyOP8AyYG7+is277aLD5U7b/8AsPP+nlv+l3W1x/IP/wDTA5OP/wBiW3PbfBrDT89Y8FzRFNNtdpWoxzJFnMdlF4l0R6K8WtjtSbX/ACjL9iB7HyPIfRba/mfqa7NRl+xA9fPO8n92XvPTo/16oOCPSlljpz2m+lb/AOopne/kcL+mVot7o/TBd6nUT+C6zb0q9CeOGacY05x8VuxfhJG7gWiMiPqVZtieQSl7TwZDS7j1kHRk96cFwz1x/wAjGPi+0qW2/CvCrDg4vP8AEupnbz80jTd9iNqdf2O2ipa1oVx6qtD2alOXGnXh105rri/enxR250T7f6Nt9oKvtPl6m7pJRvLKcs1KE329sXjhLk/HKXDFCEKtGNWnxjNZRmNktoNV2V1631nRrqVrdUfOFSPXCa+VF9nmsNHHyeLGWNx5beHz5w36Z8P0DRHUaP0R9JGk7f6N623xa6nbpK8spSzKk3ykvnQfVLyeHwN5KW1ZrOpeopkrkr1VYPbHZnTNqNInp+pUcp8adSPCdKXVKL6n9j5M5c262a1bZLWXp+o03KEsuhcRjinXj2rsa649XesM6/MPtbs5pe0+jVdL1W3VWjPjGSeJ05LlKL6pLt8uRv4/InFP9K71H02vJruPLjve4BS7+RnukLZHVdjdb+B3i9baVW3a3SjiNZdndJda81wNdUsvuLqlovG4eLy4bYrdNoVXJtczZujvbnU9jdU9bTlK50+q07m1zhS/Oj2T+/GH2rVN7wRHxuD+8zfHF41Jiy3w36qS7M2Z1zTtodIo6ppdzGvb1lwa5p9cWuprrRk8HJXR3tjqOxmr/CrXNazqtK6tXLCqLtj2TXU/J93UWzGvabtHpFHVNKuVWt6q8JQkucZLqkutFJyME4rf09r6d6hXk018srjgUfgtt8M+GeopfCfV+q9durf3M53c88Z44Kw5HOswMjLJYHPHpYP/AHX0TH/Rq/7dM8WbeUe0+lgv919E/wC7V/2qZ4o+Ze8P9uHhvVf+xL7k2eueinBfzz1l4x/UIf6w8iUknxXA9e9FOcZbZ6zh/wDwEOH/AIhnmftSelx/s1ecemjvrpNuKdOMpzqWloowXFyeamEvM9+6PtKs+hzoHhK/hGNaxs5Xt/jnVuZrLjnre84wXgjWdodkP53eltSqXVDf07RNOttQrtr2ZVE5qjF/W9rwgzH+mPtQ1R03Yuznl1X8Ovkn8lNqlB+Mk5Y/MRXRb3OmkPR3p7MXyS581PUrrVdUu9UvajqXV3WlXrSzznJ5ePuXkUFU8ijCO5FRxwSwfMm0i0rHTGnnr/qncrtTXDq7+4639FXZJaHsE9cuaW7e63NV+K4xoLhSj5rMvrnLXRls9X2x6QdM2bin6m4qKVxJc4UI+1UfuWF3tH6A2lCja2tK3oU406VKChCEVhRilhJdxwc3L2isLf0nj95ySqoAFa9AAAAAABSu5blrVn82Df2FUs9Xk46ZcY5uDivF8P3gVdPjuWNCOOVOK+wpW2XqN3Ls3I/Y3+8uoLdgorklgtdP41Lup86u17kl+4MLtmk1252V3w/H3k4rzrbv7jdZNRi2+S4mk2i9ZZ6UuuvdUpvzk6j+4lVC/wAN3XIkLkCLZAAAAAAAEASAQBIAAAAAAAAAAAgkAAPIAAPIAAAAAAAAAAAAAAAgkAAAAAAAAAAAAAAAgkAAAAAAAAAAAAAAAAAAAAAAAAAAAHEgkAAAAAAAAAAAAAAAAAQSAAAAAAAAAAAAAAAQSAAAAAAAAAAAAAAACOIJx3gAAALLXqLuNEvaK5zoTS8d14L0iSTi0+KfUGJ8NYsKylqul3S5V6U4Z+lBTX7JtCNMtZeosLFy4Ozuo05Z6lGo6b+xm5LkZshjWlP2dXqxz+MoxkvJtP70Rq3Cz9b10qkJ+6Sb+zJNyt3ULWp1S36fvW8v2Sre0vXWlajz34OPvRhPSsiS3sKjrWVGrnjKCb9xceYZARx7Q/ECy13UrTRtGvNVvqipWtnQnXrS7IxTb+4/P3arVrvafaLUNor5NV7+vKtKLedxPhGHhGKS8jpf0vNrfgGztlsja1cV9Tl6+6SfFW9Nrg/pTwvCMjluMcJRzw7y14GLUdUvPercjdopCacWkS5brw5YbWeZVUerkZzo/wBl6u1222l7P00/V3NXNxJfIoR9qpL3LC72jsyW6Y2qMVZveIh0j6LOzn8gdH8tevYqndazL4Q3LhuW8U1TXcmt6f1zmTpk2yntr0k6nr0Judk5/BrFdlvDKi12bz3p/WOmfSf2qp7HdFb0bTHG3u9WS0+1jDh6uil+EkuzEPZ8ZI41pxUcJRwkcXGrNrTklb8u/RSMUK2d7lzMlshTc9dhN8oU6kv7uP3mLWeCXA2Xo9oKpqVzNr4lu175L+B3T4U+adVdDeivUgtX2goNpTlQoSS62k6if3r3mo+kLYV7XpX1CtWptQuqNGrSb5SioKDfvi0WPR9tLHY/bax1dzcbXf8AUXiXXRnhSf1XiX1T2z0jNmYa7sStds4Kd3pSdZSjxc6DX4RcOxJS+r3nDFvazxM+Jd+GkcjhTFfMOZopLr4MqQlhZ5FDeSSSeeCIUm3hZRb7UGp29f8ARp2rWn7TXWzlzUUbfUfwltnqrRXFfWiv7vedIp5XB5OELS5u7C8o31nUdO5t6sa1Ga+TOLyn7ztHYHaC32o2VsNbtuEbmknOCfxJrhOPk00UfqGHpt1w9X6Nyuqntz8M/nqR5H6UuwUtuOjC7+B0fW6vpeb2x3V7U3Fe3TX0o5WO1RPXd1d5DinzRxUtNbRMLvJSL11L8uaMVJRa5SWUu4r04qPDB6X6SGxMNielG8o21D1em6kne2WF7MVJ+3TXZuzzw6k4nmmXnh7z0OK8XpEw8tmpNLzWWX0K9jTq/Baj9ib9h9kv8zLzSl1YNQzjDTxh5TXb2m1aVcK9s1VePWR9moux9vmT248tNd4ZLZvXtY2a1y31nRbuVreUPizSzGS64yXyovrT+9JnYnQ50l6Zt9pD3VG01e3ivhdk5Zcerfg/lQfU+rk+Jxoqaec8UXeiahqOhavb6tpF1Us722lvUqsOrtTXyovk0+DOXkcWMsbjy6eHz7YLanw/QBcVw5EnnPQ10n6ftzpqtrn1dnrlCGbm1UuE1y9ZTzzg/euT6m/RUlgpbUmk6l6vFlrlr1VYrajQNN2j0etpeq28a9vVXhKD6pRfVJdTOVukPYrU9i9Y+DXSlXsqsn8Fu0sRqL5suyaXNdfNd3X7T7TG7SaHpu0Gk19M1S2jcW1ZYlF80+pp9TXNNcjdx+ROKf6cPP8AT6cmu48uMY4cU88Hy4H0mlw4m2dJexOpbF6mqdTfuNNrSfwW6xz/ADJdk/sfNdaWoyeFzLzHeL13DxmbDfDfpvD7jLDNm6P9s9R2O1n4Zat1rSq0rq1bwqq7V2TXU/JmpueFwKdSTcfjYyL44vGpQxZbYbxeku09l9f03aPRaGq6XcKtb1o5XbF9cZLqknwaMnntZyJ0W7cXuxOtb/4SvpVw0ru3jxfZ6yC+cvtXDsOsNHv7TVNOoX9jcQuLevBTp1IPKkn1lHyOPOK39Pben8+vJp38rxPjzJZCRJzQsnPPpX/8L6J/3av+3TPE5Ps957T6WUt3WNE4/wDw1f8AbpnisXxyX/C/ah4f1T/sy+KreOw9Y9E2T/ntrSTx/ufD/WHlcoKX8T1n0UaGNtdaf/YIf6wlzI1ik9N754iHQlZaZpH8o65WjToOdJVLy4fzKcXjL7Es+9nCm2u0tbaza3VNfud6Mr6u504yfGFJezTj5RS88nRHpdbWS0bYqjsza1d281qbjVSfGNtDDn+k92Pg5HJ9JvOW2jh4OPX6pW/quabapC8qceotquViOG23gqb5f7P6NebRa9p2hWC/rN9cwo0+vdzzk+6Ky/IsLTERtUUrNrRV0P6G2xqttL1Dba8otVr+XwSybXKhB+1JfSmv7iOisIx+zmkWehaFY6PYU1TtrOhCjSj3RWPeZEoMt+u0y9jxsUYscVAAQbwAMAQ8gYfaALPVVv06FL59eCfgnvP7i98yzuPb1K2hzUFOo/dhftMC7WcFppKzYxn11JSqe+Tf7ypqFR0bGtVXONNteOOB92kPVWtKl8yCj7kBba9Wdvo17WXOFCbXjuvBgrGio6ppFql+JU5vj82nu/7RldqpZ0v1C5161Kl5Oaz9iZaaTD1u0lSeOFC1+2c/4QJR4arf8mwrkgFyBFtARxAEkcScEYAAY4k47wI4kgAAAAAAAAAARxJAEcSQAAAAAAAAAAAAAACCQBBIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA1K/t3L+WbOOVLflOH14KS/vJmzafXV1Y29zHlVpxmvNZMPqcfU7QqfJXFt9tOX8J/YXOycv9xo0Hzt6k6L7lGTS+zBKfDVXtZd6r7FtGtnHqqkZ+SfH7Gy6xnifFzTVahOk+U4uL80U9PqOrZUpy+Nu4l4rg/tItilpXs0alHP4mtKGO7OV9jRe4LGzedTvd34mYZ+lu8fs3S+BAU61SFKlKpUkowim228JJdbKh456WO2f81+jWpp9tV3L/W5/AqWHiUabWas/KPs+MkSpWb2iIa82T26TZzd0r7VvbLpG1bW4yk7WUlRs88lbwyovu3uMvrGtR4vlzLSlKXYl1LwLqlLL4ov6RFaxV43Nab3m0q8IN4xlPPM6Z9EvZGNpo97tjd0/wANft29pvL4tCD9pr6U1/cRzvs1pdzr2uWGi2Cbub+vGhB4zu5fGT7orL8EdfdJ2r23Rr0P1YaVilVt7aFhpsev1jW7F97SzN/RZx8u8zqkfKw9MwxG81vEOa/SX2lhtd0jXVK3qb9jpObK2afCUk81ZLxkt3PZBHls6eOGUX8l7Cbbba5t8X3soyjnhg7MWOKViHFmzTkvNpWm5jHDzNr6PY/8IT7oQ+9muSi88EbRsTHc027ny36yXuj/AJmbObPP6GSvIxc8yxJJ5wzpH0etq6W0mxFTQ76cat3paVvUjPj6yg0/Vt9vBOL+j3nNN3JveRkuirayexnSBZ6pVnu2NWXwa9XU6Mmva+q8S8E+05c+Prr/APG303POLJqfEq3SZs9PZHba80VxatlL1tpJ/KoS+L7uMX9EwMZHR/pJ7Lfy3sjS2jsKcal3peaknFZ37eXx+PXjhLwT7Tm2jHrT4dvadXFy+5Tujz8HtZO3iVdcVyPZfRe2ohYarc7K3dTdpXubi0T5Kql7cV4xSf1X2njcefDgi4s7u4067oahZVHSurWpGrRmuqUXlPw/cS5GL3aTDVwuRODLFndBDfcYXYfX7fabZaw1q3aUbmkpSin8Sa4Sj5STXkZs85ManUveUvF6xaHj3pWbDS2v6Na15ZUXU1XRXK8td1e1OCX4Wn9aKyl2xicNUG54lnKfFeB+ok4qUWnxTRwL09bDPYbpL1Cwt6Thpt5J3lg0vZVKbe9BfQllY7N3tLHgZe/RKp9Sw6/XDz6MHjqXiX2kXTsruM3xpS9mrHtXb5FuknjHFH1hPgWmtKW0dUabo4LCcZKUWsprrXaUppcmWGzV56ym7Ko/ait6nn5vWvIykodwiXDaJrOn1pl/eaTqNvqWnXNS1u7ae/Rq03hxl+9dqfBrgdY9CvSnZba2a0+/9Xaa7QhmrRTxGvFflKeertjzXhhnJijl4aK+nXNzp2oUb6xrVbe4oTVSlWpy3ZQkutf+uPI5uRx4yx/bs4XqFuPbv4d9eAPLuhLpRttsbNaXqkqdvr1CGZwXswuYr8pBffHq8D1BPPWUl6TSdS9fhzVzVi1Vjr+kafrmlV9M1O2hcWteO7OEl9q7Guaa4pnKfSjsPqOxOr+rqesuNMryfwS7xz/0c+pTXua4rrS67Mfr+j6frulV9M1S1hc2teO7OE19q7Guaa4pm7j8icU/05OfwKcqv9uJHJPkm2fUMt+Zt3SdsFe7D6vuSlO40q4k3a3L5vr9XP8APX2pZXWlqGVweeBe4r1yV3DxefBbDfps+8LeTa4rij0noV6Q3stqkdJ1Oq3o11PjJ8rWo/l/QfX2c+3PmMqiwsLBSnUl8nrGXDXJXplLi574MkXq7up1IVKanCSlGSymnlNE72TwT0b+kKpVjDYzWq7dWnH/AHNrTfGcEsui31uK5dsfA96jiSPPZcc4rdMvc8Xk15FOqHOfpbtvWND4/wDw1f8Abpni8W+HDie4+lhST1XQ3/2euv71M8TnGOS74UbxQ8h6pOuRMS+qc+K6z2H0V2o7Yay/+wU/2zxrGHwM5s9tVU2U2f2nr21Rwvr+zp2Vq484ucnvTT6sRy134NnKpN8c1hDgZIxZovLGdPG08dseky/1OjUc7O3bs7Pri6VNv2l9KTlLwaNBk2uX2leU1uqKXBLCLeo8t8TXTHFaxEOjJlnJebSidSceOWmuR0T6GuyXwzU9Q21u6WadonZWLa51Gk6s14Jxj5yOdbehcXt1Ss7Sk6tzXqRpUaa5znJpRXm2foN0ZbL0NjthdK2dt91u0oJVZr8pVfGcvOTbOTmZOmvSsPTMHXk65+GzLkAMlW9IAAAAAAAAFpb/AITUbmpjhBRpr7396LqTSi2+SLXSk/gaqyWJVpOo/N5X2YDCNV9uFKh11a0Y+Se8/sTLwsp/hNWpx6qNJzfjJ4X2KReghg9o571/p9DsnUrv6sML7ZonZdb9fUrj51dU0+6EUvvci21OSqbRVZN+zb20I57HKTk/sjEvtk4OOg29SXxq+9Xf15OX7yXiGqO9mWQAItwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAEZAkEADDbTRUKmn3P8AZ3Hq5PsU4uP37pT2al6vUtRtn8p07iP1o7r+2H2l3tNSdXQrvcWZ04eth9KHtL7UYuxrRhr9lWi/YuaU6We3gpx+xSJR4aZ7XbP1FjSqxtXeRqfFpt1l9GXH78l6i0u7R17mjU30oRf4SOPjpPKXvRCJbU6XQnRtI+sX4Wo3UqfSfF+7l5F2EDLKJS3VxZw56S+1kdsOk66VvWU7DSM2Nrh8JST/AAs14y4Z7Io6q6ddrlsX0banq1Koo31SKtbJdbrVOEX9XjLwizhFW+MZblxzlvi32lhwcW5myl9Vz6iKQqU4YfMrJtcU8Hwk1zKlvSrXVxTtbWm6tzVlGnSpri5zk8RivFtFlM6hRRWbTp756HuzXwzVtR2wuqeaVmnZ2ba4OpJJ1JLwjux+tIxnpVbZR1rbals3a1N6z0ZfhcPhK4muP6McLxcj2expWPRJ0LJT3ZPS7JzqY4evuZcX+lUlhdzOMri4ub2+rXt5UdW5uKsq1ao3xlOTbk/e2cGCPcyzeVnyr+zgjFHy+5YwopvhwKbjw7D7fLgfDfYmyx2p47KM8rxybZspHGhKXLfrTln3L9xq8km+ztNx0Gm4aDaxx8aDl75MhZqzz2UrnnkxdzS328pPg1hoy9dbr8jGXElxwYiGqkzHh1F6NO1MNqejp6LqMo1r3Sf6nXjPj6yi0/VyfbmKcX3xZ4N0kbP1Nkdtb/Q8NW8JqpaN/KoSy4ce7jF98WW/QttdPYrpCs9QrVHDTbqXwS/zyVOT9mf1JYeezePfvSU2Qjq+y9LaS0pqV3pSbq7q4zt5fH/R4S8FLtOSlvYzf1K/vX8Xxdx5q5s4tc+BLi5JrLS7hGKjHnwXDIi0W8Q87PaXtfov7Uu01W72SuqmKV1m5s8vlUS/CQXisS8pHRKOF9G1G50nWLTVLKW7dWlWNak+rei+T7nxT7mztTZTWrTaHZ6y1iylmjdUlUis8Yvri+9PKfeij9Qw9F+qPl670ble5j6J8wyjPH/Ss2Le0/RvU1Ozoes1PQ27uior2p0sfhYLxit7HbBHr+T4rJTpuMopp8Gn2HDjv02iYW+akXpNZfmcqicI9eYrBMXl5SeTdOnfYyWw/SPf6ZQp7mn3D+F2D6lRm3mC+jLej4JdppsFwxjh2l/jv1xEvJ5qe3bpfdtOdCvCtSlu1IvMX3m52laF5awrwwt5e0vmvrRp6XBcDLaBdO3ufVVJJUqrw89Uupk57OPNXqhnXDCyU5rKwXU4tcH1lP1fEzEuNTtLi7sryjfWVzUtrmhNTpVacsSpyXJpnVnQl0nW+2NjHTdTnSt9eoQzUguEbiK/KQX3x6vA5W3eeCvp11Xsb6je2dxUtrihNVKdWnLEoSXJp/8ArPI0cjjxlj+1hwedbjX/AKd4oHnPQr0kUNtNKdpfOnR1u0gvhFKPBVY8lVgvmvrXU+HY36JvZ5FHek0nUvZ4c1ctYtVjNqtD07aPQ7jSdUoqrb144a5Si+qUX1ST4pnHXSBs9qWx201fRtQTml7dvXUcRr0s8JrsfU11Puwztd8Vx4mkdLmwtptvs1Oynu0b+g3VsbhrjTqdj/Nlya8+aR08TkzitqfCv9T4McinVEd4ci72VyI4vm3hH1dWlzY6hX0+8oSoXNvUlTq05c4TTw1/n4ERT8j0FJ6o3Dx969M6lVtKtW1uKdzbVJ0a9KcalKpF4lCSeVJd6Z190SbY0dsdlaV3Jwhf0PwN7SXyaiXxkvmyXFe7qZx9lLtRtXRjtlW2M2qt9RjKUrKolSvqS479LPxkuuUXxXmus5ObxvcpuPMLD0vm/h8up8S9G9LFpahobbX4iv8AtUzw2pUy8ZPZfSrvKN5LZ28tqsatvXta1SnUi8qUW6TTXkzxJvJngxrFEIeqatyJmFfeb7jEbUzcbShBPhOo5PyWP3mQlNrk8GB2mr5ubeln4lLPm2/4I6bz2cmGv6mOk5IpSeSomurJSmknx44WcLmzRPaNu2I3Ons3oh7KS1/pCr69d0t6x0NKdNtcJXEk1D9Fb0u57p2UuXA889H7Yz+ZXRnp+n16W5qF1/XL7PNVaiT3X9FbsfqnoRRcjJ13mXq+Fh9rHEPoEccEml1gAAAAAAMgWmqSfwV0ovEq0lTXdnm/dkuYxUYKKWElhFpNqtqsY59m3hvP6UuC+xP3lXUazt7OrVjxkliC7ZPgl72gwpacvWVrm5fy6m5H6MeH35LxlKzoq3taVFPO5FLPa+ti8rRt7WrcTaUaUJTl4JZB4hqWo1pTttWuYfHq1Z06b70lSj9uTbbSlC3tadCC9mnBRXglg1KxpSlQ0e0n8erWhUqeSdWX2pG5IlZCkfKeoDiCLYAAAAAAAAAjJIAEZJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAABxAAAAABxAAAAAAAAAADiABA4kgAAAAAAAAAAAAAADiABBIAAAAAAAAAAAAOsAjiA4gkAB1AAfNSMZxcZJNNYaNFjJ22mW1Vt72n14qT7FTnuSf6OTezVbu0hO+1WwlwjWaqR+jUhuv+9FslVryR8topttH1gstAru50a0ryeZzox3/AKWMP7cl8R0nHgDYNe6RNp7TY7YzVNorzEoWVBzjDOHUnyhD60ml5mYjc6hi9orG5c2+lttc9Z2xpbLWtRO00eG9WafCVzNZ/uwaXjKR4ospJdfIXWpXOo6hd399VdW6uridevN/KqTeZfayYLLPQ8fHGOkQ8fyss5cs2lVjTyj1j0WtjXrnSFPXrmlvWWiQVSOVwlcTyoL6q3pdz3TyqHsPe6kstHaPRPoNp0c9E9OpqTVCsqMtQ1Kb5xm470l9WKUfqnPzcnTXUeZdHpuHrvNreIeVel1tROre6bsda1PYp/1y8w+cuKpwf96X6J4HuvuybLtjqdfaPaW9128TjWu6squ7n4ifCMfqxSXkYSdNp8s5JYKdFIhx8rP7uWZWrWOGT53eOCvKOFwKbWGzdEtEPhpKLbXUbvaQVLTLWnn4tCP3GmUbeveVo2tpRnXrz4RpwWX49y7zdKznRSo1YOEoRUceCwRtZozsdecW8ZMXXiZevut8DG3DSTxgzDXWVhcpOlKm4pprDTOuvR32qpbY9GsdL1Gcbi90yHwG7jPj6ynu4pzfapQ4PtcZHIlzLHI3DoK2yexnSRY3VxW9Xpl/iyvsvCUZP2Kj+jLHHscjRysfXTceYXHpmf2smp8SuekbZ2rsjtZfaJNSdGlLetZP5dGXGD8uMX3pmvp54tvB0p6Tmyv8p7MUNpbSkpXOlv8ADtLjK3k/a/ReJdy3jmtprrR08PN7uOPuHP6lxvYyz9SqRlKLPbvRf2rdG/vNkruotytm5s8v5X5SC8ViXlI8M3nnw7CtpOp3mjarbarp8/V3VrVjWpPqcl1Puayn3MnycUZccw1cLPODLFod35IS7TFbG63a7R7NWGt2cs0byjGoln4r64vvTyvIy+MHmJrNZmHu6Wi9YtDxT0tNi1tJ0fS1q0ob+o6E3cx3V7U6DX4aHuSl4w7zjqDWV1p8V4H6WXNGnWoyp1IRlCSalFrKafUzgnpb2JqbE9I1/okINWUpfCLCT5O3m3ur6rzD6veWfBy/+JUnqmDX64alCOX1lanT4cE2VYU0uGCrGPAsplR9UM9pFf4Vbbk3mrTWH3rqZdSju95gLKu7WvGrHq4Ndq6zYXONSKnB5i1lPtMbcmSupUH7ilPPF5K8s8eGT4lF9ROJa4To+ranoWr22raVcyt7u2nv0pYyuxxkuuLXBrrR2L0XbaWG22zNLU7Xdp3EX6u7t85lRqpcV4Pmn1p+JxrUpt5RsnRdtXebEbUUtUo79SyqYp31CP5Wl2r86PNea6zj5fHjJXceVt6bzpwXis+HakePWHFPqLfSr211LTbe/sq0K1tcU41KVSDypRaymXRSdOpexiYtG4eHekjsHC4tpbZ6ZQzXt4KOoQgvxlJcqnjHr/N+ic+yaWW+eTvC4pU61GdKrCM6c4uMoyWVJPmmjj3ph2OqbE7V1bSnGX8l3O9WsZv5meNPL64N48HFlx6fyN/47PMes8HU+7T/APWnvkfE5NPg8M+HUTWeKyim55feW7z0VllNR1q8vtB03SLhqdHTZVfg8m/ajCo4tw8E48OzODFzeerJ858gvcYrSI7Q2WmZncplybNX1yTnq1bL+Luw9yX7zbqMMyjnk3xNMvJutcVKz+XUlL3sjlrpvwSiL4c+R6J6OeyL2x6VrCnXpes0/S8X922uEt1/g4Pxnjh2RZ5rN7qTzwzxfcdpeihsb/Nro3p6tdUtzUNcau6m8uMaOPwUf0XveM2V/Ly9GPS29Pwe5l39PYEiUTzBSPUAAAAAAAAB81JRhBzk0lFZb7EfRZanmoqdnHnXliXdBcZfw8wJ0qMnbu4msTuJOq89SfxV5JIi7/DX1vb4zGGa014cI/a8+Rd8EuxFppidX114/wAtL2PoLhH38X5hheIxO1s2tDrUo/GuHGgvryUX9jZlzA7TyU73TbZf2s67XdCLS+2UTMeWLzqFPS4qttBHHxba3lLznJJfZB+82IwWy0d641G5+dXVKL7oRX+05GdFvLFPAADCYAAAAAAACCfMYAAAAAAAAAAAAAOIAAABxAAAAAAAAAAAAAOI4gYAAYAAAAAAAAAAAAAQSAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAMDAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAwAAAAAAAAAAAAAAAAAAGAAAAAAAADC6zZ309ToXNhClJypSpVZVJYjDinGWOcse1w7+ozQeAxMbWmlWcbGxhaxqSqbrk3KXNuTbfDq4tl2AGUM5d9NHbDfvNL2Gs6vCLV9f7r8qUH/el5ROmNWv7XS9LutSvaqo2trRlWqzfKMIptv3I/PnbTWK+1O1uobR3ikq1/cSrKLefVw5Qh9WKivI7OFi677+lZ6nn9vH0/bD0opJpLDLmMkuMuCXaU9zEs4ZFaW5SlJrgll+BczOoeb8vS/R92Y/nZ0l2FtWp+ssdPxfXb5xcYNerg/pTxw7FI9t9KjaR22z9psraVcVtQmqt1h8VQg+X1pY8osqei7srHZno0jrl7FUr3WsXlWU+G5QSfqovu3fa8Zs8J6Q9pqu1m22pa25N29Wp6u1Wfi0Y8IeGeMvGTKzfvZtz4hZZb/huL0R5lhKkW2+ZQlBvPDLK+83wf2FahRdSWFDPa88F3s79w873YuvDg8dXWXelbP3F8qdzdzdlZS4qbjmpVX+jj1r854XjyMzZwsLRqpKELquuKc4/g4Pui/jPvfDuK9a8lXqOdSbnOXNyeWyNu/hOMuoX2nRtLCg7bTaPwejL473t6pV+nPr8FhdxWrqjcUtytHeXJdsfBmIVdLkfULldvma+lotM2na01WxrWyc6adWl85LjFd6/eYOvJyXLgbhSuOOFLiWWo6Vb3Oalu40ar6vkS/gzZWUqXiO0tQuIosbiMJYUoprDTT68mQ1OEraq6VaEoTXVL712mJuJtvK5GztMO3Fve4dn+jttfQ286LYWWpyjcX2nxen6hCpxdSO7iM2utShjPepHOXSFoVbZPbHUNAqbzhb1M285fLoy4wl7uD70zHejptxLYjpRoSu63q9J1VqyvW37MMv8HUf0ZPDfZKR7/6U+yT1DZ6htbZUd650z2LrdXGdvJ8X9SXHwcjhwz7GfXxK55MfiuPE/MOdotPOVnPYRUbcWl2cD4g3lprGOsqKOctrkXHl53xL3X0Utr3Cpd7HXk8KWbqxcn1/lYL7JLxkdEReeJwns/qNzoWt2Ws2Txc2VaNaC+djnF9zWV5nbuzuqWut6JZ6rZT37e7oxq033NZw+9cih5+Hov1R4l6v0fle7j6JnvDIHjvpT7Ffzh2Jjr1lS3tQ0XerPdXtTt3j1sfJJTX0X2nsR81YQqUpU6kVKMk04tZTRxY7zS0TC1z4oy0msvzsilx3WmmxhZNp6YNl57F7f3+iQg42UpfCLF9tCed1d+68x+r3mq0/aWUX9ckWrEvGZcM47zEpaT5GT0Wvztpvnlwz9xYRj2FanFqalHKaeUw1XjcM5KD4cj5lDjhFS1k7iiprn8pdjKvq1j94izkntOltGnw5n04p4eM4eV3FVx7j5cXHHBI2RKO3tPo2bbu1u3sbqVbFCs5VNOlJ/FnxcqXnxkvrLsOg8nCELiva14V7erOlWpzVSnUg8OE08pp9qZ1x0O7a0dttlKd5KUIahbNUb6kn8Wol8ZL5slxXmuoqebg6Z648PWej83rp7VvMN2bSNI6ZNi6W2+x1fTouFK/o/h7GtL5FVLgn+bJey+59xu+EQ4JnFS80tFoXWXHGSk1lwDc0riyu6tne0ZULmjN06tOfOE08Si/BnzltZx5nuvpQ7BeouYba6ZQ/BzcaWpRiviy5Qq/dF/V7zw1U2vjeaPTcTLGanU8RzcE8fJ0y+Es4x1lWKwlwIwuSIc8ZOzWnFM7fdSapUKtV/Ipyl7kzSpvEVl9RsmvXMaWkXDXOUVBebX+ZqE6ykpJ8mu05s1u+nXgpOtto6MNlqm2/SDo+zcIydC4resu5L5FvD2qj7sr2V3yR+htrSp0LeFGjCMKdOKjCMeCSXBJHOvoSbIxt9B1PbW6pL12oVPglm3zVCm/aa+lPK+ojpBJJcEee5mTrvr6et9Pw+3j39i4kgHKsAAAAAAwAAIfIsrD8PXq3rXsy9il9Bdfm/swfWpSlJQtKbxUrvdyucY/Kfu4eLRc04Rp04wgkoxWEl1IMLbU5N0o21NtTrvcTXVH5T937i6pxjCEYRWIxWEl1Is7P+s3dW7fxI5pUvBP2n5vh5F8CBmtalNVNoq0n8S2t4wz2OTcpfZGJsj5GlXdSVez1KvF+3d3E6VNp9slRj92TNUMk9tNh2TpuGg21SSanXTryz2zbl+8yp8UacaVKFOCxGEVFLuR9mE4jUAADIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAwABBOEAAAAAAAAAAAAAAAAAAwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAB8t7qbYEvgant10i7HbFU09odct7WtJZhbRbnWmu6Ecyx34weS+kb05VNnLutshshWg9XUcXt7hSjZpr4kVydVrjx4R73wXLFS6uru6rXd5XrXNzWk5Va9abnUnJ83KT4tnbg4c5O9vCs5XqEY+1fLoPps6dtK2u2Ur7M7OWl/Qp3dSKuLi5jGG/STy4xim37TSXHHDPaeFc3heJaQbzxRcQlx5lrhxVxRqqg5Oa+e27LhReUsGzdGOyE9studN0PdzQrVPWXbj8m3hxnx6s8I+Mka3ScW1xSOovRF2UjYbNXm1lzT/DalN0bVtcqFN4bX0p58oxNfLyxTH2T4OCcuWN+IZ/0hdcWz3R9LSbFxo3GpL4JRhDhuUkvwjXYt32frI5ZpQ3I7uMLqXceldOu1Ntre3F5XqV4xsdPzaUJN+y8P22u1uWV4JHkOqa369uFpF06fXNr2n/AANPGxTFNoc/L7uaax4hkri7oUJJSk3L5kefn2Hwr6U4qOd2PzVyNdjV3ef3lSNylyydXS5Jx6hsUbjHyitC649prsbp4wVYXhnUNU42wRuEypCvH5xgYXi58VgqQuk0sNjUITjZ+N0kuZMbt54ZZg43LbxxPuFfxIzCPtspeRoXlB0rmCnHq7Y96fUavrGj3NopVqTdegvlJe1FfnL95mY3KXMqfC+SWcvrQiNJ0vajQ7tKVNp4aksHa3o5bXUekPoojY6pKNzf2MHp2ownx9bHdxGb+nDGe9SOUdS0SjfJ1aDjRqvqS9iXiurxRtno67Tz6Puk+3hqE5UNN1VRs73feIxbf4KrnlhSeG+yUjRycU2p1R5hc8DkV6umflT212audlNrb/QbhTcbap+BnL8pRfGEvdwfemYuKS4HSfpRbKK+2fobU2tPNfTvwdy4rjK3k+f1ZcfCUjmtNcOPHPM7eDljLj38uD1LjTgzTHxKtBrKylweUe/+i9tapRu9kLurxhm5ssv5Lf4SC8G1LH5z7Dnt1McfcXWz2t32g7R2OtWLarWlVVIrPCePjRfc45XmS5mGMuOYR9PzzgzRZ3i3hEPijG7L6vaa9oNlq9jUVS3u6MatN9zXJ965PvMpjqPLWrMTp7qlovWJh4r6Vmx71vYmOv2lDevtFk6st1cZ27x6xeWFL6r7TlWiuCzx7H3H6HXNvSuKE6NWEZ06kXGcZLKaaw00cRdJWyUtjttr/RN2SoQn62zk/lUJtuPjjjF98Sw4eTcdMqD1bBNf8kNZp0m+0uKdPHBrzPuEOOF1FWMPE73npuq6fP1M91v2J8JfxMnOHDD5GJwsYX3GRtK2/S3JNuUFjj1oNeSN9ySa45KMslWfFcWylKPWSiWuIUqiysGydE21lfYna2nqK3pafXSo31KKy5U88JJdcot5XdldZr7XWbP0T6NHXOkDRrCrCMqPwn11VPjmFNOeH4tJeZHNqaTt28Obe9Xpdg28lUoxqRk2pJNdXBlQRWI4JKDUbe8rvXdbapZ2moadcWN7RhXtrinKnVpzWVOLWGn5HD209np+nbRajY6Veu+sbe4nTt6/z4J8OPXjllcHjPWdO+kdtTV2c2Bna2dV073VanwSnKLxKEGm6kl9VYz1OSOTZNeyksKPBdRd+k4rd7/DzPruas2ikR3RUfjgoyzIrN97PiUM9fEuZh56ssDtTUSsKcM8J1s+5P8Aia6k3HnhmY2vf4e3pJ8oSm/N4/cYaPIr8v8AyWuHtSHpHQd0t6x0ba1TpzqVrvZyvUze2OXLczwdWlnlNc2uUvHDO8NH1Kx1jSrbVNNuadzZ3VKNWhWpvMZwkspryPzKST5rOGdUehNtrOvbansJeV3P4LH4dp6k+VOUsVILuUnGS+myt5mDt1wufT+VPV7cumUAgVq6AAAAAAicowi5SaSSy2ySxvH8Jrqyj8RJSrtfN6o+f3ZBKdPi61Sd9UTzVWKafyaa5e/n7uw+tSqT3I21F4rV3uxa+SvlS8l9uC5bjCDk2oxSy31JFnp6depK+mmvWLdpJ/Jp9Xv5+7sDC7oUoUaMKVOO7CCUYrsSPsAMrXVLlWmnXF0+Ko0pT9ybNW063araPYS4yjJVKn1Itt/pOJm9qpZ0+naf9Jrwpv6Ke9L+7FlposPXbRV6qXC3tlFfSqSy/shH3ko7Q1W72bEuSJIJItoAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAIAkEEgAAAAAAAAAAAAAAEEgAQSAAAAAAAAAAAAEACQAABHEkAef8AT9tu9g+ji91a3lH+UazVrYRaz+GnnEsdaik5fV7zfzk/05dYnW2j2f0BS/BW9rO7lFdcpzUE/JQl72buPj68kRLm5eT28Uy8DulWrXNS5uK061erNzrVJvMqk28uTfW2z5UeKTfAqZy+PI+kuPFI9D0xHaHkpvMzuXzDnxKsc43sEbr3uCPvepUqbqVpKEF1saYZLZPSLraLaXTtBtN5Vr+4jQjJL4ifxpeEY5fkd722i07LZOOgaPXlp1OjZ/BbWrCKk6OIbsZJdbXPxPzw0zbDVtB1SOo7O3lTTrqnGUIXEYQlUxJYfxk0vvMy+mfpWfFbdaqvqUf8Bwcrj3yT2W3Cy0xVnqh0Lfei/QvKqqVtvdQlhYinZQwv7xQj6Ktkv+fN8/8A8CH+I8Bj0ydKr/5+at7qX+Ar0+mXpSTw9utWflS/wEIxcjxEtvXxPmr3j/2V7LGP58Xv6hD/ABEf+ytY4/48X36hD/EeHLpl6T+H/vxqvupf4Cf6YulB/wDPnVV9Wl/gM+1yf5Me7xP4vb//AGWbKPLbi+/Uaf8AiPpei3Zri9t75/8A4MP4nhsumDpRa4bdat+jS/wHx/TB0p4x/PrVf0aX+Ax7XI+0JycX+L3f/wBl+zX/AD3vf1GH+I+16MNosY21vf1GH+I8G/pf6UevbrVf0aX+A+l0v9KD5bdap+jS/wABj2+R9sTfifxe9w9GW0S/46Xv6lD+JUXoz2a/55Xv6nD+J4HHpf6UP+vOqv6tL/ASul/pP4P+fGq+6l/gHt8j7Ovifxe+f+zRZ/8AXG9/U4fxEfRptc/8cbz9Th/E8Gj0vdJzSztxqv6NL/Afa6X+kvC/999U91L/AAD2+R9sdXD/AIvfqXo42tNJLa+7/U4fxPnUPRtsLy0lb3G1d1KL5P4HDK8OJ4HLpe6TMcNt9V91L/AU5dLfSdJ8NudW/wD1/wCAe3yP5JRfhxO4q7e0jR3bbGW2zmrXktY9XZq1r3FWnuyuI7u63JJvi1zOMNt9DudltrdQ0G5cn8Fq4pTf5Si+MJeaaz3pmOXSz0mcM7cat5+r/wABidY2u2g1y+hd6/qVbVKtOHq1OsoqajnOE4pdbfPPNm3h0vht3+Wnn5acisa8wySllLLEllpotrC5o3X4qWWucXwki9hHjyLiNWhRW3SXv/olbVOdrfbIXdTPqm7uycuuEn+EgvCTT+s+w6BOGtkdaudmtotP1yzTdSzq+s3E8esg+E4Pxi2vcdtaNqFrq2lWupWVVVLa6pRq0prrjJZR571HB7WTfxL1fo/K93F0z5hdnj/pO7IrVtlae0dpSzeaRmVXdXGdvL46+rwl3JS7T2ApXVGlcW9ShWhGpSqRcZwksqSaw012HDS/RaJWXIxRlxzWXBs1iWP/AEwnhZwZnpG2fuNlNt9R0Sal6mlPftZP5VCXGD78LMX3xZhYqT480XdbdUbh4fLhnHaay+k88CpRk4VFJLl3cxGPiipCm85WSTTOl0/ainFcyPVvqyfdvFfEk3z4FZxwuKaMw1yt1T5YwepejRbQq7e3FeSy6FhUcezMqkF92TzJrjwPUfRlqqG297Tb41LCTXlUh/E1cmf8cuz0yP8AZrt0ciOPaSiG8FHL3jmL0vtRqy2x0PTsv1dCynWSzw3pz3fugeKqo8cT2H0vrecdu9GvMP1dbTpU08dcKjb/AG0eNxi+w9T6dH+CNPFepd+Rbaop+J9xe82uXDmUnw5tH1Skt7D7cHfrsr9fTWtpmpavNf2cIw+zP7zFvGeXEutXret1W7qZwnVa8lw/cWby+RwX1MrGkarCJS6kb96NWp1dN6edmZxk1G4q1LWp2OM6U1j3pPyNAkuHBm9ejpZVNQ6dtk6VLP4O7lXl3Rp0pyf3HPyP25dfFj/LGn6Ex5IDHAHn3qYSAAC7AD4q1IUqcqlSSjCKy2+SQFK+rqhRzGO/Uk92nD50nyX/AK6hY2/weliUt6pN71SfzpPn5FGyhOvWd9Wi1lYowa+LHtfe/wCBVv6/weklCO/Vm92lD50v4db7gwoXjd5dKxi36uKUq7XZ1Q8+vu8S/SS5FvY23waiouW/Uk3KpPrlJ82XAEggxeuam7SEbe3Sq3tbPqqb5LtnLsiv8ubHkmdLHVa3wrXFTi807Om1L/6k8faor+8XOyVNuyr3zWPhlZ1I/QWIw96in5mHjbSrTho9vVnOrWzUu63XGDftSfZKXFL/ACNvpU4U6caVOKjCCUYxXJJdRKe3ZCnedvokAi2AAAAAAAAAIAEgAAQSAAAAAAAAAABAEgAAAAAAAAAAAAABHmBIAAAAAAAAAAAAAAAAAAAgkAAAAAAAAAAAAAAAAAAAAAAAAAAABGCQAAAAAAAAAAAADgABBIAgkAAAAAAAAAAAABHmSAI8yeAwAAAAAAAAAAAAAACASAIJAAAAA1wOLPTZjUpdK9rUmnuVNJouD8KtRM7TOZfTr2Zr3GhaJtZbUnJWVaVndSS+LTqtOEn3KccfXOniW6ckOLnU6sUuaqcuCLiliT4lnCShDem0kuvqLSvqNT4lvmEfn9b8Owvol5mce57Mvc3lK2juv26nVFPl4swl7Xq3VTeqyz2Lqj4FL1j6xnPMlMladKlKOHxfvCb5lSUU1ho+ZQeO5EdNsWFh8T6Unnn1lOTUHhyivMmM45xvwx4kJmEtTK4jLhjJUUslopRz+Mh7z7jUXz4e8zFoQmkrpceHYfXey3VRcvWQ959qafKcP0huEemVZLiEmfCnH+0h+kTlccVIY+kY7fZ0y+4dzPpYxwwU9+POM4LHeFOOfxkPHeMdmNT9KylhdhG8U1KL5zh+kSpQf5SGfpA7/SrnPWiU8NFPejnG/D9I+lKOONSHvMbNS+03nkfUVwwfCnFfLg/M+t6Ofj08/SMf/rHdcUJOlJTg2n1NdRm7DVU8Quv/ADEvvX8DXfWLP42HvJVZL8rBPxNlMnS05MU3+G7b8PVqpCUZRa4STymdBeivtjG70272Qu6v4azzcWeXxlRk/aivoyfukuw5Jp39ahLNGvFJvjFv2ZeRtnR5tjLZ/azTtbt5qFe0qqU6e9wq03wnDPfFvzwQ5U1z49fLfwYvxssWjw/QBM+ZcVjJb6VfW2pabb39pUjUt7ilGrSmuUoyWUy6SPN6mJexieqNvF/Sc2S/lLZyltNa0t660vMa+FxlbyftfovEu5bxzpQpyaW9z8Tu68t6N1a1ba4pxq0qsHCpCSypRaw0zjvbbZqtsrtTe6LUT9VRqZt5vnOjLjB9/D2X3xZZcTJuOmXmvWeNNP8AJVr1OkusrQgk+HPxKiWXwR9Ri0+w7nnt7fMVx5lfGVldfM+YwXZxRWUUk+JjekbaUsY8DbOhjU4aV0k6VOpPdp3E520n9OL3f7yianJPj1PrKKlWo1YXFCe5WpTjOnJdUk00/eheOqsw28bJOPJFvp3BF70chrJhtiNao7Q7L2Gr0WsXFFSml8ma4Sj5STXkZoorV1OpfQcd4vWLQ8k9J3ZOeu7Dw1W1pOd1o1R3GIrjKi1iol4LEvqnK0nFY3eK+8/QGtCNSnKE4qUZLDTXBo5D6dejW82N1arq2l0KlXZ64m5RlFN/A5N/i5dkMv2ZeT6s3XpfLin+Oyg9Y4U2n3aPNZPLxyPic1BOb5RTk/BH3SjJpb8cPuPnV1GnpF1UfNUml58P3l9eN12oKzqYhoMarnUlJvjKTk/NlbeWeZQnFRfDkfEptcU+PcVszqVlrfhdycXwfWdIehHsTUra5qm3d1Scbe3pOwsW18abadWS8Eoxz3yXUeK9EWwOu9I21FPSdKpSp2tOSd9fOGadrT7c9c2vix6/DLP0H2S0DTdl9nLHQdIoepsrKkqVKOct9sm+uTeW31tsr+byI6eiFn6fxZ6uuzKJd5IBUrwAIbAlvBj0/wCUK2cf1SnL/wA2S/2V9rIrznf1JW1JtW8XitUT+N+Yv3vyL6EYUqahBKMYrCS4JIMeSrUhSpSqVJKMIrLb6i1s6U6tV3txFxnJYpQf5OH8X1+SPil/X66rP/etOWaa/tZL5Xgur39hfoHkA4GI1fVJU6srKwUal1hOcpLMKCfXLtfZHm+5cQTOn1rOqq0at7eCrXc1mNPPCC+dN9UftfUYCCr/AAuVC2fwvU7hKVWpPlCPVKWPiwXHEVz97PunSrTuJ2Ont1rybU7i4q8VTz8qfa8fFgvsRsuk6bb6bbeqo70pSe9Vqz4zqS65Sf8A6S5LgT3FWrU3lGjabS063cIylVrVJb9atL41SXa+zsS6lwL4Ag2xGgABkAAAAACCQBAJAAAAAAAAAAAAAAAIJAEEgAAAAAAAAAAOAAEEgCCQAAAAAAAAAAAAAAAMAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAOscQAAAAAAAAAAAADK7QBi9qtC0/aXZ6/0HVaPrrK+oyo1Y9eH1rsaeGn1NIymSMiJ13hiYiY1L85OlXY/WNhNsLjQdZhOSjmVpcbuIXNHqnHv6mup8DUZOL5JNH6PdJuwezvSBoEtH2hs/WwTcqFem92tbz+dCXU+7k+tM486R/R2292UuKlxo9tLaXSk24VbOP9Yiurfpc2++G95Fxx+bW0asouRwZpbdPDyNY6iU+sr3lnd2NV0r+wu7SquEoXFGVOS8pI+KNOpcVFC2o1K03wUKcXNt+COzrr524ei+9aI4ZlNmtC1PaXWLXRdFtpXmoXc1ChSh29cpPqilxcnwSNw2A6DukXa64pyjotXRbBtOV5qkHSSXbGn8eXuS7zr/oa6KdnOjXTJ09OhK71OvFK61GtFesqfmxXKEM/JXnl8Tkz82tI1Hl1cfgWvbdvC06P+hTYnZ/ZGx0rVdA0nWb+nDeury5tI1JVasuMsOSyop8EupJGdXRb0c/9SNn/ANQp/wADcFgkqJyWnvtfRipEa00/+i7o65fzI2f/AFCn/Af0YdHf/UnZ/wDUKf8AA3ADrt9ntU+mn/0Y9Hif/ErZ/wDUKf8AAf0Y9HuP+JOgfqFP+BuAHXf7Pap9NP8A6Mej3H/ErQP1Cn/Af0ZdH3/UrQP1Cn/A2/rDHXb7Pap9NR/oy6PurYvQP1Cn/Af0Z9H3/UvQP1Cn/A28Drt9ntU+mox6NOj5ctjNB/UKf8D6XRr0frlsboP6hT/gbYB7lvs9qn01T+jbYD/qboP6hT/gP6NtgH/zN0H9Qp/wNrA67fZ7VPpqn9G2wP8A1N0H9Rp/wH9G+wP/AFN0H9Rp/wADawOu32e1T6ap/RxsF1bG6D+oU/4B9HGwWOOx2hfqNP8AgbWlgDrt9ns0+mo/0bbA5/4naF+o0/4Ero12AX/MzQf1Cn/A20gddvs9mn0t9NsbPTbGjY2FtStbWjFQpUaUVGEI9iS5IuQfLb7CLZCcnkfpIbMfyjoFLaO1p5utMyq+Fxnbyftfov2vDePW8Mp3dvRubapb16calKpFwnCSypJrDTJY7TS23PycMZsU0lxSo57CoormZbb7Q6uy21t3o0lJUact+2k+O/Rllwflxi++LMPTk2XFbbjbwObHOO81l94a5I+ort4s+oLLK0IJcWsmGnahOllb3vPn1aw00X8ILDZbVac4yajFvHLgTrJ3nw9N6A9sKekatLZy+qqNnfT3raTfCnWfOPhL713nQSeePUcW2tneX9xG3sbavc120407eDnNPt4cvE6n6La+1FTZelDayz+D3tP2YzdRSnVhjhKaXxZdvHvK7lUrE9UPXejcjJant3jw2spXVtQurepb3NGnWo1IuM6c4qUZJ8GmnzRVTTByRPdezETGpeIbd+j/AKVfzqXeyl7/ACRWk23a1Yudu3+bj2oeWV2I8c206FOk2hYVLW00NajmcVvWt3TlFpPLeJuL+w7SwMceR24+fmpHTtX5PTMF7dWtOCbD0fule+movZqVsn8u4vaEEvHEm/sPSNgPRNuHcQututoYypJ5dlpjlmXdKrJL3RivE6vBC/MyXbMfBx0YbZHZjQ9lNFo6Ps/ptvp9jR+LSpRxl9cpPnKT628tmYBPmcs953LsiIiNQLkB1Hy5RSy3jAZfXAsK1Spe1ZW9vJxpReKtaL/ux7+19XiRKpUv26dvJwtk8Tqx4OfdHu/O93aXtGnCjTjTpwUYRWEkuCDHl80adOjSjSpRUIxWIpdRZTb1Co6UMq0i8VJJ/jX81d3a/LtPupKV/N0qUnG2TxUqLg5/mxfZ2vyXde04QpwjCEYxjFYSiuCQExioxUYrCSwkuol8CJNRWXwSNY1PUq2puVCyqSpWSyqlxF4lV7VTfVHtn7u0zEbYtaKwudU1apWqzstNnhxe7WuUk1TfXGPbP7F19hj7KjVupSsdKl6ulCT+EXb9rdl1pN/HqdrfBdeXwI0uynqkI0rXNtpcPZdWHsuqvm0+yPbPr6u02u1t6FrbwoW9KNKlTW7CEVhJGZ1HaGusTbvKnptjb6faxt7aG7BPLbeZSb5yk+bb7S5AIt3gAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAcQAZHLqAEkPxA4gT3Dv6wAIa7eIwnzSJAFC4tbW5W7cW9Kql1TgpfefNtY2Vrn4LZ29HPP1dNRz7i5BncsdMPlwjnOOPaN1Y4I+gRZfDXiSmlzJ48zWto619SvVC4r1LfT5pKnUoPd9rsnLmuPLHB8iURtG1tNlTT5DKzjPE1uw1WvaVo22pz3oye7RusYjJ/Nn82Xfyfc+BsMGmuvImNFbRL7WesjqKN3bRuaEqU51IJ8VKEnFos6VtbUpxo3UGpSeI1FOSjP7eD7vcYSZNPuHAtXp9p1U5p9qqS/iQ517P4+9XofPSzOHiutd64hheA+KVSnUpxqU5xnGXFOLymfYZAAAYDHUABGO1DmwDHEYJA+JyjCLlLkuPIwmobYbNafUcL7WbS1kuqvPc+8zryfFSjSqx3alOM12SWUI18o23rs1j+kbYVLjtZo363D+JTqdJmwUOe1Omy+jV3vuMnf7JbMX7fw3Z/SrhvrqWkJP3tGFuOirYCvJyezFlTb/st6n+y0Tjo+Wi3vx40t7rpf2AoZxrqqvspW9WX3RMLedO2xtJ4oU9Vuf/AKdrhf3mjKVehnYGpnGl16efm3lX/EUP6ENgc5en3b//ADav8TZHtOW8cy3jTyjpU230TbapZ1bTSbq2ubZyiq9aUVvU3zi0s9aTXHhx7TTadKLWUzpK26G9gKLT/kSVRr+0uqr/ANoy9h0d7F2TTobOafldc6Sn+1k3xyaVjUKvL6RnzX6ry5dp0HOSjTlGcn8mPF/YZrS9kdq79x+B6HfTjLlOVF04++eEdT2emWFlHdtLK3t49lKmor7C6UUuojPLn4hsx+gVj/lZz5o3RBtTdbstQubKwh1pydWa8lhfab5oPQ/s5ZqNTU6lzqdRc1UluU/0Y/vbPSMYWB1mm2e8rHB6Xgxd9bWWlaTpmlUFQ06wtrSn82jTUU/cXsUkPcSapnflYVrFe0QAAwkAAAAOIAjrBb3V3CjJU4RlVry+LShzfe+xd7Aq3FanQpOrVmoQjzbLH1da/e9XjKla9VJ8JT+l2Lu9/YVqFpOVSNxeSVSqviwj8Sn4dr739hWua1G2pOrVmoRXDvb7Eut9wYfWYUopLdjCK8EkWe9U1DhByhadclwdXuXZHv6yFQrXslUuoOFDOY0OuXfP/D78l+lux4cB4PJThGnCMIRUYxWEksJI+LmvSt6M69epCnTgt6U5PCS7WyhqN/bafbuvc1d2PBJJZlJ9UYrm2+w1u/up3LjeaonSoRkvg9ovae91OSXx6nYlwXe+JmI2ja8QraldvUqc5127bTIrecZvddZds/mw/N5vr7D6sNLqaru1LunKhpy+JQa3ZV+xzXyYfm9fX2F1pulVLqpC81OGFF71G1byoPqlPqlL7F3viZ1chvXhGK77yiEIwjGMUoxSwkuCRPAlgw2iAHEAMkccDiA4kgjxAl8AAAAAAAABw5AAAQ8kgCMjiOIEgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAwAAIJAAAAAAAAAAAACCeBAAAASAAAAAAAAAAABHAAAOAEgAAAAAAAAAAAABBIAEEgCCQACAAAAAACAJIJ4EACnXo069KdKrCM4TW7KMllNdjRVAGrXunz0yE4qnK70ySxKElvzorsfzofau9cmn3lXTIRe/O80uSzCUXv1KK7n8uH2rv6towjCX2k1LepO60uKe896rat4jN9bh82X2PrxzJRO+0tc113hl7etRuKMK1CpGpSmsxnB5Ul3M+qtOFWEqdSCnCSw01lM1i0rTozndaX7L3/AOsWlT2U5dfD5E+/k+vtM9puo299CTpNxqQeKlKaxOm+9fv5PqMTGmYtEoxXs3w37i37Oc6f+Jfb4lzRq061NVKU4zg+Uk+B954ci0rWj9Y61tP1NZ8ZcMxn9Jfv5mEirauFSVa0mqVRvMotexPxXb3oUb2PrFRuYO3rPkpP2ZfRfX9/cKV4vWKhcw9RWfxU3mM/ovr8OZXq0qVam6daEZwlzjJZQFRPqHDrLL1Vza8beTr0l+SqP2l9GX7n7ypbXtCtP1e86dVc6U1uyXl1+KDK5fYR95PDAXP/ACAJLrGCcDAEYC8CQAAAAAZAeRHmA8AOAwODHABgY4jrJ8gIwnxJ4BDrAhEgAAAAAPmUoxi3JpJdb6gJyU61SnSg6lSpGEIrLlJ4SLWV7Ou9ywp+ufXUlwpx8+vyPqnZRc1Vu6nwiquK3liMfox6vHiwPj1txdLFsnQov8tOPtSX5sX979zLm2t6NvFqmuMuMpSeZSfa31n1WqU6MHUqTjCCWXKXBIs/W3N5wt96hQf5aUfal9FPl4v3BiVe4u4wqepoxdau1lU08YXbJ9S/9cSLe1frFcXM1Wr9Tx7NPuiurx5lS2oUbanu0o4y8tvi5Ptb62VnwXAMnIx2sapSslGlGDr3NRZp0Iv2pd77I9rf28i31DV5TqTtNN3KtWL3alaXGnRfZ+dL81ebRiqUKjualtp8fhd9PDuK9V5UOx1GvsgvsXElEfbXa/xClWnU+EwrXebvUKuY0KNPlFdagnyj2zf8EZrRdIlRqq+1CUa17hqKj8Sin8mCf2y5vuXAutJ0qhYRlUzKtdVPxtea9qfd3RXVFcF9pfmJn6K0+ZSADDYEE5HmBABPACASAAAAAAAAAAA4AQASBBIAAAAAAAAAAAAQCeBAAAkAAAAAAAAAAABBJAAngRwJAEEgAAAAAAAAAAAAA4ACCQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAYAAxuq6VTu5fCKM/g93FYjVis7y+bJfKj3dXVgwFxGbu6dK8jOw1CGVRrU3wn27knwku2D93WbiUL20t7y3lQuaUatOXOMvvXY+8zFtIWpvww9jrU6M42+rKNKTe7C4isUqj7H8yXc+HYzPJprga1e6dd2MZLdnqFk1hprerQXevyi+3xKWm161rSVTS68bq1zj4POfxe6Enxj9GXDwM9MfCMWmO0tnrUadam6dWEZxlzUllFr6q5tfxEnXo/2VSXtL6Mnz8H7xp2qWt65U4SlTrRXt0ai3akfFdnesovSLZGpULa6o124puNRfGpzWJLy/efVxb0LiG7VpxmlxWea8H1C4tqNxFKrBNr4r5OPg1xRbuN7b/EkrqmvkzeJrwfJ+ePECfU3lv8AiKqr018is/aXhL+OfEmnf0t9U7iM7eo+CjVWE/B8n7yaF9QqVFSk5Uqr/J1Fuy8u3yLicYTg4TipRfNNZQH1ldpOSz+Aqnxta1S3/NXGH6L/AHYI9bfUfxlvGuvnUXh/ov8AiDa9BZw1G1bUak3Ql82rFw+/gy6jKMo5TTT60Db6AAZAAAAAAAAAAAADYAFOtXo0Y71arCmvzpJFt8PjU4W1CtX74x3Y+94+wMbXvDkU69alQg6lapCnFdcngttzUKzzOpSto9lNb8ve+C9x90rG3pzVRwdSqvylR70vt5eQFNXlWvwsreUk/wApU9mHl1v/ANcQrBVXv3tWVy+e41imvq9fnku5ThTi5Tkopc23hItXfOt7NlRlXfz292mvPr8sgXS3YxSSSS9yLSV7Kq3Cyp+vlydRvFNefX4LJMbOVbEr2p63/RxWKa8uvzLyMVFKMUklwSXUBaUrLemq13P19RPMU1iEPCP73ll3yDeFxMNea0qkpUdLhG5qJ4dVvFGD75fKfdHPkNbYmYhkr68tbG3lXuq0KVNcMyfX1Jdr7ka7qF9c6hGW/KdhYJZknLcq1F2yfyI93PvXItqnC9i60qupam1mEIpLcT+auVOPe+L7WZWw0J1JxuNWnCtNPehbw/E032vPx5d78kiUREeWuZtfwstNtq9/ThTsouy0+KwqyhuymuynHqX5z8l1mx2Nnb2VvGha0o0qa44XW+tt9bfayulhYJMTO2ytdAAMJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAB3GL1HRqNxVlc285Wl0+dSmuE/px5S+/saMoAxMbajqMHR3Y6xbqluP8HeUW1CL7VL41N+PDvZdW+o6jYpet/3StsZU44jWS/Zn5YfibHKMZRakk01hpmFutAhByqaXVdnJ8XSxvUZP6HV9XHmS3E+WuazHhfadqllfxbtq6lKPx4NOM4eMXxReLDNOvoeqnF6rZu2nD4l1Sk9xeFRcY/WwvEvbS91K3jFwq09RoNcN9qNTHdJezLzS8TE1+iMnxLYa9CjXg4VqcKkeySyWvwOrR42lzOKX5Or7cf4r3lK01qyrVI0aspWteXKlXjuN+D5S8mzJZMJ9pWnwqvS/wB82s0vn0vbj7uf2FWhdW9d4pVoTkucU+K8UVyjXtqFf8bRhNrk2uK8GGVScITTjOKkn1NFq9NtU80oSovtpTcPsXAfA50/973den+bJ78ft4/aM6jT5wt667m4P7c/eA+D3cPxV9Jrsq01L7sMb+ow50reqvzZuL9zT+8hXs4vFazuafeoqa/u5PqOpWTeHcQg+yfsv7QI+F1or8JYXC+i4y/eSr+n8qjcw8aEv3Ir069Gp+Lqwn9GSZUAtP5StOupKP0qcl+4fyjZf9Ih9pdcCcICz/lOy6q8X4Jk/wAoW3U6svo0pP8AcXXDuJ4dwO6z+Hp/EtbqX/hNffgn4RdS/F2M1/8AUqRj92S6ysFKpc21L8ZcUofSmkBRa1Gfyraiu5Ob/cR8Cqz/AB97cT7oNQX2cftJepWn5Ocqz/0UHP7kPhdxP8TY1nnrqSUF+9/YDs+6Fja0Zb0KMN/5zWZe98S4eFxfItNzUKjzKtRoLshFzfveF9g/k+jJ5uJ1bh/6WWV+iuH2ATO/tozcITdafzaSc3544LzPhy1Cvwp06dtD51T25+5cF7y8hCFOChCMYxXJJYR9ZAsqenUnJTuJTuZp5zVeUvCPJe4vEkuCR8XFxRt6LrV6tOlTjzlOSil5sw9xr3rFjTbWdx2Vaj9XS97WZeSfiIjbEzEM22kYm8122hOVGzjK9rx4NUmtyL/Om+C8OL7jCX9aVTH8rXzmp/FtqacYS7lBZlPzz4FzZ2OpXUIwpUIabark6kE6mPzYLhHzz4EunXlrm8z2h8ajXnVgp6vdRVKbxG2pZUJvs+dUfdy7iva2OoXsIpRemWiWFwXrmuxLlBe99yMrp2kWVlU9dCEqtw1h16r3qjXj1LuWEZBDf0lFPta6dYWthRdO1pKCk8ylnMpvtk3xb72XQBFsAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAgCR5kAASAAAAAAAAAAAAAEDgBPuBBIEAkAAAAAAAAAAQAJIAAAkAEAAAAAAAAAQBIIAEjzAAAAAAAAAAhxTTTS4mJudBtJTdWzlUsar4uVB4jJ98H7L92TL4IEdmJiJa3c2mp0IOnc2dHUbd83RSUsd9OXB+T8i2s60KdR0tMv6trUXF2tZNpfUn7SX0WjbeBb31laXtP1d1b060erejlrwfU/AztGafTF0dauKHs6hZvC51bbM4+cfjLyyZSxvrS9pudpc06yXNRfFeK5oxVfQ7ijx06+mkuVG5TqR8pfGXvZir+nKlUU9T0yrRlHlc0W5pfXjiUV4ozqJR6rR5bkmn1kmqWN/ewpqpZ31K+o9Ua7y/BVI/vTMjR2gt4tRv6Fayl86a3qb+vHgvPBjplKLwzKSEoRksNJ+J8Ua1KtTVSjUhUhLlKLyn5lQimt6llZzy52tFvt3EfH8nWi+JTlD6M5L7mXYMmlr8Cp/JrXK/8aX8SPgS/6Vdf+ay7AY0tPgUf+k3T/wDFY+AUX8arcS8a8v4l0SDULVafZ5y6Kl9JuX3n3Ttban8ShSj4QSKw4BnRhdwwuYyslnqGqWNhj4Vc06cn8WGcyl4RXF+QYmYhekN4ZgK2tXlbhY2Xq4v8rdPd81Be0/PBjL2tTq1FS1LUK13VlytqSaT/APDhxa+k2Z6ZQnJHw2G61qwoVHRhUdzXXOlQW/JeOOEfNosq17qtwm4ep0+lzbeKlTH7Mf7xb2drqdakoWtlR02h1Osk5eVOPBeb8jIUNBtW1K+qVb+a44rv2F4QWI+9Mz2hj9VmFp/Bq1dSoUrjV7mP5T8Yov6TxCPkZGlpepXT3ru5haU3+Tt/aqec2sLyXmZ6EIQgoQioxXJJYSPoxtmKfay0/S7GwzK3oJVH8apJuU5eMnxZeJYDBhOIiEggcAykAAAAAAAAAAACAJBAAEgARgkAAAAAAAAAACABPAEcCQIBIAAAAAAAAAAgASCOBIDgCCQAAAAAAAAABAEggkAAAHAAAAAAAAAAAAQSABBIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACPIkAAAAAAAAAAAAAAAAAB5AAAAAAAAAAAAAAAEAkAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAIaySAMZe6Hp11Udb1Pqa7/LUHuT82ufnkx9XStUtcuhWpX0Pm1fwdT9JLdfmkbGGZiZhGaRLS0re2r5qK50e4k/jfi4yfis05GVpX2q2yXrFQvqfU1+Cqfvi/sM7UpwqQcKkVKLWGmspmJraBaJudjUq2E3xxQfsPxg8x9yRnf2j0zHhVttbsatSNGtOdpWfBU7iO434Pk/JsyWUazdWmrUKcoV7WjqNB83RxGWO+EuD8n5FjaVaVOr6qxvbjT6v/R55S/8ALn/s4GjrmPLdMkmvUdW1K34XNrSu4rnO3luS/Ql+6RdQ2i0zH4apVt5fNrUZRf3YfkY1LMXiWXI5dRiKmvUpr+pWl1dN8nuerj754+zJY3V7qlWDnWurfT6PX6r25JfTkkl+iNSTeIbBc3FC2pOrcVqdKC+VOSivtMbV1uNRYsLWrddk5fg6f6T4vyTMFa+qrVlUsrS41OsuVeT3kv8AxJ8F9Uy9HTNUuON3d07WD+Rbx3p+c5LHuiZ1EI9U28LTULm7lT9ZqGpwtKT+RRfq0+7fftPywW9jb1Zt/wAl6ZP2udesnTUvFyW/L3eZsNjpFhaVPWwob9f+2qtzqP6z4ryL/A6teCMe/LBW+g1avtahezkuulb5pw83nefvXgZWysLOyhuWltSox69yOM+L6y5BjcpxWIMEEgwkcgAAAAAAAAAAAAAAAB5AACCQBBIAAAAAAAAAAAAQPIkAAAAAAAAAAAAAADyAAAjyJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAoXdpbXdL1Vzb0q0PmzimvtK4BphK2z1KOfgN5cWv5jfrKfulxXk0W70zWoPEZWFVfOzOm/d7X3mxgz1ShNIa/S0nVKn4+9t6EetUabnL9KXD7C7ttB06nNVK1OV1VXKdxLfx4J8F5JGVA3LMUiERjFJJJJIkAwkAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADIAAAAOAAAAAAAAAAAAAAAABAEgAAAAAAAAAAAOAAEACQCAJAAAAAAAAAAAAgCQAAAAAAAAAAAAADJAEgAAAAAAAAAAAAAAAAgkAAAAAAAAAAAAGQAAHAAGAAAAAAAAAAAAAAAAAAAAAAAAAAAAADIAAAAAAAAAAAAAAAAAAgkAAQBIAAAAAAAAAAAEASCOBIBAAAAAAAAAAAAAAIJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAB5EEgAAAAAAAAAAAAAAADiBAJ8gAAAAAAAAAAAAAAQCfIAAAAAAAAAAAABBIAgkACOskAAAAAAAAAAAAQBAAE+Q8gAAAAAAAAAAAAAAQT5DyAAAAAAAAAAAAAABA8iePYA8gAAAAAAAAAAAAAAAQCQAAAAAAAAAAAAAACCfIAAAAAAAAAAAAAAAgkeQAAAAAAAAAAAAAAI8iQAAAAAAAAAAAAAAAAAAAAZAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADrAAAAAAAAAAAAAAAAAAAAAPIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAGAAyhlAAMrtHmAAyMgAMoZQAAAAMrtGQAGQAAAAAcAAAyAAHAAAAAGUAAGUMgAMoAABlAAAAAyu0eYADKAADIAADgAAAADK7RkABkAAAAA4DKAADzAAAABkZAAZHmAAGUAA4DIADgAAHmAAAygAHAZAAAAAAAGRlAABlAABwAAZQygAA8wAAygAGRwAADKAAZXaOAADK7RlAAAAAygAAyu0AAOAAAZQ4AAMrtAADIAAZQ4AABnvAADKAADIADKAADzAADIAAAABkZQAAAAAAAyOAAAcAAAAADK7QAAyAA4DgAAyhkABlAAAOAADgOAAAAAAAAAADKGUAAyPMABw7RwAAcBlAAMjIADgMoAAMoAAMgAM95DfBgf/Z" style="width:32px;height:32px;object-fit:contain;" alt="NX"></span><span>Nexo</span></a>
        <a class="btn" href="/">Volver al chat</a>
      </div>
      <h1>Ayuda a que Nexo pueda seguir creciendo</h1>
      <p class="lead">
        Esta IA corre en un PC local y cada modo avanzado consume mucha GPU, memoria y tiempo de procesamiento.
        Las donaciones ayudan a mejorar componentes para mantener respuestas mas rapidas, modo codigo estable y mas capacidad para usuarios.
      </p>
      <div class="actions">
        {% if donate_ready %}
          <a class="btn primary" href="/donate/go" rel="noopener">Donar ahora</a>
        {% else %}
          <span class="btn primary disabled">Donacion pendiente de configurar</span>
        {% endif %}
        <a class="btn" href="/register">Crear cuenta</a>
      </div>
    </header>

    <section class="grid">
      <article class="panel">
        <h2>Por que hace falta apoyo</h2>
        <p>
          Nexo combina chat rapido, modo combinado, modo codigo, archivos, memoria e IA local. La GTX 1080 Ti todavia es muy capaz,
          pero los modelos grandes cargan mucha VRAM y el sistema puede saturarse cuando hay varias tareas o respuestas largas.
        </p>
        <p>
          El objetivo es mejorar hardware para sostener mas contexto, menos esperas y una experiencia mas fluida sin perder privacidad local.
        </p>
        {% if not donate_ready %}
          <div class="note warning">
            Configura un enlace externo con <strong>NEXO_DONATE_URL</strong> o <strong>donate_url</strong> en settings para activar el boton de donacion.
          </div>
        {% else %}
          <div class="note">
            Gracias por apoyar el desarrollo. Cada mejora de hardware ayuda directamente a que Nexo responda mejor.
          </div>
        {% endif %}
      </article>

      <aside class="panel">
        <h2>Equipo actual</h2>
        <ul>
          <li><span>CPU</span><strong>Intel Core i7-9700K</strong></li>
          <li><span>RAM</span><strong>32 GB RAM</strong></li>
          <li><span>GPU</span><strong>NVIDIA GTX 1080 Ti 11 GB V-RAM</strong></li>
          <li><span>SSD</span><strong>SSD SATA 500 GB</strong></li>
          <li><span>HDD</span><strong>TDisco Duro HDD 2 TB</strong></li>
          <li><span>Sistema</span><strong>Windows 10 64 bits x64</strong></li>
        </ul>
      </aside>
    </section>
  </main>
</body>
</html>
"""

ADMIN_HTML = r"""
<!doctype html>
<html lang="es">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>Admin - Nexo</title>
  <style>
    :root {
      color-scheme: dark;
      --bg: #161616; --panel: #202020; --panel-2: #292929;
      --text: #f4f4f4; --muted: #b8b8b8; --line: #3a3a3a;
      --accent: #19c37d; --danger: #ff9a9a;
      --gold: #f0c040; --gold-bg: rgba(240,192,64,.10); --gold-border: rgba(240,192,64,.35);
    }
    * { box-sizing: border-box; margin: 0; padding: 0; }
    body { background: var(--bg); color: var(--text); font-family: Inter,ui-sans-serif,system-ui,-apple-system,"Segoe UI",sans-serif; min-height: 100vh; }
    .page { width: min(1100px, calc(100vw - 32px)); margin: 0 auto; padding: 36px 0 60px; }
    .topbar { display: flex; align-items: center; justify-content: space-between; gap: 12px; padding-bottom: 20px; border-bottom: 1px solid var(--line); margin-bottom: 28px; }
    .brand { display: flex; align-items: center; gap: 10px; color: var(--text); text-decoration: none; font-weight: 800; }
    .mark { display: grid; place-items: center; width: 32px; height: 32px; border-radius: 6px; background: #f2f2f2; color: #111; font-size: 13px; font-weight: 800; }
    h1 { font-size: 22px; font-weight: 800; }
    .badge-admin { display: inline-block; background: var(--gold-bg); border: 1px solid var(--gold-border); color: var(--gold); border-radius: 4px; padding: 2px 8px; font-size: 11px; font-weight: 700; letter-spacing: .05em; margin-left: 8px; }
    .btn { display: inline-flex; align-items: center; justify-content: center; gap: 6px; min-height: 38px; border: 1px solid var(--line); border-radius: 6px; background: var(--panel-2); color: var(--text); padding: 0 14px; cursor: pointer; font-size: 13px; font-weight: 600; text-decoration: none; }
    .btn:hover { background: #323232; }
    .btn.danger { background: rgba(255,154,154,.12); border-color: rgba(255,154,154,.3); color: var(--danger); }
    .stats { display: flex; gap: 14px; flex-wrap: wrap; margin-bottom: 24px; }
    .stat { background: var(--panel); border: 1px solid var(--line); border-radius: 8px; padding: 14px 20px; min-width: 140px; }
    .stat-val { font-size: 28px; font-weight: 800; }
    .stat-lbl { color: var(--muted); font-size: 12px; margin-top: 2px; }
    .stat-val.gold { color: var(--gold); }
    .legend { display: flex; align-items: center; gap: 8px; margin-bottom: 14px; font-size: 13px; color: var(--muted); }
    .legend-dot { width: 12px; height: 12px; border-radius: 3px; background: var(--gold-bg); border: 1px solid var(--gold-border); }
    table { width: 100%; border-collapse: collapse; font-size: 13px; }
    thead th { text-align: left; padding: 10px 14px; color: var(--muted); font-weight: 600; border-bottom: 1px solid var(--line); font-size: 12px; letter-spacing: .04em; text-transform: uppercase; }
    tbody tr { border-bottom: 1px solid var(--line); transition: background .12s; }
    tbody tr:hover { background: rgba(255,255,255,.03); }
    tbody tr.golden { background: var(--gold-bg); border-bottom-color: rgba(240,192,64,.15); }
    tbody tr.golden:hover { background: rgba(240,192,64,.15); }
    td { padding: 11px 14px; vertical-align: middle; }
    .order { font-weight: 700; color: var(--muted); font-size: 12px; width: 40px; }
    td.golden-text { color: var(--gold); font-weight: 700; }
    .username-cell { font-weight: 600; }
    .crown { margin-right: 5px; font-size: 14px; }
    .plan-badge { display: inline-block; padding: 2px 8px; border-radius: 4px; font-size: 11px; font-weight: 700; }
    .plan-gratis { background: rgba(184,184,184,.12); color: #b8b8b8; }
    .plan-beta_tester { background: rgba(45,212,191,.12); color: #2dd4bf; }
    .plan-developer { background: rgba(99,102,241,.15); color: #a5b4fc; }
    .plan-select { background: var(--panel-2); color: var(--text); border: 1px solid var(--line); border-radius: 5px; padding: 5px 8px; font-size: 12px; cursor: pointer; }
    .apply-btn { background: var(--accent); border: none; border-radius: 5px; color: #06140e; padding: 5px 12px; font-size: 12px; font-weight: 700; cursor: pointer; }
    .apply-btn:hover { opacity: .85; }
    .apply-btn:disabled { opacity: .4; cursor: not-allowed; }
    .date-cell { color: var(--muted); font-size: 11px; }
    .msg { padding: 10px 14px; border-radius: 6px; font-size: 13px; margin-bottom: 16px; display: none; }
    .msg.ok { background: rgba(25,195,125,.12); border: 1px solid rgba(25,195,125,.3); color: #6ee7b7; }
    .msg.err { background: rgba(255,154,154,.12); border: 1px solid rgba(255,154,154,.3); color: var(--danger); }
    @media (max-width: 700px) { table { font-size: 11px; } td, thead th { padding: 8px; } }
  </style>
</head>
<body>
  <div class="page">
    <div class="topbar">
      <div style="display:flex;align-items:center;gap:14px">
        <a class="brand" href="/"><span class="mark" style="background:transparent;display:inline-flex;align-items:center;justify-content:center;"><img src="data:image/png;base64,/9j/4AAQSkZJRgABAQAAAQABAAD/4gHYSUNDX1BST0ZJTEUAAQEAAAHIAAAAAAQwAABtbnRyUkdCIFhZWiAH4AABAAEAAAAAAABhY3NwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQAA9tYAAQAAAADTLQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAlkZXNjAAAA8AAAACRyWFlaAAABFAAAABRnWFlaAAABKAAAABRiWFlaAAABPAAAABR3dHB0AAABUAAAABRyVFJDAAABZAAAAChnVFJDAAABZAAAAChiVFJDAAABZAAAAChjcHJ0AAABjAAAADxtbHVjAAAAAAAAAAEAAAAMZW5VUwAAAAgAAAAcAHMAUgBHAEJYWVogAAAAAAAAb6IAADj1AAADkFhZWiAAAAAAAABimQAAt4UAABjaWFlaIAAAAAAAACSgAAAPhAAAts9YWVogAAAAAAAA9tYAAQAAAADTLXBhcmEAAAAAAAQAAAACZmYAAPKnAAANWQAAE9AAAApbAAAAAAAAAABtbHVjAAAAAAAAAAEAAAAMZW5VUwAAACAAAAAcAEcAbwBvAGcAbABlACAASQBuAGMALgAgADIAMAAxADb/2wBDAAUDBAQEAwUEBAQFBQUGBwwIBwcHBw8LCwkMEQ8SEhEPERETFhwXExQaFRERGCEYGh0dHx8fExciJCIeJBweHx7/2wBDAQUFBQcGBw4ICA4eFBEUHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh7/wAARCAMABYADASIAAhEBAxEB/8QAHQABAAEFAQEBAAAAAAAAAAAAAAEDBAUGCAcCCf/EAF0QAAIBAwEEBgQICQgDDQgDAQABAgMEEQUGEiExB0FRYXGBCBMikRQyQlJygqGxIzNDYpKUssHRFRYXJFOi0uFjs8IlNDU2RUZUVmR0g5PwGCZEVXN1hKNlw/Hi/8QAGwEBAAIDAQEAAAAAAAAAAAAAAAIFAQMEBgf/xAAtEQEAAgIBAwQABgIDAQEAAAAAAQIDEQQSITEFE0FRFBUiM1JhIyQyNEJxgf/aAAwDAQACEQMRAD8A7LAAAYAAEYJADAwAAwMIAAMAAAAAGAAAwAAwMAAMAAAMAAMDAAAYAADAADAwAAwMAAMDAADAAAAAAMAABgABggkAMDAAAYAAAABgYAAAAAMAAQTgAAAAGBgAAAAAAAYGAAAAAYGEAAGAAAwAAwAAGEMAAMDAAAYAADAADAwAAwMAAAAAAADAwAAwAAGBhAAMAABgYAAYGEAAAADAAAAAAAAAwAAGO4ABgYAABIABgYAAAABgYAAgnAAAAAMDAADAwAAAADAwAAwAAAwAAwMIABgjBIADAADAwAAwMAAAAAGAAGBgAAMAAMDAADAx3AAMDAAEYJwAAwMAAMIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAYAAAAAAAAAAAAAAAYAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAD5AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAwAACAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAABDJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAABxAAcQAAAAAAAAAAAADkAAAAAAAAAAAAAAAAAAHEAAAAAAAAAAAAAI4k8QBHEkAAAAAAAAAAAAAAAAABxAAAAAAAAAAAAAAAAAAAAAAAAAAAAABxAAAAAAAAAAAAAAAA4gAQSAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACOIAEkEgCCQAAAAAAAAAAAAcQAAHEAAAAAAAAAAAAAAAAAAAAAAAAAAGAAAAAAAAAAAAAAAAAAwAAAAAAAAAAAAAAAAMAAMAAAAAAAAAAAAAHAAABgAAAAAAAAAAAAAAADAAAAAAAAAAAAAMAABgAABgAAAAAAAAAAAAGCCQAwQkSAAAAAAAAAABDYE4HMpV69GhTdStVhSgucpySSMdLaHSE2o3iq466UJTXvimNMdUMtgGOtNb0u6moUb6i5vlCUt2XufEyCawCJiUgAMgAAAAAAABGCQAwAAAAAAAAAABHWSAAAAgkAAAAAAAAAAAAAAADAAAAAAAAAAAAAAAA4AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAGAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAHxUqQpwc6kowjFZcpPCRiqu0FnKThZRq381wfqFmC8ZvEftM6Rm0Qy+UU7m4oW9J1a9WnSgucpyUUveYGrealXi5Va9DT6S5+r9ueO+csRXuZZ0FZVavrLS1uNUrr8q/wiT+nP2V5GelH3N+GWq6/QmsWFvXvX86Ed2n+lLC92SxuLzVK0XOveULCiuaorLXjOfD7C5p6bqdzxr1qNnH5tJetn+lJYXuZdW+hadSmqlSlK5qrlUuJOo14Z4LySG4hHV7NepU7S4qqdtaXGqVl+Vlma/Tn7K8jJQs9anFYpWVCPVGVWUmvckjYlFJJJYSJE2IxR8tYurS+cHG80yldQ/0U1P8AuzS+xlrRqUaNRUrK+ubCr1UJ5S/Qn/sm4FK5tre6pOlcUadWD5xnFNfaIse1rxLC0dW1Kg925tad3FfLt3uT/Ql/Ev7XWtOuKipOt6ms/wAlWi6cvJPn5FCehUqfGxuq9t+Y36yn+jLl5NFvc299Gk6d3YUb6j1+pay/GE/3NjtJHVHln8onJqdq7elUVLT76vY1f+j1M4/8ufV9HBf09V1Cg926so3MV8u2eJecJfubHTKUZI+WdBj7HWLC7mqVOuo1v7KonCa+q+JfriRTiYnwkABkAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAgASAPMAAAAAAAAAAAAIJAAgkAAAAAAAAAAAABAEggnzAAAAAAAAAAAAAAABAEgEAOJJBIAAAABkABlFveXlrZ0fXXVenRpr5U5YyDa4bIzgwtbWa9fhp1nJxfKtc5pw8o/GfuXiY2+nGc1DVdRq15S5W1LMFLwhHMpebZmKoTePhmrzW9Ptqro+udeuvyVCO/PzS5eeDH19R1a5T9RTpWFP51T8LU9y9le9nxa2d9UpqlZWFHT7fqlWST8qcf3teBe0dAtpYd9VrXz+bUaVP9BYXvyZ7QjPVZgJfBrqvicrnWLiL5fjIxfgsU4+ZlaOnatcJKcqFjS5JJetqL7or7TP06VOlCNOlCNOEeCjFYS8j7HWRi+2Lt9CsYSVS4jO8qLlO4lv48F8VeSMpGKjFJLCXUggR22RWI8JBGUTkMgIcornJLzPl1afXUh7wPvqB8KpB8px96JynxT+0D6I5kNdhOe8Cjd2tvdU/V3NCnWh82cU19pjKuh+rWbC8r26/s5v1lP3S4ryaM0DMTMIzWJarf2t5ubl/ptO8pr5dFb+O/cfFeTZRsa9SEnHTdUqRcedvXzUUfGMsTj7zb8Itb7TrO9ildW9Orjk2vaj4PmvIz1Nc4teGNo63cUeGoWE1Fc6tv+Ej5x+MvczKWN/Z30HO0uadZLnuy4rxXNeZjKui3NDjYX0mlypXK315SXtLzyY69oqFRVNS02pRnHlc0cyS+vH2l5pDUSRa1fLbQa1Z3d9TpqpaXtK/odUaz4+VSP70y/o67bJqF9TqWM28L1y9h+E17P2pmNSnF4llgfKkpJNPOeKJy+8wmkEJ5JAABAAAAAAAAAACAJBBIAAAAAAAAAAAAAABBIAEEgAAAAAAAAAAAAIAkAAAAAAAAAAAAAAAAAjzAkAAAAAAAAAAAAAA4gAAAAAAAAAAAAAAAIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAIwBJBIAAAAAAAAAAAAAAAAAgkYAAAAAAAAAAAAAAAIJAEEgAAAAAAAAACCQAIJAEDBIADAAAAAAG8Ftf31rY0HXu68KUM4Tk+LfYlzb7kDelyWWo6nZ2CXwisozl8SnFb05+EVxZirnUNQveFunYW7/KVIp1pLui+EF3vL7kY6yiqlSS0m2ld1JPFS6qTe43+dUeXLwjnyJRX7apyfEMhcalqF2mqMY2FL59RKdXHh8WPnnwLC09VWuHUsaFXUrlcHcSlvJf+JLgvCPuMpabPwniep1neS5+q3d2ivq/K+s35GbpwhThGEIqMYrCSWEh1RHhiKWnyw1HSbuu96+vPVRfOla5j75v2n5bpk7Kws7KDja29Olnm4rjLxfN+ZccC3rXttSn6uVRSqf2cFvS9yMbmWyKxC4wgWnrb2t+KoRoRfyqzy/0V/EfAvWf75uK1btjndj7l+/JjTO1Wtd21KW7UrQUuqOePuKfwuU/xNrXqd7juL+9grUaFGgt2jRp01+bHBUlKMYtyaSXW2BaKWoTxinb0V+dJyf2Y+8K3vJ/jL1rup00vvyfUtQs08KvGb7IZk/sPn4ZOX4qzuZ97io/e0GOx8BT+PdXU/wDxWvuwP5Ns38anKf0pyf3seuv5fFtKUPp1v4JhR1GXy7WHhGUv3oHZK06xX/wlHzifSsLL/otD9BHx6m/fO8pLwo/5j1F5133/AOpBl9uwsnztaP6CPn+TbHOVbU0+5YIVC8/6cn40UT6q/i+FzRl40X/iAj+T7ZfE9bD6NaS/ePgU48ad7cx8ZKX3pjOoxfxbWa8ZR/iPX3kX7djvf/Tqp/fgMdj1V9F+zdU5rsnS4+9P9w9bew+Pawmu2nU/c0vvHw6EfxtG4pfSpNr3rKKtG8tazxTuKcn2KSz7gPj4dSjwrQq0e+cGl71w+0r0qtOrHepTjOPbF5R9PBQq2dtUlvujFT+dH2Ze9cQyuBhFp6i6pL8Bdb6+bWW99q4/eQ7upSf9ZtqkF8+n7cfs4/YDajd6NY16jrRhK3rPj62hLck/HHCXmmWNex1O3hKOKWoUWsNcKdTHh8WX90zdG4o1o71KpGce2LyVFxMxMwjNIlp9s40azhp11VsKy4u1qR9n/wAt9XfFoylHXZUPZ1S2dFL/AOIpZnS8+uPmsd5lL6xtL2l6u6oQqxXFZXGL7U+afejE3GmXtqnKzq/Cqa/JV5Ymvoz6/CXvJbiUNWr4Zq3rUa9KNWjVhUpyWYyjJNNdzKvmahbKEbqTs51dNvfjToyjje73DlJfnR95lLfW5UWqWqUlbZeI3EXmjJ+PyX3S97MTVKuSJ8s0SRGSkk0+ZJFsAAAAAAAcADI8yQBBIAAAAAAAAAAAAAABBKAAAAAAAAAAAAAAAIJAAgkAAAAAAAAAAAOAAgkAQSAAAAAAAAAAAAAAAMAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQ5Jcy21G/trC3de5qbkc7qSWZSk+Silxb7ka5qV1WvaM6upT+BWEeLoOeJSX+kkv2F5t8jMRMoWvEMjea26k5UNLhCvNPEq8vxNN+K+M+5ebRjI4V88KtqWpYw3wzTT/u04/a+8uLGyvL+EFGM9OsUsRSju1prsS/Jx/vfRNgsrS2srdULWjGlTTziPW+tt9b72Z8IxE27yxdvoruMVNVnGt1q2p5VFePXPz4dxmI04QgowioxSwklhJH1yLWpeJ1HStoOvVXBqL9mP0pdXhz7iO5lsiIhc5S58C1nfKb3LSnK4l2x4QXjLl7siNpKs1K9qet6/Vx4U15dfn7i6jGMYpJKKS4JLkDvKzVtc13m7uHGP9lRbivOXN/YXNvQo0IblGlCmuyKwUKl/Tc3Ttoyuai4NU+S8ZckfPqLyuv6xX9TB/k6HPzk+PuwGFevc0KCXrasIt8k3xfgubKLurirwtrSbXzqr3I+7i/sK1va0LfLpUoxk+cucn4t8WVwz3WSoXlX8bdqmuuNGGPtef3Ex021T3pwdaXbVk5/eXgBp8whCEcQiorsSwTgkBkI5kgAkAAGAAAAADBSrW9CssVaNOf0oplUAWnwCnHjQq1qD/Mnle55R841ClylRuYrqa3Je/ivuL0BjSzV9CHC4pVLd9s17P6S4FzCUJxUoSUovk08o+mk+ZbTsaO+50t6hN85Unu58Vyfmgd017KhVn6zdcKnVUg92XvXPzKf9dt/m3UF4RqL9z+wlyvbf49NXMO2Hsz9z4Mq211QuMqnP2l8aEliS8U+IEW93RrS3FJxqLi6c1iS8ivz4Mp3FvSrpKpBSxyfJrwfUUVG5t/it3FNdTwprz5PzB3h931la3tL1d1RjUiuKzzi+1Pmn3oxFzY3tnFum5X1tjDhLHrYrz4TXjh+Jmbe4pV09yXtL40WsSj4rqKw3piaxLVbCVS2j6zR60ZUU8TtKragn2Rzxpvu5dyM3puqW97KVLEqNzBZnQqLE49/eu9ZRGpaXRu5+vhKVvcpYVanzfdJcpLuf2GBvYSjVp0NUp+oqqX4C6pNxi5fmy5xl+a+feS7S1zM0bcga/Z6xWs5Ro6s06TeIXiWI+E18l9/J93Iz8WmsriYmNNlbRZJGOJIMJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAGQ3jmgJ5GL1bVo2s1bW1P4ReSWY0k8KK+dN/Jj9r6ky31LValWtUstMcXUg92tcNZhRfYl8qfdyXX2PH2lObqzstMhv1t7Nxc1XvKEu2b+VPHKK5deFglENdrfEPie/G7jUruV9qdVP1cILG6uvdXKEO2T597wjK6bou7VheajKFe5jxpxS/BUfop83+c+PhyL3TNPoWNJqDlUqz41a0+M6j7W/uS4LqLzkJt9FcffciWEUbq5pW8E6kuMniMUsyk+xLrKNe5q1KkqFnFTnF4nUl8Sn49r7l54KlrZ06MnVlJ1a0liVWfPwXYu5EU/8A4oqlc3nGu5UKL/JQl7TX50ly8F7y8o06dGmqdKEYQXBRisJFK6u6NthSblOXxKcFmUvBfvKCoXF3xu5eqpdVGnLn9KXX4Lh4gfdS9UqrpWtN3FRcHuvEY/Sl1eHF9x8qyqV3vX1b1i/sYezTXj1y8/cXdKlClTjTpwjCMeSisJH2DX2+acIU4KEIKMVwSSwkfQAZAAAAAAAAAAAAAAAAAAAAAAAAAAAKFza0LjHrKabXxZJ4lHwa4orgCyxd2vJu6p9jaVRefJ/YVra5o3Capy9qPxoNYlHxXUV2s8y3ubWlXxKScZx+LOLxKPgwxrRc21Ku1J5jUj8WpB4kvP8AdyKLuK9pwu479LqrQXL6S6vFcPAOtcWj/rKdWl/bQjxX0o/vX2F3CcKsFOElOMlwaeUwJhOM4qUWmmsprkz4uaFK4oypVqcKlOaxKMllNeBbytqlvJ1LJpLnKi3iEvD5r+z7ytbXNOvvJb0Kkfj05LEo+P8AEHnywl1p9xYKXqYzu7FrEqT9upTXdn48e58fHkUNPr1tPpxq6fJ3unvj6iLzKmv9G31L5j8scjaOZidR0uSqzvNOcadeXGpTk8U63j2S/OXnklE78tc013hkLG7t722jcW1WNSnLk19zXU+4rmp0aklc1LmxzbXkWvhNvV4KfdNLr7JrPmjPaVqdC/hJRUqVenwq0Z/Hpvv7U+prgzEwlW8SvgAYTAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADAAAAAAAAAAAAAAAAAEAcABIAAAAAAAAAAAACCQBHADgSBBIAAAAAAAAHUADBAADgABIAAAAAAAAAAAACBwAAAkCCQACAAAAAAAAIJKVetSo0p1atSNOEE5SlJ4UUutgfcmksvka1qOpVdTk7bT6sqVmm1UuoPEqvbGm+ztn7u1Rf3MtVhKVSbt9KS3nGb3ZV186fzafdzfXhcH9Weny1RKVWEqOncFGnjdlcLvXyYd3N9eFwcojXdqtabdoUdOtZ3sI0LD+rWFP2XWp8HPtjT/AHz92XxWyWVrQtLaFvb0406cFiMYr/1nxKsKcIQUIRUYxWEksJLuPm5q07elKrUmoQjzbMTO0q1ir7nKMIucpKMUstt8EWWat/8AEcqNr1zXCdTw7I9/N93MilRqXklWuouFFPNOg+vvn393V48rytVp0acqlWahCKy23wRhLymlTp0aUadOKhCKwklwRaVLmrcSdKx3Wk8Sry4wj3L5z+z7iFGtf8ailRtXyhynU8exd3Pt7C9hCNOChCKjFLCilhIHlQtLSlbtzWalWXx6k3mUvPs7lwLhcsEgMo8xwHAcAIZPVkDgA4doJCQEYQJAAciG12nzKpCKy5JGPLEzEeX2DDahtRs/YZV5rFjRkvkyrR3vdnJgb7pO2St1+DvKty+yjQm/taSJxS0/DRfl4aebQ3cg8wuumTRqcsUdK1Or3uMI/fIsZ9N1ipYWhXX1q0US9qzT+Zcf+T10Hk1HpnoVGlHZ65l9CvF/uL+h0rQq4xsrrrT66dHfMe3ZOOdht4l6UEaTZ9IdtXaVTZ3aSj3y02bX2ZM9YbRadd4Sjd0ZPhu1rSrTf2xIzEw3UzUv4ZjBJ805RnHejxTPow2gAAABgAQSBAwOAAjCLSdpKlUlVspKnJ8ZU38Sfl1PvX2l5hEoC2tbqFZulOLpVorMqcua712rvRN1awr4mpOnVh8SpHmv4ruJubanXS3sqUXmE48JRfcz4j8Pgt1/B6uPlNuLflhhhFC4nGqre7ioVX8WS+LU8Ox933l2WVeneV6bp1KVq4v8+Xv5cGVLGF1Tpbl1UhOSfsyjza7+8C11nTKd6o1ac3Ruqa/B1orLXc18qL60/sfE15+td5ClXzZanRTdKpDjGpHr3c/Hh2xfFd3Bm5tcC0v9PtL+l6q7owrQTUlGS5NdaMxbXlC2PfeFDQNTeo29RzpblWjLcnKPGnN9sJda+7kZM+KdOFKEYQioxisKKWEkfZhOPHcIJI4BkAAE4RBIAAAAAAAAAEEgCAOBIEDgCQI4EgAAAAAAAAACAAAJYAAAAAAAAABgAACAABIAYAAAAAAAAAAEEgAAAAAAAAAAAAAAAAAAgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAZAAAFpquoW+nWvr67fF7sIRWZVJPlGK62GJnSpfXdvZW07m5qRp0o85P7Eu1vsNYv7mV6/hmpr4PZ05J0raXNvPBzS5yzjEFnHe+Xxc1atStC/1PLqb27bWtP2txvqivlVH1y5LjjCy3ltH0mfrYX+oqLuI5dKknmFBPs7Z9svJYXOWtNM2m86h86fptW8nC61Gm4Uk96jay7eqVTtl2R5LvfLOYCPivWp0aUqtWajCKy2yO9t0REQXFanQpOrUliK7uL7kutlrb29StVV1drDXGlS5qn3vtl93UTbUqlesru5i44/E0n8hdr/Of2cu0r3dxC2pb88tt7sYrnN9SXeAuq9O2pesqPhnCS4uT7EutltQoVbipG4vEluvNOjzUO99svu6u0+rW3qOr8Ku8Os/ixXGNJdi7+1/uLzCxwQEjgFwGAyAgngAA6gAAYXIARkPgjxfpU6V9a0HW6ug2OkOwqRWY3d1ifrI9U6cU8Nd7b70SpSbzqHPyORTj06rPYby8tbOhK4uq9KhSisynUmoxXi2aHrXS3sxaTlR0+pV1SrH/o8fYz9N4Xuyc/6zreq63W9dquo3N7POV62eYx8IrgvJFrb1HSqxqccLmu1HdXh6j9TzXK9fvPbFGnrer9Kuu3rcLCnb6fB8sL1k/e+H2Gm6vqurao277VL25z8mdZ7v6K4fYY2LSSlFpp9aPr1hKuOsfCkzc/kZZ/VZTjTdNtxe75I+ncSjBrebwJ1E+GCjOO9F8OD4GyIhz9cz5e1dG+y+xmv7NWupT0uFW4x6u4jUrTnu1I8JcG8cefg0bza7K7N2sUqGh6fTxyxbxz78HinQlr70nan+S7iri01LEI55RrJey/Nez44Og4LgivzdVbae39K9nNhiemNwpUbS1orFKhTgvzYpFXcj81H0iOs07lbRSseIFGPYhhdg8xlDaWoSBnvHmGQAZwAAIYEgAAAAAAAAZCADHHII8wJaRGExwznIz1ASQwaJ0sdIVnsbYRo0acLvVq8W6FtvYUVy359kc+b5LralWs2nUNWbNXFWbWbhqWp6fptrK61C8oWtCHxqlaahFebNI1Hpl2DtJuFPU6t41z+DW85r34SfvObNqNb1jaa+d9rV9VvJ5e5CTxTp90Icl95iFCMeUV5nfj4Uf+peezeu23/jh0/Q6cdhqlTdq1tQtov5dS0lj+7k3LZvavZ7aOm6mi6taXqj8aNOot6PjHmvNHFzUWsOK9x92nrLW5heWletbXFN5hVozcJxfc1xJ24Ea/TLXj9dyRP647O6U8rgSeGdEPS9VrXNvoG11aPrarVO21DG6pyfBQqLkpPqkuD68dfuSaa4Ffkx2xzqXouNyacinVVIANboAAAAAAAAAAAAAAAAAAAAAAAAEAAAAAAAAAAAHAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAcQAAAAAAAAAAAAAAAGAAABAMx2sapTsKcIQputdVsqjRTw5Y5tvqiut9XjhBiZ0+9W1Gjp9KLmpVKtR7tKjD49SXYv3vkus16tKr8KhdXcXc6hVzGhQpPhBdcYZ5L503/BHzmtG7U5YvtWuY8Ir2Ywjnq+ZTXbzb7WbBpGmQslKtVn6+8qpetrNY+rFfJiupeby+JLw1d7yp6RpfqKnwu7lGreSjjMfiUl82Hd2vm+vqSyoWOoPkRltiIjw+ZSUIuUmkkstvkWNCLvq0bqomqEHmjBr4z+e/3e/wAFb+v13RXG2pSxVa5VJL5Pguv3dpe1J06NGVSclGEVlt8kgeXxc3EKFJ1JvhySXFt9SS62UbOhUnV+GXSXrmsQhnKpLs8e1/uPi0pyua6vK8XFL8TTkvir5z/Of2LzL3wAldhJBIZD5bwTy7S21G+tdNsa99fXFK2tqEHUq1aklGFOKWXJt8EkuseWJnXeVdz4ZawattZ0jbEbKzdPX9pNNsay4+olWUqv6Ecy+w5i6bPSG1TaK9uNE2Huqum6LBunO/p+xcXfa4PnTp9jXtPu5HhSowqVJVJOcqk3mU3Jtyfa3zZ2YuJNo3Kr5HqMUnVXcFb0jOjGE3GlqGoXEV8unp9XD96Rt+h9I+x+rqn6jWaNGpUScadynRlx5L2sH5/6fbesuaNFSl+EnGCWX1tHrNV09582uXHijong112lW5PWsmO3jbtOnOFSClCSlFrKafM+snJOxe3Ov7J3Mf5PuHWss+3Z1pN05L835j71w7UzpPYLazTNrtHjf6fNxnF7tehP49GXZJfc+TOLNx7Y1rwvUsfJ7eJbFlmp9J2xNjtloMraoo0r6hmpZ3GONOfY+2L5NfvSNuRBpraazuHdlxVy1mtocZX1vcWF7Wsrui6NzQqOnVg+cZJ8V/mUfWRXA9S9JjZ6VpqNptPa08U7nFtdYX5RL2JPxScfKJ47Cq285LjFk667eC5nE9nLNfhnNOu8p0H1cYfwLxZfX5Gv0ZyjOMoNxknlMz1vUjVoxqLhnn3PsMWq4bx0+H11n1lLhwPiTSXI+XLtwRa1WlUlCtCpSk4ThJShJc4yXFNeDOmtgNfhtHsxa6j7KrY9XcQXyKkeEl+9dzRzBGaXZk9A6DtpI6XtPLSbiri21LChl8I1kuH6SyvFI08inVXa89D5fs5uifEvf0QE01lcgyumdPbw1fbrbrRNjatjDWfhMfh0pRoSp0t6LlFJuLfU8PguvD7DBvpf2WXyNQz/AN2f8S/6ZNkYbabB3+j7kHdKPrrKUvkV4cY+T4xfdJnGekbS6lp0/gl9TqVqdKThOlUeKlNrg0n3PhhnZgw1yV38qP1Hl8nBf9Hh14ul7Zd4e5qH6s/4n1/S3sx1w1D9Wf8AE590a8stUoestK280vapy4Tj4r96Mh6rHW/JmyeNWPKkn17k1nT3T+lvZjONzUP1Zh9Ley/XDUOH/Zn/ABPCHBdUpe8pyjjhl+8fhqM/n/Ie9Ppc2WSzuah+rP8AiR/S9sv8zUf1Z/xPA3z6/efMufN+8fhqH5/yXvr6X9luuGo/qz/iQ+mHZVfI1H9Wf8TwCXn5s+GvHl2mfwtD8/5DoF9MOyuPxeo/qz/iQ+mLZX+z1L9Wf8Tn3Da6yfNiOLRn8+5LoD+mPZT+y1L9Vf8AEPpk2U/s9T/VX/E59kl1cT4lHL68eJL8JRn895DoR9M2yf8AZan+qv8AifL6Z9k/7HVP1X/M56lD86TPmS75e8z+Eoz+e8iXQz6adklzpap+q/5ny+mzZFfktU/VH/E54nDhjL95SlDvb8zP4OjP55yHRL6cNkM/itW/VP8AMj+nHZBPjS1b9U/zOc5QS7fefO7h837x+DozHrXIdG1unLZL1M/U0NTnU3W4RdvuqT6lnPDxPAta1W91rVrnVdSqupd3U3OpLPCPZFfmpcEu4xjiu8mL6sYwbcOCuOdw5uXz8vJrqz7qcfApSgn1cSq2kuSPnK5tHTEK6NwtpLHBIKo4lWe6lmTSXazHXddy4U8xXW+tmYjTZHdWuLqMU485fcdO+jvt5ParZyrpmo1vWappbjTnJ/GrUmvYm+18Gn3rPWcoT8De/Rx1WppnS/YUVJqlqFKpa1Fng/Z34/bD7Tl5WOL02t/S8s4csRHiXY6JPmPI+imewAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAIAkAAAAAAAAAAAAAAAAEEgAABHmSgAAAAAAAAAAIAEggkCASAAAAAAAAAABAEggIA+QD4k5ADrA7wLXVr2lp9jVuqqclBcIx5zk3hRXe20vM1NTuvhClKEa+rXrwo59iCXVnqpwzx7W+1mU2hqO41e2tOcLem7ma/Pb3Yf7b8kfWylsqjr6pNZlWk6VF/NpQbSx4vMvNdhOO0babfqtpkNH02FhSk3J1rirh160l7U3+6K6lyRkACDbEaCzv61RyjaW8sVqq4y/s4dcv4d/mV7uvC3oTq1M7sV1c2+pLvZR0+jOEZVq+PX1nvT/ADeyK7l/HtArW9KnQoxpU1uwisItGlf3Hba0ZeVSa/cvv8D7vqk6lWNlQlu1Ki3pyXyIdb8XyXv6i6o04UqUadOKjCCxFLqQH0lgYWSQGUE9QAEcjlj0yekGvcahT6PtMryja0VCtqrg/wAZJ4lTovuSxNrrzHsOo7urChbzrVHuwpxcpPsSWWfm7tFrdbaLaPU9duW3V1G7qXLy+SlJuK8EsLyOvh44tfc/Ct9SzTTHqPlZSpLfb3ePaRutdxXS4lRQTLfTzfX9rjZam6mu2keqNRzfkmz0Wl7Sw+t9ZpexVu3qs6vVTot+baRukGlgRHZx57bsrwguvBs2wm09xsltBR1Kk827ap3dNflKXX5rmv8AM1ZT44RFWSkkurK6zXkrFo004cl8WSLVl2pZ3FK7taVzQmqlKrBThJcpJrKaKuFzNE6CNRnqHRzYQqSbnaOdq2+yEsR/u7pvZRXr020+hcfJ7mOL/bWek/Qo7RbDarpaipValCUqPdUj7UP7yRx7b7soqWMZ4+HcdzzWYNdxxPtBbKx2j1azXBW9/XpLuSqSS+w7uFbzVQevYo7XUU8cmX+n3Lp1N2T9mfDwfaYxPrTPpTec5LCY28xMbbJLPXwKcpLPEtLK6dWhiTzOCw89a6mTUm314NGu6HTL7nU5ooO4q06katGpKnUg1KE0+MZJ5TXemfM5Z4lKos8ORLW0qfpnbq7oy2lp7UbI2WqJx9fKPq7iK+TVjwkvfxXc0bPnPWc5ej5tB/JG0c9FuZqFpqTSp55Rrrl+kuHikdHLGCpz4+i73vpvJ9/DEocU44ZyR6TexMNA29/ly0pKNjrWarwuEbiP4xfWTUvFyOuTTOmPZKO2Owd7pcIr4bTXwiyk/k1oJuK8Hxi+6TJcfJ7d4T53H97FMfLiq2qytq0KtGcqdSHGM4PDXmbdo21sWlR1WK7PhFOPFfSj+9e406pJb8oypyptSw4SWGmuDT78opzk+SLyYraHiMmDqnUvV1WoVKUa1CrCrSkvZlB5T8ylOfY+Z5fZalfabWdS0rOKfxoS4wl4r9/M9j6Ldmbzb/RKuoaff2VrUoVfU17eopOUHhNPh1NcvPsOXJHR3lCnBvedUYni+SPnD7Gz0mHQ1tDvZ/lTTMfRn/Arrod13HHU9Nf1Zmn3qOj8o5U+KvL91tdwcVFcD0+XQ5rv/wA0039GZ8S6Gtef/Kum+6Zn36fZ+T8r+LzBvCyvApTXHKZ6j/QxtB/800z3VP4HzLoW2gfLVdL/AEahmM+P7Z/KOV/F5c8rGT5lJcng9Rl0KbQvlqulv6tQpy6Edo3y1bS/0ahn8Rj+2fynlfxeYuSa4NHxLnzPT30IbS9WraV+jU/gR/QftJ/830r9Gp/AzHJx/ZHpPK/i8weOR8tcFwyepf0I7Sder6V+jU/gfcehLaF8Hqul/o1CX4nH9s/lXJ/i8mlFPOCnKGM4Z69/Qfr7/wCVtM/QqEPoN19/8r6Z5wqGPxOP7Tj0rlfxePNNLkfDzza9x65U6CdpG/Z1jSvONQovoG2nef8AdvSln8yoZ/E4/s/K+T81eU7/AHcCnVuKdNe005Y4RzxMr0k7P3GxmsU9HuNQtby7lR9bVdCMkqSbxFPPW8N//wCmpOTk8ybcu1nRS0WjcOa2CaW6bLmvWnV+M+9JckUZNdqKTn2+8b6Zs2x0a8JkuODYuh+hOp0ubNerftK/jJ47FCTf2Guppy5nqvov6HPUukiWquD9RpVtKbljh6yotyK/R335GjkTFccy7eDWb56xDqqK4I+iOGSSie2AAAAAAAgCQQABJHAkAOsAAAAAAAAAACABIIJAAAAAAAAAAAAAQBIIAEkEgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAMgAAABHkSAAAAAAAAAAAAAAAB5AAAAAAAAAAAAAI8iQAAAAAAAAADBRu68Le2q16nCFODnLwSyBq19cS9Zq9/Hi4ydOn37kd1L9NyNm0y2jZ6db2seVKlGHjhYNTsqcqtppVpUj+Eua8KlVe+tL7VjzN0RmWunnaSM46ieBa6jXnSoJUknWqy3KSfzn1+CWX5GGxQX9d1BvGbe2lw7J1P/APlfa+4uruvC2t5VZpvHKK5yfUl3tk2lCFtbQow4qK5vm31t97Zaw/rmoOpzoWzxHslU635Lh4t9gYVtPoSpU5VKuHXqveqvv6ku5ci6ADIAAAAAxO2Uak9lNWhSz6yVlWUcdvq5YPzYtOFvTx8yP3H6dV4xnTcJLMXwa7UfnRt9s9V2U201jZ+tBr4Fdzp08r41Nvepy84OLO7hWjcwp/VazqJYyl2YLiCy8lFLHeVocOrzLPagmG0bEwSo3dZrm4wz73+82DHBcDEbJrc0ne+fVk/dw/cZinlpfuJfCuyT+qSKecYYmt5JYxxRWhFbyWCqoxjHfnhRjxbfYa5a4ncxEPf/AEbISWxN05Z3ZahUa/Rh+89SNP6HtKnpOwGm0asNytWg7ionzTm95J+CaXkbgUmWd3l9C4VJpgrEolyZxl0izj/P/aPdxj+U6/7bOy68406M5yaSim2+44P1XVJanrWoajnhd3lauvCU5NfYzq4MfqmVV67O8cQrupFM+XW4lop8Ez6UuPAs5eX6dL23uZUa0ZpvC5rtRl3OMoqSbaayn2mvpvuz1GS0uq2nRm+fGP8AAhaqNl5jxCcVxwGuGeRTly5kYavKvRryo1IVaM5U6tOanTmucZJ5TXg0jrHo92hpbT7KWeqxaVWcNyvBfIqx4SXv5dzRyJOe7jjwPS/R02s/k7aqts/c1VG21L2qOXwVeK/2or3xRz8rF1V3Hwu/ReV7OXot4l0h5CSzFoJ8ESVT2flx76TOx72b29qana0nHT9Z3riDS4Qrr8ZHzypfWfYeXx4xznKO1+m/Y5bZbBXlhRhvX9v/AFmyfX62KeI/WWY+ZxlCgowSaafY1xXcy44mXrrr6eT9Uwe1k38StnTb6keh+j1tVPZTpFoU7mpu6bqija3WX7MJN/g5vwk8eEmaO4rd5EwSw12rDwdWTFF66lwYc84rxaHW3pC7T7YbFbJUNptmI2ta0tayjqVKrQc5KnLCjUjxXKWE1+dnhg8g07p7211Cj661vNIqR618DalHxW9wPaOiTXLPpE6KVZ6uo3VVUZ6dqdOXy3u7rb+lFqXn3HEe2uj6hsHt3qmztSrUhc6dcOFOquDq0nxpz71KLT8clfxcdJmaWju9Bzb5ZpXJit2l71Lpw2+i/wDfGlfqb/xnzPpx2+T4XGlfqT/xHjGj7UUbhqlqUY0JvlWivwb8V8n7vA2aNLejGUeKkspp8Gjp9jHHwpL8zk0nvaW/Ppz6QFyr6T+pS/xlN9O3SF/b6T+pS/xmiK3y848iHbp8MD2Mf01/mPI/k3l9PHSH/b6R+pS/xHxLp56RE+FfSP1J/wCI0Z2i57uT5drFcGn4j2MX0lHqOf8Ak3eXT30jJfj9I/Un/iPh9PvSRxxcaPn/ALi/8Ro7tFjlgpytI8jPsYvpn8xzfyb0+nzpIb/31o6//Bf+Mj+n3pIT/wB9aP8AqL/xGhytFjOSjUtmuozGDF9Ec/PP/p6Iun/pF67rSP1F/wCM+ZdP/SNnhd6T+oP/ABHnPwdyluxjl9hUp2lGDUquJ4+SuC/zHsYvpP8AH5/5PSrXpw6TblOcL/SYU1xlN2HBf3jo/ZjUtUsuju21na+4pK9haO6vZQperjTWHLG71NRwn3pnM/QpoENqdvLGwnTUrK0/rd0kvZ3INbsX9KTivDJ656UG1EbDZq32YtquLjUnv10nxjQg03+lLC8FI4s1KzkilYWvF5GSMFsuSf8A4542s1a52k2jv9evMqte1nU3PmQ5Qh5RSXkYrlwZWXBJPmfL4vJZ1rERqHnbZZtaZn5UJPtT5lObSWeJVrOMX7Tx48jKbK7IbS7W3Co6BpFxdxziVdx3KEPpVHw8ll9xC9op5luxUtl7VhhrWNxdXtGytKFS4ua9RU6NKmsyqTfKKR2h0K7FR2I2NpWVdxlqVzL199NcV6xr4qfzYrCXm+swPQx0P6fsVu6rqdSnqOvSjj1yj+Dtk+caSf2yfF9y4Hq3BIq+TyPc/THh6X07g+z+u3lIAORbAAAADrAEEgCCQAAAAAAAAAAAAAACB5EgAAAAAAAAAAAAAAEEgB5AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACORh9r6mNFqUIvErqcLdeE5JP7MmZZr200lV1TT7bi9z1lxLyjuR+2f2GYjujfw+NKirjaNSxmFpbNrHLeqSwvsg/ebIuRg9k6akr67/tblwi/zYJR+9S95nBbyxTwhllbL4TfTunxp0s0qXj8qXv4eT7SpqdadKhu0uFapJU6f0n1+Sy/IrW9KFChClD4sIpIwkpX9aVKkoUseuqy3Kfi+vwSy/IqWlCFvbwpU87sVzfNvrb7y2tf6zf1blv8AB0s0qXe/lP38PJl8CAABkAAAAAQznT0vujqvqVlT270W3lVubKl6rU6VNZlOgnmNVLrcMvP5r/NOi0j5qQjUi4yWU1hp8ieO80tuGnPijLSay/NqHf18fFFTw595656TmwuzmyW1tlU0Cq6E9ThUr1dOjH2KCTS34P5MZN4UeXB4wuB5XRt5TaiubwkXOO/XXcPKcjH7NprLcNDoqlpdvHk9xN+fEyMMJIp0qKpQUG+SS9xVptLBu8KW87mVSm3k2/ov0K12h2wstPvKlONtHNerCT41lDD3F25eM9yZqCnHzK9rqdzpl3QvrGrKjc29RVKU18mS+9dq61lGrJEzWYhPi2rTNW1o7OzYRUIKMUkkuCPo1ro72mobV7MW2qU92FZrcuKSf4qqvjR8OtdzRsmOBRWiYtqX0TFet6RavhgOkSlqlzsPrVtoai9Sq2VWFspPGZuLS49vZ34OFaEnCXqnTlTcPZcZLDi08NNdTP0InBSjg5k9Jjo4emajLbTR6DVndTS1GnFcKVV8FV8JPg/zsPrZ2cPLFZ1Kp9W41sleuPh45GXDgVItrmijCMoxWXl9pVjPHUy41vu8rb6XMH3o+1VcZb0Xh88lrvr/ADJ9Yt3kJiENM7RufXUlNcJfK7mRObZh7O7VKth8IS4P+JkKj78GqY1KE10VJ56ylSuKtpd0by2qSpV6ElUpTXOM4tNP3ktrsx3nxKMZdSaMTG4SrM1mJh2T0cbS0NrNkbHWqOFKtTxWgn+LqR4Tj5NPywbIc2+jHtOtN2iutmbmoo22oP1ttl8I14rjFfSivfDvOkimzU6LzD3fAz+9hiUNJo5M9IzZNbN7bzv7alu2Gr71xTwuEKufwsfe1L6z7DrQ0jps2V/nbsHeWdGmpX9t/WbLt9bFP2frLMfMlx8nt3iUfUOPGfDMfMONJJcUuRSct1PGCZ1U5bu649z4Pw8T4msrJd9W/Dx3RqdS9G9HTbF7NbfU7G5q7mn6xu2tXL9mNXP4KXvbh9ZdhsPpubDu50uw6QrCjmvY7tpqW6udCUvYm/ozePCfceH3G9Gm5QcoyjxTi+Kfau87F6NtXsuk7omVvq8IXE7i3np+qUvz93dk+7eTUl9JHBnicd4vC84GSMlJxS4Qoe1HEsJ5wZvQ9WvdMe5Taq22cyoT5fVfyWNqtmb7ZPavU9nL7edawuHSc5LHrIc4TXdKLi/MoUaa4LHHtO2urxtw566may9D0a9tNUpp28t2qlmVGfCcf4rvRkfgnW4nndo3SqQqU3KE48Yyi8NPuZuWi7SRlGNHU8Z5KvFftJfejVasx4VmSkxO4ZCVulyWCnKgl1YMolCrBTpyjKDWVJPKZDt21zZDq01RMsS6Cl1EStufAybt2fM1CjSnXqzhSo0/j1JyxCPi+3u5szFmYiZ8MZ8EzyiW2oUqFovw7frH+Sj8bz+avEo6ptLGf4LSVKEeTuJrE39BfJXe+PgYZTk17WW3xbfFtm+sT8tsRMeV1VrOfBKMIfMjy/zKU293OCnGRsPR3s5U2u2ysNCpuXq7ipmvJPjCjHjUfu4LvaMXmKxtuxVte8Vj5e/+jJs2tE2EqbQ3kVTudWfr96XDct45VNcep8ZfWPCOkraqe1u3Wo6upN20p+qtE+qjDhH38ZeMjoP0idep7LdHMdF03doXGpL4FQhDh6uil7bXco4j9ZHKkIJOPDG6cnGibWm8rT1C3t1jDC4k+I3steJTnJYGn0LrUtVtNMsI793d14UKMcc5yeF5dp3TfphT0xzeYiHvHoxbD2epU7zanWLKjc0Yy+DWMK1NSjlfHqJPhnOIp90joalQpUoRhSpxhCKwoxWEvIxexmh2uzey+n6HacaVpRjT3uub+VJ97eX5mYKLNknJaZe34nGrhxxGu4kGC21G9trC1lc3VVU6a4Lrcn1JLm2+xGp1rkFK1r07m2pXFJt06sFOLa6msoqgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADIAAcCAJAAAAAAAAAAAAAAAAA6ggAAAAAAAAAGRwAAgASGQSAAAAAAAAgABHXkCQQ/EASyMgeYDK7TVtQrReuahdTxuWtKFL3J1JftR9xtD4cTSY/1nTG+OdSuu35NSp/gRKrXf6bPs3byttDtKc17bpqc/pS9p/a2ZHPEiKwuHI+LmrChb1K8/iwi5PyI+U4jstI/1jVJz5wto7se+cuL9yx72VdSrzoWrdLjVm1CmvznwXu5+ROnUZUbSCqfjZZnU+k+L+8pca+q8eMLaGfry/gv2hoXNpRhb21OhDO7BYy+b7yq2OCAZBkjrHmBIIyEAJI8AgJKderTo0p1ak4whBOUpSeEkubZ99Z5B6VW2C2b6Op6Xb1HC+1uTtIbrw40sZqy/R9nxmiVKza0Q1ZskY6TaXOPS7tRLbHpB1LXIycrZyVGyT6reHCLx1bzzP6xgNJTqanbQecSqIs4vjnGOGEZbZqKer0pL5ClL7C9x16YiIeL5OSbzNpbZNttvtZTb4CU97/I+G+XE2yqYhKbfWfXFrDkj4z7ipHjHiyGk5bl0RbX1Nk9p4/CKmNLvHGld55QfyanlnD7n3I6jozVWEakZKUWspp8zjGlTi3xXimj3zoG2v8Ahtj/ADZ1Ctm6tIZtZyfGpRXye9x5eGOxldy8P/qHpPRef39m0vVi11bT7PVdMuNNv6ELi1uaUqVWnJcJRaw0XSBXx2enmImNS4k6VdlLvYfauro9xv1LWS9bZ3Evy1HPDP50eT9/Jo1b1meLOy+mjYa2262TqWS3KWpWzdawryXxKmPiv82S4PyfNI4quqd3Y31ewv6E7e6t6jpV6U/jQmnhxLriZ/cpqfLyPqPD9nJuPErp1CPWSzj9xQUkuOT6Ul2nVpWaVlMyVjX9ZS3Jv2oLh3ow+8u1lSlWdKanF8U8+Jia7YtXbOcU+LPmUspnzGpCpSU4yynx/wAj4nKKWcmvTXpd2F7c2N3RvLSo6dxQqRq0pr5M4vKfvOydgdorfanZOw1q3cV6+n+Fgn+LqLhOPlJM4mq1fYfF4wevei9tf/J+v3Gyt5Wxb6ivX2m8+Ea6XtR+tFZ8Y95ycvF1V3C59H5XtZOifEumG8FOolJPjkmMk45WcMlpFVD1vmHHvpC7IrZjpArXdtRcbDVt66o4XCNTP4WC82pfW7jz2KTOxunTZBbW7CXVC3pb+oWf9as8Li5xTzD60cx8WjjyjHKznPHh3IueLl66PI+pcecOXfxKnOnlcEj0/wBG/at7N7c09Nuqu5p+sbtvPL4Rrfk5eeXH6y7Dzfda4pce0+op78GnKEovKlF8U+prvN+XH11mHHx804skWh7d6ZexsZ2Vjt1ZUfbobtnqDiudOT/BzfhJuOfz12HNVOSeDuTYfUbDpQ6JZWuqxjUldW07HUYLnGoo4k12N5U12ZRxVtJod9sztHf6BqKfwrT7h0KjxjfS4xmu6UWpLuZz8PJ2nHPmFv6hii0Rlr4l9UOovqMd7+JY20stF9QlFdZ1zCjtDMaVdXVk06M96m+dOXGL/gbXp19QvEoxzCr/AGcufl2mn22HFF5W9iyuKi4OFOUk/CJptjizTaImWT2j2i0zRnKjOSvL5cFbUpr2H/pJ8o/RWZeHM8+1XVtQ1e4Va+qpxi/wdGC3adP6Mf3vi+tssKUIxhFY6itHHA2UxRV1VrFY1CrTk4/KzllzTqrGHz7SzTWcH3l9uCfgmsSyFOcW8HSfombLK20a+2uuaW7UvZO2tG+qjB+1JfSmv7iOaNDsL3WdodO0bT03c31eNvT/ADXJ8ZPuSy33JnZPSFq9p0cdEkqOmYpToW8LDT11+sa3YvvaScn4M4eXeZ1SPlZ+mYq0mc1vh4X0/bS/zj6QbiNCanZ6d/VKDT4Np5qS85cPCKPOqifYfcqm/wAc5eObecspuT5Y49R1YqRWsQqeTmtlyTeVvVUkuvB7B6KWyf8AKW1V1tPdUs2+lx9Vb5XB15ri/qwf988mjGc5RhCEqk5NRjFLjKT5LxbO1OijZanshsLp2juMfhEYetupL5VafGb8m8LuSObmZOmuln6RgnLk6p8Q2iHsxxxPpPiRjHWM4KiIesfWcmE1mGNfsa0/bi6NSME1ncmmnld7jn3GZ3jDbUycLa2uVw9Rcwk33Se4/sl9hKEb+FbZduOmytm+NtWnR8k8x/utGVMHoNVx1a8ovgqsIV148YS/Zj7zOCfJSdwkAGEgAAAAAAIygJHIhtDqAAe8cAJAAAAAAAAAyMgABwAAAAAAAAAAAAAAAAAAgkAAAAAAAAAAAAGQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQSBBIAAAAAAAAAAAAAABBPkAAAAAAAAAAAAAAACASAAAAAAAAAAAAAAAQSAA8gAIJAGP2hryttEvKsH7apSjD6TWI/a0YXTqEVq+m2ceMLanOp+jFQX7T9xkNqJb8LO1/trmLl9GCc/viijs5H1usX9x1U4U6K8eM5ftRJR2hqt3s2BFjqX4atbWi5VJ78/ox4/fuovnyLG1/C6nc1uqmlRj98vvS8iMNkruclCnKcniMU233FvpUGrX1s01UrydWXdnkvJYXkfOqtyoQtVzrzVP6vOX2Jl5HguWAykgZJAgkeQAAAAAHwAiTwjiL0idq3tf0j3tS3nv6fpjdjaYeVLdf4Sa8Z5WeyMTqDp52vexvRvqOo0Kijf3C+CWPb66plJr6K3pfVOIbdYpqKy0uHE7+Fj3PVKk9W5GoikPuMWuRl9msq7qVOHs08Z8WY5x4cjMbOUsQrVH1yUSziO7zmW36WZjLi/E+4rOFjifMUuwqw9lLCJSr+z6jT4rK6itFY5xwek7F7I/zq6J72dpSj/KllqFWpbPlv+xDNPPZJfakzzd+shNwqQcZJ4lGSw01wafY0aaZYvMw6M/GvirW8+JVIcOovdKvrjTL6hfWVV0bm3mqlOa6pLt7U+TXWiyjJdhG828JcOslaImNS5sd7VtFodXdH+1FrtXs7R1KjuwrL2LijnjSqLmvDrXc0bDzOUujna652Q1+N0nOdhWxC8orrh1SS+dH7eKOodLvra/sqN5a1oV6FaCnTqReVKL5MpuRinHb+nu/TefXk0iJ8rtxT5nhfpI9E9XaKnPazZq33tYoU0rq2guN5Tjycf8ASRXL5y4dSPdOYccriasWSaTuHbyMFc1emX52RbjJxeVx6+DXiVYvjg6O9Ifoid58I2v2VtHK8w56hY0o8a666kF8/tXyufPnzZv8E3jjx4F9gzVy128jyuJbBfUq2/hciHV48iipPHMb/Wb3LpktOusS9TLhGT9nuZd1JPPDjxMBvtdz5mWtLj19Deb9tcJfxNVoRtXXdVk+1Hzb3V1YXlvfWNR0bq3rQq0aifGM4vKfvR81JFGcnngY1Expik9Ntw7m6PNpLXazZDT9ctd1K4pL1kE/xdRcJw8pJo2A5d9FHbH+TtpbnY+8q7tvqKdez3nwVeK9uK+lBZ+ozqFcuZR58ft3mHtuFyPfxRKJxyuByL087JPZjbyvVt6W5p+p711btLhGWfwkF4SeUuyS7Drx8jQ+nLZL+dmwt1St6W9qFn/WrPC4ucU8w+tHMfFrsJcbJ7dmv1Hje9inXmHICiml1hLhyKmU4Jx+K1lZPiXDBeRMT3eMmJidS9T9GvataJtvLSLmpuWes4pLL4RuIp7j+ssx8d0yPpj7FxU7DbuypcPZstR3V1N/gqj824N98ew8XjVrUpqpQnKlWhJThOLw4STTUl3p8Tr3ZS+0/pW6I/V6lCEo39tK0v6cfydZLEmuxp4lHxRX56ziyRkh6DgZIz4Jw28uK6O7vezyLyhnPBMnWdJvdD2h1DRtQji6s7idGpw4Np/GXc1hruZNDHYd0TFo3Cmy1mtpiWStZPcXAudVq+r0G9knj8BJe/h+8tbf4iwj52nnu7N3HbJwgvOSMR5c3mzVIyS8D7TWMot4N4w+PeVqb7UbXZMKvPqRE5uEN7CeOoLHgVaNpXvalO0tacqtxXnGlRpxXGc5NKMV4tpEbdo2VrNpiHufofbLO+12/wBsr2lmjYp2dk2uDqyWak14Raj9ZlH0otq46ztlS2ftKilaaPH8Lh8JXE1l/oxwvFyPaLG0sOiboYjB7s3pdk5VGvy9xLi/0qkseZx7WrXN1dV7y7qurc3FWVWtN85Tk8yfvbK7DHu5JvK05lvw+GMUeZSpuCHrMvLKc5Y4pFvXqqPPhwznwO/elNFdvV/Ru2b/AJx9IsbqvTcrHR1G6qZXB1eVKPvzL6h11wPMvRw2Tlsv0eW87qluajqcvht1lcY7y9iH1Y44drZ6aUvIydd5ew9PwRhwxH2e8pXVaFvQnWmm4wWXhZb7ku0q8S0ufw97St18Wn+GqeXxV7+P1TQ7pXFCpGtRhWjxjOKkvBlprts7vSLq3ivanSko/Sxw+0q6d7Eatv8A2VVpeD9pfY8eRcy5Bie8NV0+5T1HS73lG4hKk/rxU19sftNrRpFTNrptVR56fdOX1YVN5f3GbrB5imuRKzXj+n0ACLaAAACOJPkBAJHHsAgkAAEAAAAAAAAABAJAEEoAAAAAAAAAAAR5ASQBx7ABIAAAAAAAAAAAAAABBIAAAAAAAAAAAAAAwACAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADA1/WKnrdfo0uqhaym/Gckl9kZFxsnBPT6tzj/fFxUqeW9ux+yKMTfXO5d6zfc/VNU4+FOGcfpSZsWi2/wTSLS1xxpUYRfikskp8NVY3ba5rSUKcpyeIxWX5FvpUHCypynnfqZqS8ZPP7xqz/AKlKkudWUaX6TSf2ZLpcFhLBFtWjaq6slzVCln60v8l9pdos9L9v4TX4/hK0kvCPsr7i9DEAADIAAAAAEN4XIeBrvSTtNbbIbFantDcpNWlBypwb/GVHwhDzk0vMzEbnSF7RSs2lzT6V21v8t7c0dnrWe/Z6LHFTHKVxNJy/RjhdzcjyChFQilEm4ubi91C5vL2o6tzXrSrVqjfGc5Pek/e2VYQT59ZeYqRSsQ8dyc05Mk2lUhxXIz2jU3CwUl8ubf7jCxjCK3pPGFk3PUtFudDr0tMvFu16dvRqTjjDg6kFPdfet7DNsTG9OLNE9O/haRUmuwr04PC4nxHCZUU1FYZOY7OKZdB+jMktjL+Oc/7o1P2IGuekDsU7O6ltXptL+rVpJX8IL8XN8FV8Hyffh9bM36MdXe2S1Fdmoz/YieqXltRvbSrbXNKFWjVi4VITWVKLWGmikvknFlmXs8fGryuFFJccL7j7ikbL0lbG3OyG0Do01KpplxmdnUfUuunJ/Oj9qw+01qHFcOss6Xi8bh5DPgthvNLJdPea5cOJ6H0R7dVNmryOl6pVctIrS9mT4u2k/lfQfWurn2nn64M+4tZy1xI5McXjUnG5N+Pki9XYtGpTrUo1KcozhJJxlF5TXaj6PAOibpFloVWnous1W9Lb3aNaXH4M31P8z9nw5e+0akKtONSElKMllNPKaKjLinHOnveFzacqm6+UyWUc8+kN0LyvXcbW7G2n9e41L7T6S/3x21Ka/tO2PyvHn0P1hrJjFltjncN2fj1zV6bPzkbcZOL5ru6yVJZ5nT3pDdDK1mNxtZsjapaqk53tlTWFeLrnBdVX9rx58sTqbkt2W9GSeGmsNPrTXNMvMGeuWv8AbyvK4lsFtfC4c1w4n3a3Pqayln2Xwku4s3NvkQ5S7Ebphya22KTTw08p8u8ozaRZ6XXcoujLnHjHw7C7fPiiMNU10q6Zd3Om6ta6rZVPVXdpUjVoT+bOLyn4dXg2d0bAbR2u1eyWn67aYUbqknOGfxdRcJwfepJryOD3NLqR7b6Jm2isNoLvY69rKNvqGbmy3nwjWS9uC+lFKX1X2nFzcXVXqj4W/pHJ9vJ0T4l1CfMlmIi01nPAniVG3qfMOQunvZZ7Mbe1nb09zT9SzdW2Fwi2/wAJBeEnnwkjQop+XgdddO2yMtqNh6/wWjv6jYP4TapLjJpe1BfSjlY7cHJKXLDfHivAueJk66aeP9U4/s5d/EvlRS5I9X9GbauGhbVS2fuam7Z6thQzyjcJez+kvZ8VE8ql2ZKbqVrdxrW85U61OUZU5xfGMk8prvTWTflpGSkw5OLyJw5YtD2b0utlI215Zba2dLEa+LO/cV8r8lN+WY5+ieF2rUlzOwtIuLDpV6IpUL3dXw+1lb3SivxNePBtdjUkpLyOQa2n3mj6td6TqFP1d3Z15UK8erfi8ZXc+a7mjm41+3RPwsPUscTPu18SvrbKS+wtts57mhUofPrx+xNl7bRct3gYrb6TVrY0sc5zl7kl+864jupqd7w1ylwgnnzKiZRg8R5cj6jLhjHE3adi430+PYey+iXsote26q6/d0lKy0SKlDPFSuZ5UP0Y70u57p4o5xWetrjhLi/A7c6IdAtejPogpS1VxoV40J6jqc31VJR3pL6qSj9U4eZk6a9MfKw9Pwxa/XPiGg+lhtWq13Y7IWtTMKKV3eJPnJ8KUH9sv0TwRpLgi/2m1m71/aTUNavU1WvazrNZ+Kn8WP1YpLyMa231k8GPopEOLmZpzZZtKKjfNI2fod2Vnth0j6ZptSG/ZUJfC71YyvVQae6/pS3Y+bNYqNbq49Z1N6LOyUdI2LqbQXFLF3rElODfONvHKprz9qXmiHJyRSjb6dgnNlj6h6/ShuxwlgqYzzCXYSUsf29hEajQyz0178atZ8JzqyUl2YeEvci8ZZWf4O+uqPU5Rqr6yw/ti/eZJSn6vVO6tS+2L/hL7C75otNR9h29f+zrRT8Jez+9F3ngYGt31CD1PUrWS9m4pwqY+lFwf7KMrs7cO40KyqzeZujFT+klh/amWOvJ09Xs6y/KUqlJ+KxJfdIqbJS3bW7tv7C7ml4SxNftE/hrr2szS5AhEkW0AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQAAAAAAAAAAAAAAAAAAAAAAACAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAIk8LLJLHX6/wXRL65T407ecl4qLDE+GrU18I0ukmuN9dKT71Uq7z/um7I1Owobt/o1phYpZm19Cnu/fJG2ErIY1pdPfvrSl1KUqj8lj75Ir3VRUbarVfKEHL3LJQh7erzf9nQS/Sbf+yNY42Eqa/KyjT/SkkRTVNNpulYUIP4ygs+L4suAuC4AMgAAAAAAQ+IE5OY/TG2rdxf6bsba1E4UGr2+SfOTyqUH4Lel5xOjtc1K00bRrvVb+qqVrZ0Z1q038mMU2/uOCdp9ZudpNpNQ1++T9ffXDrOLedxPhGPhGKivI7OHi677n4VXquf28fTHyxkIbrazkr0pYaeRuJnzOL3Gk9144PsLWezy8d+z0PoO2ajtb0gWFnVhv2Vni8u8rKcINbsX9KW6sdmTZenBtdJ+rpdXqV/8Aqgek+i3sotE2C/ly4p4vdZarcecaEeFNeazL6x5Z0318dK+tx680f9TA5MOTrzz/AEsOZgjFw4+5alhrtPiTeeGRGeeecn0odvWd9vDz0R3e/wDouxb2T1N//wAhL9iJ7ElwPIfReTWympr/APkJfsRPXzz/ACP3Je/9NjXHqxG1uz9htJolfS7+nmnUWYTS9qnNcpRfU1/lyZy9tVod/s3rFbS9Rp7tSHGE0vZqw6px7n9jyjrg1TpI2Nstr9GdtWao3dLMrW4S405dj7YvrX70jZx83tzqfDn9U9Ojk06q/wDKHLblxwmTFt9Z96xpmoaJq1bTNSoOhcUXiUW8prqkn1xfUynT48+KLWJiY3DxN8dqT0yrRWVwZ6Z0Ubf1NBlDSNZrSnpbeKNWXF23c+2H3eHLzSnhLgVoSTazjl1mrJji8als4vKvxrxarr6jVp1qUatKcZwkk4yTymu1H2l1ngPRd0hT2eq09J1apKekyeKdR5btm/vh3dXge929alXowq0pxnCaUoyi8pp8mmVWTFNJ7vecLm05VN18vtngfpGdCsNo6dfavZG1jT12K37q1hiMb5Lm11Kr2P5XJ9TPfSMIjjyWx23Dfmw1y16bPzWrwr0LidC4oVKNWnJwqU6kXGUZJ4cWnxTXJoRb6zrj0i+hmG1lGrtNszRhT2gpQ/DUItRjfQXU+pVEuCl18n1NciV1Wt61S3uKdSlWpScalOccShJPDjJPimutF3x80ZY/t5jlcS2C39KsKsqUlOLaaeUZONwqtJVIcE+a7GYF1OHNsrWNx6uq4Sb3Jvjx5PtOiYcdqbZWU8rhyFhfXOmahb6nYVnRu7SrGvRqL5M4vK8v3ZKTafcfMkpRafJ8DEx1RpGm6zuHfnRttLa7X7FaZtDaYjG7oKU6af4uouE4PwkmvI2M5c9ELbGOn6zdbF3lXdoX+bmw3nwVaK/CQX0opSX0Zdp1EuKKDPj9u8w9jw80ZcUSNJrHUcjdOuy/82Nu7hUKe5Yahm6tsL2Ytv8ACQXhJ5x2SR10eddPeyD2o2GrztaW/qWn5urTC4yaXtQ+tHK8cE+Nl9u7R6nxvew9vMOTcYR8y48PtKdCTnTTy2nxy+wq7nHjzLqJ28XManUvVvRn2r/kra6rs7dT3bTVVmll8I3EVw/SimvGMS/9KjY1WmsWe2dnTxSvN21vt1cFUS/Bzfik4t/mx7Tx+zrVbS8o3NrVlSuKM41aVRc4Ti04teaR13p1fTOlLoqca27GGoWzp1lHi6FePPHfGaTXgu04M8TiyRePC+4M15OC2GfMeHI1DMWkYLb6rvXFnTa4Royfvl/kbNqNjc6XrF1pd5FQurWrKlWS5KUXh47uteJpm3NXe1pQXyKEF78v953Y53O4Utcc1y9M/DD+saPqM+JRWF4nxOWN2O8lx5tm2Z07Yrvw9L9HLZeW2HSvp9KtT39P0vF/dtrg9x/g4P6U8cOtRke7eljtYrTQrLZK1q4r6hNVrpJ8VQg+Cf0pY/RZceipslT2S6LFruoRVG91pfD7iU+Dp0En6qL7lD2vGbOfukfaartftxqOvTcnRrVdy1i/k0I8ILuyvaffJlZH+bLv4hY5rfh+P0x5lim23lrmfLXPmfLqcePIb65o74UTL7F7PV9q9rNM2ftm/wCuVlGrOPyKS41JeUU/PB3bp9rQsrKjZ21ONKhQpxp04RXCMYrCS8jwL0Q9lnG01HbG6pYlXbs7LK+RF5qSXjJJfUZ0GuRUczL1319PV+lcf2sXVPmUgA5FqMs6v4PVaElyqU5QfisNfvLws9TzF21X5leOfB+z+8MS+9Tg6mn14x+NuNx8VxX2lajNVKMKkeUoprzJaTi01lMttJ/4PpQfFwzD9FtfuMfB8rLaaGLa2r9VG6ptvuk9x/tFDQZKnrl5Rzwq0KdVLvTlF/7Je7Swc9Cvd3jKNGU4+MfaX2oxVlUxtDY1o/FrUqtPxzuzX7LJx4a7drNnAQItoAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACCQAAAAAAAAAAAAAAAAAI8gAJAAAAAAAAAAAAAQABIAAAAAAAAAAAAACABJBI8gIMTte86DWpLnWlTpY+lOK+5mXMNtS80bGn8+8p5+rmX+yI8o28LbS4qrtGp9VK1k/Oc1/gNhfIwGzizq1/Pnu06MPsk/3mwdRm3linhZ2XtXt7PsnGHuin+8ah7Ve0p9tbL8ot/uQ0virmfzrif2cP3C446nax+bCpL7l+8wkvFyAAZCMkgAQCQIbIfVzPopXNSnQozrVZxhThFylJvCSXNthiZ1G3g3pf7X/ANnLPZC0qf1jU5euuknxVvB8E/pTwvCMjmem00uDTM50q7S1Ns9v9V13elK2q1fU2mfk0IcIeGeMvGTMBDgXXGp7dHkedm97JK+inyWOBsGwGzNTa7a3TtBo53LurmvNfk6EeNR+5NLvaNbpzxhHSvokbNqjpOobW3FPE7ubtbVtcqUH7bXjPh9QcjL0U7IcDj+7miJ8Pb7W0o2dlRtbanGlRowVOnCPBRilhJeRyb04Rb6Wdb4/Ko/6mB7r0b7YS2r6TNuaFCtKemaRK1sbVJ+zKcfXeumvr5jnrUEeKdONFf0q62038aj/qYHJwYmMk7WXrUR7ETHhpcVw5lWEuGGUcNdTZMXwLiXkoju6I9F552T1L/7hL9iJ68eP+i487Jany/4Ql+xA9gPPcj9yXv/AE3/AK9UPwJZINEu5pvSVsLYbXadl7tvqNFP4NcpcV+bLti+zq5o5x1bTLzR9RradqNtO2uaLxOEvskn1xfU0dgGpdJGxVjtbpmJbtDUKKfwa5S4xfzZdsX1rzR1YORNJ1PhS+p+l15FZvT/AJOY+GFjCJUlg+9esb/RdUraZqdvK2uaL9qL5NdTi+uL6mWsZNotImLRuHi8mO1LdNo7rlVG+Rv3RX0g3GztxDTNVnKrpM5ey+Lds31rth2rq5rsPOoy7HgqxeY4zjvIXxxeNS2cfk341+qrsK0uaF1b069vVhVpVIqUJwllST5NMq5Zzn0W7f1tmLiOm6jUnW0ecuC5ytm3xku2PavNdj6Gs7m3u7ancW9WFWlUipQnB5Uk+TTKrLinHL3fB51OVTceVXOVyPDfSK6E7fbGhW2l2ZpQtto6cc1KSajTv4r5MnyVTHBS6+T4Ya9yzjkQ1nnghjyWpO4deXFXLXps/M67t7qxu6tnfW1a1uqM3TrUasd2dOa5xknyZ8N5SyjtX0gehqz27s5a3okKNrtLQhiM37MLyK5U6j+d1Rn1cnw5cX6naXmm6hX0/ULaraXdvN061CrHdnTkucWi74/IjLH9vN8rizht/S+sq/rKOJcZx4P+J9zmv8zD0biVGsprivlJdhfSqp4cXwfFPtOiXBNe68stXvdJ1Oy1PTqvqr2zrwr0J9SnF5We7qfczv7YDaa02t2R03aCyaVO9oKbh1058pwffGSa8j88ZZb4HS3oSa1e1P5w7O1ZudpQVK8op/InNyjNLue5F+Oe04Odi3TqWnpmWa36Pt0v1HzUW9Fxzg+0g0U0RMPRTG405G6a9kYbLbdXXwemoWF/m7tkuUMvFSHlLiu6SNHlFZPbfSxlu6poi/7NcftUzxOWHgv+Lbqxxt4f1HHGPPMQpyxFZPafRO1i5/lfXdF3m7T1NO6jH5s8uLfmt39E8Wq8ermeseidHG2uuL/sEH/fMcuInHKfptpjPGmE9IRQo9MV/wCrgo+ttaFSWOuW61n3RXuPEdpqiqa/d8M7sow90Ue3+khGX9MNZYxmwof7R4DqVZVNWvZ81K4n9+CXGn9MJ56/57SnhnijauibYue3XSHpOguMnaTq+vvmvk28OM/De4R8ZI1Pe/BuXLgdeehzsnHTtirja+7oqN3rEtyg2uMbam2lj6Ut596UTPLy+3j/ALdXBwzky9/EMp6Tm0a2c6Pf5DsZKlc6rm2hGHDcoRX4RrsWMR+scpU/ipLC4e49E6c9qVtX0hX9xTmp2Vm3Z2mHlOEG96S+lLL8Ejz6quPBmvjY+mm5aOfmjJlmI8QSfDLZV0azutW1i00mwpqpd3tWNCjH8+Twm+5c33FtVfVk9q9EXZKWo7U3u191Szb6dH4PatrhKvNe3JfRg8fX7iWbJ0UmWviYPdyRDpHYrQbTZjZbTtBsl+AsqEaUX1ya+NJ97eW/EzIXIgpJnc7eyrWKxqEggGEklnrH/B1d8cxjv+5p/uLwo3sXO0rR+dTkvsAqQaayustdOe7K6p/Mry+3Ev3lWxk52VCfzqcX9hTteGoXce1wl/dx+4MK9emqtCdOS4Ti4vzRqGnTaoaJWb9qFanCXnCUH9rNyfI0p/g7FdXwe+/ZuP4Eq+Gu/mG7LkCFyJItoAAABAEgEAGASAAAAAAAAAAGQABAEggkCCQAAAAAAAAAABAD3kkEgB5AAAAAAAAAAAAAAZAE+QAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAHEAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADiAAAAAGD2lebzS4f6ecseFKf8TOGD2h/4T0zudZ/3BXyjfwnZlL4TqU0vy8Y+6nH+Jm+ow2zHLUH/ANsl+zEzJmfJTws9I/3lvfOqVH/fYlx1iC7LeX2yX8CNE46ZRfPO8/7zJ/5Y8Lf/AGjDK8AAZAAAAAA8m9KLa3+b3R3U0y2q7t/rUnaUsPjGnjNWX6Ps+MkesSeEcUekDtb/ADs6Tb2dGpv2GmN2Frh8Huv8JPznlZ7Io6ONj67uD1DP7WKfuXn9KnBU4YWFjkfe5hZfA+1jOT7jGPWs54F3p5Dqnaroel3mt67ZaNYx3ri9rwoUuHJyfxvBLi+5M6y6WtftOiDoOuKumNU61raxsNMi+cq8luxl3tcZv6LPPvRO2QV1rl9tdcU26Vina2mVwdWSzOS8ItL6zPP/AE4NsKmt7cWuyVrU3rHQ6e/XSfCV1USfH6MMLxnIrsk+7l6fiHoOHSMOHrt5l6F6DrktK1aU5OUqlvbTlJ85NyrZbMN04TT6VNb48pUf9TAzXoWR3NM1Jdlnaf8A9rNW6cZNdLWuLLxv0ev/AEMDZgiPflyc6/VxIj+2sPi+BDT5op0pZSy8PJcQ4rkWUvOQ6B9Fn/ilqf8A9wl+xA9iPIfRcS/mjqX/ANwl+xA9ePO8j9yXvvTv+vUKNtdW1xOtChXp1ZUanq6qjJNwnhPdfY8NPHeiszjHpJ6SNoOjz0m9przR6+/b1alurmyqSfqrheop8GuqXZJcV3rgRx45yTqHRmzRiiJl2cHxNW6Ntt9F252bo6zpFfKl7NejJr1lvUxxhNdTXua4o2lPJrtExOpTpeLxuGo9JGxGn7YaX6mslRvaSbtrqMcypvsfbF9a/ecza3ouo6BqtXTNUoSoXNJ8fmzj1Ti+uL7ffxOxzWtvtj9N2t0p213H1VzTy7a5ivapS/fF9a6/HDXTg5E4+0+FT6l6XXkR11/5OVVh8sH3vYWC/wBp9E1LZzVaum6nQ9VVg8qS+LUj1Ti+tP7OT4mL3uziWtbRaNw8blxWx26bK6e9JcFlcjeOjHb662Xu42V26lfSZy9qC4yoN85Q7u1ea489BjL/AP0rRmRyY4vGpS4/Ivx7xakuwNNvbbUbKleWVxTr0K0VKnUg8qSfWi54nNHRpt3d7J3it6ynX0ipLNWiuMqbfy4fvXX4nRulX9nqVjSvrG4p17etFSp1IPKkiozYZxy9zwOfTlU7eV15njnpE9DNn0h6fLVtH9VZbTW9PFGs/ZhdRXKlVx9kua8OB7JhZDXAhS80ncO3JjrkjVn5h6tYajomr3Gj6vZ1rO/tpblehVWJQl+9daa4NFTT6jeaMuK5xf7juHp86HNK6R9MV5bersdorWDVreYwqi/squOcex84viutPibXNK1TZ3WLjSNZsaljqFrPdq0ai4p9TXU4vmmuDRc8bkRljv5ed5nFnDP9PvEcP7zoL0JGv52bSpf9Bt/25nPUavrIqXLqfczoD0IZb22O0qT/AOT7f/WTJcz9mWn0/wDfh1hkELmChiXq3O/pZv8A3X0Ps+DV/wBqmeKvPBo9q9LJZ1jRP+7V/wBumeKtPCLziftw8R6pP+xZ9Y4cfuPXfRRgv5562/8AsFP9s8hbeD1/0UX/AO+Ws/8AcIf6wcqP8Us+l/8AZqw3pF0s9MNWTxhWFB/bM5onUcqtSo+c5t/adM+krUVLpQvqvzNKpv7JnMzi3TT7jPF/bh0Z+3IsyexWi321m1mlbM2Ckq+o3MaLmvydPnOf1YqT8jt/pg2gs+jnoj+A6Tu29WVCGmaZTi+MW47qkvoxTlntS7Tx/wBCDYrer6tt7e0niOdP09yXPGHVmvPdjnukYf0ntq3tB0kz0mhU3rHRIu3ik+Dryw6kvL2Y/VZzZJ97Lr4h37/D4Jn5l5zSqLc3UsJPEc9hDeXjJRUmlhPgPWYWXxRYR4eftEzO1aNGtcXNGhb03UrVZqFOmuc5N4il4s7o6LtlaGxuxGm6DRUXUoU964ml+MrS4zl+k3juwc1ei/sutoekNarXp79losFcNtcHXllU15e1Lxijrxciq5uXdumHpfR+P0065+UgA4V0AAARJZTXaSALTR3nTLbtVNL3Cnw1aqvnUYP7ZEaN/wAG0U+pNfayY4/lZ99BftMMLpml3/s22rxXyK9WX3SN16jTtRSzrcfz6j//AFRJUQyNwi8xTXWSUrR71rSl2wT+wqkU48AADIAAHEAAAAAAAAAAAAA4gABxAAAAAAAAAAAAAAAAAADiAAAAAAAAAAA4gAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAMFtFw1TTH31l/cM6YTaRL4bpkv9LUj76Uv4Ga+UL+H3sv8AF1Bf9sn+zEzD5GF2Zf4XUo9l0pe+nAzRifLNPCy0L/gqh4P72T/yy/8Au/8AtDRuGnwj82c17psmXDWI8OdB/ZJfxDK7AAZAAAAIbwBovTptf/M3o71DUKFVQ1Cuvgtin/bTTSf1VmX1Thqk5QkuMpdrlxz2vzPZvSx2qet7f0Nn7ao5Weiw3amOUriok5fox3V4uR496vL4Fvw8XTTf28v6nyPcydPxC4hPrwXdpCrWq0aNtB1K1WooU4Li5ybxGK8W0WD9lZbxhdh6v6L2y/8ALu3q1m5hvWWiwVfiuEq8sqmvJKUvFI35cnRWZcXHwe7kiHQNmtO6JOhiVa53ZR0mxdWs1w9fcS4tLvlUlheKOBteurvV7y91O/qOte3lWdxcT+dUm3KT97OkPTW223paZsHaVcqKWoagk+xtUYPz3pY/Nic3Uo+t3YL5bS97Ofh4+03n5WnqGbWqV8Q6y9EO1dvR1ik0vZtrSOPKoaN06QlHpZ1zs36P+pgel+i0krnaCK6oWy/1p5/07Jf0r63nnmj/AKmAwT/sTDk5Vf8ARrZotMuKUmljqLdPqx9pKqOKLKVBru6O9FuWdkdS/wDuMv2IHr+Txz0V23sfqTf/AMxl+xA9jXE85yf3Je89Nn/Xq+Jt44H5+elKp/097TSTw1O3aa/+hTP0GwcC+lJTS6ddpnj5Vv8A6imdPAjeRr9S7Y9sb0R7ba3sTrlLWtIqKSeIXlrN4p3EOe6+x9al1Puyjuno72z0TbfZ6lq+jV96Pxa1CfCpQn1wmup/Y1xR+c+nV40K633+Dlwkv3m87C7a61sLtFT1nRK6Uvi16Em/VXFP5k0vsfNff28riReOqvlWcPnzhv028P0DRDeDUOjHb/Rtvdn4anpknTrQxG6tKjXrLaePiy7U+alya80ttznkyktWazqXo8eSt46qtb2/2S03a3SHZ3sXCtDMre4ivboy7V2rtXJ+5rmDafQtU2a1mppeqUXCpHjCcfiVY/Pg+tfauTOw8ZRgNt9ktL2s0iVhqFNqS9qjXhwqUZfOi/vXJnTg5E451PhU+o+mRyY6q+XJybzho+4yeTK7W7N6lstrFTTdShjnKlVivYrR+dH966jD55NcUWtbReNw8dlxWxWmtlRyzwyzb+jXbu92S1FUarnX0utLNa3Ty4P58O/tXX4mkuSXVy7wpZfUZvii8alnBmvgvFqS7I0XVLLWNNoahp9xC4tq0d6nUg+DX7vDqL45V6ONudQ2N1F4UrjTK0s3Ftn+/Dsl9j9zXTWgaxp+uaXR1HTbmFxbVlmM4/amuprrXUU+bBOOXt+B6hXlU/tkOo836ceirSOkjQ/adOy1y1i3Y36jlx6/Vzx8am31dXNd/o/BkOKxjJpraazuHfekXrqX5q7V6XrGye0lxoG0FlKyvaLxKD4xmnynB/Ki+pr7+B7l6DUnLbTabj/yfQ/1kz3bpr6KtC6S9n1aXq+C6nbpysL+EczoS7H86D64+7DWTyP0Rtl9Y2P6T9rtA1+0+DX1vY0OK4wqw9ZPFSD64vqfinhplhfkxlxTE+VVThzizxMeHTyPoArYXDnn0r451jROv+rV/wBumeKyXJI9r9K541jReX+9a/7dM8WeHjgXnD/bh4f1X/sypy4Ya5nq/oo1H/PbWef/AAfD/WHk9RZXxcHqvooJ/wA9ta4f8nw/1hLlftyj6dOs8MF6Vl16rpA1R8v9yKSXDtcl+88D0izvdT1Wy0zTqfrr28rwtqEMZ3pze6vvR7P6YlZ0+kC/SXGWn2sffOX8CPQp2Qlre3d3tdeUt6y0OHqrdyXCV1UjjP1YN+c4mqmT28O1pXDOXkTv7dG6rPTuiHoWhb2W4/5Ls429snw9fcS4JtfnTbk/M4vulUqXM69arKrWqSc6s5PLnNvMpPvbye4elrtgtQ2ls9k7Srm303Fa6UXwdea9lP6MHn6/ceHylnr4GOJj1Xqn5Q9Uz9WSKV8QpTe7zLW4qqFOUs8EssvJYfcbl0EbHLbHpHsLS5pqpYWf9cvVjhKEGt2D+lLCx2ZOjJborMuLj09y8Vh076O2yb2V6NbGF1T3NQ1D+u3eVxUppbsX9GO6vFM9IKdKO7HGCoUV7dUzL2WKnRSKgAItgAAABD4JvsAtNG/4Npef3smP/C0v/oL9pjR1jTLfvpp+8U8PVqv5tGC97kGF2afqL9vXH+dP/UxNwNM1CWaGtz+dVqpfoRiZqhkbbYrFnQX+jj9xWPmkt2lGPYkj6MJx4AAGQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAPIeQAAAAAAAAAAAAQSABAAkAAAAAAAAAAAQMgSCCfIAAAAAAAAAAAABAEgAB1GF2oXDT6nzL2C/SjKP7zNGI2uWNFnV/salOr+jOLf2GY8o28KOzz3dT1GHzvVVPfHH+yZ18jX9Kap7Q1afL1lrn9GbX+0bB1C3linhaaXwp1ofNr1Ptk3+8iv7OqWz+dTqR/Zf7ibLhdXkP9Ipe+K/gxfYjc2dTsq7r84tfwMJLsBckAyAEASa30lbT2ux2xOq7R3WJRs6DlCDf4yo+EIfWk0vM2M5c9NXa1XFxp2xNrUzCg4318ov5TeKUH/el+ibcNOu8Q5uXljFimXh13e3V/e1r69qeturmrKtXm/lVJPMn72IotqbfYXNPi1wwy+jtGnjskzM7TPDg2/k8zsPoh2ftejnoije6w1bVPUz1LU5y5we7vbr+jFKOO1M5/wCgLZaO1XSPaULikp2Vg1eXOVwlGDW5F9uZ44diZ6l6Ze1y0/ZOx2OtKuLjV5+tuknxjbUmm0/pT3V3pSK/kz13jHC49PxRjxTms5a251292s2x1XaS+TjWv7h1VB/k4YShD6sVFeRY6TT9bqtpT+dcQXj7SJqwjKbcVxk8viXWzNPf2hsuxVN5+CTZ3RXpjSvy5Orcy6t9Fif9e2iWfkW331TQenieeljW02vyOP8AyYG7+is277aLD5U7b/8AsPP+nlv+l3W1x/IP/wDTA5OP/wBiW3PbfBrDT89Y8FzRFNNtdpWoxzJFnMdlF4l0R6K8WtjtSbX/ACjL9iB7HyPIfRba/mfqa7NRl+xA9fPO8n92XvPTo/16oOCPSlljpz2m+lb/AOopne/kcL+mVot7o/TBd6nUT+C6zb0q9CeOGacY05x8VuxfhJG7gWiMiPqVZtieQSl7TwZDS7j1kHRk96cFwz1x/wAjGPi+0qW2/CvCrDg4vP8AEupnbz80jTd9iNqdf2O2ipa1oVx6qtD2alOXGnXh105rri/enxR250T7f6Nt9oKvtPl6m7pJRvLKcs1KE329sXjhLk/HKXDFCEKtGNWnxjNZRmNktoNV2V1631nRrqVrdUfOFSPXCa+VF9nmsNHHyeLGWNx5beHz5w36Z8P0DRHUaP0R9JGk7f6N623xa6nbpK8spSzKk3ykvnQfVLyeHwN5KW1ZrOpeopkrkr1VYPbHZnTNqNInp+pUcp8adSPCdKXVKL6n9j5M5c262a1bZLWXp+o03KEsuhcRjinXj2rsa649XesM6/MPtbs5pe0+jVdL1W3VWjPjGSeJ05LlKL6pLt8uRv4/InFP9K71H02vJruPLjve4BS7+RnukLZHVdjdb+B3i9baVW3a3SjiNZdndJda81wNdUsvuLqlovG4eLy4bYrdNoVXJtczZujvbnU9jdU9bTlK50+q07m1zhS/Oj2T+/GH2rVN7wRHxuD+8zfHF41Jiy3w36qS7M2Z1zTtodIo6ppdzGvb1lwa5p9cWuprrRk8HJXR3tjqOxmr/CrXNazqtK6tXLCqLtj2TXU/J93UWzGvabtHpFHVNKuVWt6q8JQkucZLqkutFJyME4rf09r6d6hXk018srjgUfgtt8M+GeopfCfV+q9durf3M53c88Z44Kw5HOswMjLJYHPHpYP/AHX0TH/Rq/7dM8WbeUe0+lgv919E/wC7V/2qZ4o+Ze8P9uHhvVf+xL7k2eueinBfzz1l4x/UIf6w8iUknxXA9e9FOcZbZ6zh/wDwEOH/AIhnmftSelx/s1ecemjvrpNuKdOMpzqWloowXFyeamEvM9+6PtKs+hzoHhK/hGNaxs5Xt/jnVuZrLjnre84wXgjWdodkP53eltSqXVDf07RNOttQrtr2ZVE5qjF/W9rwgzH+mPtQ1R03Yuznl1X8Ovkn8lNqlB+Mk5Y/MRXRb3OmkPR3p7MXyS581PUrrVdUu9UvajqXV3WlXrSzznJ5ePuXkUFU8ijCO5FRxwSwfMm0i0rHTGnnr/qncrtTXDq7+4639FXZJaHsE9cuaW7e63NV+K4xoLhSj5rMvrnLXRls9X2x6QdM2bin6m4qKVxJc4UI+1UfuWF3tH6A2lCja2tK3oU406VKChCEVhRilhJdxwc3L2isLf0nj95ySqoAFa9AAAAAABSu5blrVn82Df2FUs9Xk46ZcY5uDivF8P3gVdPjuWNCOOVOK+wpW2XqN3Ls3I/Y3+8uoLdgorklgtdP41Lup86u17kl+4MLtmk1252V3w/H3k4rzrbv7jdZNRi2+S4mk2i9ZZ6UuuvdUpvzk6j+4lVC/wAN3XIkLkCLZAAAAAAAEASAQBIAAAAAAAAAAAgkAAPIAAPIAAAAAAAAAAAAAAAgkAAAAAAAAAAAAAAAgkAAAAAAAAAAAAAAAAAAAAAAAAAAAHEgkAAAAAAAAAAAAAAAAAQSAAAAAAAAAAAAAAAQSAAAAAAAAAAAAAAACOIJx3gAAALLXqLuNEvaK5zoTS8d14L0iSTi0+KfUGJ8NYsKylqul3S5V6U4Z+lBTX7JtCNMtZeosLFy4Ozuo05Z6lGo6b+xm5LkZshjWlP2dXqxz+MoxkvJtP70Rq3Cz9b10qkJ+6Sb+zJNyt3ULWp1S36fvW8v2Sre0vXWlajz34OPvRhPSsiS3sKjrWVGrnjKCb9xceYZARx7Q/ECy13UrTRtGvNVvqipWtnQnXrS7IxTb+4/P3arVrvafaLUNor5NV7+vKtKLedxPhGHhGKS8jpf0vNrfgGztlsja1cV9Tl6+6SfFW9Nrg/pTwvCMjluMcJRzw7y14GLUdUvPercjdopCacWkS5brw5YbWeZVUerkZzo/wBl6u1222l7P00/V3NXNxJfIoR9qpL3LC72jsyW6Y2qMVZveIh0j6LOzn8gdH8tevYqndazL4Q3LhuW8U1TXcmt6f1zmTpk2yntr0k6nr0Judk5/BrFdlvDKi12bz3p/WOmfSf2qp7HdFb0bTHG3u9WS0+1jDh6uil+EkuzEPZ8ZI41pxUcJRwkcXGrNrTklb8u/RSMUK2d7lzMlshTc9dhN8oU6kv7uP3mLWeCXA2Xo9oKpqVzNr4lu175L+B3T4U+adVdDeivUgtX2goNpTlQoSS62k6if3r3mo+kLYV7XpX1CtWptQuqNGrSb5SioKDfvi0WPR9tLHY/bax1dzcbXf8AUXiXXRnhSf1XiX1T2z0jNmYa7sStds4Kd3pSdZSjxc6DX4RcOxJS+r3nDFvazxM+Jd+GkcjhTFfMOZopLr4MqQlhZ5FDeSSSeeCIUm3hZRb7UGp29f8ARp2rWn7TXWzlzUUbfUfwltnqrRXFfWiv7vedIp5XB5OELS5u7C8o31nUdO5t6sa1Ga+TOLyn7ztHYHaC32o2VsNbtuEbmknOCfxJrhOPk00UfqGHpt1w9X6Nyuqntz8M/nqR5H6UuwUtuOjC7+B0fW6vpeb2x3V7U3Fe3TX0o5WO1RPXd1d5DinzRxUtNbRMLvJSL11L8uaMVJRa5SWUu4r04qPDB6X6SGxMNielG8o21D1em6kne2WF7MVJ+3TXZuzzw6k4nmmXnh7z0OK8XpEw8tmpNLzWWX0K9jTq/Baj9ib9h9kv8zLzSl1YNQzjDTxh5TXb2m1aVcK9s1VePWR9moux9vmT248tNd4ZLZvXtY2a1y31nRbuVreUPizSzGS64yXyovrT+9JnYnQ50l6Zt9pD3VG01e3ivhdk5Zcerfg/lQfU+rk+Jxoqaec8UXeiahqOhavb6tpF1Us722lvUqsOrtTXyovk0+DOXkcWMsbjy6eHz7YLanw/QBcVw5EnnPQ10n6ftzpqtrn1dnrlCGbm1UuE1y9ZTzzg/euT6m/RUlgpbUmk6l6vFlrlr1VYrajQNN2j0etpeq28a9vVXhKD6pRfVJdTOVukPYrU9i9Y+DXSlXsqsn8Fu0sRqL5suyaXNdfNd3X7T7TG7SaHpu0Gk19M1S2jcW1ZYlF80+pp9TXNNcjdx+ROKf6cPP8AT6cmu48uMY4cU88Hy4H0mlw4m2dJexOpbF6mqdTfuNNrSfwW6xz/ADJdk/sfNdaWoyeFzLzHeL13DxmbDfDfpvD7jLDNm6P9s9R2O1n4Zat1rSq0rq1bwqq7V2TXU/JmpueFwKdSTcfjYyL44vGpQxZbYbxeku09l9f03aPRaGq6XcKtb1o5XbF9cZLqknwaMnntZyJ0W7cXuxOtb/4SvpVw0ru3jxfZ6yC+cvtXDsOsNHv7TVNOoX9jcQuLevBTp1IPKkn1lHyOPOK39Pben8+vJp38rxPjzJZCRJzQsnPPpX/8L6J/3av+3TPE5Ps957T6WUt3WNE4/wDw1f8AbpnisXxyX/C/ah4f1T/sy+KreOw9Y9E2T/ntrSTx/ufD/WHlcoKX8T1n0UaGNtdaf/YIf6wlzI1ik9N754iHQlZaZpH8o65WjToOdJVLy4fzKcXjL7Es+9nCm2u0tbaza3VNfud6Mr6u504yfGFJezTj5RS88nRHpdbWS0bYqjsza1d281qbjVSfGNtDDn+k92Pg5HJ9JvOW2jh4OPX6pW/quabapC8qceotquViOG23gqb5f7P6NebRa9p2hWC/rN9cwo0+vdzzk+6Ky/IsLTERtUUrNrRV0P6G2xqttL1Dba8otVr+XwSybXKhB+1JfSmv7iOisIx+zmkWehaFY6PYU1TtrOhCjSj3RWPeZEoMt+u0y9jxsUYscVAAQbwAMAQ8gYfaALPVVv06FL59eCfgnvP7i98yzuPb1K2hzUFOo/dhftMC7WcFppKzYxn11JSqe+Tf7ypqFR0bGtVXONNteOOB92kPVWtKl8yCj7kBba9Wdvo17WXOFCbXjuvBgrGio6ppFql+JU5vj82nu/7RldqpZ0v1C5161Kl5Oaz9iZaaTD1u0lSeOFC1+2c/4QJR4arf8mwrkgFyBFtARxAEkcScEYAAY4k47wI4kgAAAAAAAAAARxJAEcSQAAAAAAAAAAAAAACCQBBIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA1K/t3L+WbOOVLflOH14KS/vJmzafXV1Y29zHlVpxmvNZMPqcfU7QqfJXFt9tOX8J/YXOycv9xo0Hzt6k6L7lGTS+zBKfDVXtZd6r7FtGtnHqqkZ+SfH7Gy6xnifFzTVahOk+U4uL80U9PqOrZUpy+Nu4l4rg/tItilpXs0alHP4mtKGO7OV9jRe4LGzedTvd34mYZ+lu8fs3S+BAU61SFKlKpUkowim228JJdbKh456WO2f81+jWpp9tV3L/W5/AqWHiUabWas/KPs+MkSpWb2iIa82T26TZzd0r7VvbLpG1bW4yk7WUlRs88lbwyovu3uMvrGtR4vlzLSlKXYl1LwLqlLL4ov6RFaxV43Nab3m0q8IN4xlPPM6Z9EvZGNpo97tjd0/wANft29pvL4tCD9pr6U1/cRzvs1pdzr2uWGi2Cbub+vGhB4zu5fGT7orL8EdfdJ2r23Rr0P1YaVilVt7aFhpsev1jW7F97SzN/RZx8u8zqkfKw9MwxG81vEOa/SX2lhtd0jXVK3qb9jpObK2afCUk81ZLxkt3PZBHls6eOGUX8l7Cbbba5t8X3soyjnhg7MWOKViHFmzTkvNpWm5jHDzNr6PY/8IT7oQ+9muSi88EbRsTHc027ny36yXuj/AJmbObPP6GSvIxc8yxJJ5wzpH0etq6W0mxFTQ76cat3paVvUjPj6yg0/Vt9vBOL+j3nNN3JveRkuirayexnSBZ6pVnu2NWXwa9XU6Mmva+q8S8E+05c+Prr/APG303POLJqfEq3SZs9PZHba80VxatlL1tpJ/KoS+L7uMX9EwMZHR/pJ7Lfy3sjS2jsKcal3peaknFZ37eXx+PXjhLwT7Tm2jHrT4dvadXFy+5Tujz8HtZO3iVdcVyPZfRe2ohYarc7K3dTdpXubi0T5Kql7cV4xSf1X2njcefDgi4s7u4067oahZVHSurWpGrRmuqUXlPw/cS5GL3aTDVwuRODLFndBDfcYXYfX7fabZaw1q3aUbmkpSin8Sa4Sj5STXkZs85ManUveUvF6xaHj3pWbDS2v6Na15ZUXU1XRXK8td1e1OCX4Wn9aKyl2xicNUG54lnKfFeB+ok4qUWnxTRwL09bDPYbpL1Cwt6Thpt5J3lg0vZVKbe9BfQllY7N3tLHgZe/RKp9Sw6/XDz6MHjqXiX2kXTsruM3xpS9mrHtXb5FuknjHFH1hPgWmtKW0dUabo4LCcZKUWsprrXaUppcmWGzV56ym7Ko/ait6nn5vWvIykodwiXDaJrOn1pl/eaTqNvqWnXNS1u7ae/Rq03hxl+9dqfBrgdY9CvSnZba2a0+/9Xaa7QhmrRTxGvFflKeertjzXhhnJijl4aK+nXNzp2oUb6xrVbe4oTVSlWpy3ZQkutf+uPI5uRx4yx/bs4XqFuPbv4d9eAPLuhLpRttsbNaXqkqdvr1CGZwXswuYr8pBffHq8D1BPPWUl6TSdS9fhzVzVi1Vjr+kafrmlV9M1O2hcWteO7OEl9q7Guaa4pnKfSjsPqOxOr+rqesuNMryfwS7xz/0c+pTXua4rrS67Mfr+j6frulV9M1S1hc2teO7OE19q7Guaa4pm7j8icU/05OfwKcqv9uJHJPkm2fUMt+Zt3SdsFe7D6vuSlO40q4k3a3L5vr9XP8APX2pZXWlqGVweeBe4r1yV3DxefBbDfps+8LeTa4rij0noV6Q3stqkdJ1Oq3o11PjJ8rWo/l/QfX2c+3PmMqiwsLBSnUl8nrGXDXJXplLi574MkXq7up1IVKanCSlGSymnlNE72TwT0b+kKpVjDYzWq7dWnH/AHNrTfGcEsui31uK5dsfA96jiSPPZcc4rdMvc8Xk15FOqHOfpbtvWND4/wDw1f8Abpni8W+HDie4+lhST1XQ3/2euv71M8TnGOS74UbxQ8h6pOuRMS+qc+K6z2H0V2o7Yay/+wU/2zxrGHwM5s9tVU2U2f2nr21Rwvr+zp2Vq484ucnvTT6sRy134NnKpN8c1hDgZIxZovLGdPG08dseky/1OjUc7O3bs7Pri6VNv2l9KTlLwaNBk2uX2leU1uqKXBLCLeo8t8TXTHFaxEOjJlnJebSidSceOWmuR0T6GuyXwzU9Q21u6WadonZWLa51Gk6s14Jxj5yOdbehcXt1Ss7Sk6tzXqRpUaa5znJpRXm2foN0ZbL0NjthdK2dt91u0oJVZr8pVfGcvOTbOTmZOmvSsPTMHXk65+GzLkAMlW9IAAAAAAAAFpb/AITUbmpjhBRpr7396LqTSi2+SLXSk/gaqyWJVpOo/N5X2YDCNV9uFKh11a0Y+Se8/sTLwsp/hNWpx6qNJzfjJ4X2KReghg9o571/p9DsnUrv6sML7ZonZdb9fUrj51dU0+6EUvvci21OSqbRVZN+zb20I57HKTk/sjEvtk4OOg29SXxq+9Xf15OX7yXiGqO9mWQAItwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAEZAkEADDbTRUKmn3P8AZ3Hq5PsU4uP37pT2al6vUtRtn8p07iP1o7r+2H2l3tNSdXQrvcWZ04eth9KHtL7UYuxrRhr9lWi/YuaU6We3gpx+xSJR4aZ7XbP1FjSqxtXeRqfFpt1l9GXH78l6i0u7R17mjU30oRf4SOPjpPKXvRCJbU6XQnRtI+sX4Wo3UqfSfF+7l5F2EDLKJS3VxZw56S+1kdsOk66VvWU7DSM2Nrh8JST/AAs14y4Z7Io6q6ddrlsX0banq1Koo31SKtbJdbrVOEX9XjLwizhFW+MZblxzlvi32lhwcW5myl9Vz6iKQqU4YfMrJtcU8Hwk1zKlvSrXVxTtbWm6tzVlGnSpri5zk8RivFtFlM6hRRWbTp756HuzXwzVtR2wuqeaVmnZ2ba4OpJJ1JLwjux+tIxnpVbZR1rbals3a1N6z0ZfhcPhK4muP6McLxcj2expWPRJ0LJT3ZPS7JzqY4evuZcX+lUlhdzOMri4ub2+rXt5UdW5uKsq1ao3xlOTbk/e2cGCPcyzeVnyr+zgjFHy+5YwopvhwKbjw7D7fLgfDfYmyx2p47KM8rxybZspHGhKXLfrTln3L9xq8km+ztNx0Gm4aDaxx8aDl75MhZqzz2UrnnkxdzS328pPg1hoy9dbr8jGXElxwYiGqkzHh1F6NO1MNqejp6LqMo1r3Sf6nXjPj6yi0/VyfbmKcX3xZ4N0kbP1Nkdtb/Q8NW8JqpaN/KoSy4ce7jF98WW/QttdPYrpCs9QrVHDTbqXwS/zyVOT9mf1JYeezePfvSU2Qjq+y9LaS0pqV3pSbq7q4zt5fH/R4S8FLtOSlvYzf1K/vX8Xxdx5q5s4tc+BLi5JrLS7hGKjHnwXDIi0W8Q87PaXtfov7Uu01W72SuqmKV1m5s8vlUS/CQXisS8pHRKOF9G1G50nWLTVLKW7dWlWNak+rei+T7nxT7mztTZTWrTaHZ6y1iylmjdUlUis8Yvri+9PKfeij9Qw9F+qPl670ble5j6J8wyjPH/Ss2Le0/RvU1Ozoes1PQ27uior2p0sfhYLxit7HbBHr+T4rJTpuMopp8Gn2HDjv02iYW+akXpNZfmcqicI9eYrBMXl5SeTdOnfYyWw/SPf6ZQp7mn3D+F2D6lRm3mC+jLej4JdppsFwxjh2l/jv1xEvJ5qe3bpfdtOdCvCtSlu1IvMX3m52laF5awrwwt5e0vmvrRp6XBcDLaBdO3ufVVJJUqrw89Uupk57OPNXqhnXDCyU5rKwXU4tcH1lP1fEzEuNTtLi7sryjfWVzUtrmhNTpVacsSpyXJpnVnQl0nW+2NjHTdTnSt9eoQzUguEbiK/KQX3x6vA5W3eeCvp11Xsb6je2dxUtrihNVKdWnLEoSXJp/8ArPI0cjjxlj+1hwedbjX/AKd4oHnPQr0kUNtNKdpfOnR1u0gvhFKPBVY8lVgvmvrXU+HY36JvZ5FHek0nUvZ4c1ctYtVjNqtD07aPQ7jSdUoqrb144a5Si+qUX1ST4pnHXSBs9qWx201fRtQTml7dvXUcRr0s8JrsfU11Puwztd8Vx4mkdLmwtptvs1Oynu0b+g3VsbhrjTqdj/Nlya8+aR08TkzitqfCv9T4McinVEd4ci72VyI4vm3hH1dWlzY6hX0+8oSoXNvUlTq05c4TTw1/n4ERT8j0FJ6o3Dx969M6lVtKtW1uKdzbVJ0a9KcalKpF4lCSeVJd6Z190SbY0dsdlaV3Jwhf0PwN7SXyaiXxkvmyXFe7qZx9lLtRtXRjtlW2M2qt9RjKUrKolSvqS479LPxkuuUXxXmus5ObxvcpuPMLD0vm/h8up8S9G9LFpahobbX4iv8AtUzw2pUy8ZPZfSrvKN5LZ28tqsatvXta1SnUi8qUW6TTXkzxJvJngxrFEIeqatyJmFfeb7jEbUzcbShBPhOo5PyWP3mQlNrk8GB2mr5ubeln4lLPm2/4I6bz2cmGv6mOk5IpSeSomurJSmknx44WcLmzRPaNu2I3Ons3oh7KS1/pCr69d0t6x0NKdNtcJXEk1D9Fb0u57p2UuXA889H7Yz+ZXRnp+n16W5qF1/XL7PNVaiT3X9FbsfqnoRRcjJ13mXq+Fh9rHEPoEccEml1gAAAAAAMgWmqSfwV0ovEq0lTXdnm/dkuYxUYKKWElhFpNqtqsY59m3hvP6UuC+xP3lXUazt7OrVjxkliC7ZPgl72gwpacvWVrm5fy6m5H6MeH35LxlKzoq3taVFPO5FLPa+ti8rRt7WrcTaUaUJTl4JZB4hqWo1pTttWuYfHq1Z06b70lSj9uTbbSlC3tadCC9mnBRXglg1KxpSlQ0e0n8erWhUqeSdWX2pG5IlZCkfKeoDiCLYAAAAAAAAAjJIAEZJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAABxAAAAABxAAAAAAAAAADiABA4kgAAAAAAAAAAAAAADiABBIAAAAAAAAAAAAOsAjiA4gkAB1AAfNSMZxcZJNNYaNFjJ22mW1Vt72n14qT7FTnuSf6OTezVbu0hO+1WwlwjWaqR+jUhuv+9FslVryR8topttH1gstAru50a0ryeZzox3/AKWMP7cl8R0nHgDYNe6RNp7TY7YzVNorzEoWVBzjDOHUnyhD60ml5mYjc6hi9orG5c2+lttc9Z2xpbLWtRO00eG9WafCVzNZ/uwaXjKR4ospJdfIXWpXOo6hd399VdW6uridevN/KqTeZfayYLLPQ8fHGOkQ8fyss5cs2lVjTyj1j0WtjXrnSFPXrmlvWWiQVSOVwlcTyoL6q3pdz3TyqHsPe6kstHaPRPoNp0c9E9OpqTVCsqMtQ1Kb5xm470l9WKUfqnPzcnTXUeZdHpuHrvNreIeVel1tROre6bsda1PYp/1y8w+cuKpwf96X6J4HuvuybLtjqdfaPaW9128TjWu6squ7n4ifCMfqxSXkYSdNp8s5JYKdFIhx8rP7uWZWrWOGT53eOCvKOFwKbWGzdEtEPhpKLbXUbvaQVLTLWnn4tCP3GmUbeveVo2tpRnXrz4RpwWX49y7zdKznRSo1YOEoRUceCwRtZozsdecW8ZMXXiZevut8DG3DSTxgzDXWVhcpOlKm4pprDTOuvR32qpbY9GsdL1Gcbi90yHwG7jPj6ynu4pzfapQ4PtcZHIlzLHI3DoK2yexnSRY3VxW9Xpl/iyvsvCUZP2Kj+jLHHscjRysfXTceYXHpmf2smp8SuekbZ2rsjtZfaJNSdGlLetZP5dGXGD8uMX3pmvp54tvB0p6Tmyv8p7MUNpbSkpXOlv8ADtLjK3k/a/ReJdy3jmtprrR08PN7uOPuHP6lxvYyz9SqRlKLPbvRf2rdG/vNkruotytm5s8v5X5SC8ViXlI8M3nnw7CtpOp3mjarbarp8/V3VrVjWpPqcl1Puayn3MnycUZccw1cLPODLFod35IS7TFbG63a7R7NWGt2cs0byjGoln4r64vvTyvIy+MHmJrNZmHu6Wi9YtDxT0tNi1tJ0fS1q0ob+o6E3cx3V7U6DX4aHuSl4w7zjqDWV1p8V4H6WXNGnWoyp1IRlCSalFrKafUzgnpb2JqbE9I1/okINWUpfCLCT5O3m3ur6rzD6veWfBy/+JUnqmDX64alCOX1lanT4cE2VYU0uGCrGPAsplR9UM9pFf4Vbbk3mrTWH3rqZdSju95gLKu7WvGrHq4Ndq6zYXONSKnB5i1lPtMbcmSupUH7ilPPF5K8s8eGT4lF9ROJa4To+ranoWr22raVcyt7u2nv0pYyuxxkuuLXBrrR2L0XbaWG22zNLU7Xdp3EX6u7t85lRqpcV4Pmn1p+JxrUpt5RsnRdtXebEbUUtUo79SyqYp31CP5Wl2r86PNea6zj5fHjJXceVt6bzpwXis+HakePWHFPqLfSr211LTbe/sq0K1tcU41KVSDypRaymXRSdOpexiYtG4eHekjsHC4tpbZ6ZQzXt4KOoQgvxlJcqnjHr/N+ic+yaWW+eTvC4pU61GdKrCM6c4uMoyWVJPmmjj3ph2OqbE7V1bSnGX8l3O9WsZv5meNPL64N48HFlx6fyN/47PMes8HU+7T/APWnvkfE5NPg8M+HUTWeKyim55feW7z0VllNR1q8vtB03SLhqdHTZVfg8m/ajCo4tw8E48OzODFzeerJ858gvcYrSI7Q2WmZncplybNX1yTnq1bL+Luw9yX7zbqMMyjnk3xNMvJutcVKz+XUlL3sjlrpvwSiL4c+R6J6OeyL2x6VrCnXpes0/S8X922uEt1/g4Pxnjh2RZ5rN7qTzwzxfcdpeihsb/Nro3p6tdUtzUNcau6m8uMaOPwUf0XveM2V/Ly9GPS29Pwe5l39PYEiUTzBSPUAAAAAAAAB81JRhBzk0lFZb7EfRZanmoqdnHnXliXdBcZfw8wJ0qMnbu4msTuJOq89SfxV5JIi7/DX1vb4zGGa014cI/a8+Rd8EuxFppidX114/wAtL2PoLhH38X5hheIxO1s2tDrUo/GuHGgvryUX9jZlzA7TyU73TbZf2s67XdCLS+2UTMeWLzqFPS4qttBHHxba3lLznJJfZB+82IwWy0d641G5+dXVKL7oRX+05GdFvLFPAADCYAAAAAAACCfMYAAAAAAAAAAAAAOIAAABxAAAAAAAAAAAAAOI4gYAAYAAAAAAAAAAAAAQSAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAMDAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAwAAAAAAAAAAAAAAAAAAGAAAAAAAADC6zZ309ToXNhClJypSpVZVJYjDinGWOcse1w7+ozQeAxMbWmlWcbGxhaxqSqbrk3KXNuTbfDq4tl2AGUM5d9NHbDfvNL2Gs6vCLV9f7r8qUH/el5ROmNWv7XS9LutSvaqo2trRlWqzfKMIptv3I/PnbTWK+1O1uobR3ikq1/cSrKLefVw5Qh9WKivI7OFi677+lZ6nn9vH0/bD0opJpLDLmMkuMuCXaU9zEs4ZFaW5SlJrgll+BczOoeb8vS/R92Y/nZ0l2FtWp+ssdPxfXb5xcYNerg/pTxw7FI9t9KjaR22z9psraVcVtQmqt1h8VQg+X1pY8osqei7srHZno0jrl7FUr3WsXlWU+G5QSfqovu3fa8Zs8J6Q9pqu1m22pa25N29Wp6u1Wfi0Y8IeGeMvGTKzfvZtz4hZZb/huL0R5lhKkW2+ZQlBvPDLK+83wf2FahRdSWFDPa88F3s79w873YuvDg8dXWXelbP3F8qdzdzdlZS4qbjmpVX+jj1r854XjyMzZwsLRqpKELquuKc4/g4Pui/jPvfDuK9a8lXqOdSbnOXNyeWyNu/hOMuoX2nRtLCg7bTaPwejL473t6pV+nPr8FhdxWrqjcUtytHeXJdsfBmIVdLkfULldvma+lotM2na01WxrWyc6adWl85LjFd6/eYOvJyXLgbhSuOOFLiWWo6Vb3Oalu40ar6vkS/gzZWUqXiO0tQuIosbiMJYUoprDTT68mQ1OEraq6VaEoTXVL712mJuJtvK5GztMO3Fve4dn+jttfQ286LYWWpyjcX2nxen6hCpxdSO7iM2utShjPepHOXSFoVbZPbHUNAqbzhb1M285fLoy4wl7uD70zHejptxLYjpRoSu63q9J1VqyvW37MMv8HUf0ZPDfZKR7/6U+yT1DZ6htbZUd650z2LrdXGdvJ8X9SXHwcjhwz7GfXxK55MfiuPE/MOdotPOVnPYRUbcWl2cD4g3lprGOsqKOctrkXHl53xL3X0Utr3Cpd7HXk8KWbqxcn1/lYL7JLxkdEReeJwns/qNzoWt2Ws2Txc2VaNaC+djnF9zWV5nbuzuqWut6JZ6rZT37e7oxq033NZw+9cih5+Hov1R4l6v0fle7j6JnvDIHjvpT7Ffzh2Jjr1lS3tQ0XerPdXtTt3j1sfJJTX0X2nsR81YQqUpU6kVKMk04tZTRxY7zS0TC1z4oy0msvzsilx3WmmxhZNp6YNl57F7f3+iQg42UpfCLF9tCed1d+68x+r3mq0/aWUX9ckWrEvGZcM47zEpaT5GT0Wvztpvnlwz9xYRj2FanFqalHKaeUw1XjcM5KD4cj5lDjhFS1k7iiprn8pdjKvq1j94izkntOltGnw5n04p4eM4eV3FVx7j5cXHHBI2RKO3tPo2bbu1u3sbqVbFCs5VNOlJ/FnxcqXnxkvrLsOg8nCELiva14V7erOlWpzVSnUg8OE08pp9qZ1x0O7a0dttlKd5KUIahbNUb6kn8Wol8ZL5slxXmuoqebg6Z648PWej83rp7VvMN2bSNI6ZNi6W2+x1fTouFK/o/h7GtL5FVLgn+bJey+59xu+EQ4JnFS80tFoXWXHGSk1lwDc0riyu6tne0ZULmjN06tOfOE08Si/BnzltZx5nuvpQ7BeouYba6ZQ/BzcaWpRiviy5Qq/dF/V7zw1U2vjeaPTcTLGanU8RzcE8fJ0y+Es4x1lWKwlwIwuSIc8ZOzWnFM7fdSapUKtV/Ipyl7kzSpvEVl9RsmvXMaWkXDXOUVBebX+ZqE6ykpJ8mu05s1u+nXgpOtto6MNlqm2/SDo+zcIydC4resu5L5FvD2qj7sr2V3yR+htrSp0LeFGjCMKdOKjCMeCSXBJHOvoSbIxt9B1PbW6pL12oVPglm3zVCm/aa+lPK+ojpBJJcEee5mTrvr6et9Pw+3j39i4kgHKsAAAAAAwAAIfIsrD8PXq3rXsy9il9Bdfm/swfWpSlJQtKbxUrvdyucY/Kfu4eLRc04Rp04wgkoxWEl1IMLbU5N0o21NtTrvcTXVH5T937i6pxjCEYRWIxWEl1Is7P+s3dW7fxI5pUvBP2n5vh5F8CBmtalNVNoq0n8S2t4wz2OTcpfZGJsj5GlXdSVez1KvF+3d3E6VNp9slRj92TNUMk9tNh2TpuGg21SSanXTryz2zbl+8yp8UacaVKFOCxGEVFLuR9mE4jUAADIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAwABBOEAAAAAAAAAAAAAAAAAAwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAB8t7qbYEvgant10i7HbFU09odct7WtJZhbRbnWmu6Ecyx34weS+kb05VNnLutshshWg9XUcXt7hSjZpr4kVydVrjx4R73wXLFS6uru6rXd5XrXNzWk5Va9abnUnJ83KT4tnbg4c5O9vCs5XqEY+1fLoPps6dtK2u2Ur7M7OWl/Qp3dSKuLi5jGG/STy4xim37TSXHHDPaeFc3heJaQbzxRcQlx5lrhxVxRqqg5Oa+e27LhReUsGzdGOyE9studN0PdzQrVPWXbj8m3hxnx6s8I+Mka3ScW1xSOovRF2UjYbNXm1lzT/DalN0bVtcqFN4bX0p58oxNfLyxTH2T4OCcuWN+IZ/0hdcWz3R9LSbFxo3GpL4JRhDhuUkvwjXYt32frI5ZpQ3I7uMLqXceldOu1Ntre3F5XqV4xsdPzaUJN+y8P22u1uWV4JHkOqa369uFpF06fXNr2n/AANPGxTFNoc/L7uaax4hkri7oUJJSk3L5kefn2Hwr6U4qOd2PzVyNdjV3ef3lSNylyydXS5Jx6hsUbjHyitC649prsbp4wVYXhnUNU42wRuEypCvH5xgYXi58VgqQuk0sNjUITjZ+N0kuZMbt54ZZg43LbxxPuFfxIzCPtspeRoXlB0rmCnHq7Y96fUavrGj3NopVqTdegvlJe1FfnL95mY3KXMqfC+SWcvrQiNJ0vajQ7tKVNp4aksHa3o5bXUekPoojY6pKNzf2MHp2ownx9bHdxGb+nDGe9SOUdS0SjfJ1aDjRqvqS9iXiurxRtno67Tz6Puk+3hqE5UNN1VRs73feIxbf4KrnlhSeG+yUjRycU2p1R5hc8DkV6umflT212audlNrb/QbhTcbap+BnL8pRfGEvdwfemYuKS4HSfpRbKK+2fobU2tPNfTvwdy4rjK3k+f1ZcfCUjmtNcOPHPM7eDljLj38uD1LjTgzTHxKtBrKylweUe/+i9tapRu9kLurxhm5ssv5Lf4SC8G1LH5z7Dnt1McfcXWz2t32g7R2OtWLarWlVVIrPCePjRfc45XmS5mGMuOYR9PzzgzRZ3i3hEPijG7L6vaa9oNlq9jUVS3u6MatN9zXJ965PvMpjqPLWrMTp7qlovWJh4r6Vmx71vYmOv2lDevtFk6st1cZ27x6xeWFL6r7TlWiuCzx7H3H6HXNvSuKE6NWEZ06kXGcZLKaaw00cRdJWyUtjttr/RN2SoQn62zk/lUJtuPjjjF98Sw4eTcdMqD1bBNf8kNZp0m+0uKdPHBrzPuEOOF1FWMPE73npuq6fP1M91v2J8JfxMnOHDD5GJwsYX3GRtK2/S3JNuUFjj1oNeSN9ySa45KMslWfFcWylKPWSiWuIUqiysGydE21lfYna2nqK3pafXSo31KKy5U88JJdcot5XdldZr7XWbP0T6NHXOkDRrCrCMqPwn11VPjmFNOeH4tJeZHNqaTt28Obe9Xpdg28lUoxqRk2pJNdXBlQRWI4JKDUbe8rvXdbapZ2moadcWN7RhXtrinKnVpzWVOLWGn5HD209np+nbRajY6Veu+sbe4nTt6/z4J8OPXjllcHjPWdO+kdtTV2c2Bna2dV073VanwSnKLxKEGm6kl9VYz1OSOTZNeyksKPBdRd+k4rd7/DzPruas2ikR3RUfjgoyzIrN97PiUM9fEuZh56ssDtTUSsKcM8J1s+5P8Aia6k3HnhmY2vf4e3pJ8oSm/N4/cYaPIr8v8AyWuHtSHpHQd0t6x0ba1TpzqVrvZyvUze2OXLczwdWlnlNc2uUvHDO8NH1Kx1jSrbVNNuadzZ3VKNWhWpvMZwkspryPzKST5rOGdUehNtrOvbansJeV3P4LH4dp6k+VOUsVILuUnGS+myt5mDt1wufT+VPV7cumUAgVq6AAAAAAicowi5SaSSy2ySxvH8Jrqyj8RJSrtfN6o+f3ZBKdPi61Sd9UTzVWKafyaa5e/n7uw+tSqT3I21F4rV3uxa+SvlS8l9uC5bjCDk2oxSy31JFnp6depK+mmvWLdpJ/Jp9Xv5+7sDC7oUoUaMKVOO7CCUYrsSPsAMrXVLlWmnXF0+Ko0pT9ybNW063araPYS4yjJVKn1Itt/pOJm9qpZ0+naf9Jrwpv6Ke9L+7FlposPXbRV6qXC3tlFfSqSy/shH3ko7Q1W72bEuSJIJItoAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAIAkEEgAAAAAAAAAAAAAAEEgAQSAAAAAAAAAAAAEACQAABHEkAef8AT9tu9g+ji91a3lH+UazVrYRaz+GnnEsdaik5fV7zfzk/05dYnW2j2f0BS/BW9rO7lFdcpzUE/JQl72buPj68kRLm5eT28Uy8DulWrXNS5uK061erNzrVJvMqk28uTfW2z5UeKTfAqZy+PI+kuPFI9D0xHaHkpvMzuXzDnxKsc43sEbr3uCPvepUqbqVpKEF1saYZLZPSLraLaXTtBtN5Vr+4jQjJL4ifxpeEY5fkd722i07LZOOgaPXlp1OjZ/BbWrCKk6OIbsZJdbXPxPzw0zbDVtB1SOo7O3lTTrqnGUIXEYQlUxJYfxk0vvMy+mfpWfFbdaqvqUf8Bwcrj3yT2W3Cy0xVnqh0Lfei/QvKqqVtvdQlhYinZQwv7xQj6Ktkv+fN8/8A8CH+I8Bj0ydKr/5+at7qX+Ar0+mXpSTw9utWflS/wEIxcjxEtvXxPmr3j/2V7LGP58Xv6hD/ABEf+ytY4/48X36hD/EeHLpl6T+H/vxqvupf4Cf6YulB/wDPnVV9Wl/gM+1yf5Me7xP4vb//AGWbKPLbi+/Uaf8AiPpei3Zri9t75/8A4MP4nhsumDpRa4bdat+jS/wHx/TB0p4x/PrVf0aX+Ax7XI+0JycX+L3f/wBl+zX/AD3vf1GH+I+16MNosY21vf1GH+I8G/pf6UevbrVf0aX+A+l0v9KD5bdap+jS/wABj2+R9sTfifxe9w9GW0S/46Xv6lD+JUXoz2a/55Xv6nD+J4HHpf6UP+vOqv6tL/ASul/pP4P+fGq+6l/gHt8j7Ovifxe+f+zRZ/8AXG9/U4fxEfRptc/8cbz9Th/E8Gj0vdJzSztxqv6NL/Afa6X+kvC/999U91L/AAD2+R9sdXD/AIvfqXo42tNJLa+7/U4fxPnUPRtsLy0lb3G1d1KL5P4HDK8OJ4HLpe6TMcNt9V91L/AU5dLfSdJ8NudW/wD1/wCAe3yP5JRfhxO4q7e0jR3bbGW2zmrXktY9XZq1r3FWnuyuI7u63JJvi1zOMNt9DudltrdQ0G5cn8Fq4pTf5Si+MJeaaz3pmOXSz0mcM7cat5+r/wABidY2u2g1y+hd6/qVbVKtOHq1OsoqajnOE4pdbfPPNm3h0vht3+Wnn5acisa8wySllLLEllpotrC5o3X4qWWucXwki9hHjyLiNWhRW3SXv/olbVOdrfbIXdTPqm7uycuuEn+EgvCTT+s+w6BOGtkdaudmtotP1yzTdSzq+s3E8esg+E4Pxi2vcdtaNqFrq2lWupWVVVLa6pRq0prrjJZR571HB7WTfxL1fo/K93F0z5hdnj/pO7IrVtlae0dpSzeaRmVXdXGdvL46+rwl3JS7T2ApXVGlcW9ShWhGpSqRcZwksqSaw012HDS/RaJWXIxRlxzWXBs1iWP/AEwnhZwZnpG2fuNlNt9R0Sal6mlPftZP5VCXGD78LMX3xZhYqT480XdbdUbh4fLhnHaay+k88CpRk4VFJLl3cxGPiipCm85WSTTOl0/ainFcyPVvqyfdvFfEk3z4FZxwuKaMw1yt1T5YwepejRbQq7e3FeSy6FhUcezMqkF92TzJrjwPUfRlqqG297Tb41LCTXlUh/E1cmf8cuz0yP8AZrt0ciOPaSiG8FHL3jmL0vtRqy2x0PTsv1dCynWSzw3pz3fugeKqo8cT2H0vrecdu9GvMP1dbTpU08dcKjb/AG0eNxi+w9T6dH+CNPFepd+Rbaop+J9xe82uXDmUnw5tH1Skt7D7cHfrsr9fTWtpmpavNf2cIw+zP7zFvGeXEutXret1W7qZwnVa8lw/cWby+RwX1MrGkarCJS6kb96NWp1dN6edmZxk1G4q1LWp2OM6U1j3pPyNAkuHBm9ejpZVNQ6dtk6VLP4O7lXl3Rp0pyf3HPyP25dfFj/LGn6Ex5IDHAHn3qYSAAC7AD4q1IUqcqlSSjCKy2+SQFK+rqhRzGO/Uk92nD50nyX/AK6hY2/weliUt6pN71SfzpPn5FGyhOvWd9Wi1lYowa+LHtfe/wCBVv6/weklCO/Vm92lD50v4db7gwoXjd5dKxi36uKUq7XZ1Q8+vu8S/SS5FvY23waiouW/Uk3KpPrlJ82XAEggxeuam7SEbe3Sq3tbPqqb5LtnLsiv8ubHkmdLHVa3wrXFTi807Om1L/6k8faor+8XOyVNuyr3zWPhlZ1I/QWIw96in5mHjbSrTho9vVnOrWzUu63XGDftSfZKXFL/ACNvpU4U6caVOKjCCUYxXJJdRKe3ZCnedvokAi2AAAAAAAAAIAEgAAQSAAAAAAAAAABAEgAAAAAAAAAAAAABHmBIAAAAAAAAAAAAAAAAAAAgkAAAAAAAAAAAAAAAAAAAAAAAAAAABGCQAAAAAAAAAAAADgABBIAgkAAAAAAAAAAAABHmSAI8yeAwAAAAAAAAAAAAAACASAIJAAAAA1wOLPTZjUpdK9rUmnuVNJouD8KtRM7TOZfTr2Zr3GhaJtZbUnJWVaVndSS+LTqtOEn3KccfXOniW6ckOLnU6sUuaqcuCLiliT4lnCShDem0kuvqLSvqNT4lvmEfn9b8Owvol5mce57Mvc3lK2juv26nVFPl4swl7Xq3VTeqyz2Lqj4FL1j6xnPMlMladKlKOHxfvCb5lSUU1ho+ZQeO5EdNsWFh8T6Unnn1lOTUHhyivMmM45xvwx4kJmEtTK4jLhjJUUslopRz+Mh7z7jUXz4e8zFoQmkrpceHYfXey3VRcvWQ959qafKcP0huEemVZLiEmfCnH+0h+kTlccVIY+kY7fZ0y+4dzPpYxwwU9+POM4LHeFOOfxkPHeMdmNT9KylhdhG8U1KL5zh+kSpQf5SGfpA7/SrnPWiU8NFPejnG/D9I+lKOONSHvMbNS+03nkfUVwwfCnFfLg/M+t6Ofj08/SMf/rHdcUJOlJTg2n1NdRm7DVU8Quv/ADEvvX8DXfWLP42HvJVZL8rBPxNlMnS05MU3+G7b8PVqpCUZRa4STymdBeivtjG70272Qu6v4azzcWeXxlRk/aivoyfukuw5Jp39ahLNGvFJvjFv2ZeRtnR5tjLZ/azTtbt5qFe0qqU6e9wq03wnDPfFvzwQ5U1z49fLfwYvxssWjw/QBM+ZcVjJb6VfW2pabb39pUjUt7ilGrSmuUoyWUy6SPN6mJexieqNvF/Sc2S/lLZyltNa0t660vMa+FxlbyftfovEu5bxzpQpyaW9z8Tu68t6N1a1ba4pxq0qsHCpCSypRaw0zjvbbZqtsrtTe6LUT9VRqZt5vnOjLjB9/D2X3xZZcTJuOmXmvWeNNP8AJVr1OkusrQgk+HPxKiWXwR9Ri0+w7nnt7fMVx5lfGVldfM+YwXZxRWUUk+JjekbaUsY8DbOhjU4aV0k6VOpPdp3E520n9OL3f7yianJPj1PrKKlWo1YXFCe5WpTjOnJdUk00/eheOqsw28bJOPJFvp3BF70chrJhtiNao7Q7L2Gr0WsXFFSml8ma4Sj5STXkZoorV1OpfQcd4vWLQ8k9J3ZOeu7Dw1W1pOd1o1R3GIrjKi1iol4LEvqnK0nFY3eK+8/QGtCNSnKE4qUZLDTXBo5D6dejW82N1arq2l0KlXZ64m5RlFN/A5N/i5dkMv2ZeT6s3XpfLin+Oyg9Y4U2n3aPNZPLxyPic1BOb5RTk/BH3SjJpb8cPuPnV1GnpF1UfNUml58P3l9eN12oKzqYhoMarnUlJvjKTk/NlbeWeZQnFRfDkfEptcU+PcVszqVlrfhdycXwfWdIehHsTUra5qm3d1Scbe3pOwsW18abadWS8Eoxz3yXUeK9EWwOu9I21FPSdKpSp2tOSd9fOGadrT7c9c2vix6/DLP0H2S0DTdl9nLHQdIoepsrKkqVKOct9sm+uTeW31tsr+byI6eiFn6fxZ6uuzKJd5IBUrwAIbAlvBj0/wCUK2cf1SnL/wA2S/2V9rIrznf1JW1JtW8XitUT+N+Yv3vyL6EYUqahBKMYrCS4JIMeSrUhSpSqVJKMIrLb6i1s6U6tV3txFxnJYpQf5OH8X1+SPil/X66rP/etOWaa/tZL5Xgur39hfoHkA4GI1fVJU6srKwUal1hOcpLMKCfXLtfZHm+5cQTOn1rOqq0at7eCrXc1mNPPCC+dN9UftfUYCCr/AAuVC2fwvU7hKVWpPlCPVKWPiwXHEVz97PunSrTuJ2Ont1rybU7i4q8VTz8qfa8fFgvsRsuk6bb6bbeqo70pSe9Vqz4zqS65Sf8A6S5LgT3FWrU3lGjabS063cIylVrVJb9atL41SXa+zsS6lwL4Ag2xGgABkAAAAACCQBAJAAAAAAAAAAAAAAAIJAEEgAAAAAAAAAAOAAEEgCCQAAAAAAAAAAAAAAAMAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAOscQAAAAAAAAAAAADK7QBi9qtC0/aXZ6/0HVaPrrK+oyo1Y9eH1rsaeGn1NIymSMiJ13hiYiY1L85OlXY/WNhNsLjQdZhOSjmVpcbuIXNHqnHv6mup8DUZOL5JNH6PdJuwezvSBoEtH2hs/WwTcqFem92tbz+dCXU+7k+tM486R/R2292UuKlxo9tLaXSk24VbOP9Yiurfpc2++G95Fxx+bW0asouRwZpbdPDyNY6iU+sr3lnd2NV0r+wu7SquEoXFGVOS8pI+KNOpcVFC2o1K03wUKcXNt+COzrr524ei+9aI4ZlNmtC1PaXWLXRdFtpXmoXc1ChSh29cpPqilxcnwSNw2A6DukXa64pyjotXRbBtOV5qkHSSXbGn8eXuS7zr/oa6KdnOjXTJ09OhK71OvFK61GtFesqfmxXKEM/JXnl8Tkz82tI1Hl1cfgWvbdvC06P+hTYnZ/ZGx0rVdA0nWb+nDeury5tI1JVasuMsOSyop8EupJGdXRb0c/9SNn/ANQp/wADcFgkqJyWnvtfRipEa00/+i7o65fzI2f/AFCn/Af0YdHf/UnZ/wDUKf8AA3ADrt9ntU+mn/0Y9Hif/ErZ/wDUKf8AAf0Y9HuP+JOgfqFP+BuAHXf7Pap9NP8A6Mej3H/ErQP1Cn/Af0ZdH3/UrQP1Cn/A2/rDHXb7Pap9NR/oy6PurYvQP1Cn/Af0Z9H3/UvQP1Cn/A28Drt9ntU+mox6NOj5ctjNB/UKf8D6XRr0frlsboP6hT/gbYB7lvs9qn01T+jbYD/qboP6hT/gP6NtgH/zN0H9Qp/wNrA67fZ7VPpqn9G2wP8A1N0H9Rp/wH9G+wP/AFN0H9Rp/wADawOu32e1T6ap/RxsF1bG6D+oU/4B9HGwWOOx2hfqNP8AgbWlgDrt9ns0+mo/0bbA5/4naF+o0/4Ero12AX/MzQf1Cn/A20gddvs9mn0t9NsbPTbGjY2FtStbWjFQpUaUVGEI9iS5IuQfLb7CLZCcnkfpIbMfyjoFLaO1p5utMyq+Fxnbyftfov2vDePW8Mp3dvRubapb16calKpFwnCSypJrDTJY7TS23PycMZsU0lxSo57CoormZbb7Q6uy21t3o0lJUact+2k+O/Rllwflxi++LMPTk2XFbbjbwObHOO81l94a5I+ort4s+oLLK0IJcWsmGnahOllb3vPn1aw00X8ILDZbVac4yajFvHLgTrJ3nw9N6A9sKekatLZy+qqNnfT3raTfCnWfOPhL713nQSeePUcW2tneX9xG3sbavc120407eDnNPt4cvE6n6La+1FTZelDayz+D3tP2YzdRSnVhjhKaXxZdvHvK7lUrE9UPXejcjJant3jw2spXVtQurepb3NGnWo1IuM6c4qUZJ8GmnzRVTTByRPdezETGpeIbd+j/AKVfzqXeyl7/ACRWk23a1Yudu3+bj2oeWV2I8c206FOk2hYVLW00NajmcVvWt3TlFpPLeJuL+w7SwMceR24+fmpHTtX5PTMF7dWtOCbD0fule+movZqVsn8u4vaEEvHEm/sPSNgPRNuHcQututoYypJ5dlpjlmXdKrJL3RivE6vBC/MyXbMfBx0YbZHZjQ9lNFo6Ps/ptvp9jR+LSpRxl9cpPnKT628tmYBPmcs953LsiIiNQLkB1Hy5RSy3jAZfXAsK1Spe1ZW9vJxpReKtaL/ux7+19XiRKpUv26dvJwtk8Tqx4OfdHu/O93aXtGnCjTjTpwUYRWEkuCDHl80adOjSjSpRUIxWIpdRZTb1Co6UMq0i8VJJ/jX81d3a/LtPupKV/N0qUnG2TxUqLg5/mxfZ2vyXde04QpwjCEYxjFYSiuCQExioxUYrCSwkuol8CJNRWXwSNY1PUq2puVCyqSpWSyqlxF4lV7VTfVHtn7u0zEbYtaKwudU1apWqzstNnhxe7WuUk1TfXGPbP7F19hj7KjVupSsdKl6ulCT+EXb9rdl1pN/HqdrfBdeXwI0uynqkI0rXNtpcPZdWHsuqvm0+yPbPr6u02u1t6FrbwoW9KNKlTW7CEVhJGZ1HaGusTbvKnptjb6faxt7aG7BPLbeZSb5yk+bb7S5AIt3gAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAcQAZHLqAEkPxA4gT3Dv6wAIa7eIwnzSJAFC4tbW5W7cW9Kql1TgpfefNtY2Vrn4LZ29HPP1dNRz7i5BncsdMPlwjnOOPaN1Y4I+gRZfDXiSmlzJ48zWto619SvVC4r1LfT5pKnUoPd9rsnLmuPLHB8iURtG1tNlTT5DKzjPE1uw1WvaVo22pz3oye7RusYjJ/Nn82Xfyfc+BsMGmuvImNFbRL7WesjqKN3bRuaEqU51IJ8VKEnFos6VtbUpxo3UGpSeI1FOSjP7eD7vcYSZNPuHAtXp9p1U5p9qqS/iQ517P4+9XofPSzOHiutd64hheA+KVSnUpxqU5xnGXFOLymfYZAAAYDHUABGO1DmwDHEYJA+JyjCLlLkuPIwmobYbNafUcL7WbS1kuqvPc+8zryfFSjSqx3alOM12SWUI18o23rs1j+kbYVLjtZo363D+JTqdJmwUOe1Omy+jV3vuMnf7JbMX7fw3Z/SrhvrqWkJP3tGFuOirYCvJyezFlTb/st6n+y0Tjo+Wi3vx40t7rpf2AoZxrqqvspW9WX3RMLedO2xtJ4oU9Vuf/AKdrhf3mjKVehnYGpnGl16efm3lX/EUP6ENgc5en3b//ADav8TZHtOW8cy3jTyjpU230TbapZ1bTSbq2ubZyiq9aUVvU3zi0s9aTXHhx7TTadKLWUzpK26G9gKLT/kSVRr+0uqr/ANoy9h0d7F2TTobOafldc6Sn+1k3xyaVjUKvL6RnzX6ry5dp0HOSjTlGcn8mPF/YZrS9kdq79x+B6HfTjLlOVF04++eEdT2emWFlHdtLK3t49lKmor7C6UUuojPLn4hsx+gVj/lZz5o3RBtTdbstQubKwh1pydWa8lhfab5oPQ/s5ZqNTU6lzqdRc1UluU/0Y/vbPSMYWB1mm2e8rHB6Xgxd9bWWlaTpmlUFQ06wtrSn82jTUU/cXsUkPcSapnflYVrFe0QAAwkAAAAOIAjrBb3V3CjJU4RlVry+LShzfe+xd7Aq3FanQpOrVmoQjzbLH1da/e9XjKla9VJ8JT+l2Lu9/YVqFpOVSNxeSVSqviwj8Sn4dr739hWua1G2pOrVmoRXDvb7Eut9wYfWYUopLdjCK8EkWe9U1DhByhadclwdXuXZHv6yFQrXslUuoOFDOY0OuXfP/D78l+lux4cB4PJThGnCMIRUYxWEksJI+LmvSt6M69epCnTgt6U5PCS7WyhqN/bafbuvc1d2PBJJZlJ9UYrm2+w1u/up3LjeaonSoRkvg9ovae91OSXx6nYlwXe+JmI2ja8QraldvUqc5127bTIrecZvddZds/mw/N5vr7D6sNLqaru1LunKhpy+JQa3ZV+xzXyYfm9fX2F1pulVLqpC81OGFF71G1byoPqlPqlL7F3viZ1chvXhGK77yiEIwjGMUoxSwkuCRPAlgw2iAHEAMkccDiA4kgjxAl8AAAAAAAABw5AAAQ8kgCMjiOIEgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAwAAIJAAAAAAAAAAAACCeBAAAASAAAAAAAAAAABHAAAOAEgAAAAAAAAAAAABBIAEEgCCQACAAAAAACAJIJ4EACnXo069KdKrCM4TW7KMllNdjRVAGrXunz0yE4qnK70ySxKElvzorsfzofau9cmn3lXTIRe/O80uSzCUXv1KK7n8uH2rv6towjCX2k1LepO60uKe896rat4jN9bh82X2PrxzJRO+0tc113hl7etRuKMK1CpGpSmsxnB5Ul3M+qtOFWEqdSCnCSw01lM1i0rTozndaX7L3/AOsWlT2U5dfD5E+/k+vtM9puo299CTpNxqQeKlKaxOm+9fv5PqMTGmYtEoxXs3w37i37Oc6f+Jfb4lzRq061NVKU4zg+Uk+B954ci0rWj9Y61tP1NZ8ZcMxn9Jfv5mEirauFSVa0mqVRvMotexPxXb3oUb2PrFRuYO3rPkpP2ZfRfX9/cKV4vWKhcw9RWfxU3mM/ovr8OZXq0qVam6daEZwlzjJZQFRPqHDrLL1Vza8beTr0l+SqP2l9GX7n7ypbXtCtP1e86dVc6U1uyXl1+KDK5fYR95PDAXP/ACAJLrGCcDAEYC8CQAAAAAZAeRHmA8AOAwODHABgY4jrJ8gIwnxJ4BDrAhEgAAAAAPmUoxi3JpJdb6gJyU61SnSg6lSpGEIrLlJ4SLWV7Ou9ywp+ufXUlwpx8+vyPqnZRc1Vu6nwiquK3liMfox6vHiwPj1txdLFsnQov8tOPtSX5sX979zLm2t6NvFqmuMuMpSeZSfa31n1WqU6MHUqTjCCWXKXBIs/W3N5wt96hQf5aUfal9FPl4v3BiVe4u4wqepoxdau1lU08YXbJ9S/9cSLe1frFcXM1Wr9Tx7NPuiurx5lS2oUbanu0o4y8tvi5Ptb62VnwXAMnIx2sapSslGlGDr3NRZp0Iv2pd77I9rf28i31DV5TqTtNN3KtWL3alaXGnRfZ+dL81ebRiqUKjualtp8fhd9PDuK9V5UOx1GvsgvsXElEfbXa/xClWnU+EwrXebvUKuY0KNPlFdagnyj2zf8EZrRdIlRqq+1CUa17hqKj8Sin8mCf2y5vuXAutJ0qhYRlUzKtdVPxtea9qfd3RXVFcF9pfmJn6K0+ZSADDYEE5HmBABPACASAAAAAAAAAAA4AQASBBIAAAAAAAAAAAAQCeBAAAkAAAAAAAAAAABBJAAngRwJAEEgAAAAAAAAAAAAA4ACCQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAYAAxuq6VTu5fCKM/g93FYjVis7y+bJfKj3dXVgwFxGbu6dK8jOw1CGVRrU3wn27knwku2D93WbiUL20t7y3lQuaUatOXOMvvXY+8zFtIWpvww9jrU6M42+rKNKTe7C4isUqj7H8yXc+HYzPJprga1e6dd2MZLdnqFk1hprerQXevyi+3xKWm161rSVTS68bq1zj4POfxe6Enxj9GXDwM9MfCMWmO0tnrUadam6dWEZxlzUllFr6q5tfxEnXo/2VSXtL6Mnz8H7xp2qWt65U4SlTrRXt0ai3akfFdnesovSLZGpULa6o124puNRfGpzWJLy/efVxb0LiG7VpxmlxWea8H1C4tqNxFKrBNr4r5OPg1xRbuN7b/EkrqmvkzeJrwfJ+ePECfU3lv8AiKqr018is/aXhL+OfEmnf0t9U7iM7eo+CjVWE/B8n7yaF9QqVFSk5Uqr/J1Fuy8u3yLicYTg4TipRfNNZQH1ldpOSz+Aqnxta1S3/NXGH6L/AHYI9bfUfxlvGuvnUXh/ov8AiDa9BZw1G1bUak3Ql82rFw+/gy6jKMo5TTT60Db6AAZAAAAAAAAAAAADYAFOtXo0Y71arCmvzpJFt8PjU4W1CtX74x3Y+94+wMbXvDkU69alQg6lapCnFdcngttzUKzzOpSto9lNb8ve+C9x90rG3pzVRwdSqvylR70vt5eQFNXlWvwsreUk/wApU9mHl1v/ANcQrBVXv3tWVy+e41imvq9fnku5ThTi5Tkopc23hItXfOt7NlRlXfz292mvPr8sgXS3YxSSSS9yLSV7Kq3Cyp+vlydRvFNefX4LJMbOVbEr2p63/RxWKa8uvzLyMVFKMUklwSXUBaUrLemq13P19RPMU1iEPCP73ll3yDeFxMNea0qkpUdLhG5qJ4dVvFGD75fKfdHPkNbYmYhkr68tbG3lXuq0KVNcMyfX1Jdr7ka7qF9c6hGW/KdhYJZknLcq1F2yfyI93PvXItqnC9i60qupam1mEIpLcT+auVOPe+L7WZWw0J1JxuNWnCtNPehbw/E032vPx5d78kiUREeWuZtfwstNtq9/ThTsouy0+KwqyhuymuynHqX5z8l1mx2Nnb2VvGha0o0qa44XW+tt9bfayulhYJMTO2ytdAAMJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAB3GL1HRqNxVlc285Wl0+dSmuE/px5S+/saMoAxMbajqMHR3Y6xbqluP8HeUW1CL7VL41N+PDvZdW+o6jYpet/3StsZU44jWS/Zn5YfibHKMZRakk01hpmFutAhByqaXVdnJ8XSxvUZP6HV9XHmS3E+WuazHhfadqllfxbtq6lKPx4NOM4eMXxReLDNOvoeqnF6rZu2nD4l1Sk9xeFRcY/WwvEvbS91K3jFwq09RoNcN9qNTHdJezLzS8TE1+iMnxLYa9CjXg4VqcKkeySyWvwOrR42lzOKX5Or7cf4r3lK01qyrVI0aspWteXKlXjuN+D5S8mzJZMJ9pWnwqvS/wB82s0vn0vbj7uf2FWhdW9d4pVoTkucU+K8UVyjXtqFf8bRhNrk2uK8GGVScITTjOKkn1NFq9NtU80oSovtpTcPsXAfA50/973den+bJ78ft4/aM6jT5wt667m4P7c/eA+D3cPxV9Jrsq01L7sMb+ow50reqvzZuL9zT+8hXs4vFazuafeoqa/u5PqOpWTeHcQg+yfsv7QI+F1or8JYXC+i4y/eSr+n8qjcw8aEv3Ir069Gp+Lqwn9GSZUAtP5StOupKP0qcl+4fyjZf9Ih9pdcCcICz/lOy6q8X4Jk/wAoW3U6svo0pP8AcXXDuJ4dwO6z+Hp/EtbqX/hNffgn4RdS/F2M1/8AUqRj92S6ysFKpc21L8ZcUofSmkBRa1Gfyraiu5Ob/cR8Cqz/AB97cT7oNQX2cftJepWn5Ocqz/0UHP7kPhdxP8TY1nnrqSUF+9/YDs+6Fja0Zb0KMN/5zWZe98S4eFxfItNzUKjzKtRoLshFzfveF9g/k+jJ5uJ1bh/6WWV+iuH2ATO/tozcITdafzaSc3544LzPhy1Cvwp06dtD51T25+5cF7y8hCFOChCMYxXJJYR9ZAsqenUnJTuJTuZp5zVeUvCPJe4vEkuCR8XFxRt6LrV6tOlTjzlOSil5sw9xr3rFjTbWdx2Vaj9XS97WZeSfiIjbEzEM22kYm8122hOVGzjK9rx4NUmtyL/Om+C8OL7jCX9aVTH8rXzmp/FtqacYS7lBZlPzz4FzZ2OpXUIwpUIabark6kE6mPzYLhHzz4EunXlrm8z2h8ajXnVgp6vdRVKbxG2pZUJvs+dUfdy7iva2OoXsIpRemWiWFwXrmuxLlBe99yMrp2kWVlU9dCEqtw1h16r3qjXj1LuWEZBDf0lFPta6dYWthRdO1pKCk8ylnMpvtk3xb72XQBFsAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAgCR5kAASAAAAAAAAAAAAAEDgBPuBBIEAkAAAAAAAAAAQAJIAAAkAEAAAAAAAAAQBIIAEjzAAAAAAAAAAhxTTTS4mJudBtJTdWzlUsar4uVB4jJ98H7L92TL4IEdmJiJa3c2mp0IOnc2dHUbd83RSUsd9OXB+T8i2s60KdR0tMv6trUXF2tZNpfUn7SX0WjbeBb31laXtP1d1b060erejlrwfU/AztGafTF0dauKHs6hZvC51bbM4+cfjLyyZSxvrS9pudpc06yXNRfFeK5oxVfQ7ijx06+mkuVG5TqR8pfGXvZir+nKlUU9T0yrRlHlc0W5pfXjiUV4ozqJR6rR5bkmn1kmqWN/ewpqpZ31K+o9Ua7y/BVI/vTMjR2gt4tRv6Fayl86a3qb+vHgvPBjplKLwzKSEoRksNJ+J8Ua1KtTVSjUhUhLlKLyn5lQimt6llZzy52tFvt3EfH8nWi+JTlD6M5L7mXYMmlr8Cp/JrXK/8aX8SPgS/6Vdf+ay7AY0tPgUf+k3T/wDFY+AUX8arcS8a8v4l0SDULVafZ5y6Kl9JuX3n3Ttban8ShSj4QSKw4BnRhdwwuYyslnqGqWNhj4Vc06cn8WGcyl4RXF+QYmYhekN4ZgK2tXlbhY2Xq4v8rdPd81Be0/PBjL2tTq1FS1LUK13VlytqSaT/APDhxa+k2Z6ZQnJHw2G61qwoVHRhUdzXXOlQW/JeOOEfNosq17qtwm4ep0+lzbeKlTH7Mf7xb2drqdakoWtlR02h1Osk5eVOPBeb8jIUNBtW1K+qVb+a44rv2F4QWI+9Mz2hj9VmFp/Bq1dSoUrjV7mP5T8Yov6TxCPkZGlpepXT3ru5haU3+Tt/aqec2sLyXmZ6EIQgoQioxXJJYSPoxtmKfay0/S7GwzK3oJVH8apJuU5eMnxZeJYDBhOIiEggcAykAAAAAAAAAAACAJBAAEgARgkAAAAAAAAAACABPAEcCQIBIAAAAAAAAAAgASCOBIDgCCQAAAAAAAAABAEggkAAAHAAAAAAAAAAAAQSABBIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACPIkAAAAAAAAAAAAAAAAAB5AAAAAAAAAAAAAAAEAkAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAIaySAMZe6Hp11Udb1Pqa7/LUHuT82ufnkx9XStUtcuhWpX0Pm1fwdT9JLdfmkbGGZiZhGaRLS0re2r5qK50e4k/jfi4yfis05GVpX2q2yXrFQvqfU1+Cqfvi/sM7UpwqQcKkVKLWGmspmJraBaJudjUq2E3xxQfsPxg8x9yRnf2j0zHhVttbsatSNGtOdpWfBU7iO434Pk/JsyWUazdWmrUKcoV7WjqNB83RxGWO+EuD8n5FjaVaVOr6qxvbjT6v/R55S/8ALn/s4GjrmPLdMkmvUdW1K34XNrSu4rnO3luS/Ql+6RdQ2i0zH4apVt5fNrUZRf3YfkY1LMXiWXI5dRiKmvUpr+pWl1dN8nuerj754+zJY3V7qlWDnWurfT6PX6r25JfTkkl+iNSTeIbBc3FC2pOrcVqdKC+VOSivtMbV1uNRYsLWrddk5fg6f6T4vyTMFa+qrVlUsrS41OsuVeT3kv8AxJ8F9Uy9HTNUuON3d07WD+Rbx3p+c5LHuiZ1EI9U28LTULm7lT9ZqGpwtKT+RRfq0+7fftPywW9jb1Zt/wAl6ZP2udesnTUvFyW/L3eZsNjpFhaVPWwob9f+2qtzqP6z4ryL/A6teCMe/LBW+g1avtahezkuulb5pw83nefvXgZWysLOyhuWltSox69yOM+L6y5BjcpxWIMEEgwkcgAAAAAAAAAAAAAAAB5AACCQBBIAAAAAAAAAAAAQPIkAAAAAAAAAAAAAADyAAAjyJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAoXdpbXdL1Vzb0q0PmzimvtK4BphK2z1KOfgN5cWv5jfrKfulxXk0W70zWoPEZWFVfOzOm/d7X3mxgz1ShNIa/S0nVKn4+9t6EetUabnL9KXD7C7ttB06nNVK1OV1VXKdxLfx4J8F5JGVA3LMUiERjFJJJJIkAwkAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADIAAAAOAAAAAAAAAAAAAAAABAEgAAAAAAAAAAAOAAEACQCAJAAAAAAAAAAAAgCQAAAAAAAAAAAAADJAEgAAAAAAAAAAAAAAAAgkAAAAAAAAAAAAGQAAHAAGAAAAAAAAAAAAAAAAAAAAAAAAAAAAADIAAAAAAAAAAAAAAAAAAgkAAQBIAAAAAAAAAAAEASCOBIBAAAAAAAAAAAAAAIJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAB5EEgAAAAAAAAAAAAAAADiBAJ8gAAAAAAAAAAAAAAQCfIAAAAAAAAAAAABBIAgkACOskAAAAAAAAAAAAQBAAE+Q8gAAAAAAAAAAAAAAQT5DyAAAAAAAAAAAAAABA8iePYA8gAAAAAAAAAAAAAAAQCQAAAAAAAAAAAAAACCfIAAAAAAAAAAAAAAAgkeQAAAAAAAAAAAAAAI8iQAAAAAAAAAAAAAAAAAAAAZAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADrAAAAAAAAAAAAAAAAAAAAAPIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAGAAyhlAAMrtHmAAyMgAMoZQAAAAMrtGQAGQAAAAAcAAAyAAHAAAAAGUAAGUMgAMoAABlAAAAAyu0eYADKAADIAADgAAAADK7RkABkAAAAA4DKAADzAAAABkZAAZHmAAGUAA4DIADgAAHmAAAygAHAZAAAAAAAGRlAABlAABwAAZQygAA8wAAygAGRwAADKAAZXaOAADK7RlAAAAAygAAyu0AAOAAAZQ4AAMrtAADIAAZQ4AABnvAADKAADIADKAADzAADIAAAABkZQAAAAAAAyOAAAcAAAAADK7QAAyAA4DgAAyhkABlAAAOAADgOAAAAAAAAAADKGUAAyPMABw7RwAAcBlAAMjIADgMoAAMoAAMgAM95DfBgf/Z" style="width:32px;height:32px;object-fit:contain;" alt="NX"></span><span>Nexo</span></a>
        <h1>Panel Admin <span class="badge-admin">ADMIN</span></h1>
      </div>
      <a class="btn" href="/">← Volver al chat</a>
    </div>

    <div id="msg" class="msg"></div>

    <div class="stats" id="statsBar">
      <div class="stat"><div class="stat-val" id="statTotal">—</div><div class="stat-lbl">Usuarios totales</div></div>
      <div class="stat"><div class="stat-val gold" id="statGold">—</div><div class="stat-lbl">Primeros 30 (dorados)</div></div>
      <div class="stat"><div class="stat-val" id="statDev">—</div><div class="stat-lbl">Plan Developer</div></div>
      <div class="stat"><div class="stat-val" id="statBeta">—</div><div class="stat-lbl">Plan Beta Tester</div></div>
    </div>

    <div class="legend"><div class="legend-dot"></div> Los primeros 30 usuarios registrados aparecen en dorado</div>

    <table>
      <thead>
        <tr>
          <th>#</th>
          <th>Usuario</th>
          <th>Plan actual</th>
          <th>Registrado</th>
          <th>Cambiar plan</th>
        </tr>
      </thead>
      <tbody id="usersBody"></tbody>
    </table>
  </div>

  <script>
    const PLANS = [
      {key:'gratis', label:'Plan Gratis'},
      {key:'beta_tester', label:'Plan BETA Tester'},
      {key:'developer', label:'Plan Developer'},
    ];

    function planBadge(key) {
      const p = PLANS.find(p=>p.key===key) || {label: key};
      return `<span class="plan-badge plan-${key}">${p.label}</span>`;
    }

    function fmtDate(iso) {
      if (!iso) return '—';
      try { return new Date(iso).toLocaleDateString('es-ES', {day:'2-digit',month:'short',year:'numeric'}); }
      catch { return iso.slice(0,10); }
    }

    function showMsg(text, type='ok') {
      const el = document.getElementById('msg');
      el.textContent = text;
      el.className = 'msg ' + type;
      el.style.display = 'block';
      setTimeout(() => el.style.display='none', 3500);
    }

    async function loadUsers() {
      const res = await fetch('/api/admin/users');
      if (!res.ok) { showMsg('Error cargando usuarios', 'err'); return; }
      const users = await res.json();

      document.getElementById('statTotal').textContent = users.length;
      document.getElementById('statGold').textContent = Math.min(30, users.length);
      document.getElementById('statDev').textContent = users.filter(u=>u.plan==='developer').length;
      document.getElementById('statBeta').textContent = users.filter(u=>u.plan==='beta_tester').length;

      const tbody = document.getElementById('usersBody');
      tbody.innerHTML = '';
      users.forEach(u => {
        const isGold = u.registration_order <= 30;
        const tr = document.createElement('tr');
        if (isGold) tr.className = 'golden';

        const selectId = 'sel_' + u.id;
        const btnId = 'btn_' + u.id;
        const planOptions = PLANS.map(p =>
          `<option value="${p.key}" ${p.key===u.plan?'selected':''}>${p.label}</option>`
        ).join('');

        tr.innerHTML = `
          <td class="order ${isGold?'golden-text':''}">${u.registration_order}</td>
          <td class="username-cell ${isGold?'golden-text':''}">${isGold?'<span class="crown">👑</span>':''}${u.username}</td>
          <td>${planBadge(u.plan)}</td>
          <td class="date-cell">${fmtDate(u.created_at)}</td>
          <td style="display:flex;gap:8px;align-items:center">
            <select class="plan-select" id="${selectId}">${planOptions}</select>
            <button class="apply-btn" id="${btnId}" onclick="applyPlan('${u.id}','${selectId}','${btnId}',this.closest('tr'))">Aplicar</button>
          </td>
        `;
        tbody.appendChild(tr);
      });
    }

    async function applyPlan(userId, selectId, btnId, tr) {
      const sel = document.getElementById(selectId);
      const btn = document.getElementById(btnId);
      const newPlan = sel.value;
      btn.disabled = true;
      btn.textContent = '...';
      try {
        const res = await fetch(`/api/admin/users/${userId}/plan`, {
          method: 'POST',
          headers: {'Content-Type':'application/json'},
          body: JSON.stringify({plan: newPlan})
        });
        const data = await res.json();
        if (!res.ok) throw new Error(data.error || 'Error');
        showMsg(`Plan actualizado correctamente`, 'ok');
        // Actualizar badge en la fila
        const planBadgeCell = tr.querySelectorAll('td')[2];
        if (planBadgeCell) planBadgeCell.innerHTML = planBadge(newPlan);
        document.getElementById('statDev').textContent = document.querySelectorAll('.plan-developer').length;
        document.getElementById('statBeta').textContent = document.querySelectorAll('.plan-beta_tester').length;
      } catch(e) {
        showMsg(e.message, 'err');
      } finally {
        btn.disabled = false;
        btn.textContent = 'Aplicar';
      }
    }

    loadUsers();

    // ═══════════════════════════════════════════════════
    //  NEXO MEJORAS — GPU Monitor, Voice, Personality
    // ═══════════════════════════════════════════════════

    // --- Prism.js dinámico ---
    (function loadPrism() {
      if (window.Prism) return;
      const s = document.createElement('script');
      s.src = 'https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/prism.min.js';
      s.onload = () => {
        const langs = ['python','javascript','typescript','bash','sql','json',
                        'css','html','cpp','csharp','go','rust','java','lua'];
        const base = 'https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/components/';
        langs.forEach(l => {
          const sc = document.createElement('script');
          sc.src = `${base}prism-${l}.min.js`;
          document.head.appendChild(sc);
        });
      };
      document.head.appendChild(s);
    })();

    // --- GPU / System Monitor ---
    const gpuWidget   = document.getElementById('gpuWidget');
    const gpuBar      = document.getElementById('gpuBar');
    const cpuBar      = document.getElementById('cpuBar');
    const gpuLoadLbl  = document.getElementById('gpuLoadLabel');
    const vramText    = document.getElementById('vramText');
    const cpuLabel    = document.getElementById('cpuLabel');
    const queueWidget = document.getElementById('queueWidget');
    const queueCount  = document.getElementById('queueCount');
    const gpuNameEl   = document.getElementById('gpuName');

    async function pollSystemStats() {
      try {
        const r = await fetch('/api/system-stats');
        if (!r.ok) return;
        const d = await r.json();
        if (!gpuWidget) return;
        gpuWidget.style.display = 'block';
        if (d.gpu_name && gpuNameEl) gpuNameEl.textContent = '💻 ' + d.gpu_name.replace('NVIDIA ','');
        if (d.gpu_load != null) {
          gpuBar.style.width = d.gpu_load + '%';
          gpuBar.style.background = d.gpu_load > 90 ? 'var(--danger)' : d.gpu_load > 70 ? '#f0c040' : 'var(--accent)';
          gpuLoadLbl.textContent = d.gpu_load + '%';
        }
        if (d.vram_used != null) {
          vramText.textContent = `VRAM ${d.vram_used}MB / ${d.vram_total}MB`;
        }
        if (d.cpu != null) {
          cpuBar.style.width = d.cpu + '%';
          cpuLabel.textContent = 'CPU ' + d.cpu + '%';
        }
        const q = d.queue_size || 0;
        if (queueWidget) {
          queueWidget.style.display = q > 0 ? 'block' : 'none';
          if (queueCount) queueCount.textContent = q;
        }
      } catch(e) {}
    }
    pollSystemStats();
    setInterval(pollSystemStats, 5000);

    // --- Voice-to-Text (Web Speech API) ---
    const micBtn = document.getElementById('micBtn');
    if (micBtn) {
      const SpeechRecognition = window.SpeechRecognition || window.webkitSpeechRecognition;
      if (SpeechRecognition) {
        const recog = new SpeechRecognition();
        recog.lang = 'es-ES';
        recog.interimResults = true;
        recog.maxAlternatives = 1;
        let isListening = false;
        recog.onresult = (e) => {
          const transcript = Array.from(e.results).map(r => r[0].transcript).join('');
          els.prompt.value = transcript;
          els.prompt.style.height = 'auto';
          els.prompt.style.height = Math.min(180, els.prompt.scrollHeight) + 'px';
        };
        recog.onend = () => {
          isListening = false;
          micBtn.textContent = '🎤';
          micBtn.style.color = '';
        };
        recog.onerror = (e) => {
          isListening = false;
          micBtn.textContent = '🎤';
          micBtn.style.color = '';
          if (e.error !== 'no-speech') setStatus('Mic: ' + e.error);
        };
        micBtn.addEventListener('click', () => {
          if (isListening) { recog.stop(); return; }
          isListening = true;
          micBtn.textContent = '🔴';
          micBtn.style.color = 'var(--danger)';
          recog.start();
        });
      } else {
        micBtn.title = 'Tu navegador no soporta Web Speech API';
        micBtn.style.opacity = '0.4';
      }
    }

    // --- Personality selector ---
    const personalitySelect = document.getElementById('personalitySelect');
    // El valor se lee en sendMessage() y se incluye en el payload.
    // Patch del sendMessage existente para inyectar personality:
    const _origSend = window._nexoSendOverride || null;
    // Interceptamos el fetch de /api/chat/stream mediante monkey-patch del JSON payload
    const _origFetch = window.fetch.bind(window);
    window.fetch = function(url, opts, ...rest) {
      if (typeof url === 'string' && url.includes('/api/chat/stream') && opts && opts.body) {
        try {
          const parsed = JSON.parse(opts.body);
          if (parsed && personalitySelect) {
            parsed.personality = personalitySelect.value || 'normal';
            opts = { ...opts, body: JSON.stringify(parsed) };
          }
        } catch(e) {}
      }
      return _origFetch(url, opts, ...rest);
    };

    // --- Rate limit feedback ---
    // El error 429 ya llega como JSON {error:'...',retry_after:N}.
    // El handler existente de errores en el fetch lo mostrará via setStatus().
    // No se necesita código extra aquí.

    // --- Scroll forzado al cargar chat ---
    _autoScroll = true;
  </script>
</body>
</html>
"""

PLANES_HTML = r"""
<!doctype html>
<html lang="es">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>Planes - Nexo</title>
  <style>
    :root {
      color-scheme: dark;
      --bg: #161616; --panel: #202020; --panel-2: #292929;
      --text: #f4f4f4; --muted: #b8b8b8; --line: #3a3a3a;
      --accent: #19c37d; --accent-2: #2dd4bf; --danger: #ff9a9a;
    }
    * { box-sizing: border-box; margin: 0; padding: 0; }
    body { background: var(--bg); color: var(--text); font-family: Inter,ui-sans-serif,system-ui,-apple-system,"Segoe UI",sans-serif; min-height: 100vh; }
    .page { width: min(980px, calc(100vw - 32px)); margin: 0 auto; padding: 42px 0 60px; }
    .topbar { display: flex; justify-content: space-between; align-items: center; padding-bottom: 24px; border-bottom: 1px solid var(--line); margin-bottom: 36px; }
    .brand { display: flex; align-items: center; gap: 10px; color: var(--text); text-decoration: none; font-weight: 800; }
    .mark { display: grid; place-items: center; width: 32px; height: 32px; border-radius: 6px; background: #f2f2f2; color: #111; font-size: 13px; font-weight: 800; }
    .btn { display: inline-flex; align-items: center; justify-content: center; gap: 6px; min-height: 42px; border: 1px solid var(--line); border-radius: 6px; background: var(--panel-2); color: var(--text); padding: 0 16px; cursor: pointer; font-size: 14px; font-weight: 600; text-decoration: none; }
    .btn:hover { background: #323232; }
    .btn.primary { background: var(--accent); border-color: var(--accent); color: #06140e; }
    .btn.primary:hover { opacity: .88; }
    h1 { font-size: clamp(28px,5vw,46px); font-weight: 800; line-height: 1.08; margin-bottom: 10px; }
    .subtitle { color: var(--muted); font-size: 17px; line-height: 1.6; max-width: 620px; margin-bottom: 40px; }
    .current-plan { display: inline-flex; align-items: center; gap: 8px; background: rgba(25,195,125,.1); border: 1px solid rgba(25,195,125,.3); border-radius: 6px; padding: 8px 14px; font-size: 13px; color: #6ee7b7; margin-bottom: 32px; }
    .grid { display: grid; grid-template-columns: repeat(auto-fit, minmax(260px, 1fr)); gap: 16px; margin-bottom: 40px; }
    .card { background: var(--panel); border: 1px solid var(--line); border-radius: 12px; padding: 24px; display: flex; flex-direction: column; gap: 16px; position: relative; transition: border-color .15s; }
    .card:hover { border-color: #555; }
    .card.current { border-color: var(--accent); }
    .card.highlighted { border-color: var(--accent-2); box-shadow: 0 0 0 1px var(--accent-2) inset; }
    .card-badge { position: absolute; top: -1px; right: 16px; background: var(--accent-2); color: #06140e; font-size: 10px; font-weight: 800; padding: 3px 10px; border-radius: 0 0 6px 6px; letter-spacing: .06em; }
    .card-name { font-size: 18px; font-weight: 800; }
    .card-price { display: flex; align-items: baseline; gap: 6px; }
    .price-val { font-size: 36px; font-weight: 800; }
    .price-period { color: var(--muted); font-size: 14px; }
    .features { display: grid; gap: 8px; }
    .feat { display: flex; align-items: center; gap: 8px; font-size: 13px; color: var(--muted); }
    .feat-icon { color: var(--accent); font-size: 14px; flex-shrink: 0; }
    .feat strong { color: var(--text); }
    .card-action { margin-top: auto; }
    .tag-current { display: inline-flex; align-items: center; justify-content: center; width: 100%; height: 42px; border-radius: 6px; background: rgba(25,195,125,.1); border: 1px solid rgba(25,195,125,.3); color: #6ee7b7; font-size: 13px; font-weight: 700; }
    .contact-box { background: var(--panel); border: 1px solid var(--line); border-radius: 10px; padding: 24px; display: flex; flex-direction: column; gap: 12px; }
    .contact-box h2 { font-size: 18px; font-weight: 700; }
    .contact-box p { color: var(--muted); font-size: 14px; line-height: 1.6; }
    .contact-methods { display: flex; flex-wrap: wrap; gap: 10px; margin-top: 4px; }
    @media (max-width: 600px) { .grid { grid-template-columns: 1fr; } }
  </style>
</head>
<body>
  <div class="page">
    <div class="topbar">
      <a class="brand" href="/"><span class="mark" style="background:transparent;display:inline-flex;align-items:center;justify-content:center;"><img src="data:image/png;base64,/9j/4AAQSkZJRgABAQAAAQABAAD/4gHYSUNDX1BST0ZJTEUAAQEAAAHIAAAAAAQwAABtbnRyUkdCIFhZWiAH4AABAAEAAAAAAABhY3NwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQAA9tYAAQAAAADTLQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAlkZXNjAAAA8AAAACRyWFlaAAABFAAAABRnWFlaAAABKAAAABRiWFlaAAABPAAAABR3dHB0AAABUAAAABRyVFJDAAABZAAAAChnVFJDAAABZAAAAChiVFJDAAABZAAAAChjcHJ0AAABjAAAADxtbHVjAAAAAAAAAAEAAAAMZW5VUwAAAAgAAAAcAHMAUgBHAEJYWVogAAAAAAAAb6IAADj1AAADkFhZWiAAAAAAAABimQAAt4UAABjaWFlaIAAAAAAAACSgAAAPhAAAts9YWVogAAAAAAAA9tYAAQAAAADTLXBhcmEAAAAAAAQAAAACZmYAAPKnAAANWQAAE9AAAApbAAAAAAAAAABtbHVjAAAAAAAAAAEAAAAMZW5VUwAAACAAAAAcAEcAbwBvAGcAbABlACAASQBuAGMALgAgADIAMAAxADb/2wBDAAUDBAQEAwUEBAQFBQUGBwwIBwcHBw8LCwkMEQ8SEhEPERETFhwXExQaFRERGCEYGh0dHx8fExciJCIeJBweHx7/2wBDAQUFBQcGBw4ICA4eFBEUHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh7/wAARCAMABYADASIAAhEBAxEB/8QAHQABAAEFAQEBAAAAAAAAAAAAAAEDBAUGCAcCCf/EAF0QAAIBAwEEBgQICQgDDQgDAQABAgMEEQUGEiExB0FRYXGBCBMikRQyQlJygqGxIzNDYpKUssHRFRYXJFOi0uFjs8IlNDU2RUZUVmR0g5PwGCZEVXN1hKNlw/Hi/8QAGwEBAAIDAQEAAAAAAAAAAAAAAAIFAQMEBgf/xAAtEQEAAgIBAwQABgIDAQEAAAAAAQIDEQQSITEFE0FRFBUiM1JhIyQyNEJxgf/aAAwDAQACEQMRAD8A7LAAAYAAEYJADAwAAwMIAAMAAAAAGAAAwAAwMAAMAAAMAAMDAAAYAADAADAwAAwMAAMDAADAAAAAAMAABgABggkAMDAAAYAAAABgYAAAAAMAAQTgAAAAGBgAAAAAAAYGAAAAAYGEAAGAAAwAAwAAGEMAAMDAAAYAADAADAwAAwMAAAAAAADAwAAwAAGBhAAMAABgYAAYGEAAAADAAAAAAAAAwAAGO4ABgYAABIABgYAAAABgYAAgnAAAAAMDAADAwAAAADAwAAwAAAwAAwMIABgjBIADAADAwAAwMAAAAAGAAGBgAAMAAMDAADAx3AAMDAAEYJwAAwMAAMIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAYAAAAAAAAAAAAAAAYAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAD5AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAwAACAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAABDJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAABxAAcQAAAAAAAAAAAADkAAAAAAAAAAAAAAAAAAHEAAAAAAAAAAAAAI4k8QBHEkAAAAAAAAAAAAAAAAABxAAAAAAAAAAAAAAAAAAAAAAAAAAAAABxAAAAAAAAAAAAAAAA4gAQSAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACOIAEkEgCCQAAAAAAAAAAAAcQAAHEAAAAAAAAAAAAAAAAAAAAAAAAAAGAAAAAAAAAAAAAAAAAAwAAAAAAAAAAAAAAAAMAAMAAAAAAAAAAAAAHAAABgAAAAAAAAAAAAAAADAAAAAAAAAAAAAMAABgAABgAAAAAAAAAAAAGCCQAwQkSAAAAAAAAAABDYE4HMpV69GhTdStVhSgucpySSMdLaHSE2o3iq466UJTXvimNMdUMtgGOtNb0u6moUb6i5vlCUt2XufEyCawCJiUgAMgAAAAAAABGCQAwAAAAAAAAAABHWSAAAAgkAAAAAAAAAAAAAAADAAAAAAAAAAAAAAAA4AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAGAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAHxUqQpwc6kowjFZcpPCRiqu0FnKThZRq381wfqFmC8ZvEftM6Rm0Qy+UU7m4oW9J1a9WnSgucpyUUveYGrealXi5Va9DT6S5+r9ueO+csRXuZZ0FZVavrLS1uNUrr8q/wiT+nP2V5GelH3N+GWq6/QmsWFvXvX86Ed2n+lLC92SxuLzVK0XOveULCiuaorLXjOfD7C5p6bqdzxr1qNnH5tJetn+lJYXuZdW+hadSmqlSlK5qrlUuJOo14Z4LySG4hHV7NepU7S4qqdtaXGqVl+Vlma/Tn7K8jJQs9anFYpWVCPVGVWUmvckjYlFJJJYSJE2IxR8tYurS+cHG80yldQ/0U1P8AuzS+xlrRqUaNRUrK+ubCr1UJ5S/Qn/sm4FK5tre6pOlcUadWD5xnFNfaIse1rxLC0dW1Kg925tad3FfLt3uT/Ql/Ev7XWtOuKipOt6ms/wAlWi6cvJPn5FCehUqfGxuq9t+Y36yn+jLl5NFvc299Gk6d3YUb6j1+pay/GE/3NjtJHVHln8onJqdq7elUVLT76vY1f+j1M4/8ufV9HBf09V1Cg926so3MV8u2eJecJfubHTKUZI+WdBj7HWLC7mqVOuo1v7KonCa+q+JfriRTiYnwkABkAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAgASAPMAAAAAAAAAAAAIJAAgkAAAAAAAAAAAABAEggnzAAAAAAAAAAAAAAABAEgEAOJJBIAAAABkABlFveXlrZ0fXXVenRpr5U5YyDa4bIzgwtbWa9fhp1nJxfKtc5pw8o/GfuXiY2+nGc1DVdRq15S5W1LMFLwhHMpebZmKoTePhmrzW9Ptqro+udeuvyVCO/PzS5eeDH19R1a5T9RTpWFP51T8LU9y9le9nxa2d9UpqlZWFHT7fqlWST8qcf3teBe0dAtpYd9VrXz+bUaVP9BYXvyZ7QjPVZgJfBrqvicrnWLiL5fjIxfgsU4+ZlaOnatcJKcqFjS5JJetqL7or7TP06VOlCNOlCNOEeCjFYS8j7HWRi+2Lt9CsYSVS4jO8qLlO4lv48F8VeSMpGKjFJLCXUggR22RWI8JBGUTkMgIcornJLzPl1afXUh7wPvqB8KpB8px96JynxT+0D6I5kNdhOe8Cjd2tvdU/V3NCnWh82cU19pjKuh+rWbC8r26/s5v1lP3S4ryaM0DMTMIzWJarf2t5ubl/ptO8pr5dFb+O/cfFeTZRsa9SEnHTdUqRcedvXzUUfGMsTj7zb8Itb7TrO9ildW9Orjk2vaj4PmvIz1Nc4teGNo63cUeGoWE1Fc6tv+Ej5x+MvczKWN/Z30HO0uadZLnuy4rxXNeZjKui3NDjYX0mlypXK315SXtLzyY69oqFRVNS02pRnHlc0cyS+vH2l5pDUSRa1fLbQa1Z3d9TpqpaXtK/odUaz4+VSP70y/o67bJqF9TqWM28L1y9h+E17P2pmNSnF4llgfKkpJNPOeKJy+8wmkEJ5JAABAAAAAAAAAACAJBBIAAAAAAAAAAAAAABBIAEEgAAAAAAAAAAAAIAkAAAAAAAAAAAAAAAAAjzAkAAAAAAAAAAAAAA4gAAAAAAAAAAAAAAAIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAIwBJBIAAAAAAAAAAAAAAAAAgkYAAAAAAAAAAAAAAAIJAEEgAAAAAAAAACCQAIJAEDBIADAAAAAAG8Ftf31rY0HXu68KUM4Tk+LfYlzb7kDelyWWo6nZ2CXwisozl8SnFb05+EVxZirnUNQveFunYW7/KVIp1pLui+EF3vL7kY6yiqlSS0m2ld1JPFS6qTe43+dUeXLwjnyJRX7apyfEMhcalqF2mqMY2FL59RKdXHh8WPnnwLC09VWuHUsaFXUrlcHcSlvJf+JLgvCPuMpabPwniep1neS5+q3d2ivq/K+s35GbpwhThGEIqMYrCSWEh1RHhiKWnyw1HSbuu96+vPVRfOla5j75v2n5bpk7Kws7KDja29Olnm4rjLxfN+ZccC3rXttSn6uVRSqf2cFvS9yMbmWyKxC4wgWnrb2t+KoRoRfyqzy/0V/EfAvWf75uK1btjndj7l+/JjTO1Wtd21KW7UrQUuqOePuKfwuU/xNrXqd7juL+9grUaFGgt2jRp01+bHBUlKMYtyaSXW2BaKWoTxinb0V+dJyf2Y+8K3vJ/jL1rup00vvyfUtQs08KvGb7IZk/sPn4ZOX4qzuZ97io/e0GOx8BT+PdXU/wDxWvuwP5Ns38anKf0pyf3seuv5fFtKUPp1v4JhR1GXy7WHhGUv3oHZK06xX/wlHzifSsLL/otD9BHx6m/fO8pLwo/5j1F5133/AOpBl9uwsnztaP6CPn+TbHOVbU0+5YIVC8/6cn40UT6q/i+FzRl40X/iAj+T7ZfE9bD6NaS/ePgU48ad7cx8ZKX3pjOoxfxbWa8ZR/iPX3kX7djvf/Tqp/fgMdj1V9F+zdU5rsnS4+9P9w9bew+Pawmu2nU/c0vvHw6EfxtG4pfSpNr3rKKtG8tazxTuKcn2KSz7gPj4dSjwrQq0e+cGl71w+0r0qtOrHepTjOPbF5R9PBQq2dtUlvujFT+dH2Ze9cQyuBhFp6i6pL8Bdb6+bWW99q4/eQ7upSf9ZtqkF8+n7cfs4/YDajd6NY16jrRhK3rPj62hLck/HHCXmmWNex1O3hKOKWoUWsNcKdTHh8WX90zdG4o1o71KpGce2LyVFxMxMwjNIlp9s40azhp11VsKy4u1qR9n/wAt9XfFoylHXZUPZ1S2dFL/AOIpZnS8+uPmsd5lL6xtL2l6u6oQqxXFZXGL7U+afejE3GmXtqnKzq/Cqa/JV5Ymvoz6/CXvJbiUNWr4Zq3rUa9KNWjVhUpyWYyjJNNdzKvmahbKEbqTs51dNvfjToyjje73DlJfnR95lLfW5UWqWqUlbZeI3EXmjJ+PyX3S97MTVKuSJ8s0SRGSkk0+ZJFsAAAAAAAcADI8yQBBIAAAAAAAAAAAAAABBKAAAAAAAAAAAAAAAIJAAgkAAAAAAAAAAAOAAgkAQSAAAAAAAAAAAAAAAMAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQ5Jcy21G/trC3de5qbkc7qSWZSk+Silxb7ka5qV1WvaM6upT+BWEeLoOeJSX+kkv2F5t8jMRMoWvEMjea26k5UNLhCvNPEq8vxNN+K+M+5ebRjI4V88KtqWpYw3wzTT/u04/a+8uLGyvL+EFGM9OsUsRSju1prsS/Jx/vfRNgsrS2srdULWjGlTTziPW+tt9b72Z8IxE27yxdvoruMVNVnGt1q2p5VFePXPz4dxmI04QgowioxSwklhJH1yLWpeJ1HStoOvVXBqL9mP0pdXhz7iO5lsiIhc5S58C1nfKb3LSnK4l2x4QXjLl7siNpKs1K9qet6/Vx4U15dfn7i6jGMYpJKKS4JLkDvKzVtc13m7uHGP9lRbivOXN/YXNvQo0IblGlCmuyKwUKl/Tc3Ttoyuai4NU+S8ZckfPqLyuv6xX9TB/k6HPzk+PuwGFevc0KCXrasIt8k3xfgubKLurirwtrSbXzqr3I+7i/sK1va0LfLpUoxk+cucn4t8WVwz3WSoXlX8bdqmuuNGGPtef3Ex021T3pwdaXbVk5/eXgBp8whCEcQiorsSwTgkBkI5kgAkAAGAAAAADBSrW9CssVaNOf0oplUAWnwCnHjQq1qD/Mnle55R841ClylRuYrqa3Je/ivuL0BjSzV9CHC4pVLd9s17P6S4FzCUJxUoSUovk08o+mk+ZbTsaO+50t6hN85Unu58Vyfmgd017KhVn6zdcKnVUg92XvXPzKf9dt/m3UF4RqL9z+wlyvbf49NXMO2Hsz9z4Mq211QuMqnP2l8aEliS8U+IEW93RrS3FJxqLi6c1iS8ivz4Mp3FvSrpKpBSxyfJrwfUUVG5t/it3FNdTwprz5PzB3h931la3tL1d1RjUiuKzzi+1Pmn3oxFzY3tnFum5X1tjDhLHrYrz4TXjh+Jmbe4pV09yXtL40WsSj4rqKw3piaxLVbCVS2j6zR60ZUU8TtKragn2Rzxpvu5dyM3puqW97KVLEqNzBZnQqLE49/eu9ZRGpaXRu5+vhKVvcpYVanzfdJcpLuf2GBvYSjVp0NUp+oqqX4C6pNxi5fmy5xl+a+feS7S1zM0bcga/Z6xWs5Ro6s06TeIXiWI+E18l9/J93Iz8WmsriYmNNlbRZJGOJIMJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAGQ3jmgJ5GL1bVo2s1bW1P4ReSWY0k8KK+dN/Jj9r6ky31LValWtUstMcXUg92tcNZhRfYl8qfdyXX2PH2lObqzstMhv1t7Nxc1XvKEu2b+VPHKK5deFglENdrfEPie/G7jUruV9qdVP1cILG6uvdXKEO2T597wjK6bou7VheajKFe5jxpxS/BUfop83+c+PhyL3TNPoWNJqDlUqz41a0+M6j7W/uS4LqLzkJt9FcffciWEUbq5pW8E6kuMniMUsyk+xLrKNe5q1KkqFnFTnF4nUl8Sn49r7l54KlrZ06MnVlJ1a0liVWfPwXYu5EU/8A4oqlc3nGu5UKL/JQl7TX50ly8F7y8o06dGmqdKEYQXBRisJFK6u6NthSblOXxKcFmUvBfvKCoXF3xu5eqpdVGnLn9KXX4Lh4gfdS9UqrpWtN3FRcHuvEY/Sl1eHF9x8qyqV3vX1b1i/sYezTXj1y8/cXdKlClTjTpwjCMeSisJH2DX2+acIU4KEIKMVwSSwkfQAZAAAAAAAAAAAAAAAAAAAAAAAAAAAKFza0LjHrKabXxZJ4lHwa4orgCyxd2vJu6p9jaVRefJ/YVra5o3Capy9qPxoNYlHxXUV2s8y3ubWlXxKScZx+LOLxKPgwxrRc21Ku1J5jUj8WpB4kvP8AdyKLuK9pwu479LqrQXL6S6vFcPAOtcWj/rKdWl/bQjxX0o/vX2F3CcKsFOElOMlwaeUwJhOM4qUWmmsprkz4uaFK4oypVqcKlOaxKMllNeBbytqlvJ1LJpLnKi3iEvD5r+z7ytbXNOvvJb0Kkfj05LEo+P8AEHnywl1p9xYKXqYzu7FrEqT9upTXdn48e58fHkUNPr1tPpxq6fJ3unvj6iLzKmv9G31L5j8scjaOZidR0uSqzvNOcadeXGpTk8U63j2S/OXnklE78tc013hkLG7t722jcW1WNSnLk19zXU+4rmp0aklc1LmxzbXkWvhNvV4KfdNLr7JrPmjPaVqdC/hJRUqVenwq0Z/Hpvv7U+prgzEwlW8SvgAYTAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADAAAAAAAAAAAAAAAAAEAcABIAAAAAAAAAAAACCQBHADgSBBIAAAAAAAAHUADBAADgABIAAAAAAAAAAAACBwAAAkCCQACAAAAAAAAIJKVetSo0p1atSNOEE5SlJ4UUutgfcmksvka1qOpVdTk7bT6sqVmm1UuoPEqvbGm+ztn7u1Rf3MtVhKVSbt9KS3nGb3ZV186fzafdzfXhcH9Weny1RKVWEqOncFGnjdlcLvXyYd3N9eFwcojXdqtabdoUdOtZ3sI0LD+rWFP2XWp8HPtjT/AHz92XxWyWVrQtLaFvb0406cFiMYr/1nxKsKcIQUIRUYxWEksJLuPm5q07elKrUmoQjzbMTO0q1ir7nKMIucpKMUstt8EWWat/8AEcqNr1zXCdTw7I9/N93MilRqXklWuouFFPNOg+vvn393V48rytVp0acqlWahCKy23wRhLymlTp0aUadOKhCKwklwRaVLmrcSdKx3Wk8Sry4wj3L5z+z7iFGtf8ailRtXyhynU8exd3Pt7C9hCNOChCKjFLCilhIHlQtLSlbtzWalWXx6k3mUvPs7lwLhcsEgMo8xwHAcAIZPVkDgA4doJCQEYQJAAciG12nzKpCKy5JGPLEzEeX2DDahtRs/YZV5rFjRkvkyrR3vdnJgb7pO2St1+DvKty+yjQm/taSJxS0/DRfl4aebQ3cg8wuumTRqcsUdK1Or3uMI/fIsZ9N1ipYWhXX1q0US9qzT+Zcf+T10Hk1HpnoVGlHZ65l9CvF/uL+h0rQq4xsrrrT66dHfMe3ZOOdht4l6UEaTZ9IdtXaVTZ3aSj3y02bX2ZM9YbRadd4Sjd0ZPhu1rSrTf2xIzEw3UzUv4ZjBJ805RnHejxTPow2gAAABgAQSBAwOAAjCLSdpKlUlVspKnJ8ZU38Sfl1PvX2l5hEoC2tbqFZulOLpVorMqcua712rvRN1awr4mpOnVh8SpHmv4ruJubanXS3sqUXmE48JRfcz4j8Pgt1/B6uPlNuLflhhhFC4nGqre7ioVX8WS+LU8Ox933l2WVeneV6bp1KVq4v8+Xv5cGVLGF1Tpbl1UhOSfsyjza7+8C11nTKd6o1ac3Ruqa/B1orLXc18qL60/sfE15+td5ClXzZanRTdKpDjGpHr3c/Hh2xfFd3Bm5tcC0v9PtL+l6q7owrQTUlGS5NdaMxbXlC2PfeFDQNTeo29RzpblWjLcnKPGnN9sJda+7kZM+KdOFKEYQioxisKKWEkfZhOPHcIJI4BkAAE4RBIAAAAAAAAAEEgCAOBIEDgCQI4EgAAAAAAAAACAAAJYAAAAAAAAABgAACAABIAYAAAAAAAAAAEEgAAAAAAAAAAAAAAAAAAgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAZAAAFpquoW+nWvr67fF7sIRWZVJPlGK62GJnSpfXdvZW07m5qRp0o85P7Eu1vsNYv7mV6/hmpr4PZ05J0raXNvPBzS5yzjEFnHe+Xxc1atStC/1PLqb27bWtP2txvqivlVH1y5LjjCy3ltH0mfrYX+oqLuI5dKknmFBPs7Z9svJYXOWtNM2m86h86fptW8nC61Gm4Uk96jay7eqVTtl2R5LvfLOYCPivWp0aUqtWajCKy2yO9t0REQXFanQpOrUliK7uL7kutlrb29StVV1drDXGlS5qn3vtl93UTbUqlesru5i44/E0n8hdr/Of2cu0r3dxC2pb88tt7sYrnN9SXeAuq9O2pesqPhnCS4uT7EutltQoVbipG4vEluvNOjzUO99svu6u0+rW3qOr8Ku8Os/ixXGNJdi7+1/uLzCxwQEjgFwGAyAgngAA6gAAYXIARkPgjxfpU6V9a0HW6ug2OkOwqRWY3d1ifrI9U6cU8Nd7b70SpSbzqHPyORTj06rPYby8tbOhK4uq9KhSisynUmoxXi2aHrXS3sxaTlR0+pV1SrH/o8fYz9N4Xuyc/6zreq63W9dquo3N7POV62eYx8IrgvJFrb1HSqxqccLmu1HdXh6j9TzXK9fvPbFGnrer9Kuu3rcLCnb6fB8sL1k/e+H2Gm6vqurao277VL25z8mdZ7v6K4fYY2LSSlFpp9aPr1hKuOsfCkzc/kZZ/VZTjTdNtxe75I+ncSjBrebwJ1E+GCjOO9F8OD4GyIhz9cz5e1dG+y+xmv7NWupT0uFW4x6u4jUrTnu1I8JcG8cefg0bza7K7N2sUqGh6fTxyxbxz78HinQlr70nan+S7iri01LEI55RrJey/Nez44Og4LgivzdVbae39K9nNhiemNwpUbS1orFKhTgvzYpFXcj81H0iOs07lbRSseIFGPYhhdg8xlDaWoSBnvHmGQAZwAAIYEgAAAAAAAAZCADHHII8wJaRGExwznIz1ASQwaJ0sdIVnsbYRo0acLvVq8W6FtvYUVy359kc+b5LralWs2nUNWbNXFWbWbhqWp6fptrK61C8oWtCHxqlaahFebNI1Hpl2DtJuFPU6t41z+DW85r34SfvObNqNb1jaa+d9rV9VvJ5e5CTxTp90Icl95iFCMeUV5nfj4Uf+peezeu23/jh0/Q6cdhqlTdq1tQtov5dS0lj+7k3LZvavZ7aOm6mi6taXqj8aNOot6PjHmvNHFzUWsOK9x92nrLW5heWletbXFN5hVozcJxfc1xJ24Ea/TLXj9dyRP647O6U8rgSeGdEPS9VrXNvoG11aPrarVO21DG6pyfBQqLkpPqkuD68dfuSaa4Ffkx2xzqXouNyacinVVIANboAAAAAAAAAAAAAAAAAAAAAAAAEAAAAAAAAAAAHAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAcQAAAAAAAAAAAAAAAGAAABAMx2sapTsKcIQputdVsqjRTw5Y5tvqiut9XjhBiZ0+9W1Gjp9KLmpVKtR7tKjD49SXYv3vkus16tKr8KhdXcXc6hVzGhQpPhBdcYZ5L503/BHzmtG7U5YvtWuY8Ir2Ywjnq+ZTXbzb7WbBpGmQslKtVn6+8qpetrNY+rFfJiupeby+JLw1d7yp6RpfqKnwu7lGreSjjMfiUl82Hd2vm+vqSyoWOoPkRltiIjw+ZSUIuUmkkstvkWNCLvq0bqomqEHmjBr4z+e/3e/wAFb+v13RXG2pSxVa5VJL5Pguv3dpe1J06NGVSclGEVlt8kgeXxc3EKFJ1JvhySXFt9SS62UbOhUnV+GXSXrmsQhnKpLs8e1/uPi0pyua6vK8XFL8TTkvir5z/Of2LzL3wAldhJBIZD5bwTy7S21G+tdNsa99fXFK2tqEHUq1aklGFOKWXJt8EkuseWJnXeVdz4ZawattZ0jbEbKzdPX9pNNsay4+olWUqv6Ecy+w5i6bPSG1TaK9uNE2Huqum6LBunO/p+xcXfa4PnTp9jXtPu5HhSowqVJVJOcqk3mU3Jtyfa3zZ2YuJNo3Kr5HqMUnVXcFb0jOjGE3GlqGoXEV8unp9XD96Rt+h9I+x+rqn6jWaNGpUScadynRlx5L2sH5/6fbesuaNFSl+EnGCWX1tHrNV09582uXHijong112lW5PWsmO3jbtOnOFSClCSlFrKafM+snJOxe3Ov7J3Mf5PuHWss+3Z1pN05L835j71w7UzpPYLazTNrtHjf6fNxnF7tehP49GXZJfc+TOLNx7Y1rwvUsfJ7eJbFlmp9J2xNjtloMraoo0r6hmpZ3GONOfY+2L5NfvSNuRBpraazuHdlxVy1mtocZX1vcWF7Wsrui6NzQqOnVg+cZJ8V/mUfWRXA9S9JjZ6VpqNptPa08U7nFtdYX5RL2JPxScfKJ47Cq285LjFk667eC5nE9nLNfhnNOu8p0H1cYfwLxZfX5Gv0ZyjOMoNxknlMz1vUjVoxqLhnn3PsMWq4bx0+H11n1lLhwPiTSXI+XLtwRa1WlUlCtCpSk4ThJShJc4yXFNeDOmtgNfhtHsxa6j7KrY9XcQXyKkeEl+9dzRzBGaXZk9A6DtpI6XtPLSbiri21LChl8I1kuH6SyvFI08inVXa89D5fs5uifEvf0QE01lcgyumdPbw1fbrbrRNjatjDWfhMfh0pRoSp0t6LlFJuLfU8PguvD7DBvpf2WXyNQz/AN2f8S/6ZNkYbabB3+j7kHdKPrrKUvkV4cY+T4xfdJnGekbS6lp0/gl9TqVqdKThOlUeKlNrg0n3PhhnZgw1yV38qP1Hl8nBf9Hh14ul7Zd4e5qH6s/4n1/S3sx1w1D9Wf8AE590a8stUoestK280vapy4Tj4r96Mh6rHW/JmyeNWPKkn17k1nT3T+lvZjONzUP1Zh9Ley/XDUOH/Zn/ABPCHBdUpe8pyjjhl+8fhqM/n/Ie9Ppc2WSzuah+rP8AiR/S9sv8zUf1Z/xPA3z6/efMufN+8fhqH5/yXvr6X9luuGo/qz/iQ+mHZVfI1H9Wf8TwCXn5s+GvHl2mfwtD8/5DoF9MOyuPxeo/qz/iQ+mLZX+z1L9Wf8Tn3Da6yfNiOLRn8+5LoD+mPZT+y1L9Vf8AEPpk2U/s9T/VX/E59kl1cT4lHL68eJL8JRn895DoR9M2yf8AZan+qv8AifL6Z9k/7HVP1X/M56lD86TPmS75e8z+Eoz+e8iXQz6adklzpap+q/5ny+mzZFfktU/VH/E54nDhjL95SlDvb8zP4OjP55yHRL6cNkM/itW/VP8AMj+nHZBPjS1b9U/zOc5QS7fefO7h837x+DozHrXIdG1unLZL1M/U0NTnU3W4RdvuqT6lnPDxPAta1W91rVrnVdSqupd3U3OpLPCPZFfmpcEu4xjiu8mL6sYwbcOCuOdw5uXz8vJrqz7qcfApSgn1cSq2kuSPnK5tHTEK6NwtpLHBIKo4lWe6lmTSXazHXddy4U8xXW+tmYjTZHdWuLqMU485fcdO+jvt5ParZyrpmo1vWappbjTnJ/GrUmvYm+18Gn3rPWcoT8De/Rx1WppnS/YUVJqlqFKpa1Fng/Z34/bD7Tl5WOL02t/S8s4csRHiXY6JPmPI+imewAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAIAkAAAAAAAAAAAAAAAAEEgAABHmSgAAAAAAAAAAIAEggkCASAAAAAAAAAABAEggIA+QD4k5ADrA7wLXVr2lp9jVuqqclBcIx5zk3hRXe20vM1NTuvhClKEa+rXrwo59iCXVnqpwzx7W+1mU2hqO41e2tOcLem7ma/Pb3Yf7b8kfWylsqjr6pNZlWk6VF/NpQbSx4vMvNdhOO0babfqtpkNH02FhSk3J1rirh160l7U3+6K6lyRkACDbEaCzv61RyjaW8sVqq4y/s4dcv4d/mV7uvC3oTq1M7sV1c2+pLvZR0+jOEZVq+PX1nvT/ADeyK7l/HtArW9KnQoxpU1uwisItGlf3Hba0ZeVSa/cvv8D7vqk6lWNlQlu1Ki3pyXyIdb8XyXv6i6o04UqUadOKjCCxFLqQH0lgYWSQGUE9QAEcjlj0yekGvcahT6PtMryja0VCtqrg/wAZJ4lTovuSxNrrzHsOo7urChbzrVHuwpxcpPsSWWfm7tFrdbaLaPU9duW3V1G7qXLy+SlJuK8EsLyOvh44tfc/Ct9SzTTHqPlZSpLfb3ePaRutdxXS4lRQTLfTzfX9rjZam6mu2keqNRzfkmz0Wl7Sw+t9ZpexVu3qs6vVTot+baRukGlgRHZx57bsrwguvBs2wm09xsltBR1Kk827ap3dNflKXX5rmv8AM1ZT44RFWSkkurK6zXkrFo004cl8WSLVl2pZ3FK7taVzQmqlKrBThJcpJrKaKuFzNE6CNRnqHRzYQqSbnaOdq2+yEsR/u7pvZRXr020+hcfJ7mOL/bWek/Qo7RbDarpaipValCUqPdUj7UP7yRx7b7soqWMZ4+HcdzzWYNdxxPtBbKx2j1azXBW9/XpLuSqSS+w7uFbzVQevYo7XUU8cmX+n3Lp1N2T9mfDwfaYxPrTPpTec5LCY28xMbbJLPXwKcpLPEtLK6dWhiTzOCw89a6mTUm314NGu6HTL7nU5ooO4q06katGpKnUg1KE0+MZJ5TXemfM5Z4lKos8ORLW0qfpnbq7oy2lp7UbI2WqJx9fKPq7iK+TVjwkvfxXc0bPnPWc5ej5tB/JG0c9FuZqFpqTSp55Rrrl+kuHikdHLGCpz4+i73vpvJ9/DEocU44ZyR6TexMNA29/ly0pKNjrWarwuEbiP4xfWTUvFyOuTTOmPZKO2Owd7pcIr4bTXwiyk/k1oJuK8Hxi+6TJcfJ7d4T53H97FMfLiq2qytq0KtGcqdSHGM4PDXmbdo21sWlR1WK7PhFOPFfSj+9e406pJb8oypyptSw4SWGmuDT78opzk+SLyYraHiMmDqnUvV1WoVKUa1CrCrSkvZlB5T8ylOfY+Z5fZalfabWdS0rOKfxoS4wl4r9/M9j6Ldmbzb/RKuoaff2VrUoVfU17eopOUHhNPh1NcvPsOXJHR3lCnBvedUYni+SPnD7Gz0mHQ1tDvZ/lTTMfRn/Arrod13HHU9Nf1Zmn3qOj8o5U+KvL91tdwcVFcD0+XQ5rv/wA0039GZ8S6Gtef/Kum+6Zn36fZ+T8r+LzBvCyvApTXHKZ6j/QxtB/800z3VP4HzLoW2gfLVdL/AEahmM+P7Z/KOV/F5c8rGT5lJcng9Rl0KbQvlqulv6tQpy6Edo3y1bS/0ahn8Rj+2fynlfxeYuSa4NHxLnzPT30IbS9WraV+jU/gR/QftJ/830r9Gp/AzHJx/ZHpPK/i8weOR8tcFwyepf0I7Sder6V+jU/gfcehLaF8Hqul/o1CX4nH9s/lXJ/i8mlFPOCnKGM4Z69/Qfr7/wCVtM/QqEPoN19/8r6Z5wqGPxOP7Tj0rlfxePNNLkfDzza9x65U6CdpG/Z1jSvONQovoG2nef8AdvSln8yoZ/E4/s/K+T81eU7/AHcCnVuKdNe005Y4RzxMr0k7P3GxmsU9HuNQtby7lR9bVdCMkqSbxFPPW8N//wCmpOTk8ybcu1nRS0WjcOa2CaW6bLmvWnV+M+9JckUZNdqKTn2+8b6Zs2x0a8JkuODYuh+hOp0ubNerftK/jJ47FCTf2Guppy5nqvov6HPUukiWquD9RpVtKbljh6yotyK/R335GjkTFccy7eDWb56xDqqK4I+iOGSSie2AAAAAAAgCQQABJHAkAOsAAAAAAAAAACABIIJAAAAAAAAAAAAAQBIIAEkEgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAMgAAABHkSAAAAAAAAAAAAAAAB5AAAAAAAAAAAAAI8iQAAAAAAAAADBRu68Le2q16nCFODnLwSyBq19cS9Zq9/Hi4ydOn37kd1L9NyNm0y2jZ6db2seVKlGHjhYNTsqcqtppVpUj+Eua8KlVe+tL7VjzN0RmWunnaSM46ieBa6jXnSoJUknWqy3KSfzn1+CWX5GGxQX9d1BvGbe2lw7J1P/APlfa+4uruvC2t5VZpvHKK5yfUl3tk2lCFtbQow4qK5vm31t97Zaw/rmoOpzoWzxHslU635Lh4t9gYVtPoSpU5VKuHXqveqvv6ku5ci6ADIAAAAAxO2Uak9lNWhSz6yVlWUcdvq5YPzYtOFvTx8yP3H6dV4xnTcJLMXwa7UfnRt9s9V2U201jZ+tBr4Fdzp08r41Nvepy84OLO7hWjcwp/VazqJYyl2YLiCy8lFLHeVocOrzLPagmG0bEwSo3dZrm4wz73+82DHBcDEbJrc0ne+fVk/dw/cZinlpfuJfCuyT+qSKecYYmt5JYxxRWhFbyWCqoxjHfnhRjxbfYa5a4ncxEPf/AEbISWxN05Z3ZahUa/Rh+89SNP6HtKnpOwGm0asNytWg7ionzTm95J+CaXkbgUmWd3l9C4VJpgrEolyZxl0izj/P/aPdxj+U6/7bOy68406M5yaSim2+44P1XVJanrWoajnhd3lauvCU5NfYzq4MfqmVV67O8cQrupFM+XW4lop8Ez6UuPAs5eX6dL23uZUa0ZpvC5rtRl3OMoqSbaayn2mvpvuz1GS0uq2nRm+fGP8AAhaqNl5jxCcVxwGuGeRTly5kYavKvRryo1IVaM5U6tOanTmucZJ5TXg0jrHo92hpbT7KWeqxaVWcNyvBfIqx4SXv5dzRyJOe7jjwPS/R02s/k7aqts/c1VG21L2qOXwVeK/2or3xRz8rF1V3Hwu/ReV7OXot4l0h5CSzFoJ8ESVT2flx76TOx72b29qana0nHT9Z3riDS4Qrr8ZHzypfWfYeXx4xznKO1+m/Y5bZbBXlhRhvX9v/AFmyfX62KeI/WWY+ZxlCgowSaafY1xXcy44mXrrr6eT9Uwe1k38StnTb6keh+j1tVPZTpFoU7mpu6bqija3WX7MJN/g5vwk8eEmaO4rd5EwSw12rDwdWTFF66lwYc84rxaHW3pC7T7YbFbJUNptmI2ta0tayjqVKrQc5KnLCjUjxXKWE1+dnhg8g07p7211Cj661vNIqR618DalHxW9wPaOiTXLPpE6KVZ6uo3VVUZ6dqdOXy3u7rb+lFqXn3HEe2uj6hsHt3qmztSrUhc6dcOFOquDq0nxpz71KLT8clfxcdJmaWju9Bzb5ZpXJit2l71Lpw2+i/wDfGlfqb/xnzPpx2+T4XGlfqT/xHjGj7UUbhqlqUY0JvlWivwb8V8n7vA2aNLejGUeKkspp8Gjp9jHHwpL8zk0nvaW/Ppz6QFyr6T+pS/xlN9O3SF/b6T+pS/xmiK3y848iHbp8MD2Mf01/mPI/k3l9PHSH/b6R+pS/xHxLp56RE+FfSP1J/wCI0Z2i57uT5drFcGn4j2MX0lHqOf8Ak3eXT30jJfj9I/Un/iPh9PvSRxxcaPn/ALi/8Ro7tFjlgpytI8jPsYvpn8xzfyb0+nzpIb/31o6//Bf+Mj+n3pIT/wB9aP8AqL/xGhytFjOSjUtmuozGDF9Ec/PP/p6Iun/pF67rSP1F/wCM+ZdP/SNnhd6T+oP/ABHnPwdyluxjl9hUp2lGDUquJ4+SuC/zHsYvpP8AH5/5PSrXpw6TblOcL/SYU1xlN2HBf3jo/ZjUtUsuju21na+4pK9haO6vZQperjTWHLG71NRwn3pnM/QpoENqdvLGwnTUrK0/rd0kvZ3INbsX9KTivDJ656UG1EbDZq32YtquLjUnv10nxjQg03+lLC8FI4s1KzkilYWvF5GSMFsuSf8A4542s1a52k2jv9evMqte1nU3PmQ5Qh5RSXkYrlwZWXBJPmfL4vJZ1rERqHnbZZtaZn5UJPtT5lObSWeJVrOMX7Tx48jKbK7IbS7W3Co6BpFxdxziVdx3KEPpVHw8ll9xC9op5luxUtl7VhhrWNxdXtGytKFS4ua9RU6NKmsyqTfKKR2h0K7FR2I2NpWVdxlqVzL199NcV6xr4qfzYrCXm+swPQx0P6fsVu6rqdSnqOvSjj1yj+Dtk+caSf2yfF9y4Hq3BIq+TyPc/THh6X07g+z+u3lIAORbAAAADrAEEgCCQAAAAAAAAAAAAAACB5EgAAAAAAAAAAAAAAEEgB5AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACORh9r6mNFqUIvErqcLdeE5JP7MmZZr200lV1TT7bi9z1lxLyjuR+2f2GYjujfw+NKirjaNSxmFpbNrHLeqSwvsg/ebIuRg9k6akr67/tblwi/zYJR+9S95nBbyxTwhllbL4TfTunxp0s0qXj8qXv4eT7SpqdadKhu0uFapJU6f0n1+Sy/IrW9KFChClD4sIpIwkpX9aVKkoUseuqy3Kfi+vwSy/IqWlCFvbwpU87sVzfNvrb7y2tf6zf1blv8AB0s0qXe/lP38PJl8CAABkAAAAAQznT0vujqvqVlT270W3lVubKl6rU6VNZlOgnmNVLrcMvP5r/NOi0j5qQjUi4yWU1hp8ieO80tuGnPijLSay/NqHf18fFFTw595656TmwuzmyW1tlU0Cq6E9ThUr1dOjH2KCTS34P5MZN4UeXB4wuB5XRt5TaiubwkXOO/XXcPKcjH7NprLcNDoqlpdvHk9xN+fEyMMJIp0qKpQUG+SS9xVptLBu8KW87mVSm3k2/ov0K12h2wstPvKlONtHNerCT41lDD3F25eM9yZqCnHzK9rqdzpl3QvrGrKjc29RVKU18mS+9dq61lGrJEzWYhPi2rTNW1o7OzYRUIKMUkkuCPo1ro72mobV7MW2qU92FZrcuKSf4qqvjR8OtdzRsmOBRWiYtqX0TFet6RavhgOkSlqlzsPrVtoai9Sq2VWFspPGZuLS49vZ34OFaEnCXqnTlTcPZcZLDi08NNdTP0InBSjg5k9Jjo4emajLbTR6DVndTS1GnFcKVV8FV8JPg/zsPrZ2cPLFZ1Kp9W41sleuPh45GXDgVItrmijCMoxWXl9pVjPHUy41vu8rb6XMH3o+1VcZb0Xh88lrvr/ADJ9Yt3kJiENM7RufXUlNcJfK7mRObZh7O7VKth8IS4P+JkKj78GqY1KE10VJ56ylSuKtpd0by2qSpV6ElUpTXOM4tNP3ktrsx3nxKMZdSaMTG4SrM1mJh2T0cbS0NrNkbHWqOFKtTxWgn+LqR4Tj5NPywbIc2+jHtOtN2iutmbmoo22oP1ttl8I14rjFfSivfDvOkimzU6LzD3fAz+9hiUNJo5M9IzZNbN7bzv7alu2Gr71xTwuEKufwsfe1L6z7DrQ0jps2V/nbsHeWdGmpX9t/WbLt9bFP2frLMfMlx8nt3iUfUOPGfDMfMONJJcUuRSct1PGCZ1U5bu649z4Pw8T4msrJd9W/Dx3RqdS9G9HTbF7NbfU7G5q7mn6xu2tXL9mNXP4KXvbh9ZdhsPpubDu50uw6QrCjmvY7tpqW6udCUvYm/ozePCfceH3G9Gm5QcoyjxTi+Kfau87F6NtXsuk7omVvq8IXE7i3np+qUvz93dk+7eTUl9JHBnicd4vC84GSMlJxS4Qoe1HEsJ5wZvQ9WvdMe5Taq22cyoT5fVfyWNqtmb7ZPavU9nL7edawuHSc5LHrIc4TXdKLi/MoUaa4LHHtO2urxtw566may9D0a9tNUpp28t2qlmVGfCcf4rvRkfgnW4nndo3SqQqU3KE48Yyi8NPuZuWi7SRlGNHU8Z5KvFftJfejVasx4VmSkxO4ZCVulyWCnKgl1YMolCrBTpyjKDWVJPKZDt21zZDq01RMsS6Cl1EStufAybt2fM1CjSnXqzhSo0/j1JyxCPi+3u5szFmYiZ8MZ8EzyiW2oUqFovw7frH+Sj8bz+avEo6ptLGf4LSVKEeTuJrE39BfJXe+PgYZTk17WW3xbfFtm+sT8tsRMeV1VrOfBKMIfMjy/zKU293OCnGRsPR3s5U2u2ysNCpuXq7ipmvJPjCjHjUfu4LvaMXmKxtuxVte8Vj5e/+jJs2tE2EqbQ3kVTudWfr96XDct45VNcep8ZfWPCOkraqe1u3Wo6upN20p+qtE+qjDhH38ZeMjoP0idep7LdHMdF03doXGpL4FQhDh6uil7bXco4j9ZHKkIJOPDG6cnGibWm8rT1C3t1jDC4k+I3steJTnJYGn0LrUtVtNMsI793d14UKMcc5yeF5dp3TfphT0xzeYiHvHoxbD2epU7zanWLKjc0Yy+DWMK1NSjlfHqJPhnOIp90joalQpUoRhSpxhCKwoxWEvIxexmh2uzey+n6HacaVpRjT3uub+VJ97eX5mYKLNknJaZe34nGrhxxGu4kGC21G9trC1lc3VVU6a4Lrcn1JLm2+xGp1rkFK1r07m2pXFJt06sFOLa6msoqgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADIAAcCAJAAAAAAAAAAAAAAAAA6ggAAAAAAAAAGRwAAgASGQSAAAAAAAAgABHXkCQQ/EASyMgeYDK7TVtQrReuahdTxuWtKFL3J1JftR9xtD4cTSY/1nTG+OdSuu35NSp/gRKrXf6bPs3byttDtKc17bpqc/pS9p/a2ZHPEiKwuHI+LmrChb1K8/iwi5PyI+U4jstI/1jVJz5wto7se+cuL9yx72VdSrzoWrdLjVm1CmvznwXu5+ROnUZUbSCqfjZZnU+k+L+8pca+q8eMLaGfry/gv2hoXNpRhb21OhDO7BYy+b7yq2OCAZBkjrHmBIIyEAJI8AgJKderTo0p1ak4whBOUpSeEkubZ99Z5B6VW2C2b6Op6Xb1HC+1uTtIbrw40sZqy/R9nxmiVKza0Q1ZskY6TaXOPS7tRLbHpB1LXIycrZyVGyT6reHCLx1bzzP6xgNJTqanbQecSqIs4vjnGOGEZbZqKer0pL5ClL7C9x16YiIeL5OSbzNpbZNttvtZTb4CU97/I+G+XE2yqYhKbfWfXFrDkj4z7ipHjHiyGk5bl0RbX1Nk9p4/CKmNLvHGld55QfyanlnD7n3I6jozVWEakZKUWspp8zjGlTi3xXimj3zoG2v8Ahtj/ADZ1Ctm6tIZtZyfGpRXye9x5eGOxldy8P/qHpPRef39m0vVi11bT7PVdMuNNv6ELi1uaUqVWnJcJRaw0XSBXx2enmImNS4k6VdlLvYfauro9xv1LWS9bZ3Evy1HPDP50eT9/Jo1b1meLOy+mjYa2262TqWS3KWpWzdawryXxKmPiv82S4PyfNI4quqd3Y31ewv6E7e6t6jpV6U/jQmnhxLriZ/cpqfLyPqPD9nJuPErp1CPWSzj9xQUkuOT6Ul2nVpWaVlMyVjX9ZS3Jv2oLh3ow+8u1lSlWdKanF8U8+Jia7YtXbOcU+LPmUspnzGpCpSU4yynx/wAj4nKKWcmvTXpd2F7c2N3RvLSo6dxQqRq0pr5M4vKfvOydgdorfanZOw1q3cV6+n+Fgn+LqLhOPlJM4mq1fYfF4wevei9tf/J+v3Gyt5Wxb6ivX2m8+Ea6XtR+tFZ8Y95ycvF1V3C59H5XtZOifEumG8FOolJPjkmMk45WcMlpFVD1vmHHvpC7IrZjpArXdtRcbDVt66o4XCNTP4WC82pfW7jz2KTOxunTZBbW7CXVC3pb+oWf9as8Li5xTzD60cx8WjjyjHKznPHh3IueLl66PI+pcecOXfxKnOnlcEj0/wBG/at7N7c09Nuqu5p+sbtvPL4Rrfk5eeXH6y7Dzfda4pce0+op78GnKEovKlF8U+prvN+XH11mHHx804skWh7d6ZexsZ2Vjt1ZUfbobtnqDiudOT/BzfhJuOfz12HNVOSeDuTYfUbDpQ6JZWuqxjUldW07HUYLnGoo4k12N5U12ZRxVtJod9sztHf6BqKfwrT7h0KjxjfS4xmu6UWpLuZz8PJ2nHPmFv6hii0Rlr4l9UOovqMd7+JY20stF9QlFdZ1zCjtDMaVdXVk06M96m+dOXGL/gbXp19QvEoxzCr/AGcufl2mn22HFF5W9iyuKi4OFOUk/CJptjizTaImWT2j2i0zRnKjOSvL5cFbUpr2H/pJ8o/RWZeHM8+1XVtQ1e4Va+qpxi/wdGC3adP6Mf3vi+tssKUIxhFY6itHHA2UxRV1VrFY1CrTk4/KzllzTqrGHz7SzTWcH3l9uCfgmsSyFOcW8HSfombLK20a+2uuaW7UvZO2tG+qjB+1JfSmv7iOaNDsL3WdodO0bT03c31eNvT/ADXJ8ZPuSy33JnZPSFq9p0cdEkqOmYpToW8LDT11+sa3YvvaScn4M4eXeZ1SPlZ+mYq0mc1vh4X0/bS/zj6QbiNCanZ6d/VKDT4Np5qS85cPCKPOqifYfcqm/wAc5eObecspuT5Y49R1YqRWsQqeTmtlyTeVvVUkuvB7B6KWyf8AKW1V1tPdUs2+lx9Vb5XB15ri/qwf988mjGc5RhCEqk5NRjFLjKT5LxbO1OijZanshsLp2juMfhEYetupL5VafGb8m8LuSObmZOmuln6RgnLk6p8Q2iHsxxxPpPiRjHWM4KiIesfWcmE1mGNfsa0/bi6NSME1ncmmnld7jn3GZ3jDbUycLa2uVw9Rcwk33Se4/sl9hKEb+FbZduOmytm+NtWnR8k8x/utGVMHoNVx1a8ovgqsIV148YS/Zj7zOCfJSdwkAGEgAAAAAAIygJHIhtDqAAe8cAJAAAAAAAAAyMgABwAAAAAAAAAAAAAAAAAAgkAAAAAAAAAAAAGQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQSBBIAAAAAAAAAAAAAABBPkAAAAAAAAAAAAAAACASAAAAAAAAAAAAAAAQSAA8gAIJAGP2hryttEvKsH7apSjD6TWI/a0YXTqEVq+m2ceMLanOp+jFQX7T9xkNqJb8LO1/trmLl9GCc/viijs5H1usX9x1U4U6K8eM5ftRJR2hqt3s2BFjqX4atbWi5VJ78/ox4/fuovnyLG1/C6nc1uqmlRj98vvS8iMNkruclCnKcniMU233FvpUGrX1s01UrydWXdnkvJYXkfOqtyoQtVzrzVP6vOX2Jl5HguWAykgZJAgkeQAAAAAHwAiTwjiL0idq3tf0j3tS3nv6fpjdjaYeVLdf4Sa8Z5WeyMTqDp52vexvRvqOo0Kijf3C+CWPb66plJr6K3pfVOIbdYpqKy0uHE7+Fj3PVKk9W5GoikPuMWuRl9msq7qVOHs08Z8WY5x4cjMbOUsQrVH1yUSziO7zmW36WZjLi/E+4rOFjifMUuwqw9lLCJSr+z6jT4rK6itFY5xwek7F7I/zq6J72dpSj/KllqFWpbPlv+xDNPPZJfakzzd+shNwqQcZJ4lGSw01wafY0aaZYvMw6M/GvirW8+JVIcOovdKvrjTL6hfWVV0bm3mqlOa6pLt7U+TXWiyjJdhG828JcOslaImNS5sd7VtFodXdH+1FrtXs7R1KjuwrL2LijnjSqLmvDrXc0bDzOUujna652Q1+N0nOdhWxC8orrh1SS+dH7eKOodLvra/sqN5a1oV6FaCnTqReVKL5MpuRinHb+nu/TefXk0iJ8rtxT5nhfpI9E9XaKnPazZq33tYoU0rq2guN5Tjycf8ASRXL5y4dSPdOYccriasWSaTuHbyMFc1emX52RbjJxeVx6+DXiVYvjg6O9Ifoid58I2v2VtHK8w56hY0o8a666kF8/tXyufPnzZv8E3jjx4F9gzVy128jyuJbBfUq2/hciHV48iipPHMb/Wb3LpktOusS9TLhGT9nuZd1JPPDjxMBvtdz5mWtLj19Deb9tcJfxNVoRtXXdVk+1Hzb3V1YXlvfWNR0bq3rQq0aifGM4vKfvR81JFGcnngY1Expik9Ntw7m6PNpLXazZDT9ctd1K4pL1kE/xdRcJw8pJo2A5d9FHbH+TtpbnY+8q7tvqKdez3nwVeK9uK+lBZ+ozqFcuZR58ft3mHtuFyPfxRKJxyuByL087JPZjbyvVt6W5p+p711btLhGWfwkF4SeUuyS7Drx8jQ+nLZL+dmwt1St6W9qFn/WrPC4ucU8w+tHMfFrsJcbJ7dmv1Hje9inXmHICiml1hLhyKmU4Jx+K1lZPiXDBeRMT3eMmJidS9T9GvataJtvLSLmpuWes4pLL4RuIp7j+ssx8d0yPpj7FxU7DbuypcPZstR3V1N/gqj824N98ew8XjVrUpqpQnKlWhJThOLw4STTUl3p8Tr3ZS+0/pW6I/V6lCEo39tK0v6cfydZLEmuxp4lHxRX56ziyRkh6DgZIz4Jw28uK6O7vezyLyhnPBMnWdJvdD2h1DRtQji6s7idGpw4Np/GXc1hruZNDHYd0TFo3Cmy1mtpiWStZPcXAudVq+r0G9knj8BJe/h+8tbf4iwj52nnu7N3HbJwgvOSMR5c3mzVIyS8D7TWMot4N4w+PeVqb7UbXZMKvPqRE5uEN7CeOoLHgVaNpXvalO0tacqtxXnGlRpxXGc5NKMV4tpEbdo2VrNpiHufofbLO+12/wBsr2lmjYp2dk2uDqyWak14Raj9ZlH0otq46ztlS2ftKilaaPH8Lh8JXE1l/oxwvFyPaLG0sOiboYjB7s3pdk5VGvy9xLi/0qkseZx7WrXN1dV7y7qurc3FWVWtN85Tk8yfvbK7DHu5JvK05lvw+GMUeZSpuCHrMvLKc5Y4pFvXqqPPhwznwO/elNFdvV/Ru2b/AJx9IsbqvTcrHR1G6qZXB1eVKPvzL6h11wPMvRw2Tlsv0eW87qluajqcvht1lcY7y9iH1Y44drZ6aUvIydd5ew9PwRhwxH2e8pXVaFvQnWmm4wWXhZb7ku0q8S0ufw97St18Wn+GqeXxV7+P1TQ7pXFCpGtRhWjxjOKkvBlprts7vSLq3ivanSko/Sxw+0q6d7Eatv8A2VVpeD9pfY8eRcy5Bie8NV0+5T1HS73lG4hKk/rxU19sftNrRpFTNrptVR56fdOX1YVN5f3GbrB5imuRKzXj+n0ACLaAAACOJPkBAJHHsAgkAAEAAAAAAAAABAJAEEoAAAAAAAAAAAR5ASQBx7ABIAAAAAAAAAAAAAABBIAAAAAAAAAAAAAAwACAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADA1/WKnrdfo0uqhaym/Gckl9kZFxsnBPT6tzj/fFxUqeW9ux+yKMTfXO5d6zfc/VNU4+FOGcfpSZsWi2/wTSLS1xxpUYRfikskp8NVY3ba5rSUKcpyeIxWX5FvpUHCypynnfqZqS8ZPP7xqz/AKlKkudWUaX6TSf2ZLpcFhLBFtWjaq6slzVCln60v8l9pdos9L9v4TX4/hK0kvCPsr7i9DEAADIAAAAAEN4XIeBrvSTtNbbIbFantDcpNWlBypwb/GVHwhDzk0vMzEbnSF7RSs2lzT6V21v8t7c0dnrWe/Z6LHFTHKVxNJy/RjhdzcjyChFQilEm4ubi91C5vL2o6tzXrSrVqjfGc5Pek/e2VYQT59ZeYqRSsQ8dyc05Mk2lUhxXIz2jU3CwUl8ubf7jCxjCK3pPGFk3PUtFudDr0tMvFu16dvRqTjjDg6kFPdfet7DNsTG9OLNE9O/haRUmuwr04PC4nxHCZUU1FYZOY7OKZdB+jMktjL+Oc/7o1P2IGuekDsU7O6ltXptL+rVpJX8IL8XN8FV8Hyffh9bM36MdXe2S1Fdmoz/YieqXltRvbSrbXNKFWjVi4VITWVKLWGmikvknFlmXs8fGryuFFJccL7j7ikbL0lbG3OyG0Do01KpplxmdnUfUuunJ/Oj9qw+01qHFcOss6Xi8bh5DPgthvNLJdPea5cOJ6H0R7dVNmryOl6pVctIrS9mT4u2k/lfQfWurn2nn64M+4tZy1xI5McXjUnG5N+Pki9XYtGpTrUo1KcozhJJxlF5TXaj6PAOibpFloVWnous1W9Lb3aNaXH4M31P8z9nw5e+0akKtONSElKMllNPKaKjLinHOnveFzacqm6+UyWUc8+kN0LyvXcbW7G2n9e41L7T6S/3x21Ka/tO2PyvHn0P1hrJjFltjncN2fj1zV6bPzkbcZOL5ru6yVJZ5nT3pDdDK1mNxtZsjapaqk53tlTWFeLrnBdVX9rx58sTqbkt2W9GSeGmsNPrTXNMvMGeuWv8AbyvK4lsFtfC4c1w4n3a3Pqayln2Xwku4s3NvkQ5S7Ebphya22KTTw08p8u8ozaRZ6XXcoujLnHjHw7C7fPiiMNU10q6Zd3Om6ta6rZVPVXdpUjVoT+bOLyn4dXg2d0bAbR2u1eyWn67aYUbqknOGfxdRcJwfepJryOD3NLqR7b6Jm2isNoLvY69rKNvqGbmy3nwjWS9uC+lFKX1X2nFzcXVXqj4W/pHJ9vJ0T4l1CfMlmIi01nPAniVG3qfMOQunvZZ7Mbe1nb09zT9SzdW2Fwi2/wAJBeEnnwkjQop+XgdddO2yMtqNh6/wWjv6jYP4TapLjJpe1BfSjlY7cHJKXLDfHivAueJk66aeP9U4/s5d/EvlRS5I9X9GbauGhbVS2fuam7Z6thQzyjcJez+kvZ8VE8ql2ZKbqVrdxrW85U61OUZU5xfGMk8prvTWTflpGSkw5OLyJw5YtD2b0utlI215Zba2dLEa+LO/cV8r8lN+WY5+ieF2rUlzOwtIuLDpV6IpUL3dXw+1lb3SivxNePBtdjUkpLyOQa2n3mj6td6TqFP1d3Z15UK8erfi8ZXc+a7mjm41+3RPwsPUscTPu18SvrbKS+wtts57mhUofPrx+xNl7bRct3gYrb6TVrY0sc5zl7kl+864jupqd7w1ylwgnnzKiZRg8R5cj6jLhjHE3adi430+PYey+iXsote26q6/d0lKy0SKlDPFSuZ5UP0Y70u57p4o5xWetrjhLi/A7c6IdAtejPogpS1VxoV40J6jqc31VJR3pL6qSj9U4eZk6a9MfKw9Pwxa/XPiGg+lhtWq13Y7IWtTMKKV3eJPnJ8KUH9sv0TwRpLgi/2m1m71/aTUNavU1WvazrNZ+Kn8WP1YpLyMa231k8GPopEOLmZpzZZtKKjfNI2fod2Vnth0j6ZptSG/ZUJfC71YyvVQae6/pS3Y+bNYqNbq49Z1N6LOyUdI2LqbQXFLF3rElODfONvHKprz9qXmiHJyRSjb6dgnNlj6h6/ShuxwlgqYzzCXYSUsf29hEajQyz0178atZ8JzqyUl2YeEvci8ZZWf4O+uqPU5Rqr6yw/ti/eZJSn6vVO6tS+2L/hL7C75otNR9h29f+zrRT8Jez+9F3ngYGt31CD1PUrWS9m4pwqY+lFwf7KMrs7cO40KyqzeZujFT+klh/amWOvJ09Xs6y/KUqlJ+KxJfdIqbJS3bW7tv7C7ml4SxNftE/hrr2szS5AhEkW0AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQAAAAAAAAAAAAAAAAAAAAAAACAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAIk8LLJLHX6/wXRL65T407ecl4qLDE+GrU18I0ukmuN9dKT71Uq7z/um7I1Owobt/o1phYpZm19Cnu/fJG2ErIY1pdPfvrSl1KUqj8lj75Ir3VRUbarVfKEHL3LJQh7erzf9nQS/Sbf+yNY42Eqa/KyjT/SkkRTVNNpulYUIP4ygs+L4suAuC4AMgAAAAAAQ+IE5OY/TG2rdxf6bsba1E4UGr2+SfOTyqUH4Lel5xOjtc1K00bRrvVb+qqVrZ0Z1q038mMU2/uOCdp9ZudpNpNQ1++T9ffXDrOLedxPhGPhGKivI7OHi677n4VXquf28fTHyxkIbrazkr0pYaeRuJnzOL3Gk9144PsLWezy8d+z0PoO2ajtb0gWFnVhv2Vni8u8rKcINbsX9KW6sdmTZenBtdJ+rpdXqV/8Aqgek+i3sotE2C/ly4p4vdZarcecaEeFNeazL6x5Z0318dK+tx680f9TA5MOTrzz/AEsOZgjFw4+5alhrtPiTeeGRGeeecn0odvWd9vDz0R3e/wDouxb2T1N//wAhL9iJ7ElwPIfReTWympr/APkJfsRPXzz/ACP3Je/9NjXHqxG1uz9htJolfS7+nmnUWYTS9qnNcpRfU1/lyZy9tVod/s3rFbS9Rp7tSHGE0vZqw6px7n9jyjrg1TpI2Nstr9GdtWao3dLMrW4S405dj7YvrX70jZx83tzqfDn9U9Ojk06q/wDKHLblxwmTFt9Z96xpmoaJq1bTNSoOhcUXiUW8prqkn1xfUynT48+KLWJiY3DxN8dqT0yrRWVwZ6Z0Ubf1NBlDSNZrSnpbeKNWXF23c+2H3eHLzSnhLgVoSTazjl1mrJji8als4vKvxrxarr6jVp1qUatKcZwkk4yTymu1H2l1ngPRd0hT2eq09J1apKekyeKdR5btm/vh3dXge929alXowq0pxnCaUoyi8pp8mmVWTFNJ7vecLm05VN18vtngfpGdCsNo6dfavZG1jT12K37q1hiMb5Lm11Kr2P5XJ9TPfSMIjjyWx23Dfmw1y16bPzWrwr0LidC4oVKNWnJwqU6kXGUZJ4cWnxTXJoRb6zrj0i+hmG1lGrtNszRhT2gpQ/DUItRjfQXU+pVEuCl18n1NciV1Wt61S3uKdSlWpScalOccShJPDjJPimutF3x80ZY/t5jlcS2C39KsKsqUlOLaaeUZONwqtJVIcE+a7GYF1OHNsrWNx6uq4Sb3Jvjx5PtOiYcdqbZWU8rhyFhfXOmahb6nYVnRu7SrGvRqL5M4vK8v3ZKTafcfMkpRafJ8DEx1RpGm6zuHfnRttLa7X7FaZtDaYjG7oKU6af4uouE4PwkmvI2M5c9ELbGOn6zdbF3lXdoX+bmw3nwVaK/CQX0opSX0Zdp1EuKKDPj9u8w9jw80ZcUSNJrHUcjdOuy/82Nu7hUKe5Yahm6tsL2Ytv8ACQXhJ5x2SR10eddPeyD2o2GrztaW/qWn5urTC4yaXtQ+tHK8cE+Nl9u7R6nxvew9vMOTcYR8y48PtKdCTnTTy2nxy+wq7nHjzLqJ28XManUvVvRn2r/kra6rs7dT3bTVVmll8I3EVw/SimvGMS/9KjY1WmsWe2dnTxSvN21vt1cFUS/Bzfik4t/mx7Tx+zrVbS8o3NrVlSuKM41aVRc4Ti04teaR13p1fTOlLoqca27GGoWzp1lHi6FePPHfGaTXgu04M8TiyRePC+4M15OC2GfMeHI1DMWkYLb6rvXFnTa4Royfvl/kbNqNjc6XrF1pd5FQurWrKlWS5KUXh47uteJpm3NXe1pQXyKEF78v953Y53O4Utcc1y9M/DD+saPqM+JRWF4nxOWN2O8lx5tm2Z07Yrvw9L9HLZeW2HSvp9KtT39P0vF/dtrg9x/g4P6U8cOtRke7eljtYrTQrLZK1q4r6hNVrpJ8VQg+Cf0pY/RZceipslT2S6LFruoRVG91pfD7iU+Dp0En6qL7lD2vGbOfukfaartftxqOvTcnRrVdy1i/k0I8ILuyvaffJlZH+bLv4hY5rfh+P0x5lim23lrmfLXPmfLqcePIb65o74UTL7F7PV9q9rNM2ftm/wCuVlGrOPyKS41JeUU/PB3bp9rQsrKjZ21ONKhQpxp04RXCMYrCS8jwL0Q9lnG01HbG6pYlXbs7LK+RF5qSXjJJfUZ0GuRUczL1319PV+lcf2sXVPmUgA5FqMs6v4PVaElyqU5QfisNfvLws9TzF21X5leOfB+z+8MS+9Tg6mn14x+NuNx8VxX2lajNVKMKkeUoprzJaTi01lMttJ/4PpQfFwzD9FtfuMfB8rLaaGLa2r9VG6ptvuk9x/tFDQZKnrl5Rzwq0KdVLvTlF/7Je7Swc9Cvd3jKNGU4+MfaX2oxVlUxtDY1o/FrUqtPxzuzX7LJx4a7drNnAQItoAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACCQAAAAAAAAAAAAAAAAAI8gAJAAAAAAAAAAAAAQABIAAAAAAAAAAAAACABJBI8gIMTte86DWpLnWlTpY+lOK+5mXMNtS80bGn8+8p5+rmX+yI8o28LbS4qrtGp9VK1k/Oc1/gNhfIwGzizq1/Pnu06MPsk/3mwdRm3linhZ2XtXt7PsnGHuin+8ah7Ve0p9tbL8ot/uQ0virmfzrif2cP3C446nax+bCpL7l+8wkvFyAAZCMkgAQCQIbIfVzPopXNSnQozrVZxhThFylJvCSXNthiZ1G3g3pf7X/ANnLPZC0qf1jU5euuknxVvB8E/pTwvCMjmem00uDTM50q7S1Ns9v9V13elK2q1fU2mfk0IcIeGeMvGTMBDgXXGp7dHkedm97JK+inyWOBsGwGzNTa7a3TtBo53LurmvNfk6EeNR+5NLvaNbpzxhHSvokbNqjpOobW3FPE7ubtbVtcqUH7bXjPh9QcjL0U7IcDj+7miJ8Pb7W0o2dlRtbanGlRowVOnCPBRilhJeRyb04Rb6Wdb4/Ko/6mB7r0b7YS2r6TNuaFCtKemaRK1sbVJ+zKcfXeumvr5jnrUEeKdONFf0q62038aj/qYHJwYmMk7WXrUR7ETHhpcVw5lWEuGGUcNdTZMXwLiXkoju6I9F552T1L/7hL9iJ68eP+i487Jany/4Ql+xA9gPPcj9yXv/AE3/AK9UPwJZINEu5pvSVsLYbXadl7tvqNFP4NcpcV+bLti+zq5o5x1bTLzR9RradqNtO2uaLxOEvskn1xfU0dgGpdJGxVjtbpmJbtDUKKfwa5S4xfzZdsX1rzR1YORNJ1PhS+p+l15FZvT/AJOY+GFjCJUlg+9esb/RdUraZqdvK2uaL9qL5NdTi+uL6mWsZNotImLRuHi8mO1LdNo7rlVG+Rv3RX0g3GztxDTNVnKrpM5ey+Lds31rth2rq5rsPOoy7HgqxeY4zjvIXxxeNS2cfk341+qrsK0uaF1b069vVhVpVIqUJwllST5NMq5Zzn0W7f1tmLiOm6jUnW0ecuC5ytm3xku2PavNdj6Gs7m3u7ancW9WFWlUipQnB5Uk+TTKrLinHL3fB51OVTceVXOVyPDfSK6E7fbGhW2l2ZpQtto6cc1KSajTv4r5MnyVTHBS6+T4Ya9yzjkQ1nnghjyWpO4deXFXLXps/M67t7qxu6tnfW1a1uqM3TrUasd2dOa5xknyZ8N5SyjtX0gehqz27s5a3okKNrtLQhiM37MLyK5U6j+d1Rn1cnw5cX6naXmm6hX0/ULaraXdvN061CrHdnTkucWi74/IjLH9vN8rizht/S+sq/rKOJcZx4P+J9zmv8zD0biVGsprivlJdhfSqp4cXwfFPtOiXBNe68stXvdJ1Oy1PTqvqr2zrwr0J9SnF5We7qfczv7YDaa02t2R03aCyaVO9oKbh1058pwffGSa8j88ZZb4HS3oSa1e1P5w7O1ZudpQVK8op/InNyjNLue5F+Oe04Odi3TqWnpmWa36Pt0v1HzUW9Fxzg+0g0U0RMPRTG405G6a9kYbLbdXXwemoWF/m7tkuUMvFSHlLiu6SNHlFZPbfSxlu6poi/7NcftUzxOWHgv+Lbqxxt4f1HHGPPMQpyxFZPafRO1i5/lfXdF3m7T1NO6jH5s8uLfmt39E8Wq8ermeseidHG2uuL/sEH/fMcuInHKfptpjPGmE9IRQo9MV/wCrgo+ttaFSWOuW61n3RXuPEdpqiqa/d8M7sow90Ue3+khGX9MNZYxmwof7R4DqVZVNWvZ81K4n9+CXGn9MJ56/57SnhnijauibYue3XSHpOguMnaTq+vvmvk28OM/De4R8ZI1Pe/BuXLgdeehzsnHTtirja+7oqN3rEtyg2uMbam2lj6Ut596UTPLy+3j/ALdXBwzky9/EMp6Tm0a2c6Pf5DsZKlc6rm2hGHDcoRX4RrsWMR+scpU/ipLC4e49E6c9qVtX0hX9xTmp2Vm3Z2mHlOEG96S+lLL8Ejz6quPBmvjY+mm5aOfmjJlmI8QSfDLZV0azutW1i00mwpqpd3tWNCjH8+Twm+5c33FtVfVk9q9EXZKWo7U3u191Szb6dH4PatrhKvNe3JfRg8fX7iWbJ0UmWviYPdyRDpHYrQbTZjZbTtBsl+AsqEaUX1ya+NJ97eW/EzIXIgpJnc7eyrWKxqEggGEklnrH/B1d8cxjv+5p/uLwo3sXO0rR+dTkvsAqQaayustdOe7K6p/Mry+3Ev3lWxk52VCfzqcX9hTteGoXce1wl/dx+4MK9emqtCdOS4Ti4vzRqGnTaoaJWb9qFanCXnCUH9rNyfI0p/g7FdXwe+/ZuP4Eq+Gu/mG7LkCFyJItoAAABAEgEAGASAAAAAAAAAAGQABAEggkCCQAAAAAAAAAABAD3kkEgB5AAAAAAAAAAAAAAZAE+QAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAHEAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADiAAAAAGD2lebzS4f6ecseFKf8TOGD2h/4T0zudZ/3BXyjfwnZlL4TqU0vy8Y+6nH+Jm+ow2zHLUH/ANsl+zEzJmfJTws9I/3lvfOqVH/fYlx1iC7LeX2yX8CNE46ZRfPO8/7zJ/5Y8Lf/AGjDK8AAZAAAAAA8m9KLa3+b3R3U0y2q7t/rUnaUsPjGnjNWX6Ps+MkesSeEcUekDtb/ADs6Tb2dGpv2GmN2Frh8Huv8JPznlZ7Io6ONj67uD1DP7WKfuXn9KnBU4YWFjkfe5hZfA+1jOT7jGPWs54F3p5Dqnaroel3mt67ZaNYx3ri9rwoUuHJyfxvBLi+5M6y6WtftOiDoOuKumNU61raxsNMi+cq8luxl3tcZv6LPPvRO2QV1rl9tdcU26Vina2mVwdWSzOS8ItL6zPP/AE4NsKmt7cWuyVrU3rHQ6e/XSfCV1USfH6MMLxnIrsk+7l6fiHoOHSMOHrt5l6F6DrktK1aU5OUqlvbTlJ85NyrZbMN04TT6VNb48pUf9TAzXoWR3NM1Jdlnaf8A9rNW6cZNdLWuLLxv0ev/AEMDZgiPflyc6/VxIj+2sPi+BDT5op0pZSy8PJcQ4rkWUvOQ6B9Fn/ilqf8A9wl+xA9iPIfRcS/mjqX/ANwl+xA9ePO8j9yXvvTv+vUKNtdW1xOtChXp1ZUanq6qjJNwnhPdfY8NPHeiszjHpJ6SNoOjz0m9przR6+/b1alurmyqSfqrheop8GuqXZJcV3rgRx45yTqHRmzRiiJl2cHxNW6Ntt9F252bo6zpFfKl7NejJr1lvUxxhNdTXua4o2lPJrtExOpTpeLxuGo9JGxGn7YaX6mslRvaSbtrqMcypvsfbF9a/ecza3ouo6BqtXTNUoSoXNJ8fmzj1Ti+uL7ffxOxzWtvtj9N2t0p213H1VzTy7a5ivapS/fF9a6/HDXTg5E4+0+FT6l6XXkR11/5OVVh8sH3vYWC/wBp9E1LZzVaum6nQ9VVg8qS+LUj1Ti+tP7OT4mL3uziWtbRaNw8blxWx26bK6e9JcFlcjeOjHb662Xu42V26lfSZy9qC4yoN85Q7u1ea489BjL/AP0rRmRyY4vGpS4/Ivx7xakuwNNvbbUbKleWVxTr0K0VKnUg8qSfWi54nNHRpt3d7J3it6ynX0ipLNWiuMqbfy4fvXX4nRulX9nqVjSvrG4p17etFSp1IPKkiozYZxy9zwOfTlU7eV15njnpE9DNn0h6fLVtH9VZbTW9PFGs/ZhdRXKlVx9kua8OB7JhZDXAhS80ncO3JjrkjVn5h6tYajomr3Gj6vZ1rO/tpblehVWJQl+9daa4NFTT6jeaMuK5xf7juHp86HNK6R9MV5bersdorWDVreYwqi/squOcex84viutPibXNK1TZ3WLjSNZsaljqFrPdq0ai4p9TXU4vmmuDRc8bkRljv5ed5nFnDP9PvEcP7zoL0JGv52bSpf9Bt/25nPUavrIqXLqfczoD0IZb22O0qT/AOT7f/WTJcz9mWn0/wDfh1hkELmChiXq3O/pZv8A3X0Ps+DV/wBqmeKvPBo9q9LJZ1jRP+7V/wBumeKtPCLziftw8R6pP+xZ9Y4cfuPXfRRgv5562/8AsFP9s8hbeD1/0UX/AO+Ws/8AcIf6wcqP8Us+l/8AZqw3pF0s9MNWTxhWFB/bM5onUcqtSo+c5t/adM+krUVLpQvqvzNKpv7JnMzi3TT7jPF/bh0Z+3IsyexWi321m1mlbM2Ckq+o3MaLmvydPnOf1YqT8jt/pg2gs+jnoj+A6Tu29WVCGmaZTi+MW47qkvoxTlntS7Tx/wBCDYrer6tt7e0niOdP09yXPGHVmvPdjnukYf0ntq3tB0kz0mhU3rHRIu3ik+Dryw6kvL2Y/VZzZJ97Lr4h37/D4Jn5l5zSqLc3UsJPEc9hDeXjJRUmlhPgPWYWXxRYR4eftEzO1aNGtcXNGhb03UrVZqFOmuc5N4il4s7o6LtlaGxuxGm6DRUXUoU964ml+MrS4zl+k3juwc1ei/sutoekNarXp79losFcNtcHXllU15e1Lxijrxciq5uXdumHpfR+P0065+UgA4V0AAARJZTXaSALTR3nTLbtVNL3Cnw1aqvnUYP7ZEaN/wAG0U+pNfayY4/lZ99BftMMLpml3/s22rxXyK9WX3SN16jTtRSzrcfz6j//AFRJUQyNwi8xTXWSUrR71rSl2wT+wqkU48AADIAAHEAAAAAAAAAAAAA4gABxAAAAAAAAAAAAAAAAAADiAAAAAAAAAAA4gAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAMFtFw1TTH31l/cM6YTaRL4bpkv9LUj76Uv4Ga+UL+H3sv8AF1Bf9sn+zEzD5GF2Zf4XUo9l0pe+nAzRifLNPCy0L/gqh4P72T/yy/8Au/8AtDRuGnwj82c17psmXDWI8OdB/ZJfxDK7AAZAAAAIbwBovTptf/M3o71DUKFVQ1Cuvgtin/bTTSf1VmX1Thqk5QkuMpdrlxz2vzPZvSx2qet7f0Nn7ao5Weiw3amOUriok5fox3V4uR496vL4Fvw8XTTf28v6nyPcydPxC4hPrwXdpCrWq0aNtB1K1WooU4Li5ybxGK8W0WD9lZbxhdh6v6L2y/8ALu3q1m5hvWWiwVfiuEq8sqmvJKUvFI35cnRWZcXHwe7kiHQNmtO6JOhiVa53ZR0mxdWs1w9fcS4tLvlUlheKOBteurvV7y91O/qOte3lWdxcT+dUm3KT97OkPTW223paZsHaVcqKWoagk+xtUYPz3pY/Nic3Uo+t3YL5bS97Ofh4+03n5WnqGbWqV8Q6y9EO1dvR1ik0vZtrSOPKoaN06QlHpZ1zs36P+pgel+i0krnaCK6oWy/1p5/07Jf0r63nnmj/AKmAwT/sTDk5Vf8ARrZotMuKUmljqLdPqx9pKqOKLKVBru6O9FuWdkdS/wDuMv2IHr+Txz0V23sfqTf/AMxl+xA9jXE85yf3Je89Nn/Xq+Jt44H5+elKp/097TSTw1O3aa/+hTP0GwcC+lJTS6ddpnj5Vv8A6imdPAjeRr9S7Y9sb0R7ba3sTrlLWtIqKSeIXlrN4p3EOe6+x9al1Puyjuno72z0TbfZ6lq+jV96Pxa1CfCpQn1wmup/Y1xR+c+nV40K633+Dlwkv3m87C7a61sLtFT1nRK6Uvi16Em/VXFP5k0vsfNff28riReOqvlWcPnzhv028P0DRDeDUOjHb/Rtvdn4anpknTrQxG6tKjXrLaePiy7U+alya80ttznkyktWazqXo8eSt46qtb2/2S03a3SHZ3sXCtDMre4ivboy7V2rtXJ+5rmDafQtU2a1mppeqUXCpHjCcfiVY/Pg+tfauTOw8ZRgNt9ktL2s0iVhqFNqS9qjXhwqUZfOi/vXJnTg5E451PhU+o+mRyY6q+XJybzho+4yeTK7W7N6lstrFTTdShjnKlVivYrR+dH966jD55NcUWtbReNw8dlxWxWmtlRyzwyzb+jXbu92S1FUarnX0utLNa3Ty4P58O/tXX4mkuSXVy7wpZfUZvii8alnBmvgvFqS7I0XVLLWNNoahp9xC4tq0d6nUg+DX7vDqL45V6ONudQ2N1F4UrjTK0s3Ftn+/Dsl9j9zXTWgaxp+uaXR1HTbmFxbVlmM4/amuprrXUU+bBOOXt+B6hXlU/tkOo836ceirSOkjQ/adOy1y1i3Y36jlx6/Vzx8am31dXNd/o/BkOKxjJpraazuHfekXrqX5q7V6XrGye0lxoG0FlKyvaLxKD4xmnynB/Ki+pr7+B7l6DUnLbTabj/yfQ/1kz3bpr6KtC6S9n1aXq+C6nbpysL+EczoS7H86D64+7DWTyP0Rtl9Y2P6T9rtA1+0+DX1vY0OK4wqw9ZPFSD64vqfinhplhfkxlxTE+VVThzizxMeHTyPoArYXDnn0r451jROv+rV/wBumeKyXJI9r9K541jReX+9a/7dM8WeHjgXnD/bh4f1X/sypy4Ya5nq/oo1H/PbWef/AAfD/WHk9RZXxcHqvooJ/wA9ta4f8nw/1hLlftyj6dOs8MF6Vl16rpA1R8v9yKSXDtcl+88D0izvdT1Wy0zTqfrr28rwtqEMZ3pze6vvR7P6YlZ0+kC/SXGWn2sffOX8CPQp2Qlre3d3tdeUt6y0OHqrdyXCV1UjjP1YN+c4mqmT28O1pXDOXkTv7dG6rPTuiHoWhb2W4/5Ls429snw9fcS4JtfnTbk/M4vulUqXM69arKrWqSc6s5PLnNvMpPvbye4elrtgtQ2ls9k7Srm303Fa6UXwdea9lP6MHn6/ceHylnr4GOJj1Xqn5Q9Uz9WSKV8QpTe7zLW4qqFOUs8EssvJYfcbl0EbHLbHpHsLS5pqpYWf9cvVjhKEGt2D+lLCx2ZOjJborMuLj09y8Vh076O2yb2V6NbGF1T3NQ1D+u3eVxUppbsX9GO6vFM9IKdKO7HGCoUV7dUzL2WKnRSKgAItgAAABD4JvsAtNG/4Npef3smP/C0v/oL9pjR1jTLfvpp+8U8PVqv5tGC97kGF2afqL9vXH+dP/UxNwNM1CWaGtz+dVqpfoRiZqhkbbYrFnQX+jj9xWPmkt2lGPYkj6MJx4AAGQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAPIeQAAAAAAAAAAAAQSABAAkAAAAAAAAAAAQMgSCCfIAAAAAAAAAAAABAEgAB1GF2oXDT6nzL2C/SjKP7zNGI2uWNFnV/salOr+jOLf2GY8o28KOzz3dT1GHzvVVPfHH+yZ18jX9Kap7Q1afL1lrn9GbX+0bB1C3linhaaXwp1ofNr1Ptk3+8iv7OqWz+dTqR/Zf7ibLhdXkP9Ipe+K/gxfYjc2dTsq7r84tfwMJLsBckAyAEASa30lbT2ux2xOq7R3WJRs6DlCDf4yo+EIfWk0vM2M5c9NXa1XFxp2xNrUzCg4318ov5TeKUH/el+ibcNOu8Q5uXljFimXh13e3V/e1r69qeturmrKtXm/lVJPMn72IotqbfYXNPi1wwy+jtGnjskzM7TPDg2/k8zsPoh2ftejnoije6w1bVPUz1LU5y5we7vbr+jFKOO1M5/wCgLZaO1XSPaULikp2Vg1eXOVwlGDW5F9uZ44diZ6l6Ze1y0/ZOx2OtKuLjV5+tuknxjbUmm0/pT3V3pSK/kz13jHC49PxRjxTms5a251292s2x1XaS+TjWv7h1VB/k4YShD6sVFeRY6TT9bqtpT+dcQXj7SJqwjKbcVxk8viXWzNPf2hsuxVN5+CTZ3RXpjSvy5Orcy6t9Fif9e2iWfkW331TQenieeljW02vyOP8AyYG7+is277aLD5U7b/8AsPP+nlv+l3W1x/IP/wDTA5OP/wBiW3PbfBrDT89Y8FzRFNNtdpWoxzJFnMdlF4l0R6K8WtjtSbX/ACjL9iB7HyPIfRba/mfqa7NRl+xA9fPO8n92XvPTo/16oOCPSlljpz2m+lb/AOopne/kcL+mVot7o/TBd6nUT+C6zb0q9CeOGacY05x8VuxfhJG7gWiMiPqVZtieQSl7TwZDS7j1kHRk96cFwz1x/wAjGPi+0qW2/CvCrDg4vP8AEupnbz80jTd9iNqdf2O2ipa1oVx6qtD2alOXGnXh105rri/enxR250T7f6Nt9oKvtPl6m7pJRvLKcs1KE329sXjhLk/HKXDFCEKtGNWnxjNZRmNktoNV2V1631nRrqVrdUfOFSPXCa+VF9nmsNHHyeLGWNx5beHz5w36Z8P0DRHUaP0R9JGk7f6N623xa6nbpK8spSzKk3ykvnQfVLyeHwN5KW1ZrOpeopkrkr1VYPbHZnTNqNInp+pUcp8adSPCdKXVKL6n9j5M5c262a1bZLWXp+o03KEsuhcRjinXj2rsa649XesM6/MPtbs5pe0+jVdL1W3VWjPjGSeJ05LlKL6pLt8uRv4/InFP9K71H02vJruPLjve4BS7+RnukLZHVdjdb+B3i9baVW3a3SjiNZdndJda81wNdUsvuLqlovG4eLy4bYrdNoVXJtczZujvbnU9jdU9bTlK50+q07m1zhS/Oj2T+/GH2rVN7wRHxuD+8zfHF41Jiy3w36qS7M2Z1zTtodIo6ppdzGvb1lwa5p9cWuprrRk8HJXR3tjqOxmr/CrXNazqtK6tXLCqLtj2TXU/J93UWzGvabtHpFHVNKuVWt6q8JQkucZLqkutFJyME4rf09r6d6hXk018srjgUfgtt8M+GeopfCfV+q9durf3M53c88Z44Kw5HOswMjLJYHPHpYP/AHX0TH/Rq/7dM8WbeUe0+lgv919E/wC7V/2qZ4o+Ze8P9uHhvVf+xL7k2eueinBfzz1l4x/UIf6w8iUknxXA9e9FOcZbZ6zh/wDwEOH/AIhnmftSelx/s1ecemjvrpNuKdOMpzqWloowXFyeamEvM9+6PtKs+hzoHhK/hGNaxs5Xt/jnVuZrLjnre84wXgjWdodkP53eltSqXVDf07RNOttQrtr2ZVE5qjF/W9rwgzH+mPtQ1R03Yuznl1X8Ovkn8lNqlB+Mk5Y/MRXRb3OmkPR3p7MXyS581PUrrVdUu9UvajqXV3WlXrSzznJ5ePuXkUFU8ijCO5FRxwSwfMm0i0rHTGnnr/qncrtTXDq7+4639FXZJaHsE9cuaW7e63NV+K4xoLhSj5rMvrnLXRls9X2x6QdM2bin6m4qKVxJc4UI+1UfuWF3tH6A2lCja2tK3oU406VKChCEVhRilhJdxwc3L2isLf0nj95ySqoAFa9AAAAAABSu5blrVn82Df2FUs9Xk46ZcY5uDivF8P3gVdPjuWNCOOVOK+wpW2XqN3Ls3I/Y3+8uoLdgorklgtdP41Lup86u17kl+4MLtmk1252V3w/H3k4rzrbv7jdZNRi2+S4mk2i9ZZ6UuuvdUpvzk6j+4lVC/wAN3XIkLkCLZAAAAAAAEASAQBIAAAAAAAAAAAgkAAPIAAPIAAAAAAAAAAAAAAAgkAAAAAAAAAAAAAAAgkAAAAAAAAAAAAAAAAAAAAAAAAAAAHEgkAAAAAAAAAAAAAAAAAQSAAAAAAAAAAAAAAAQSAAAAAAAAAAAAAAACOIJx3gAAALLXqLuNEvaK5zoTS8d14L0iSTi0+KfUGJ8NYsKylqul3S5V6U4Z+lBTX7JtCNMtZeosLFy4Ozuo05Z6lGo6b+xm5LkZshjWlP2dXqxz+MoxkvJtP70Rq3Cz9b10qkJ+6Sb+zJNyt3ULWp1S36fvW8v2Sre0vXWlajz34OPvRhPSsiS3sKjrWVGrnjKCb9xceYZARx7Q/ECy13UrTRtGvNVvqipWtnQnXrS7IxTb+4/P3arVrvafaLUNor5NV7+vKtKLedxPhGHhGKS8jpf0vNrfgGztlsja1cV9Tl6+6SfFW9Nrg/pTwvCMjluMcJRzw7y14GLUdUvPercjdopCacWkS5brw5YbWeZVUerkZzo/wBl6u1222l7P00/V3NXNxJfIoR9qpL3LC72jsyW6Y2qMVZveIh0j6LOzn8gdH8tevYqndazL4Q3LhuW8U1TXcmt6f1zmTpk2yntr0k6nr0Judk5/BrFdlvDKi12bz3p/WOmfSf2qp7HdFb0bTHG3u9WS0+1jDh6uil+EkuzEPZ8ZI41pxUcJRwkcXGrNrTklb8u/RSMUK2d7lzMlshTc9dhN8oU6kv7uP3mLWeCXA2Xo9oKpqVzNr4lu175L+B3T4U+adVdDeivUgtX2goNpTlQoSS62k6if3r3mo+kLYV7XpX1CtWptQuqNGrSb5SioKDfvi0WPR9tLHY/bax1dzcbXf8AUXiXXRnhSf1XiX1T2z0jNmYa7sStds4Kd3pSdZSjxc6DX4RcOxJS+r3nDFvazxM+Jd+GkcjhTFfMOZopLr4MqQlhZ5FDeSSSeeCIUm3hZRb7UGp29f8ARp2rWn7TXWzlzUUbfUfwltnqrRXFfWiv7vedIp5XB5OELS5u7C8o31nUdO5t6sa1Ga+TOLyn7ztHYHaC32o2VsNbtuEbmknOCfxJrhOPk00UfqGHpt1w9X6Nyuqntz8M/nqR5H6UuwUtuOjC7+B0fW6vpeb2x3V7U3Fe3TX0o5WO1RPXd1d5DinzRxUtNbRMLvJSL11L8uaMVJRa5SWUu4r04qPDB6X6SGxMNielG8o21D1em6kne2WF7MVJ+3TXZuzzw6k4nmmXnh7z0OK8XpEw8tmpNLzWWX0K9jTq/Baj9ib9h9kv8zLzSl1YNQzjDTxh5TXb2m1aVcK9s1VePWR9moux9vmT248tNd4ZLZvXtY2a1y31nRbuVreUPizSzGS64yXyovrT+9JnYnQ50l6Zt9pD3VG01e3ivhdk5Zcerfg/lQfU+rk+Jxoqaec8UXeiahqOhavb6tpF1Us722lvUqsOrtTXyovk0+DOXkcWMsbjy6eHz7YLanw/QBcVw5EnnPQ10n6ftzpqtrn1dnrlCGbm1UuE1y9ZTzzg/euT6m/RUlgpbUmk6l6vFlrlr1VYrajQNN2j0etpeq28a9vVXhKD6pRfVJdTOVukPYrU9i9Y+DXSlXsqsn8Fu0sRqL5suyaXNdfNd3X7T7TG7SaHpu0Gk19M1S2jcW1ZYlF80+pp9TXNNcjdx+ROKf6cPP8AT6cmu48uMY4cU88Hy4H0mlw4m2dJexOpbF6mqdTfuNNrSfwW6xz/ADJdk/sfNdaWoyeFzLzHeL13DxmbDfDfpvD7jLDNm6P9s9R2O1n4Zat1rSq0rq1bwqq7V2TXU/JmpueFwKdSTcfjYyL44vGpQxZbYbxeku09l9f03aPRaGq6XcKtb1o5XbF9cZLqknwaMnntZyJ0W7cXuxOtb/4SvpVw0ru3jxfZ6yC+cvtXDsOsNHv7TVNOoX9jcQuLevBTp1IPKkn1lHyOPOK39Pben8+vJp38rxPjzJZCRJzQsnPPpX/8L6J/3av+3TPE5Ps957T6WUt3WNE4/wDw1f8AbpnisXxyX/C/ah4f1T/sy+KreOw9Y9E2T/ntrSTx/ufD/WHlcoKX8T1n0UaGNtdaf/YIf6wlzI1ik9N754iHQlZaZpH8o65WjToOdJVLy4fzKcXjL7Es+9nCm2u0tbaza3VNfud6Mr6u504yfGFJezTj5RS88nRHpdbWS0bYqjsza1d281qbjVSfGNtDDn+k92Pg5HJ9JvOW2jh4OPX6pW/quabapC8qceotquViOG23gqb5f7P6NebRa9p2hWC/rN9cwo0+vdzzk+6Ky/IsLTERtUUrNrRV0P6G2xqttL1Dba8otVr+XwSybXKhB+1JfSmv7iOisIx+zmkWehaFY6PYU1TtrOhCjSj3RWPeZEoMt+u0y9jxsUYscVAAQbwAMAQ8gYfaALPVVv06FL59eCfgnvP7i98yzuPb1K2hzUFOo/dhftMC7WcFppKzYxn11JSqe+Tf7ypqFR0bGtVXONNteOOB92kPVWtKl8yCj7kBba9Wdvo17WXOFCbXjuvBgrGio6ppFql+JU5vj82nu/7RldqpZ0v1C5161Kl5Oaz9iZaaTD1u0lSeOFC1+2c/4QJR4arf8mwrkgFyBFtARxAEkcScEYAAY4k47wI4kgAAAAAAAAAARxJAEcSQAAAAAAAAAAAAAACCQBBIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA1K/t3L+WbOOVLflOH14KS/vJmzafXV1Y29zHlVpxmvNZMPqcfU7QqfJXFt9tOX8J/YXOycv9xo0Hzt6k6L7lGTS+zBKfDVXtZd6r7FtGtnHqqkZ+SfH7Gy6xnifFzTVahOk+U4uL80U9PqOrZUpy+Nu4l4rg/tItilpXs0alHP4mtKGO7OV9jRe4LGzedTvd34mYZ+lu8fs3S+BAU61SFKlKpUkowim228JJdbKh456WO2f81+jWpp9tV3L/W5/AqWHiUabWas/KPs+MkSpWb2iIa82T26TZzd0r7VvbLpG1bW4yk7WUlRs88lbwyovu3uMvrGtR4vlzLSlKXYl1LwLqlLL4ov6RFaxV43Nab3m0q8IN4xlPPM6Z9EvZGNpo97tjd0/wANft29pvL4tCD9pr6U1/cRzvs1pdzr2uWGi2Cbub+vGhB4zu5fGT7orL8EdfdJ2r23Rr0P1YaVilVt7aFhpsev1jW7F97SzN/RZx8u8zqkfKw9MwxG81vEOa/SX2lhtd0jXVK3qb9jpObK2afCUk81ZLxkt3PZBHls6eOGUX8l7Cbbba5t8X3soyjnhg7MWOKViHFmzTkvNpWm5jHDzNr6PY/8IT7oQ+9muSi88EbRsTHc027ny36yXuj/AJmbObPP6GSvIxc8yxJJ5wzpH0etq6W0mxFTQ76cat3paVvUjPj6yg0/Vt9vBOL+j3nNN3JveRkuirayexnSBZ6pVnu2NWXwa9XU6Mmva+q8S8E+05c+Prr/APG303POLJqfEq3SZs9PZHba80VxatlL1tpJ/KoS+L7uMX9EwMZHR/pJ7Lfy3sjS2jsKcal3peaknFZ37eXx+PXjhLwT7Tm2jHrT4dvadXFy+5Tujz8HtZO3iVdcVyPZfRe2ohYarc7K3dTdpXubi0T5Kql7cV4xSf1X2njcefDgi4s7u4067oahZVHSurWpGrRmuqUXlPw/cS5GL3aTDVwuRODLFndBDfcYXYfX7fabZaw1q3aUbmkpSin8Sa4Sj5STXkZs85ManUveUvF6xaHj3pWbDS2v6Na15ZUXU1XRXK8td1e1OCX4Wn9aKyl2xicNUG54lnKfFeB+ok4qUWnxTRwL09bDPYbpL1Cwt6Thpt5J3lg0vZVKbe9BfQllY7N3tLHgZe/RKp9Sw6/XDz6MHjqXiX2kXTsruM3xpS9mrHtXb5FuknjHFH1hPgWmtKW0dUabo4LCcZKUWsprrXaUppcmWGzV56ym7Ko/ait6nn5vWvIykodwiXDaJrOn1pl/eaTqNvqWnXNS1u7ae/Rq03hxl+9dqfBrgdY9CvSnZba2a0+/9Xaa7QhmrRTxGvFflKeertjzXhhnJijl4aK+nXNzp2oUb6xrVbe4oTVSlWpy3ZQkutf+uPI5uRx4yx/bs4XqFuPbv4d9eAPLuhLpRttsbNaXqkqdvr1CGZwXswuYr8pBffHq8D1BPPWUl6TSdS9fhzVzVi1Vjr+kafrmlV9M1O2hcWteO7OEl9q7Guaa4pnKfSjsPqOxOr+rqesuNMryfwS7xz/0c+pTXua4rrS67Mfr+j6frulV9M1S1hc2teO7OE19q7Guaa4pm7j8icU/05OfwKcqv9uJHJPkm2fUMt+Zt3SdsFe7D6vuSlO40q4k3a3L5vr9XP8APX2pZXWlqGVweeBe4r1yV3DxefBbDfps+8LeTa4rij0noV6Q3stqkdJ1Oq3o11PjJ8rWo/l/QfX2c+3PmMqiwsLBSnUl8nrGXDXJXplLi574MkXq7up1IVKanCSlGSymnlNE72TwT0b+kKpVjDYzWq7dWnH/AHNrTfGcEsui31uK5dsfA96jiSPPZcc4rdMvc8Xk15FOqHOfpbtvWND4/wDw1f8Abpni8W+HDie4+lhST1XQ3/2euv71M8TnGOS74UbxQ8h6pOuRMS+qc+K6z2H0V2o7Yay/+wU/2zxrGHwM5s9tVU2U2f2nr21Rwvr+zp2Vq484ucnvTT6sRy134NnKpN8c1hDgZIxZovLGdPG08dseky/1OjUc7O3bs7Pri6VNv2l9KTlLwaNBk2uX2leU1uqKXBLCLeo8t8TXTHFaxEOjJlnJebSidSceOWmuR0T6GuyXwzU9Q21u6WadonZWLa51Gk6s14Jxj5yOdbehcXt1Ss7Sk6tzXqRpUaa5znJpRXm2foN0ZbL0NjthdK2dt91u0oJVZr8pVfGcvOTbOTmZOmvSsPTMHXk65+GzLkAMlW9IAAAAAAAAFpb/AITUbmpjhBRpr7396LqTSi2+SLXSk/gaqyWJVpOo/N5X2YDCNV9uFKh11a0Y+Se8/sTLwsp/hNWpx6qNJzfjJ4X2KReghg9o571/p9DsnUrv6sML7ZonZdb9fUrj51dU0+6EUvvci21OSqbRVZN+zb20I57HKTk/sjEvtk4OOg29SXxq+9Xf15OX7yXiGqO9mWQAItwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAEZAkEADDbTRUKmn3P8AZ3Hq5PsU4uP37pT2al6vUtRtn8p07iP1o7r+2H2l3tNSdXQrvcWZ04eth9KHtL7UYuxrRhr9lWi/YuaU6We3gpx+xSJR4aZ7XbP1FjSqxtXeRqfFpt1l9GXH78l6i0u7R17mjU30oRf4SOPjpPKXvRCJbU6XQnRtI+sX4Wo3UqfSfF+7l5F2EDLKJS3VxZw56S+1kdsOk66VvWU7DSM2Nrh8JST/AAs14y4Z7Io6q6ddrlsX0banq1Koo31SKtbJdbrVOEX9XjLwizhFW+MZblxzlvi32lhwcW5myl9Vz6iKQqU4YfMrJtcU8Hwk1zKlvSrXVxTtbWm6tzVlGnSpri5zk8RivFtFlM6hRRWbTp756HuzXwzVtR2wuqeaVmnZ2ba4OpJJ1JLwjux+tIxnpVbZR1rbals3a1N6z0ZfhcPhK4muP6McLxcj2expWPRJ0LJT3ZPS7JzqY4evuZcX+lUlhdzOMri4ub2+rXt5UdW5uKsq1ao3xlOTbk/e2cGCPcyzeVnyr+zgjFHy+5YwopvhwKbjw7D7fLgfDfYmyx2p47KM8rxybZspHGhKXLfrTln3L9xq8km+ztNx0Gm4aDaxx8aDl75MhZqzz2UrnnkxdzS328pPg1hoy9dbr8jGXElxwYiGqkzHh1F6NO1MNqejp6LqMo1r3Sf6nXjPj6yi0/VyfbmKcX3xZ4N0kbP1Nkdtb/Q8NW8JqpaN/KoSy4ce7jF98WW/QttdPYrpCs9QrVHDTbqXwS/zyVOT9mf1JYeezePfvSU2Qjq+y9LaS0pqV3pSbq7q4zt5fH/R4S8FLtOSlvYzf1K/vX8Xxdx5q5s4tc+BLi5JrLS7hGKjHnwXDIi0W8Q87PaXtfov7Uu01W72SuqmKV1m5s8vlUS/CQXisS8pHRKOF9G1G50nWLTVLKW7dWlWNak+rei+T7nxT7mztTZTWrTaHZ6y1iylmjdUlUis8Yvri+9PKfeij9Qw9F+qPl670ble5j6J8wyjPH/Ss2Le0/RvU1Ozoes1PQ27uior2p0sfhYLxit7HbBHr+T4rJTpuMopp8Gn2HDjv02iYW+akXpNZfmcqicI9eYrBMXl5SeTdOnfYyWw/SPf6ZQp7mn3D+F2D6lRm3mC+jLej4JdppsFwxjh2l/jv1xEvJ5qe3bpfdtOdCvCtSlu1IvMX3m52laF5awrwwt5e0vmvrRp6XBcDLaBdO3ufVVJJUqrw89Uupk57OPNXqhnXDCyU5rKwXU4tcH1lP1fEzEuNTtLi7sryjfWVzUtrmhNTpVacsSpyXJpnVnQl0nW+2NjHTdTnSt9eoQzUguEbiK/KQX3x6vA5W3eeCvp11Xsb6je2dxUtrihNVKdWnLEoSXJp/8ArPI0cjjxlj+1hwedbjX/AKd4oHnPQr0kUNtNKdpfOnR1u0gvhFKPBVY8lVgvmvrXU+HY36JvZ5FHek0nUvZ4c1ctYtVjNqtD07aPQ7jSdUoqrb144a5Si+qUX1ST4pnHXSBs9qWx201fRtQTml7dvXUcRr0s8JrsfU11Puwztd8Vx4mkdLmwtptvs1Oynu0b+g3VsbhrjTqdj/Nlya8+aR08TkzitqfCv9T4McinVEd4ci72VyI4vm3hH1dWlzY6hX0+8oSoXNvUlTq05c4TTw1/n4ERT8j0FJ6o3Dx969M6lVtKtW1uKdzbVJ0a9KcalKpF4lCSeVJd6Z190SbY0dsdlaV3Jwhf0PwN7SXyaiXxkvmyXFe7qZx9lLtRtXRjtlW2M2qt9RjKUrKolSvqS479LPxkuuUXxXmus5ObxvcpuPMLD0vm/h8up8S9G9LFpahobbX4iv8AtUzw2pUy8ZPZfSrvKN5LZ28tqsatvXta1SnUi8qUW6TTXkzxJvJngxrFEIeqatyJmFfeb7jEbUzcbShBPhOo5PyWP3mQlNrk8GB2mr5ubeln4lLPm2/4I6bz2cmGv6mOk5IpSeSomurJSmknx44WcLmzRPaNu2I3Ons3oh7KS1/pCr69d0t6x0NKdNtcJXEk1D9Fb0u57p2UuXA889H7Yz+ZXRnp+n16W5qF1/XL7PNVaiT3X9FbsfqnoRRcjJ13mXq+Fh9rHEPoEccEml1gAAAAAAMgWmqSfwV0ovEq0lTXdnm/dkuYxUYKKWElhFpNqtqsY59m3hvP6UuC+xP3lXUazt7OrVjxkliC7ZPgl72gwpacvWVrm5fy6m5H6MeH35LxlKzoq3taVFPO5FLPa+ti8rRt7WrcTaUaUJTl4JZB4hqWo1pTttWuYfHq1Z06b70lSj9uTbbSlC3tadCC9mnBRXglg1KxpSlQ0e0n8erWhUqeSdWX2pG5IlZCkfKeoDiCLYAAAAAAAAAjJIAEZJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAABxAAAAABxAAAAAAAAAADiABA4kgAAAAAAAAAAAAAADiABBIAAAAAAAAAAAAOsAjiA4gkAB1AAfNSMZxcZJNNYaNFjJ22mW1Vt72n14qT7FTnuSf6OTezVbu0hO+1WwlwjWaqR+jUhuv+9FslVryR8topttH1gstAru50a0ryeZzox3/AKWMP7cl8R0nHgDYNe6RNp7TY7YzVNorzEoWVBzjDOHUnyhD60ml5mYjc6hi9orG5c2+lttc9Z2xpbLWtRO00eG9WafCVzNZ/uwaXjKR4ospJdfIXWpXOo6hd399VdW6uridevN/KqTeZfayYLLPQ8fHGOkQ8fyss5cs2lVjTyj1j0WtjXrnSFPXrmlvWWiQVSOVwlcTyoL6q3pdz3TyqHsPe6kstHaPRPoNp0c9E9OpqTVCsqMtQ1Kb5xm470l9WKUfqnPzcnTXUeZdHpuHrvNreIeVel1tROre6bsda1PYp/1y8w+cuKpwf96X6J4HuvuybLtjqdfaPaW9128TjWu6squ7n4ifCMfqxSXkYSdNp8s5JYKdFIhx8rP7uWZWrWOGT53eOCvKOFwKbWGzdEtEPhpKLbXUbvaQVLTLWnn4tCP3GmUbeveVo2tpRnXrz4RpwWX49y7zdKznRSo1YOEoRUceCwRtZozsdecW8ZMXXiZevut8DG3DSTxgzDXWVhcpOlKm4pprDTOuvR32qpbY9GsdL1Gcbi90yHwG7jPj6ynu4pzfapQ4PtcZHIlzLHI3DoK2yexnSRY3VxW9Xpl/iyvsvCUZP2Kj+jLHHscjRysfXTceYXHpmf2smp8SuekbZ2rsjtZfaJNSdGlLetZP5dGXGD8uMX3pmvp54tvB0p6Tmyv8p7MUNpbSkpXOlv8ADtLjK3k/a/ReJdy3jmtprrR08PN7uOPuHP6lxvYyz9SqRlKLPbvRf2rdG/vNkruotytm5s8v5X5SC8ViXlI8M3nnw7CtpOp3mjarbarp8/V3VrVjWpPqcl1Puayn3MnycUZccw1cLPODLFod35IS7TFbG63a7R7NWGt2cs0byjGoln4r64vvTyvIy+MHmJrNZmHu6Wi9YtDxT0tNi1tJ0fS1q0ob+o6E3cx3V7U6DX4aHuSl4w7zjqDWV1p8V4H6WXNGnWoyp1IRlCSalFrKafUzgnpb2JqbE9I1/okINWUpfCLCT5O3m3ur6rzD6veWfBy/+JUnqmDX64alCOX1lanT4cE2VYU0uGCrGPAsplR9UM9pFf4Vbbk3mrTWH3rqZdSju95gLKu7WvGrHq4Ndq6zYXONSKnB5i1lPtMbcmSupUH7ilPPF5K8s8eGT4lF9ROJa4To+ranoWr22raVcyt7u2nv0pYyuxxkuuLXBrrR2L0XbaWG22zNLU7Xdp3EX6u7t85lRqpcV4Pmn1p+JxrUpt5RsnRdtXebEbUUtUo79SyqYp31CP5Wl2r86PNea6zj5fHjJXceVt6bzpwXis+HakePWHFPqLfSr211LTbe/sq0K1tcU41KVSDypRaymXRSdOpexiYtG4eHekjsHC4tpbZ6ZQzXt4KOoQgvxlJcqnjHr/N+ic+yaWW+eTvC4pU61GdKrCM6c4uMoyWVJPmmjj3ph2OqbE7V1bSnGX8l3O9WsZv5meNPL64N48HFlx6fyN/47PMes8HU+7T/APWnvkfE5NPg8M+HUTWeKyim55feW7z0VllNR1q8vtB03SLhqdHTZVfg8m/ajCo4tw8E48OzODFzeerJ858gvcYrSI7Q2WmZncplybNX1yTnq1bL+Luw9yX7zbqMMyjnk3xNMvJutcVKz+XUlL3sjlrpvwSiL4c+R6J6OeyL2x6VrCnXpes0/S8X922uEt1/g4Pxnjh2RZ5rN7qTzwzxfcdpeihsb/Nro3p6tdUtzUNcau6m8uMaOPwUf0XveM2V/Ly9GPS29Pwe5l39PYEiUTzBSPUAAAAAAAAB81JRhBzk0lFZb7EfRZanmoqdnHnXliXdBcZfw8wJ0qMnbu4msTuJOq89SfxV5JIi7/DX1vb4zGGa014cI/a8+Rd8EuxFppidX114/wAtL2PoLhH38X5hheIxO1s2tDrUo/GuHGgvryUX9jZlzA7TyU73TbZf2s67XdCLS+2UTMeWLzqFPS4qttBHHxba3lLznJJfZB+82IwWy0d641G5+dXVKL7oRX+05GdFvLFPAADCYAAAAAAACCfMYAAAAAAAAAAAAAOIAAABxAAAAAAAAAAAAAOI4gYAAYAAAAAAAAAAAAAQSAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAMDAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAwAAAAAAAAAAAAAAAAAAGAAAAAAAADC6zZ309ToXNhClJypSpVZVJYjDinGWOcse1w7+ozQeAxMbWmlWcbGxhaxqSqbrk3KXNuTbfDq4tl2AGUM5d9NHbDfvNL2Gs6vCLV9f7r8qUH/el5ROmNWv7XS9LutSvaqo2trRlWqzfKMIptv3I/PnbTWK+1O1uobR3ikq1/cSrKLefVw5Qh9WKivI7OFi677+lZ6nn9vH0/bD0opJpLDLmMkuMuCXaU9zEs4ZFaW5SlJrgll+BczOoeb8vS/R92Y/nZ0l2FtWp+ssdPxfXb5xcYNerg/pTxw7FI9t9KjaR22z9psraVcVtQmqt1h8VQg+X1pY8osqei7srHZno0jrl7FUr3WsXlWU+G5QSfqovu3fa8Zs8J6Q9pqu1m22pa25N29Wp6u1Wfi0Y8IeGeMvGTKzfvZtz4hZZb/huL0R5lhKkW2+ZQlBvPDLK+83wf2FahRdSWFDPa88F3s79w873YuvDg8dXWXelbP3F8qdzdzdlZS4qbjmpVX+jj1r854XjyMzZwsLRqpKELquuKc4/g4Pui/jPvfDuK9a8lXqOdSbnOXNyeWyNu/hOMuoX2nRtLCg7bTaPwejL473t6pV+nPr8FhdxWrqjcUtytHeXJdsfBmIVdLkfULldvma+lotM2na01WxrWyc6adWl85LjFd6/eYOvJyXLgbhSuOOFLiWWo6Vb3Oalu40ar6vkS/gzZWUqXiO0tQuIosbiMJYUoprDTT68mQ1OEraq6VaEoTXVL712mJuJtvK5GztMO3Fve4dn+jttfQ286LYWWpyjcX2nxen6hCpxdSO7iM2utShjPepHOXSFoVbZPbHUNAqbzhb1M285fLoy4wl7uD70zHejptxLYjpRoSu63q9J1VqyvW37MMv8HUf0ZPDfZKR7/6U+yT1DZ6htbZUd650z2LrdXGdvJ8X9SXHwcjhwz7GfXxK55MfiuPE/MOdotPOVnPYRUbcWl2cD4g3lprGOsqKOctrkXHl53xL3X0Utr3Cpd7HXk8KWbqxcn1/lYL7JLxkdEReeJwns/qNzoWt2Ws2Txc2VaNaC+djnF9zWV5nbuzuqWut6JZ6rZT37e7oxq033NZw+9cih5+Hov1R4l6v0fle7j6JnvDIHjvpT7Ffzh2Jjr1lS3tQ0XerPdXtTt3j1sfJJTX0X2nsR81YQqUpU6kVKMk04tZTRxY7zS0TC1z4oy0msvzsilx3WmmxhZNp6YNl57F7f3+iQg42UpfCLF9tCed1d+68x+r3mq0/aWUX9ckWrEvGZcM47zEpaT5GT0Wvztpvnlwz9xYRj2FanFqalHKaeUw1XjcM5KD4cj5lDjhFS1k7iiprn8pdjKvq1j94izkntOltGnw5n04p4eM4eV3FVx7j5cXHHBI2RKO3tPo2bbu1u3sbqVbFCs5VNOlJ/FnxcqXnxkvrLsOg8nCELiva14V7erOlWpzVSnUg8OE08pp9qZ1x0O7a0dttlKd5KUIahbNUb6kn8Wol8ZL5slxXmuoqebg6Z648PWej83rp7VvMN2bSNI6ZNi6W2+x1fTouFK/o/h7GtL5FVLgn+bJey+59xu+EQ4JnFS80tFoXWXHGSk1lwDc0riyu6tne0ZULmjN06tOfOE08Si/BnzltZx5nuvpQ7BeouYba6ZQ/BzcaWpRiviy5Qq/dF/V7zw1U2vjeaPTcTLGanU8RzcE8fJ0y+Es4x1lWKwlwIwuSIc8ZOzWnFM7fdSapUKtV/Ipyl7kzSpvEVl9RsmvXMaWkXDXOUVBebX+ZqE6ykpJ8mu05s1u+nXgpOtto6MNlqm2/SDo+zcIydC4resu5L5FvD2qj7sr2V3yR+htrSp0LeFGjCMKdOKjCMeCSXBJHOvoSbIxt9B1PbW6pL12oVPglm3zVCm/aa+lPK+ojpBJJcEee5mTrvr6et9Pw+3j39i4kgHKsAAAAAAwAAIfIsrD8PXq3rXsy9il9Bdfm/swfWpSlJQtKbxUrvdyucY/Kfu4eLRc04Rp04wgkoxWEl1IMLbU5N0o21NtTrvcTXVH5T937i6pxjCEYRWIxWEl1Is7P+s3dW7fxI5pUvBP2n5vh5F8CBmtalNVNoq0n8S2t4wz2OTcpfZGJsj5GlXdSVez1KvF+3d3E6VNp9slRj92TNUMk9tNh2TpuGg21SSanXTryz2zbl+8yp8UacaVKFOCxGEVFLuR9mE4jUAADIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAwABBOEAAAAAAAAAAAAAAAAAAwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAB8t7qbYEvgant10i7HbFU09odct7WtJZhbRbnWmu6Ecyx34weS+kb05VNnLutshshWg9XUcXt7hSjZpr4kVydVrjx4R73wXLFS6uru6rXd5XrXNzWk5Va9abnUnJ83KT4tnbg4c5O9vCs5XqEY+1fLoPps6dtK2u2Ur7M7OWl/Qp3dSKuLi5jGG/STy4xim37TSXHHDPaeFc3heJaQbzxRcQlx5lrhxVxRqqg5Oa+e27LhReUsGzdGOyE9studN0PdzQrVPWXbj8m3hxnx6s8I+Mka3ScW1xSOovRF2UjYbNXm1lzT/DalN0bVtcqFN4bX0p58oxNfLyxTH2T4OCcuWN+IZ/0hdcWz3R9LSbFxo3GpL4JRhDhuUkvwjXYt32frI5ZpQ3I7uMLqXceldOu1Ntre3F5XqV4xsdPzaUJN+y8P22u1uWV4JHkOqa369uFpF06fXNr2n/AANPGxTFNoc/L7uaax4hkri7oUJJSk3L5kefn2Hwr6U4qOd2PzVyNdjV3ef3lSNylyydXS5Jx6hsUbjHyitC649prsbp4wVYXhnUNU42wRuEypCvH5xgYXi58VgqQuk0sNjUITjZ+N0kuZMbt54ZZg43LbxxPuFfxIzCPtspeRoXlB0rmCnHq7Y96fUavrGj3NopVqTdegvlJe1FfnL95mY3KXMqfC+SWcvrQiNJ0vajQ7tKVNp4aksHa3o5bXUekPoojY6pKNzf2MHp2ownx9bHdxGb+nDGe9SOUdS0SjfJ1aDjRqvqS9iXiurxRtno67Tz6Puk+3hqE5UNN1VRs73feIxbf4KrnlhSeG+yUjRycU2p1R5hc8DkV6umflT212audlNrb/QbhTcbap+BnL8pRfGEvdwfemYuKS4HSfpRbKK+2fobU2tPNfTvwdy4rjK3k+f1ZcfCUjmtNcOPHPM7eDljLj38uD1LjTgzTHxKtBrKylweUe/+i9tapRu9kLurxhm5ssv5Lf4SC8G1LH5z7Dnt1McfcXWz2t32g7R2OtWLarWlVVIrPCePjRfc45XmS5mGMuOYR9PzzgzRZ3i3hEPijG7L6vaa9oNlq9jUVS3u6MatN9zXJ965PvMpjqPLWrMTp7qlovWJh4r6Vmx71vYmOv2lDevtFk6st1cZ27x6xeWFL6r7TlWiuCzx7H3H6HXNvSuKE6NWEZ06kXGcZLKaaw00cRdJWyUtjttr/RN2SoQn62zk/lUJtuPjjjF98Sw4eTcdMqD1bBNf8kNZp0m+0uKdPHBrzPuEOOF1FWMPE73npuq6fP1M91v2J8JfxMnOHDD5GJwsYX3GRtK2/S3JNuUFjj1oNeSN9ySa45KMslWfFcWylKPWSiWuIUqiysGydE21lfYna2nqK3pafXSo31KKy5U88JJdcot5XdldZr7XWbP0T6NHXOkDRrCrCMqPwn11VPjmFNOeH4tJeZHNqaTt28Obe9Xpdg28lUoxqRk2pJNdXBlQRWI4JKDUbe8rvXdbapZ2moadcWN7RhXtrinKnVpzWVOLWGn5HD209np+nbRajY6Veu+sbe4nTt6/z4J8OPXjllcHjPWdO+kdtTV2c2Bna2dV073VanwSnKLxKEGm6kl9VYz1OSOTZNeyksKPBdRd+k4rd7/DzPruas2ikR3RUfjgoyzIrN97PiUM9fEuZh56ssDtTUSsKcM8J1s+5P8Aia6k3HnhmY2vf4e3pJ8oSm/N4/cYaPIr8v8AyWuHtSHpHQd0t6x0ba1TpzqVrvZyvUze2OXLczwdWlnlNc2uUvHDO8NH1Kx1jSrbVNNuadzZ3VKNWhWpvMZwkspryPzKST5rOGdUehNtrOvbansJeV3P4LH4dp6k+VOUsVILuUnGS+myt5mDt1wufT+VPV7cumUAgVq6AAAAAAicowi5SaSSy2ySxvH8Jrqyj8RJSrtfN6o+f3ZBKdPi61Sd9UTzVWKafyaa5e/n7uw+tSqT3I21F4rV3uxa+SvlS8l9uC5bjCDk2oxSy31JFnp6depK+mmvWLdpJ/Jp9Xv5+7sDC7oUoUaMKVOO7CCUYrsSPsAMrXVLlWmnXF0+Ko0pT9ybNW063araPYS4yjJVKn1Itt/pOJm9qpZ0+naf9Jrwpv6Ke9L+7FlposPXbRV6qXC3tlFfSqSy/shH3ko7Q1W72bEuSJIJItoAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAIAkEEgAAAAAAAAAAAAAAEEgAQSAAAAAAAAAAAAEACQAABHEkAef8AT9tu9g+ji91a3lH+UazVrYRaz+GnnEsdaik5fV7zfzk/05dYnW2j2f0BS/BW9rO7lFdcpzUE/JQl72buPj68kRLm5eT28Uy8DulWrXNS5uK061erNzrVJvMqk28uTfW2z5UeKTfAqZy+PI+kuPFI9D0xHaHkpvMzuXzDnxKsc43sEbr3uCPvepUqbqVpKEF1saYZLZPSLraLaXTtBtN5Vr+4jQjJL4ifxpeEY5fkd722i07LZOOgaPXlp1OjZ/BbWrCKk6OIbsZJdbXPxPzw0zbDVtB1SOo7O3lTTrqnGUIXEYQlUxJYfxk0vvMy+mfpWfFbdaqvqUf8Bwcrj3yT2W3Cy0xVnqh0Lfei/QvKqqVtvdQlhYinZQwv7xQj6Ktkv+fN8/8A8CH+I8Bj0ydKr/5+at7qX+Ar0+mXpSTw9utWflS/wEIxcjxEtvXxPmr3j/2V7LGP58Xv6hD/ABEf+ytY4/48X36hD/EeHLpl6T+H/vxqvupf4Cf6YulB/wDPnVV9Wl/gM+1yf5Me7xP4vb//AGWbKPLbi+/Uaf8AiPpei3Zri9t75/8A4MP4nhsumDpRa4bdat+jS/wHx/TB0p4x/PrVf0aX+Ax7XI+0JycX+L3f/wBl+zX/AD3vf1GH+I+16MNosY21vf1GH+I8G/pf6UevbrVf0aX+A+l0v9KD5bdap+jS/wABj2+R9sTfifxe9w9GW0S/46Xv6lD+JUXoz2a/55Xv6nD+J4HHpf6UP+vOqv6tL/ASul/pP4P+fGq+6l/gHt8j7Ovifxe+f+zRZ/8AXG9/U4fxEfRptc/8cbz9Th/E8Gj0vdJzSztxqv6NL/Afa6X+kvC/999U91L/AAD2+R9sdXD/AIvfqXo42tNJLa+7/U4fxPnUPRtsLy0lb3G1d1KL5P4HDK8OJ4HLpe6TMcNt9V91L/AU5dLfSdJ8NudW/wD1/wCAe3yP5JRfhxO4q7e0jR3bbGW2zmrXktY9XZq1r3FWnuyuI7u63JJvi1zOMNt9DudltrdQ0G5cn8Fq4pTf5Si+MJeaaz3pmOXSz0mcM7cat5+r/wABidY2u2g1y+hd6/qVbVKtOHq1OsoqajnOE4pdbfPPNm3h0vht3+Wnn5acisa8wySllLLEllpotrC5o3X4qWWucXwki9hHjyLiNWhRW3SXv/olbVOdrfbIXdTPqm7uycuuEn+EgvCTT+s+w6BOGtkdaudmtotP1yzTdSzq+s3E8esg+E4Pxi2vcdtaNqFrq2lWupWVVVLa6pRq0prrjJZR571HB7WTfxL1fo/K93F0z5hdnj/pO7IrVtlae0dpSzeaRmVXdXGdvL46+rwl3JS7T2ApXVGlcW9ShWhGpSqRcZwksqSaw012HDS/RaJWXIxRlxzWXBs1iWP/AEwnhZwZnpG2fuNlNt9R0Sal6mlPftZP5VCXGD78LMX3xZhYqT480XdbdUbh4fLhnHaay+k88CpRk4VFJLl3cxGPiipCm85WSTTOl0/ainFcyPVvqyfdvFfEk3z4FZxwuKaMw1yt1T5YwepejRbQq7e3FeSy6FhUcezMqkF92TzJrjwPUfRlqqG297Tb41LCTXlUh/E1cmf8cuz0yP8AZrt0ciOPaSiG8FHL3jmL0vtRqy2x0PTsv1dCynWSzw3pz3fugeKqo8cT2H0vrecdu9GvMP1dbTpU08dcKjb/AG0eNxi+w9T6dH+CNPFepd+Rbaop+J9xe82uXDmUnw5tH1Skt7D7cHfrsr9fTWtpmpavNf2cIw+zP7zFvGeXEutXret1W7qZwnVa8lw/cWby+RwX1MrGkarCJS6kb96NWp1dN6edmZxk1G4q1LWp2OM6U1j3pPyNAkuHBm9ejpZVNQ6dtk6VLP4O7lXl3Rp0pyf3HPyP25dfFj/LGn6Ex5IDHAHn3qYSAAC7AD4q1IUqcqlSSjCKy2+SQFK+rqhRzGO/Uk92nD50nyX/AK6hY2/weliUt6pN71SfzpPn5FGyhOvWd9Wi1lYowa+LHtfe/wCBVv6/weklCO/Vm92lD50v4db7gwoXjd5dKxi36uKUq7XZ1Q8+vu8S/SS5FvY23waiouW/Uk3KpPrlJ82XAEggxeuam7SEbe3Sq3tbPqqb5LtnLsiv8ubHkmdLHVa3wrXFTi807Om1L/6k8faor+8XOyVNuyr3zWPhlZ1I/QWIw96in5mHjbSrTho9vVnOrWzUu63XGDftSfZKXFL/ACNvpU4U6caVOKjCCUYxXJJdRKe3ZCnedvokAi2AAAAAAAAAIAEgAAQSAAAAAAAAAABAEgAAAAAAAAAAAAABHmBIAAAAAAAAAAAAAAAAAAAgkAAAAAAAAAAAAAAAAAAAAAAAAAAABGCQAAAAAAAAAAAADgABBIAgkAAAAAAAAAAAABHmSAI8yeAwAAAAAAAAAAAAAACASAIJAAAAA1wOLPTZjUpdK9rUmnuVNJouD8KtRM7TOZfTr2Zr3GhaJtZbUnJWVaVndSS+LTqtOEn3KccfXOniW6ckOLnU6sUuaqcuCLiliT4lnCShDem0kuvqLSvqNT4lvmEfn9b8Owvol5mce57Mvc3lK2juv26nVFPl4swl7Xq3VTeqyz2Lqj4FL1j6xnPMlMladKlKOHxfvCb5lSUU1ho+ZQeO5EdNsWFh8T6Unnn1lOTUHhyivMmM45xvwx4kJmEtTK4jLhjJUUslopRz+Mh7z7jUXz4e8zFoQmkrpceHYfXey3VRcvWQ959qafKcP0huEemVZLiEmfCnH+0h+kTlccVIY+kY7fZ0y+4dzPpYxwwU9+POM4LHeFOOfxkPHeMdmNT9KylhdhG8U1KL5zh+kSpQf5SGfpA7/SrnPWiU8NFPejnG/D9I+lKOONSHvMbNS+03nkfUVwwfCnFfLg/M+t6Ofj08/SMf/rHdcUJOlJTg2n1NdRm7DVU8Quv/ADEvvX8DXfWLP42HvJVZL8rBPxNlMnS05MU3+G7b8PVqpCUZRa4STymdBeivtjG70272Qu6v4azzcWeXxlRk/aivoyfukuw5Jp39ahLNGvFJvjFv2ZeRtnR5tjLZ/azTtbt5qFe0qqU6e9wq03wnDPfFvzwQ5U1z49fLfwYvxssWjw/QBM+ZcVjJb6VfW2pabb39pUjUt7ilGrSmuUoyWUy6SPN6mJexieqNvF/Sc2S/lLZyltNa0t660vMa+FxlbyftfovEu5bxzpQpyaW9z8Tu68t6N1a1ba4pxq0qsHCpCSypRaw0zjvbbZqtsrtTe6LUT9VRqZt5vnOjLjB9/D2X3xZZcTJuOmXmvWeNNP8AJVr1OkusrQgk+HPxKiWXwR9Ri0+w7nnt7fMVx5lfGVldfM+YwXZxRWUUk+JjekbaUsY8DbOhjU4aV0k6VOpPdp3E520n9OL3f7yianJPj1PrKKlWo1YXFCe5WpTjOnJdUk00/eheOqsw28bJOPJFvp3BF70chrJhtiNao7Q7L2Gr0WsXFFSml8ma4Sj5STXkZoorV1OpfQcd4vWLQ8k9J3ZOeu7Dw1W1pOd1o1R3GIrjKi1iol4LEvqnK0nFY3eK+8/QGtCNSnKE4qUZLDTXBo5D6dejW82N1arq2l0KlXZ64m5RlFN/A5N/i5dkMv2ZeT6s3XpfLin+Oyg9Y4U2n3aPNZPLxyPic1BOb5RTk/BH3SjJpb8cPuPnV1GnpF1UfNUml58P3l9eN12oKzqYhoMarnUlJvjKTk/NlbeWeZQnFRfDkfEptcU+PcVszqVlrfhdycXwfWdIehHsTUra5qm3d1Scbe3pOwsW18abadWS8Eoxz3yXUeK9EWwOu9I21FPSdKpSp2tOSd9fOGadrT7c9c2vix6/DLP0H2S0DTdl9nLHQdIoepsrKkqVKOct9sm+uTeW31tsr+byI6eiFn6fxZ6uuzKJd5IBUrwAIbAlvBj0/wCUK2cf1SnL/wA2S/2V9rIrznf1JW1JtW8XitUT+N+Yv3vyL6EYUqahBKMYrCS4JIMeSrUhSpSqVJKMIrLb6i1s6U6tV3txFxnJYpQf5OH8X1+SPil/X66rP/etOWaa/tZL5Xgur39hfoHkA4GI1fVJU6srKwUal1hOcpLMKCfXLtfZHm+5cQTOn1rOqq0at7eCrXc1mNPPCC+dN9UftfUYCCr/AAuVC2fwvU7hKVWpPlCPVKWPiwXHEVz97PunSrTuJ2Ont1rybU7i4q8VTz8qfa8fFgvsRsuk6bb6bbeqo70pSe9Vqz4zqS65Sf8A6S5LgT3FWrU3lGjabS063cIylVrVJb9atL41SXa+zsS6lwL4Ag2xGgABkAAAAACCQBAJAAAAAAAAAAAAAAAIJAEEgAAAAAAAAAAOAAEEgCCQAAAAAAAAAAAAAAAMAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAOscQAAAAAAAAAAAADK7QBi9qtC0/aXZ6/0HVaPrrK+oyo1Y9eH1rsaeGn1NIymSMiJ13hiYiY1L85OlXY/WNhNsLjQdZhOSjmVpcbuIXNHqnHv6mup8DUZOL5JNH6PdJuwezvSBoEtH2hs/WwTcqFem92tbz+dCXU+7k+tM486R/R2292UuKlxo9tLaXSk24VbOP9Yiurfpc2++G95Fxx+bW0asouRwZpbdPDyNY6iU+sr3lnd2NV0r+wu7SquEoXFGVOS8pI+KNOpcVFC2o1K03wUKcXNt+COzrr524ei+9aI4ZlNmtC1PaXWLXRdFtpXmoXc1ChSh29cpPqilxcnwSNw2A6DukXa64pyjotXRbBtOV5qkHSSXbGn8eXuS7zr/oa6KdnOjXTJ09OhK71OvFK61GtFesqfmxXKEM/JXnl8Tkz82tI1Hl1cfgWvbdvC06P+hTYnZ/ZGx0rVdA0nWb+nDeury5tI1JVasuMsOSyop8EupJGdXRb0c/9SNn/ANQp/wADcFgkqJyWnvtfRipEa00/+i7o65fzI2f/AFCn/Af0YdHf/UnZ/wDUKf8AA3ADrt9ntU+mn/0Y9Hif/ErZ/wDUKf8AAf0Y9HuP+JOgfqFP+BuAHXf7Pap9NP8A6Mej3H/ErQP1Cn/Af0ZdH3/UrQP1Cn/A2/rDHXb7Pap9NR/oy6PurYvQP1Cn/Af0Z9H3/UvQP1Cn/A28Drt9ntU+mox6NOj5ctjNB/UKf8D6XRr0frlsboP6hT/gbYB7lvs9qn01T+jbYD/qboP6hT/gP6NtgH/zN0H9Qp/wNrA67fZ7VPpqn9G2wP8A1N0H9Rp/wH9G+wP/AFN0H9Rp/wADawOu32e1T6ap/RxsF1bG6D+oU/4B9HGwWOOx2hfqNP8AgbWlgDrt9ns0+mo/0bbA5/4naF+o0/4Ero12AX/MzQf1Cn/A20gddvs9mn0t9NsbPTbGjY2FtStbWjFQpUaUVGEI9iS5IuQfLb7CLZCcnkfpIbMfyjoFLaO1p5utMyq+Fxnbyftfov2vDePW8Mp3dvRubapb16calKpFwnCSypJrDTJY7TS23PycMZsU0lxSo57CoormZbb7Q6uy21t3o0lJUact+2k+O/Rllwflxi++LMPTk2XFbbjbwObHOO81l94a5I+ort4s+oLLK0IJcWsmGnahOllb3vPn1aw00X8ILDZbVac4yajFvHLgTrJ3nw9N6A9sKekatLZy+qqNnfT3raTfCnWfOPhL713nQSeePUcW2tneX9xG3sbavc120407eDnNPt4cvE6n6La+1FTZelDayz+D3tP2YzdRSnVhjhKaXxZdvHvK7lUrE9UPXejcjJant3jw2spXVtQurepb3NGnWo1IuM6c4qUZJ8GmnzRVTTByRPdezETGpeIbd+j/AKVfzqXeyl7/ACRWk23a1Yudu3+bj2oeWV2I8c206FOk2hYVLW00NajmcVvWt3TlFpPLeJuL+w7SwMceR24+fmpHTtX5PTMF7dWtOCbD0fule+movZqVsn8u4vaEEvHEm/sPSNgPRNuHcQututoYypJ5dlpjlmXdKrJL3RivE6vBC/MyXbMfBx0YbZHZjQ9lNFo6Ps/ptvp9jR+LSpRxl9cpPnKT628tmYBPmcs953LsiIiNQLkB1Hy5RSy3jAZfXAsK1Spe1ZW9vJxpReKtaL/ux7+19XiRKpUv26dvJwtk8Tqx4OfdHu/O93aXtGnCjTjTpwUYRWEkuCDHl80adOjSjSpRUIxWIpdRZTb1Co6UMq0i8VJJ/jX81d3a/LtPupKV/N0qUnG2TxUqLg5/mxfZ2vyXde04QpwjCEYxjFYSiuCQExioxUYrCSwkuol8CJNRWXwSNY1PUq2puVCyqSpWSyqlxF4lV7VTfVHtn7u0zEbYtaKwudU1apWqzstNnhxe7WuUk1TfXGPbP7F19hj7KjVupSsdKl6ulCT+EXb9rdl1pN/HqdrfBdeXwI0uynqkI0rXNtpcPZdWHsuqvm0+yPbPr6u02u1t6FrbwoW9KNKlTW7CEVhJGZ1HaGusTbvKnptjb6faxt7aG7BPLbeZSb5yk+bb7S5AIt3gAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAcQAZHLqAEkPxA4gT3Dv6wAIa7eIwnzSJAFC4tbW5W7cW9Kql1TgpfefNtY2Vrn4LZ29HPP1dNRz7i5BncsdMPlwjnOOPaN1Y4I+gRZfDXiSmlzJ48zWto619SvVC4r1LfT5pKnUoPd9rsnLmuPLHB8iURtG1tNlTT5DKzjPE1uw1WvaVo22pz3oye7RusYjJ/Nn82Xfyfc+BsMGmuvImNFbRL7WesjqKN3bRuaEqU51IJ8VKEnFos6VtbUpxo3UGpSeI1FOSjP7eD7vcYSZNPuHAtXp9p1U5p9qqS/iQ517P4+9XofPSzOHiutd64hheA+KVSnUpxqU5xnGXFOLymfYZAAAYDHUABGO1DmwDHEYJA+JyjCLlLkuPIwmobYbNafUcL7WbS1kuqvPc+8zryfFSjSqx3alOM12SWUI18o23rs1j+kbYVLjtZo363D+JTqdJmwUOe1Omy+jV3vuMnf7JbMX7fw3Z/SrhvrqWkJP3tGFuOirYCvJyezFlTb/st6n+y0Tjo+Wi3vx40t7rpf2AoZxrqqvspW9WX3RMLedO2xtJ4oU9Vuf/AKdrhf3mjKVehnYGpnGl16efm3lX/EUP6ENgc5en3b//ADav8TZHtOW8cy3jTyjpU230TbapZ1bTSbq2ubZyiq9aUVvU3zi0s9aTXHhx7TTadKLWUzpK26G9gKLT/kSVRr+0uqr/ANoy9h0d7F2TTobOafldc6Sn+1k3xyaVjUKvL6RnzX6ry5dp0HOSjTlGcn8mPF/YZrS9kdq79x+B6HfTjLlOVF04++eEdT2emWFlHdtLK3t49lKmor7C6UUuojPLn4hsx+gVj/lZz5o3RBtTdbstQubKwh1pydWa8lhfab5oPQ/s5ZqNTU6lzqdRc1UluU/0Y/vbPSMYWB1mm2e8rHB6Xgxd9bWWlaTpmlUFQ06wtrSn82jTUU/cXsUkPcSapnflYVrFe0QAAwkAAAAOIAjrBb3V3CjJU4RlVry+LShzfe+xd7Aq3FanQpOrVmoQjzbLH1da/e9XjKla9VJ8JT+l2Lu9/YVqFpOVSNxeSVSqviwj8Sn4dr739hWua1G2pOrVmoRXDvb7Eut9wYfWYUopLdjCK8EkWe9U1DhByhadclwdXuXZHv6yFQrXslUuoOFDOY0OuXfP/D78l+lux4cB4PJThGnCMIRUYxWEksJI+LmvSt6M69epCnTgt6U5PCS7WyhqN/bafbuvc1d2PBJJZlJ9UYrm2+w1u/up3LjeaonSoRkvg9ovae91OSXx6nYlwXe+JmI2ja8QraldvUqc5127bTIrecZvddZds/mw/N5vr7D6sNLqaru1LunKhpy+JQa3ZV+xzXyYfm9fX2F1pulVLqpC81OGFF71G1byoPqlPqlL7F3viZ1chvXhGK77yiEIwjGMUoxSwkuCRPAlgw2iAHEAMkccDiA4kgjxAl8AAAAAAAABw5AAAQ8kgCMjiOIEgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAwAAIJAAAAAAAAAAAACCeBAAAASAAAAAAAAAAABHAAAOAEgAAAAAAAAAAAABBIAEEgCCQACAAAAAACAJIJ4EACnXo069KdKrCM4TW7KMllNdjRVAGrXunz0yE4qnK70ySxKElvzorsfzofau9cmn3lXTIRe/O80uSzCUXv1KK7n8uH2rv6towjCX2k1LepO60uKe896rat4jN9bh82X2PrxzJRO+0tc113hl7etRuKMK1CpGpSmsxnB5Ul3M+qtOFWEqdSCnCSw01lM1i0rTozndaX7L3/AOsWlT2U5dfD5E+/k+vtM9puo299CTpNxqQeKlKaxOm+9fv5PqMTGmYtEoxXs3w37i37Oc6f+Jfb4lzRq061NVKU4zg+Uk+B954ci0rWj9Y61tP1NZ8ZcMxn9Jfv5mEirauFSVa0mqVRvMotexPxXb3oUb2PrFRuYO3rPkpP2ZfRfX9/cKV4vWKhcw9RWfxU3mM/ovr8OZXq0qVam6daEZwlzjJZQFRPqHDrLL1Vza8beTr0l+SqP2l9GX7n7ypbXtCtP1e86dVc6U1uyXl1+KDK5fYR95PDAXP/ACAJLrGCcDAEYC8CQAAAAAZAeRHmA8AOAwODHABgY4jrJ8gIwnxJ4BDrAhEgAAAAAPmUoxi3JpJdb6gJyU61SnSg6lSpGEIrLlJ4SLWV7Ou9ywp+ufXUlwpx8+vyPqnZRc1Vu6nwiquK3liMfox6vHiwPj1txdLFsnQov8tOPtSX5sX979zLm2t6NvFqmuMuMpSeZSfa31n1WqU6MHUqTjCCWXKXBIs/W3N5wt96hQf5aUfal9FPl4v3BiVe4u4wqepoxdau1lU08YXbJ9S/9cSLe1frFcXM1Wr9Tx7NPuiurx5lS2oUbanu0o4y8tvi5Ptb62VnwXAMnIx2sapSslGlGDr3NRZp0Iv2pd77I9rf28i31DV5TqTtNN3KtWL3alaXGnRfZ+dL81ebRiqUKjualtp8fhd9PDuK9V5UOx1GvsgvsXElEfbXa/xClWnU+EwrXebvUKuY0KNPlFdagnyj2zf8EZrRdIlRqq+1CUa17hqKj8Sin8mCf2y5vuXAutJ0qhYRlUzKtdVPxtea9qfd3RXVFcF9pfmJn6K0+ZSADDYEE5HmBABPACASAAAAAAAAAAA4AQASBBIAAAAAAAAAAAAQCeBAAAkAAAAAAAAAAABBJAAngRwJAEEgAAAAAAAAAAAAA4ACCQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAYAAxuq6VTu5fCKM/g93FYjVis7y+bJfKj3dXVgwFxGbu6dK8jOw1CGVRrU3wn27knwku2D93WbiUL20t7y3lQuaUatOXOMvvXY+8zFtIWpvww9jrU6M42+rKNKTe7C4isUqj7H8yXc+HYzPJprga1e6dd2MZLdnqFk1hprerQXevyi+3xKWm161rSVTS68bq1zj4POfxe6Enxj9GXDwM9MfCMWmO0tnrUadam6dWEZxlzUllFr6q5tfxEnXo/2VSXtL6Mnz8H7xp2qWt65U4SlTrRXt0ai3akfFdnesovSLZGpULa6o124puNRfGpzWJLy/efVxb0LiG7VpxmlxWea8H1C4tqNxFKrBNr4r5OPg1xRbuN7b/EkrqmvkzeJrwfJ+ePECfU3lv8AiKqr018is/aXhL+OfEmnf0t9U7iM7eo+CjVWE/B8n7yaF9QqVFSk5Uqr/J1Fuy8u3yLicYTg4TipRfNNZQH1ldpOSz+Aqnxta1S3/NXGH6L/AHYI9bfUfxlvGuvnUXh/ov8AiDa9BZw1G1bUak3Ql82rFw+/gy6jKMo5TTT60Db6AAZAAAAAAAAAAAADYAFOtXo0Y71arCmvzpJFt8PjU4W1CtX74x3Y+94+wMbXvDkU69alQg6lapCnFdcngttzUKzzOpSto9lNb8ve+C9x90rG3pzVRwdSqvylR70vt5eQFNXlWvwsreUk/wApU9mHl1v/ANcQrBVXv3tWVy+e41imvq9fnku5ThTi5Tkopc23hItXfOt7NlRlXfz292mvPr8sgXS3YxSSSS9yLSV7Kq3Cyp+vlydRvFNefX4LJMbOVbEr2p63/RxWKa8uvzLyMVFKMUklwSXUBaUrLemq13P19RPMU1iEPCP73ll3yDeFxMNea0qkpUdLhG5qJ4dVvFGD75fKfdHPkNbYmYhkr68tbG3lXuq0KVNcMyfX1Jdr7ka7qF9c6hGW/KdhYJZknLcq1F2yfyI93PvXItqnC9i60qupam1mEIpLcT+auVOPe+L7WZWw0J1JxuNWnCtNPehbw/E032vPx5d78kiUREeWuZtfwstNtq9/ThTsouy0+KwqyhuymuynHqX5z8l1mx2Nnb2VvGha0o0qa44XW+tt9bfayulhYJMTO2ytdAAMJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAB3GL1HRqNxVlc285Wl0+dSmuE/px5S+/saMoAxMbajqMHR3Y6xbqluP8HeUW1CL7VL41N+PDvZdW+o6jYpet/3StsZU44jWS/Zn5YfibHKMZRakk01hpmFutAhByqaXVdnJ8XSxvUZP6HV9XHmS3E+WuazHhfadqllfxbtq6lKPx4NOM4eMXxReLDNOvoeqnF6rZu2nD4l1Sk9xeFRcY/WwvEvbS91K3jFwq09RoNcN9qNTHdJezLzS8TE1+iMnxLYa9CjXg4VqcKkeySyWvwOrR42lzOKX5Or7cf4r3lK01qyrVI0aspWteXKlXjuN+D5S8mzJZMJ9pWnwqvS/wB82s0vn0vbj7uf2FWhdW9d4pVoTkucU+K8UVyjXtqFf8bRhNrk2uK8GGVScITTjOKkn1NFq9NtU80oSovtpTcPsXAfA50/973den+bJ78ft4/aM6jT5wt667m4P7c/eA+D3cPxV9Jrsq01L7sMb+ow50reqvzZuL9zT+8hXs4vFazuafeoqa/u5PqOpWTeHcQg+yfsv7QI+F1or8JYXC+i4y/eSr+n8qjcw8aEv3Ir069Gp+Lqwn9GSZUAtP5StOupKP0qcl+4fyjZf9Ih9pdcCcICz/lOy6q8X4Jk/wAoW3U6svo0pP8AcXXDuJ4dwO6z+Hp/EtbqX/hNffgn4RdS/F2M1/8AUqRj92S6ysFKpc21L8ZcUofSmkBRa1Gfyraiu5Ob/cR8Cqz/AB97cT7oNQX2cftJepWn5Ocqz/0UHP7kPhdxP8TY1nnrqSUF+9/YDs+6Fja0Zb0KMN/5zWZe98S4eFxfItNzUKjzKtRoLshFzfveF9g/k+jJ5uJ1bh/6WWV+iuH2ATO/tozcITdafzaSc3544LzPhy1Cvwp06dtD51T25+5cF7y8hCFOChCMYxXJJYR9ZAsqenUnJTuJTuZp5zVeUvCPJe4vEkuCR8XFxRt6LrV6tOlTjzlOSil5sw9xr3rFjTbWdx2Vaj9XS97WZeSfiIjbEzEM22kYm8122hOVGzjK9rx4NUmtyL/Om+C8OL7jCX9aVTH8rXzmp/FtqacYS7lBZlPzz4FzZ2OpXUIwpUIabark6kE6mPzYLhHzz4EunXlrm8z2h8ajXnVgp6vdRVKbxG2pZUJvs+dUfdy7iva2OoXsIpRemWiWFwXrmuxLlBe99yMrp2kWVlU9dCEqtw1h16r3qjXj1LuWEZBDf0lFPta6dYWthRdO1pKCk8ylnMpvtk3xb72XQBFsAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAgCR5kAASAAAAAAAAAAAAAEDgBPuBBIEAkAAAAAAAAAAQAJIAAAkAEAAAAAAAAAQBIIAEjzAAAAAAAAAAhxTTTS4mJudBtJTdWzlUsar4uVB4jJ98H7L92TL4IEdmJiJa3c2mp0IOnc2dHUbd83RSUsd9OXB+T8i2s60KdR0tMv6trUXF2tZNpfUn7SX0WjbeBb31laXtP1d1b060erejlrwfU/AztGafTF0dauKHs6hZvC51bbM4+cfjLyyZSxvrS9pudpc06yXNRfFeK5oxVfQ7ijx06+mkuVG5TqR8pfGXvZir+nKlUU9T0yrRlHlc0W5pfXjiUV4ozqJR6rR5bkmn1kmqWN/ewpqpZ31K+o9Ua7y/BVI/vTMjR2gt4tRv6Fayl86a3qb+vHgvPBjplKLwzKSEoRksNJ+J8Ua1KtTVSjUhUhLlKLyn5lQimt6llZzy52tFvt3EfH8nWi+JTlD6M5L7mXYMmlr8Cp/JrXK/8aX8SPgS/6Vdf+ay7AY0tPgUf+k3T/wDFY+AUX8arcS8a8v4l0SDULVafZ5y6Kl9JuX3n3Ttban8ShSj4QSKw4BnRhdwwuYyslnqGqWNhj4Vc06cn8WGcyl4RXF+QYmYhekN4ZgK2tXlbhY2Xq4v8rdPd81Be0/PBjL2tTq1FS1LUK13VlytqSaT/APDhxa+k2Z6ZQnJHw2G61qwoVHRhUdzXXOlQW/JeOOEfNosq17qtwm4ep0+lzbeKlTH7Mf7xb2drqdakoWtlR02h1Osk5eVOPBeb8jIUNBtW1K+qVb+a44rv2F4QWI+9Mz2hj9VmFp/Bq1dSoUrjV7mP5T8Yov6TxCPkZGlpepXT3ru5haU3+Tt/aqec2sLyXmZ6EIQgoQioxXJJYSPoxtmKfay0/S7GwzK3oJVH8apJuU5eMnxZeJYDBhOIiEggcAykAAAAAAAAAAACAJBAAEgARgkAAAAAAAAAACABPAEcCQIBIAAAAAAAAAAgASCOBIDgCCQAAAAAAAAABAEggkAAAHAAAAAAAAAAAAQSABBIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACPIkAAAAAAAAAAAAAAAAAB5AAAAAAAAAAAAAAAEAkAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAIaySAMZe6Hp11Udb1Pqa7/LUHuT82ufnkx9XStUtcuhWpX0Pm1fwdT9JLdfmkbGGZiZhGaRLS0re2r5qK50e4k/jfi4yfis05GVpX2q2yXrFQvqfU1+Cqfvi/sM7UpwqQcKkVKLWGmspmJraBaJudjUq2E3xxQfsPxg8x9yRnf2j0zHhVttbsatSNGtOdpWfBU7iO434Pk/JsyWUazdWmrUKcoV7WjqNB83RxGWO+EuD8n5FjaVaVOr6qxvbjT6v/R55S/8ALn/s4GjrmPLdMkmvUdW1K34XNrSu4rnO3luS/Ql+6RdQ2i0zH4apVt5fNrUZRf3YfkY1LMXiWXI5dRiKmvUpr+pWl1dN8nuerj754+zJY3V7qlWDnWurfT6PX6r25JfTkkl+iNSTeIbBc3FC2pOrcVqdKC+VOSivtMbV1uNRYsLWrddk5fg6f6T4vyTMFa+qrVlUsrS41OsuVeT3kv8AxJ8F9Uy9HTNUuON3d07WD+Rbx3p+c5LHuiZ1EI9U28LTULm7lT9ZqGpwtKT+RRfq0+7fftPywW9jb1Zt/wAl6ZP2udesnTUvFyW/L3eZsNjpFhaVPWwob9f+2qtzqP6z4ryL/A6teCMe/LBW+g1avtahezkuulb5pw83nefvXgZWysLOyhuWltSox69yOM+L6y5BjcpxWIMEEgwkcgAAAAAAAAAAAAAAAB5AACCQBBIAAAAAAAAAAAAQPIkAAAAAAAAAAAAAADyAAAjyJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAoXdpbXdL1Vzb0q0PmzimvtK4BphK2z1KOfgN5cWv5jfrKfulxXk0W70zWoPEZWFVfOzOm/d7X3mxgz1ShNIa/S0nVKn4+9t6EetUabnL9KXD7C7ttB06nNVK1OV1VXKdxLfx4J8F5JGVA3LMUiERjFJJJJIkAwkAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADIAAAAOAAAAAAAAAAAAAAAABAEgAAAAAAAAAAAOAAEACQCAJAAAAAAAAAAAAgCQAAAAAAAAAAAAADJAEgAAAAAAAAAAAAAAAAgkAAAAAAAAAAAAGQAAHAAGAAAAAAAAAAAAAAAAAAAAAAAAAAAAADIAAAAAAAAAAAAAAAAAAgkAAQBIAAAAAAAAAAAEASCOBIBAAAAAAAAAAAAAAIJAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAB5EEgAAAAAAAAAAAAAAADiBAJ8gAAAAAAAAAAAAAAQCfIAAAAAAAAAAAABBIAgkACOskAAAAAAAAAAAAQBAAE+Q8gAAAAAAAAAAAAAAQT5DyAAAAAAAAAAAAAABA8iePYA8gAAAAAAAAAAAAAAAQCQAAAAAAAAAAAAAACCfIAAAAAAAAAAAAAAAgkeQAAAAAAAAAAAAAAI8iQAAAAAAAAAAAAAAAAAAAAZAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADrAAAAAAAAAAAAAAAAAAAAAPIAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAGAAyhlAAMrtHmAAyMgAMoZQAAAAMrtGQAGQAAAAAcAAAyAAHAAAAAGUAAGUMgAMoAABlAAAAAyu0eYADKAADIAADgAAAADK7RkABkAAAAA4DKAADzAAAABkZAAZHmAAGUAA4DIADgAAHmAAAygAHAZAAAAAAAGRlAABlAABwAAZQygAA8wAAygAGRwAADKAAZXaOAADK7RlAAAAAygAAyu0AAOAAAZQ4AAMrtAADIAAZQ4AABnvAADKAADIADKAADzAADIAAAABkZQAAAAAAAyOAAAcAAAAADK7QAAyAA4DgAAyhkABlAAAOAADgOAAAAAAAAAADKGUAAyPMABw7RwAAcBlAAMjIADgMoAAMoAAMgAM95DfBgf/Z" style="width:32px;height:32px;object-fit:contain;" alt="NX"></span><span>Nexo</span></a>
      <a class="btn" href="/">← Volver al chat</a>
    </div>

    <h1>Elige tu plan</h1>
    <p class="subtitle">Cada plan desbloquea más modos y prioridad de respuesta. Pago único sin suscripción.</p>

    <div id="currentPlanBox" class="current-plan" style="display:none">
      Plan actual: <strong id="currentPlanName"></strong>
    </div>

    <div class="grid" id="plansGrid"></div>

    <div class="contact-box">
      <h2>¿Cómo comprar?</h2>
      <p>
        El pago se realiza directamente por mensaje. Después de pagar, el plan se activa en menos de 24 horas.
        Escribe indicando tu <strong>nombre de usuario en Nexo</strong> y el plan que quieres.
      </p>
      <div class="contact-methods">
        {% if donate_url %}
          <a class="btn primary" href="{{ donate_url }}" target="_blank" rel="noopener">💳 Pagar ahora</a>
        {% endif %}
        <a class="btn" href="https://www.tiktok.com/@teamnexoai?is_from_webapp=1&sender_device=pc" target="_blank" rel="noopener">TikTok</a>
        <a class="btn" href="/" >💬 Escríbenos en el chat</a>
      </div>
    </div>
  </div>

  <script>
    const PLAN_FEATURES = {
      gratis: {
        name: 'Plan Gratis', price: 0,
        features: ['Modo Rápido', 'Acceso al chat', 'Memoria básica'],
        modes: ['auto', 'rapido'],
      },
      beta_tester: {
        name: 'Plan BETA Tester', price: 5,
        features: ['<strong>Modo Combinado</strong>', 'Mayor prioridad de respuesta', 'Modo Rápido incluido', 'Acceso anticipado a novedades'],
        modes: ['auto', 'rapido', 'combinado'],
        highlight: true,
        badge: 'MÁS POPULAR',
      },
      developer: {
        name: 'Plan Developer', price: 15,
        features: ['<strong>Modo Código</strong> (multi-agente)', '<strong>API Key</strong> incluida', 'Prioridad máxima', 'Todos los modos'],
        modes: ['auto', 'rapido', 'combinado', 'codigo'],
      },
    };

    async function init() {
      let currentPlan = 'gratis';
      try {
        const r = await fetch('/api/account');
        if (r.ok) {
          const d = await r.json();
          currentPlan = d.plan?.key || 'gratis';
          document.getElementById('currentPlanName').textContent = d.plan?.label || 'Plan Gratis';
          document.getElementById('currentPlanBox').style.display = '';
        }
      } catch {}

      const grid = document.getElementById('plansGrid');
      grid.innerHTML = '';
      const planKeys = ['gratis', 'beta_tester', 'developer'];
      planKeys.forEach(key => {
        const p = PLAN_FEATURES[key];
        const isCurrent = key === currentPlan;
        const priceStr = p.price === 0 ? 'Gratis' : `${p.price} €`;
        const priceHtml = p.price === 0
          ? `<span class="price-val">Gratis</span>`
          : `<span class="price-val">${p.price} €</span><span class="price-period">/ mes</span>`;
        const featHtml = p.features.map(f => `<div class="feat"><span class="feat-icon">✓</span><span>${f}</span></div>`).join('');
        const actionHtml = isCurrent
          ? `<div class="tag-current">✓ Tu plan actual</div>`
          : `<a class="btn primary card-action" href="/planes#contacto" onclick="document.querySelector('.contact-box').scrollIntoView({behavior:'smooth'});return false;">Obtener este plan</a>`;
        const badgeHtml = p.badge ? `<div class="card-badge">${p.badge}</div>` : '';
        const cardClass = `card${isCurrent?' current':''}${p.highlight?' highlighted':''}`;
        grid.innerHTML += `
          <div class="${cardClass}">
            ${badgeHtml}
            <div class="card-name">${p.name}</div>
            <div class="card-price">${priceHtml}</div>
            <div class="features">${featHtml}</div>
            ${actionHtml}
          </div>`;
      });
    }
    init();
  </script>
</body>
</html>
"""

NEURAL_HTML = r"""
<!doctype html>
<html lang="es">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>🧠 Nexo Neural</title>
  <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&family=JetBrains+Mono:wght@400;500&display=swap');
    :root {
      color-scheme: dark;
      --bg:          #080c10;
      --panel:       #0d1117;
      --panel-2:     #131920;
      --panel-3:     #1a2230;
      --text:        #e6edf3;
      --muted:       #7d8590;
      --muted-2:     #484f58;
      --line:        #21262d;
      --line-2:      #30363d;
      --accent:      #00d97e;
      --accent-dim:  #00d97e22;
      --accent-2:    #4d9eff;
      --accent-2-dim:#4d9eff18;
      --purple:      #a78bfa;
      --orange:      #f59e0b;
      --danger:      #ff6b6b;
      --node-new:    #f59e0b;
      --node-normal: #4d9eff;
      --glow-green:  0 0 20px #00d97e44;
      --glow-blue:   0 0 20px #4d9eff44;
      --radius:      10px;
      --radius-sm:   6px;
    }
    * { box-sizing: border-box; margin: 0; padding: 0; }
    body {
      background: var(--bg);
      color: var(--text);
      font-family: 'Inter', ui-sans-serif, system-ui, sans-serif;
      min-height: 100vh;
      display: flex;
      flex-direction: column;
      overflow: hidden;
    }

    /* ══ HEADER ══ */
    header {
      display: flex;
      align-items: center;
      gap: 14px;
      padding: 0 20px;
      height: 52px;
      background: var(--panel);
      border-bottom: 1px solid var(--line);
      flex-shrink: 0;
      position: relative;
      z-index: 10;
    }
    header::after {
      content: '';
      position: absolute;
      bottom: 0; left: 0; right: 0;
      height: 1px;
      background: linear-gradient(90deg, transparent, var(--accent)44, transparent);
    }
    .back-btn {
      display: flex; align-items: center; gap: 6px;
      background: transparent;
      border: 1px solid var(--line-2);
      color: var(--muted);
      border-radius: var(--radius-sm);
      padding: 5px 12px;
      font-size: 12px;
      font-family: inherit;
      cursor: pointer;
      text-decoration: none;
      transition: all .2s;
    }
    .back-btn:hover { border-color: var(--accent); color: var(--accent); background: var(--accent-dim); }
    .header-title {
      display: flex; align-items: center; gap: 10px;
      font-size: 15px; font-weight: 700;
      letter-spacing: -.2px;
    }
    .header-badge {
      font-size: 10px; font-weight: 600;
      background: var(--accent-dim);
      color: var(--accent);
      border: 1px solid var(--accent)44;
      border-radius: 4px;
      padding: 2px 7px;
      letter-spacing: .5px;
      text-transform: uppercase;
    }
    .header-spacer { flex: 1; }
    .header-clock {
      font-family: 'JetBrains Mono', monospace;
      font-size: 12px;
      color: var(--muted);
    }

    /* ══ LAYOUT ══ */
    .main-layout {
      display: flex;
      flex: 1;
      overflow: hidden;
      height: calc(100vh - 52px);
    }

    /* ══ CANVAS PANEL ══ */
    .canvas-panel {
      flex: 1;
      position: relative;
      background: radial-gradient(ellipse at 50% 50%, #0d1f1455 0%, #080c10 70%);
      overflow: hidden;
    }
    #neuralCanvas { display: block; width: 100%; height: 100%; }

    /* Scan lines overlay */
    .canvas-panel::before {
      content: '';
      position: absolute; inset: 0;
      background: repeating-linear-gradient(0deg, transparent, transparent 2px, #00000008 2px, #00000008 4px);
      pointer-events: none;
      z-index: 1;
    }

    /* Grid overlay */
    .canvas-panel::after {
      content: '';
      position: absolute; inset: 0;
      background-image:
        linear-gradient(var(--line) 1px, transparent 1px),
        linear-gradient(90deg, var(--line) 1px, transparent 1px);
      background-size: 60px 60px;
      opacity: .18;
      pointer-events: none;
      z-index: 1;
    }

    /* Central AI orb */
    .ai-orb {
      position: absolute;
      top: 50%; left: 50%;
      transform: translate(-50%, -50%);
      width: 72px; height: 72px;
      border-radius: 50%;
      background: radial-gradient(circle at 35% 35%, #00d97e60, #00d97e10 70%);
      border: 1.5px solid var(--accent);
      box-shadow: var(--glow-green), inset 0 0 20px #00d97e18;
      pointer-events: none;
      transition: box-shadow .4s;
      z-index: 5;
    }
    .ai-orb.busy {
      animation: orbSpin 1.6s linear infinite;
      box-shadow: 0 0 50px #00d97e88, 0 0 100px #00d97e22, inset 0 0 30px #00d97e30;
    }
    @keyframes orbSpin {
      from { transform: translate(-50%, -50%) rotate(0deg); }
      to   { transform: translate(-50%, -50%) rotate(360deg); }
    }
    .orb-ring {
      position: absolute;
      top: 50%; left: 50%;
      transform: translate(-50%, -50%);
      width: 100px; height: 100px;
      border-radius: 50%;
      border: 1px solid var(--accent)22;
      pointer-events: none;
      z-index: 4;
      animation: ringPulse 3s ease-in-out infinite;
    }
    .orb-ring-2 {
      width: 140px; height: 140px;
      border-color: var(--accent)11;
      animation-delay: 1s;
    }
    @keyframes ringPulse {
      0%, 100% { opacity: .4; transform: translate(-50%, -50%) scale(1); }
      50% { opacity: 1; transform: translate(-50%, -50%) scale(1.04); }
    }
    .ai-orb-label {
      position: absolute;
      top: 50%; left: 50%;
      transform: translate(-50%, calc(-50% + 52px));
      font-size: 9px;
      color: var(--accent);
      font-weight: 700;
      letter-spacing: 2px;
      pointer-events: none;
      z-index: 5;
      text-transform: uppercase;
    }

    /* Canvas corner decorations */
    .corner { position: absolute; width: 20px; height: 20px; z-index: 5; opacity: .5; pointer-events: none; }
    .corner-tl { top: 16px; left: 16px; border-top: 1.5px solid var(--accent); border-left: 1.5px solid var(--accent); }
    .corner-tr { top: 16px; right: 16px; border-top: 1.5px solid var(--accent); border-right: 1.5px solid var(--accent); }
    .corner-bl { bottom: 16px; left: 16px; border-bottom: 1.5px solid var(--accent); border-left: 1.5px solid var(--accent); }
    .corner-br { bottom: 16px; right: 16px; border-bottom: 1.5px solid var(--accent); border-right: 1.5px solid var(--accent); }

    /* ══ SIDE PANEL ══ */
    .side-panel {
      width: 300px;
      flex-shrink: 0;
      background: var(--panel);
      border-left: 1px solid var(--line);
      display: flex;
      flex-direction: column;
      gap: 0;
      overflow-y: auto;
      scrollbar-width: thin;
      scrollbar-color: var(--line-2) transparent;
    }
    .side-panel::-webkit-scrollbar { width: 4px; }
    .side-panel::-webkit-scrollbar-thumb { background: var(--line-2); border-radius: 4px; }

    .panel-section {
      padding: 14px 16px;
      border-bottom: 1px solid var(--line);
    }
    .section-title {
      font-size: 10px;
      font-weight: 700;
      text-transform: uppercase;
      letter-spacing: 1.2px;
      color: var(--muted);
      margin-bottom: 12px;
      display: flex;
      align-items: center;
      gap: 6px;
    }
    .section-title::before {
      content: '';
      width: 3px; height: 10px;
      background: var(--accent);
      border-radius: 2px;
      flex-shrink: 0;
    }

    /* ── Status ── */
    .status-row {
      display: flex;
      align-items: center;
      justify-content: space-between;
    }
    .status-left { display: flex; align-items: center; gap: 10px; }
    .status-dot {
      width: 8px; height: 8px;
      border-radius: 50%;
      background: var(--accent);
      box-shadow: 0 0 8px var(--accent);
      flex-shrink: 0;
    }
    .status-dot.busy {
      background: var(--orange);
      box-shadow: 0 0 8px var(--orange);
      animation: dotBlink .8s ease-in-out infinite;
    }
    @keyframes dotBlink { 0%,100%{opacity:1} 50%{opacity:.3} }
    #statusText { font-size: 13px; font-weight: 500; }
    .status-queue {
      font-size: 10px;
      font-family: 'JetBrains Mono', monospace;
      color: var(--muted);
      background: var(--panel-2);
      border: 1px solid var(--line-2);
      border-radius: 4px;
      padding: 2px 7px;
    }

    /* ── Metrics bars ── */
    .bar-group { display: flex; flex-direction: column; gap: 11px; }
    .bar-item label {
      display: flex;
      justify-content: space-between;
      align-items: center;
      font-size: 11px;
      color: var(--muted);
      margin-bottom: 5px;
    }
    .bar-label-name { font-weight: 500; color: var(--text); }
    .bar-val {
      font-family: 'JetBrains Mono', monospace;
      font-size: 11px;
      color: var(--accent);
      min-width: 34px;
      text-align: right;
    }
    .bar-track {
      background: var(--panel-3);
      border-radius: 99px;
      height: 5px;
      overflow: hidden;
      position: relative;
    }
    .bar-fill {
      height: 100%;
      border-radius: 99px;
      transition: width .9s cubic-bezier(.4,0,.2,1);
      position: relative;
    }
    .bar-fill::after {
      content: '';
      position: absolute;
      right: 0; top: 0; bottom: 0;
      width: 4px;
      background: rgba(255,255,255,.5);
      border-radius: 99px;
      filter: blur(1px);
    }
    .bar-fill.creativity { background: linear-gradient(90deg, #00d97e, #4d9eff); }
    .bar-fill.potential  { background: linear-gradient(90deg, #4d9eff, #a78bfa); }
    .bar-fill.learning   { background: linear-gradient(90deg, #f59e0b, #ef4444); }
    .bar-fill.gpu  { background: linear-gradient(90deg, #00d97e, #06b6d4); }
    .bar-fill.cpu  { background: linear-gradient(90deg, #4d9eff, #6366f1); }
    .bar-fill.ram  { background: linear-gradient(90deg, #a78bfa, #ec4899); }

    /* ── System cards ── */
    .sys-grid { display: grid; grid-template-columns: 1fr 1fr 1fr; gap: 6px; margin-bottom: 12px; }
    .sys-card {
      background: var(--panel-2);
      border: 1px solid var(--line-2);
      border-radius: var(--radius-sm);
      padding: 10px 6px;
      text-align: center;
      position: relative;
      overflow: hidden;
      transition: border-color .2s;
    }
    .sys-card:hover { border-color: var(--accent)66; }
    .sys-card .val {
      font-size: 18px;
      font-weight: 700;
      font-family: 'JetBrains Mono', monospace;
      color: var(--accent);
      line-height: 1;
    }
    .sys-card .lbl {
      font-size: 9px;
      color: var(--muted);
      margin-top: 4px;
      letter-spacing: .5px;
      text-transform: uppercase;
    }
    .sys-card::before {
      content: '';
      position: absolute;
      top: 0; left: 0; right: 0;
      height: 2px;
      opacity: 0;
      transition: opacity .3s;
    }
    .sys-card:nth-child(1)::before { background: linear-gradient(90deg, #00d97e, #06b6d4); }
    .sys-card:nth-child(2)::before { background: linear-gradient(90deg, #4d9eff, #6366f1); }
    .sys-card:nth-child(3)::before { background: linear-gradient(90deg, #a78bfa, #ec4899); }
    .sys-card:hover::before { opacity: 1; }

    /* ── Internet search ── */
    .input-row { display: flex; gap: 6px; }
    .styled-input {
      flex: 1;
      background: var(--panel-2);
      border: 1px solid var(--line-2);
      border-radius: var(--radius-sm);
      color: var(--text);
      padding: 8px 11px;
      font-size: 12px;
      font-family: inherit;
      outline: none;
      transition: border-color .2s, box-shadow .2s;
    }
    .styled-input::placeholder { color: var(--muted-2); }
    .styled-input:focus {
      border-color: var(--accent-2);
      box-shadow: 0 0 0 3px var(--accent-2-dim);
    }
    .btn-primary {
      background: var(--accent-2);
      border: none;
      border-radius: var(--radius-sm);
      color: #fff;
      padding: 8px 13px;
      font-size: 12px;
      font-weight: 600;
      font-family: inherit;
      cursor: pointer;
      transition: opacity .15s, transform .1s;
      white-space: nowrap;
    }
    .btn-primary:hover { opacity: .88; }
    .btn-primary:active { transform: scale(.97); }
    .btn-primary:disabled { opacity: .4; cursor: default; }
    .internet-result {
      font-size: 11.5px;
      line-height: 1.5;
      color: var(--muted);
      background: var(--panel-2);
      border: 1px solid var(--line-2);
      border-radius: var(--radius-sm);
      padding: 9px 11px;
      max-height: 80px;
      overflow-y: auto;
      display: none;
      margin-top: 8px;
      font-family: 'JetBrains Mono', monospace;
    }

    /* ── Drop zone ── */
    .drop-zone {
      border: 1.5px dashed var(--line-2);
      border-radius: var(--radius);
      padding: 20px 12px;
      text-align: center;
      font-size: 12px;
      color: var(--muted);
      cursor: pointer;
      transition: border-color .2s, background .2s;
      position: relative;
    }
    .drop-zone.drag-over {
      border-color: var(--accent-2);
      background: var(--accent-2-dim);
    }
    .drop-zone input[type=file] {
      position: absolute; inset: 0; opacity: 0; cursor: pointer; width: 100%; height: 100%;
    }
    .drop-icon { font-size: 24px; margin-bottom: 6px; display: block; }
    .drop-cta { color: var(--accent); font-weight: 600; }
    .link-row { display: flex; gap: 6px; margin-top: 10px; }
    .btn-accent {
      background: transparent;
      border: 1px solid var(--accent);
      border-radius: var(--radius-sm);
      color: var(--accent);
      padding: 8px 12px;
      font-size: 12px;
      font-weight: 700;
      font-family: inherit;
      cursor: pointer;
      transition: background .15s;
      white-space: nowrap;
    }
    .btn-accent:hover { background: var(--accent-dim); }

    /* ── Neuron counter ── */
    .neuron-counter {
      padding: 14px 16px;
      display: flex;
      align-items: center;
      justify-content: space-between;
    }
    .neuron-label { font-size: 11px; color: var(--muted); text-transform: uppercase; letter-spacing: .8px; }
    .neuron-value {
      font-family: 'JetBrains Mono', monospace;
      font-size: 28px;
      font-weight: 700;
      color: var(--node-normal);
      line-height: 1;
      text-shadow: 0 0 20px var(--node-normal)88;
      transition: color .4s;
    }
    .neuron-value.active { color: var(--accent); text-shadow: 0 0 20px var(--accent)88; }

    /* Responsive */
    @media (max-width: 760px) {
      .main-layout { flex-direction: column; overflow-y: auto; height: auto; }
      .canvas-panel { min-height: 300px; flex: none; }
      .side-panel { width: 100%; border-left: none; border-top: 1px solid var(--line); }
    }
  </style>
</head>
<body>

<header>
  <a class="back-btn" href="javascript:window.close()">← Volver al chat</a>
  <div class="header-title">
    🧠 Nexo Neural
    <span class="header-badge" id="headerBadge">Online</span>
  </div>
  <div class="header-spacer"></div>
  <div class="header-clock" id="headerClock"></div>
</header>

<div class="main-layout">

  <!-- ═══ LEFT: Neural canvas ═══ -->
  <div class="canvas-panel">
    <canvas id="neuralCanvas"></canvas>
    <div class="orb-ring"></div>
    <div class="orb-ring orb-ring-2"></div>
    <div class="ai-orb" id="aiOrb"></div>
    <div class="ai-orb-label">NEXO</div>
    <div class="corner corner-tl"></div>
    <div class="corner corner-tr"></div>
    <div class="corner corner-bl"></div>
    <div class="corner corner-br"></div>
  </div>

  <!-- ═══ RIGHT: Controls ═══ -->
  <aside class="side-panel">

    <!-- ── Status ── -->
    <div class="panel-section">
      <div class="section-title">Actividad de la IA</div>
      <div class="status-row">
        <div class="status-left">
          <div class="status-dot" id="statusDot"></div>
          <span id="statusText">Online</span>
        </div>
        <span class="status-queue" id="statusQueue">Cola: 0</span>
      </div>
    </div>

    <!-- ── Cognitive metrics ── -->
    <div class="panel-section">
      <div class="section-title">Métricas cognitivas</div>
      <div class="bar-group">
        <div class="bar-item">
          <label>
            <span class="bar-label-name">Creatividad</span>
            <span class="bar-val" id="creVal">0%</span>
          </label>
          <div class="bar-track"><div class="bar-fill creativity" id="creBar" style="width:0%"></div></div>
        </div>
        <div class="bar-item">
          <label>
            <span class="bar-label-name">Potencial</span>
            <span class="bar-val" id="potVal">0%</span>
          </label>
          <div class="bar-track"><div class="bar-fill potential" id="potBar" style="width:0%"></div></div>
        </div>
        <div class="bar-item">
          <label>
            <span class="bar-label-name">Aprendizaje</span>
            <span class="bar-val" id="leaVal">0%</span>
          </label>
          <div class="bar-track"><div class="bar-fill learning" id="leaBar" style="width:0%"></div></div>
        </div>
      </div>
    </div>

    <!-- ── System monitor ── -->
    <div class="panel-section">
      <div class="section-title">Monitor del sistema</div>
      <div class="sys-grid">
        <div class="sys-card"><div class="val" id="gpuVal">—</div><div class="lbl">GPU %</div></div>
        <div class="sys-card"><div class="val" id="cpuVal">—</div><div class="lbl">CPU %</div></div>
        <div class="sys-card"><div class="val" id="ramVal">—</div><div class="lbl">RAM %</div></div>
      </div>
      <div class="bar-group">
        <div class="bar-item">
          <label><span class="bar-label-name">GPU</span><span class="bar-val" id="gpuPct">0%</span></label>
          <div class="bar-track"><div class="bar-fill gpu" id="gpuBar" style="width:0%"></div></div>
        </div>
        <div class="bar-item">
          <label><span class="bar-label-name">CPU</span><span class="bar-val" id="cpuPct">0%</span></label>
          <div class="bar-track"><div class="bar-fill cpu" id="cpuBar" style="width:0%"></div></div>
        </div>
        <div class="bar-item">
          <label><span class="bar-label-name">RAM</span><span class="bar-val" id="ramPct">0%</span></label>
          <div class="bar-track"><div class="bar-fill ram" id="ramBar" style="width:0%"></div></div>
        </div>
      </div>
    </div>

    <!-- ── Internet search ── -->
    <div class="panel-section">
      <div class="section-title">🌐 Internet</div>
      <div class="input-row">
        <input type="text" class="styled-input" id="searchInput" placeholder="Buscar en internet...">
        <button class="btn-primary" id="searchBtn" onclick="doSearch()">Buscar</button>
      </div>
      <div class="internet-result" id="searchResult"></div>
    </div>

    <!-- ── Feed neural ── -->
    <div class="panel-section">
      <div class="section-title">📎 Alimentar la red</div>
      <div class="drop-zone" id="dropZone">
        <input type="file" id="fileInput" multiple onchange="handleFiles(this.files)">
        <span class="drop-icon">📂</span>
        <div>Arrastra archivos aquí</div>
        <div><span class="drop-cta">o haz clic para seleccionar</span></div>
      </div>
      <div class="link-row">
        <input type="text" class="styled-input" id="linkInput" placeholder="Pega un link aquí...">
        <button class="btn-accent" onclick="handleLink()">+ Link</button>
      </div>
    </div>

    <!-- ── Neuron counter ── -->
    <div class="neuron-counter">
      <div>
        <div class="neuron-label">Neuronas activas</div>
      </div>
      <div class="neuron-value" id="neuronCount">10</div>
    </div>

  </aside>
</div>

<script>
// ═══════════════════════════════════════════════
//  STATE
// ═══════════════════════════════════════════════
const canvas = document.getElementById('neuralCanvas');
const ctx    = canvas.getContext('2d');
const aiOrb  = document.getElementById('aiOrb');

let nodes = [];
let busy  = false;

const COLOR_NORMAL  = '#4d9eff';
const COLOR_NEW     = '#f59e0b';
const COLOR_CENTRAL = '#00d97e';
const LINE_COLOR    = '#4d9eff';

// ─── Clock ───
function updateClock() {
  const now = new Date();
  const hh = String(now.getHours()).padStart(2,'0');
  const mm = String(now.getMinutes()).padStart(2,'0');
  const ss = String(now.getSeconds()).padStart(2,'0');
  document.getElementById('headerClock').textContent = hh + ':' + mm + ':' + ss;
}
updateClock();
setInterval(updateClock, 1000);

// ─── Canvas resize ───
function resizeCanvas() {
  const rect = canvas.parentElement.getBoundingClientRect();
  canvas.width  = rect.width;
  canvas.height = rect.height;
}
window.addEventListener('resize', () => { resizeCanvas(); initNodes(); });
resizeCanvas();

// ─── Init nodes ───
function initNodes() {
  nodes = [];
  const cx = canvas.width / 2;
  const cy = canvas.height / 2;
  const count = Math.max(14, Math.floor(canvas.width * canvas.height / 12000));
  for (let i = 0; i < count; i++) {
    let x, y;
    do {
      x = Math.random() * canvas.width;
      y = Math.random() * canvas.height;
    } while (Math.hypot(x - cx, y - cy) < 90);
    nodes.push(makeNode(x, y, COLOR_NORMAL));
  }
}

function makeNode(x, y, color) {
  return {
    x, y,
    vx: (Math.random() - .5) * .4,
    vy: (Math.random() - .5) * .4,
    r: 3.5 + Math.random() * 3,
    color,
    newborn: color === COLOR_NEW
  };
}

// ─── Add neuron (visual only) ───
function addNeuron(label) {
  const margin = 40;
  const cx = canvas.width / 2;
  const cy = canvas.height / 2;
  let x, y;
  do {
    x = margin + Math.random() * (canvas.width  - margin * 2);
    y = margin + Math.random() * (canvas.height - margin * 2);
  } while (Math.hypot(x - cx, y - cy) < 90);
  const extra = 1 + Math.floor(Math.random() * 2);
  for (let i = 0; i <= extra; i++) {
    nodes.push(makeNode(
      x + (Math.random() - .5) * 60,
      y + (Math.random() - .5) * 60,
      COLOR_NEW
    ));
  }
  setTimeout(() => {
    nodes.forEach(n => { if (n.newborn) { n.color = COLOR_NORMAL; n.newborn = false; } });
  }, 3000);
}

// ─── Animation loop ───
function animate() {
  requestAnimationFrame(animate);
  ctx.clearRect(0, 0, canvas.width, canvas.height);

  const cx = canvas.width / 2;
  const cy = canvas.height / 2;

  for (let i = 0; i < nodes.length; i++) {
    for (let j = i + 1; j < nodes.length; j++) {
      const d = Math.hypot(nodes[i].x - nodes[j].x, nodes[i].y - nodes[j].y);
      if (d < 120) {
        ctx.beginPath();
        ctx.moveTo(nodes[i].x, nodes[i].y);
        ctx.lineTo(nodes[j].x, nodes[j].y);
        const a = Math.floor((1 - d / 120) * 55).toString(16).padStart(2,'0');
        ctx.strokeStyle = LINE_COLOR + a;
        ctx.lineWidth = .5;
        ctx.stroke();
      }
    }
    const dc = Math.hypot(nodes[i].x - cx, nodes[i].y - cy);
    if (dc < 260) {
      ctx.beginPath();
      ctx.moveTo(nodes[i].x, nodes[i].y);
      ctx.lineTo(cx, cy);
      const a = Math.floor((1 - dc / 260) * 35).toString(16).padStart(2,'0');
      ctx.strokeStyle = COLOR_CENTRAL + a;
      ctx.lineWidth = .35;
      ctx.stroke();
    }
  }

  for (const n of nodes) {
    n.x += n.vx; n.y += n.vy;
    if (n.x < 0 || n.x > canvas.width)  n.vx *= -1;
    if (n.y < 0 || n.y > canvas.height) n.vy *= -1;

    const grd = ctx.createRadialGradient(n.x, n.y, 0, n.x, n.y, n.r * 3);
    grd.addColorStop(0, n.color + 'bb');
    grd.addColorStop(1, n.color + '00');
    ctx.beginPath();
    ctx.arc(n.x, n.y, n.r * 3, 0, Math.PI * 2);
    ctx.fillStyle = grd;
    ctx.fill();

    ctx.beginPath();
    ctx.arc(n.x, n.y, n.r, 0, Math.PI * 2);
    ctx.fillStyle = n.color;
    ctx.fill();
  }
}

initNodes();
animate();

// ─── Busy state ───
function setBusy(val) {
  busy = val;
  aiOrb.classList.toggle('busy', val);
  const dot = document.getElementById('statusDot');
  const txt = document.getElementById('statusText');
  const badge = document.getElementById('headerBadge');
  dot.classList.toggle('busy', val);
  txt.textContent   = val ? 'Procesando...' : 'Online';
  badge.textContent = val ? 'Busy' : 'Online';
  badge.style.color = val ? 'var(--orange)' : '';
  badge.style.background = val ? 'var(--orange)22' : '';
  badge.style.borderColor = val ? 'var(--orange)44' : '';
}

// ─── Metric helpers ───
function boostMetrics() { /* now driven by real data from fetchStats */ }
function updateBars()   { /* now driven by real data from fetchStats */ }

// ─── Real stats from /api/neural/stats ───
async function fetchStats() {
  try {
    const r = await fetch('/api/neural/stats');
    if (!r.ok) return;
    const d = await r.json();

    const gpu = d.gpu_load  ?? 0;
    const cpu = d.cpu_load  ?? 0;
    const ram = d.ram_pct   ?? 0;
    const cre = d.creativity ?? gpu;
    const pot = d.potential  ?? cpu;
    const lea = d.learning   ?? ram;
    const nc  = d.neuron_count ?? 10;
    const isBusy = d.status === 'Procesando...';

    // System cards
    document.getElementById('gpuVal').textContent = gpu;
    document.getElementById('cpuVal').textContent = cpu;
    document.getElementById('ramVal').textContent = ram;

    // System bars
    document.getElementById('gpuBar').style.width = gpu + '%';
    document.getElementById('gpuPct').textContent = gpu + '%';
    document.getElementById('cpuBar').style.width = cpu + '%';
    document.getElementById('cpuPct').textContent = cpu + '%';
    document.getElementById('ramBar').style.width = ram + '%';
    document.getElementById('ramPct').textContent = ram + '%';

    // Cognitive bars
    document.getElementById('creBar').style.width = cre + '%';
    document.getElementById('creVal').textContent = Math.round(cre) + '%';
    document.getElementById('potBar').style.width = pot + '%';
    document.getElementById('potVal').textContent = Math.round(pot) + '%';
    document.getElementById('leaBar').style.width = lea + '%';
    document.getElementById('leaVal').textContent = Math.round(lea) + '%';

    // Status
    const dot   = document.getElementById('statusDot');
    const txt   = document.getElementById('statusText');
    const badge = document.getElementById('headerBadge');
    const queue = document.getElementById('statusQueue');
    dot.classList.toggle('busy', isBusy);
    txt.textContent   = d.status || 'Online';
    badge.textContent = isBusy ? 'Busy' : 'Online';
    badge.style.color = isBusy ? 'var(--orange)' : '';
    badge.style.background = isBusy ? 'var(--orange)22' : '';
    badge.style.borderColor = isBusy ? 'var(--orange)44' : '';
    aiOrb.classList.toggle('busy', isBusy);
    queue.textContent = 'Cola: ' + (d.queued ?? 0);

    // Neuron counter: use canvas nodes.length (real visual count), not the server value
    const ncEl = document.getElementById('neuronCount');
    ncEl.textContent = nodes.length;
    ncEl.classList.toggle('active', isBusy);

  } catch(_) {}
}
fetchStats();
setInterval(fetchStats, 5000);

// ─── Internet search ───
async function doSearch() {
  const q = document.getElementById('searchInput').value.trim();
  if (!q) return;
  const btn = document.getElementById('searchBtn');
  btn.disabled = true;
  setBusy(true);
  const resultBox = document.getElementById('searchResult');
  resultBox.style.display = 'block';
  resultBox.textContent = 'Buscando...';
  try {
    const r = await fetch('/api/neural/process', {
      method: 'POST',
      headers: {'Content-Type':'application/json'},
      body: JSON.stringify({type:'search', content: q})
    });
    const d = await r.json();
    resultBox.textContent = d.message || 'Procesado';
    for (let i = 0; i < (d.added || 2); i++) addNeuron('search');
  } catch(e) {
    resultBox.textContent = 'Error: ' + e.message;
  }
  setBusy(false);
  btn.disabled = false;
}
document.getElementById('searchInput').addEventListener('keydown', e => { if (e.key === 'Enter') doSearch(); });

// ─── File drop zone ───
const dropZone = document.getElementById('dropZone');
dropZone.addEventListener('dragover',  e => { e.preventDefault(); dropZone.classList.add('drag-over'); });
dropZone.addEventListener('dragleave', () => dropZone.classList.remove('drag-over'));
dropZone.addEventListener('drop', e => {
  e.preventDefault();
  dropZone.classList.remove('drag-over');
  handleFiles(e.dataTransfer.files);
});

async function handleFiles(files) {
  if (!files || !files.length) return;
  setBusy(true);
  for (const file of files) {
    try {
      const r = await fetch('/api/neural/process', {
        method: 'POST',
        headers: {'Content-Type':'application/json'},
        body: JSON.stringify({type:'file', content: file.name})
      });
      const d = await r.json();
      for (let i = 0; i < (d.added || 2); i++) addNeuron('file');
    } catch(_) {
      addNeuron('file');
    }
  }
  setBusy(false);
}

// ─── Link ───
async function handleLink() {
  const url = document.getElementById('linkInput').value.trim();
  if (!url) return;
  setBusy(true);
  try {
    const r = await fetch('/api/neural/process', {
      method: 'POST',
      headers: {'Content-Type':'application/json'},
      body: JSON.stringify({type:'link', content: url})
    });
    const d = await r.json();
    for (let i = 0; i < (d.added || 2); i++) addNeuron('link');
  } catch(_) {
    addNeuron('link');
  }
  document.getElementById('linkInput').value = '';
  setBusy(false);
}
</script>
</body>
</html>
"""


MAIN_HTML = r"""
<!doctype html>
<html lang="es">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>Nexo</title>
  <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/themes/prism-tomorrow.min.css">
  <style>
    :root {
      color-scheme: dark;
      --bg: #212121;
      --sidebar: #171717;
      --sidebar-2: #202020;
      --panel: #2b2b2b;
      --panel-soft: #303030;
      --text: #f4f4f4;
      --muted: #b4b4b4;
      --muted-2: #8c8c8c;
      --line: #3b3b3b;
      --accent: #19c37d;
      --accent-2: #2dd4bf;
      --danger: #ff6b6b;
      --code: #111111;
    }
    * { box-sizing: border-box; }
    html, body { height: 100%; }
    html { overflow-x: hidden; }
    body {
      margin: 0;
      background: var(--bg);
      color: var(--text);
      font-family: Inter, ui-sans-serif, system-ui, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
      overflow: hidden;
    }
    button, input, textarea, select { font: inherit; }
    button { cursor: pointer; }
    .app {
      display: grid;
      grid-template-columns: 292px 1fr;
      height: 100vh;
      min-height: 0;
    }
    body.sidebar-collapsed .app { grid-template-columns: 0 1fr; }
    .sidebar {
      display: flex;
      flex-direction: column;
      min-height: 0;
      background: var(--sidebar);
      border-right: 1px solid #111;
    }
    body.sidebar-collapsed .sidebar { display: none; }
    .brand {
      display: flex;
      align-items: center;
      gap: 10px;
      padding: 14px;
      border-bottom: 1px solid #242424;
    }
    .mark {
      display: grid;
      place-items: center;
      width: 32px;
      height: 32px;
      border-radius: 6px;
      background: #f2f2f2;
      color: #111;
      font-weight: 800;
    }
    .brand strong { font-size: 15px; }
    .brand span { display: block; color: var(--muted-2); font-size: 12px; margin-top: 1px; }
    .side-actions { padding: 12px; display: grid; gap: 8px; }
    .select-group { display: grid; gap: 4px; }
    .select-label { font-size: 11px; color: var(--muted-2); text-transform: uppercase; letter-spacing: .06em; padding: 0 2px; }
    .plan-card {
      margin: 0 12px 12px;
      border: 1px solid var(--line);
      border-radius: 8px;
      background: #141414;
      padding: 10px;
      display: grid;
      gap: 8px;
      font-size: 12px;
    }
    .plan-row {
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 8px;
    }
    .plan-row strong { font-size: 13px; }
    .plan-row span, .plan-modes { color: var(--muted-2); }
    .plan-api-btn { min-height: 34px; font-size: 12px; }
    .api-key-preview {
      display: block;
      min-width: 0;
      overflow-wrap: anywhere;
      padding: 8px;
      border-radius: 6px;
      background: #0e0e0e;
      color: #dcdcdc;
      border: 1px solid #262626;
    }
    .btn {
      display: inline-flex;
      align-items: center;
      justify-content: center;
      gap: 8px;
      min-height: 44px;
      border: 1px solid var(--line);
      border-radius: 6px;
      background: var(--sidebar-2);
      color: var(--text);
      padding: 0 12px;
    }
    .btn:hover { background: #282828; }
    .btn.primary { background: var(--accent); border-color: var(--accent); color: #07130e; font-weight: 700; }
    .btn.ghost { background: transparent; }
    .btn.danger { border-color: rgba(255, 107, 107, .45); color: #ffd1d1; }
    a.btn { text-decoration: none; }
    select {
      width: 100%;
      height: 44px;
      border: 1px solid var(--line);
      border-radius: 6px;
      background: #121212;
      color: var(--text);
      padding: 0 10px;
      outline: none;
    }
    .history-title {
      padding: 8px 14px;
      color: var(--muted-2);
      font-size: 12px;
      text-transform: uppercase;
      letter-spacing: .04em;
    }
    .chat-list {
      overflow: auto;
      padding: 0 8px 12px;
      flex: 1;
      min-height: 0;
    }
    .chat-item {
      width: 100%;
      text-align: left;
      border: 0;
      border-radius: 6px;
      background: transparent;
      color: #e8e8e8;
      padding: 10px 10px;
      margin: 2px 0;
      overflow: hidden;
      text-overflow: ellipsis;
      white-space: nowrap;
    }
    .chat-item:hover, .chat-item.active { background: #242424; }
    .chat-row {
      display: flex;
      align-items: center;
      gap: 4px;
      margin: 2px 0;
      border-radius: 6px;
    }
    .chat-row:hover, .chat-row.active { background: #242424; }
    .chat-row .chat-item {
      flex: 1;
      margin: 0;
      min-width: 0;
      border-radius: 6px 0 0 6px;
    }
    .chat-row:hover .chat-item, .chat-row.active .chat-item { background: transparent; }
    .chat-delete-btn {
      flex-shrink: 0;
      display: none;
      align-items: center;
      justify-content: center;
      width: 28px;
      height: 36px;
      border: 0;
      border-radius: 0 6px 6px 0;
      background: transparent;
      color: var(--muted);
      font-size: 14px;
      cursor: pointer;
      padding: 0;
    }
    .chat-row:hover .chat-delete-btn { display: flex; }
    .chat-delete-btn:hover { color: var(--danger); }
    .sidebar-footer {
      padding: 10px 12px;
      border-top: 1px solid #242424;
      display: grid;
      gap: 7px;
    }
    .footer-row {
      display: grid;
      grid-template-columns: 1fr 1fr;
      gap: 7px;
    }
    .footer-row-3 {
      display: grid;
      grid-template-columns: 1fr 1fr 1fr;
      gap: 6px;
    }
    .footer-row .btn, .footer-row-3 .btn { min-height: 36px; font-size: 12px; padding: 0 6px; }
    .footer-row-3 .btn { min-height: 34px; font-size: 11px; padding: 0 4px; }
    .main {
      display: flex;
      flex-direction: column;
      min-width: 0;
      min-height: 0;
      overflow: hidden;
      background: var(--bg);
      grid-column: 2;
    }
    .topbar {
      height: 58px;
      display: flex;
      align-items: center;
      justify-content: space-between;
      padding: 0 18px;
      border-bottom: 1px solid var(--line);
      background: var(--sidebar);
    }
    .topbar h2 { margin: 0; font-size: 16px; letter-spacing: 0; }
    .topbar-actions { display: flex; align-items: center; gap: 8px; }
    .mobile-menu { display: none; }
    .desktop-sidebar-toggle { display: inline-flex; }
    .scrim {
      position: fixed;
      inset: 0;
      display: none;
      background: rgba(0, 0, 0, .48);
      z-index: 9;
    }
    .messages {
      flex: 1;
      min-height: 0;
      overflow: auto;
      padding: 24px 18px 190px;
    }
    .empty {
      min-height: calc(100% - 214px);
      display: flex;
      flex-direction: column;
      align-items: center;
      justify-content: center;
      text-align: center;
      color: var(--muted);
      padding: 24px;
    }
    .empty h1 { margin: 0 0 10px; color: var(--text); font-size: clamp(28px, 5vw, 44px); letter-spacing: 0; }
    .empty p { margin: 0; max-width: 560px; line-height: 1.55; }
    .message {
      width: min(860px, 100%);
      margin: 0 auto 20px;
      display: grid;
      grid-template-columns: 34px minmax(0, 1fr);
      gap: 12px;
    }
    .avatar {
      width: 34px;
      height: 34px;
      border-radius: 6px;
      display: grid;
      place-items: center;
      font-size: 13px;
      font-weight: 800;
      background: #f2f2f2;
      color: #111;
    }
    .message.user .avatar { background: var(--accent-2); color: #06211e; }
    .bubble {
      min-width: 0;
      padding: 6px 0 0;
      line-height: 1.6;
      color: #f2f2f2;
      word-wrap: break-word;
    }
    .message.user .bubble {
      justify-self: end;
      max-width: min(680px, 100%);
      background: var(--panel-soft);
      padding: 10px 14px;
      border-radius: 8px;
    }
    .bubble p { margin: 0 0 12px; }
    .bubble p:last-child { margin-bottom: 0; }
    .bubble ul, .bubble ol { margin: 8px 0 12px 22px; padding: 0; }
    .bubble pre {
      position: relative;
      margin: 12px 0;
      border-radius: 8px;
      overflow: hidden;
      background: var(--code);
      border: 1px solid #303030;
    }
    .bubble code {
      font-family: "Cascadia Code", Consolas, "SFMono-Regular", monospace;
      font-size: 13px;
    }
    .bubble pre code {
      display: block;
      padding: 42px 14px 14px;
      overflow-x: auto;
      white-space: pre;
    }
    .code-copy {
      position: absolute;
      top: 8px;
      right: 8px;
      height: 28px;
      border: 1px solid #3f3f3f;
      border-radius: 6px;
      background: #202020;
      color: #e8e8e8;
      padding: 0 10px;
      font-size: 12px;
    }
    .attachment-list, .source-list {
      display: grid;
      gap: 8px;
      margin-top: 10px;
    }
    .attachment-item, .source-item {
      border: 1px solid var(--line);
      border-radius: 6px;
      background: rgba(255, 255, 255, .035);
      padding: 9px 10px;
      font-size: 13px;
      color: #dedede;
    }
    .attachment-item strong, .source-item strong {
      display: block;
      color: var(--text);
      font-size: 13px;
      margin-bottom: 3px;
      overflow-wrap: anywhere;
    }
    .attachment-item span, .source-item span {
      color: var(--muted-2);
      font-size: 12px;
    }
    .source-item a {
      color: var(--accent-2);
      text-decoration: none;
      overflow-wrap: anywhere;
    }
    .source-item a:hover { text-decoration: underline; }
    .status {
      width: min(860px, 100%);
      margin: 0 auto 12px;
      color: var(--muted-2);
      font-size: 13px;
    }
    .composer {
      position: fixed;
      left: 292px;
      right: 0;
      bottom: 0;
      padding: 16px 18px 20px;
      background: linear-gradient(to top, var(--bg) 72%, rgba(33, 33, 33, 0));
      transition: left .18s ease;
    }
    body.sidebar-collapsed .composer { left: 0; }
    .composer.dragging .composer-box {
      border-color: var(--accent-2);
      background: #263430;
    }
    .pending-files {
      width: min(860px, 100%);
      margin: 0 auto 8px;
      display: none;
      gap: 8px;
      grid-template-columns: repeat(auto-fit, minmax(180px, 1fr));
    }
    .pending-files.open { display: grid; }
    .pending-file {
      min-width: 0;
      display: grid;
      grid-template-columns: 42px minmax(0, 1fr) 28px;
      gap: 8px;
      align-items: center;
      border: 1px solid var(--line);
      border-radius: 6px;
      background: var(--panel);
      padding: 7px;
    }
    .pending-thumb {
      width: 42px;
      height: 42px;
      border-radius: 6px;
      object-fit: cover;
      background: #111;
      display: grid;
      place-items: center;
      color: var(--muted);
      font-size: 11px;
      border: 1px solid #383838;
    }
    .pending-file strong {
      display: block;
      overflow: hidden;
      white-space: nowrap;
      text-overflow: ellipsis;
      font-size: 13px;
    }
    .pending-file span {
      display: block;
      color: var(--muted-2);
      font-size: 12px;
      margin-top: 2px;
    }
    .remove-file, .attach {
      width: 44px;
      height: 44px;
      border: 1px solid var(--line);
      border-radius: 6px;
      background: #202020;
      color: var(--text);
    }
    .remove-file {
      width: 28px;
      height: 28px;
      color: #ffd0d0;
    }
    .composer-box {
      width: min(860px, 100%);
      margin: 0 auto;
      display: grid;
      grid-template-columns: 44px 1fr 44px 44px;
      gap: 10px;
      align-items: end;
      border: 1px solid var(--line);
      border-radius: 8px;
      background: var(--panel);
      padding: 10px;
      box-shadow: 0 18px 50px rgba(0, 0, 0, .18);
    }
    textarea {
      resize: none;
      max-height: 180px;
      min-height: 42px;
      font-size: 16px;
      border: 0;
      outline: none;
      background: transparent;
      color: var(--text);
      padding: 10px 8px;
      line-height: 1.45;
    }
    .send {
      width: 44px;
      height: 44px;
      border: 0;
      border-radius: 6px;
      background: var(--accent);
      color: #06140e;
      font-weight: 900;
    }
    .send:disabled { opacity: .5; cursor: not-allowed; }
    .bubble.streaming::after {
      content: '▋';
      display: inline-block;
      color: var(--accent);
      animation: blink .7s step-end infinite;
      margin-left: 2px;
      font-size: 14px;
      vertical-align: middle;
    }
    @keyframes blink {
      0%, 100% { opacity: 1; }
      50% { opacity: 0; }
    }
    .topbar-right { display: flex; align-items: center; gap: 10px; }
    .topbar-memory-btn { min-height: 34px; font-size: 12px; padding: 0 10px; border-color: rgba(255,255,255,.12); }
    .topbar-memory-btn:hover { border-color: var(--accent); color: var(--accent); }
    .topbar-status { display: flex; align-items: center; gap: 6px; font-size: 13px; color: var(--muted); }
    .topbar-status.busy { color: var(--accent); }
    .topbar-status.busy::before {
      content: '';
      display: inline-block;
      width: 7px;
      height: 7px;
      border-radius: 50%;
      background: var(--accent);
      animation: pulse .9s ease-in-out infinite;
    }
    @keyframes pulse {
      0%, 100% { opacity: 1; transform: scale(1); }
      50% { opacity: .4; transform: scale(.75); }
    }
    .report-form { display: grid; gap: 12px; margin-top: 8px; }
    .report-textarea {
      min-height: 130px;
      width: 100%;
      resize: vertical;
      background: #101010;
      border: 1px solid #2a2a2a;
      border-radius: 6px;
      padding: 10px 12px;
      color: #e8e8e8;
      font-size: 14px;
      line-height: 1.5;
      outline: none;
      transition: border-color .15s;
    }
    .report-textarea:focus { border-color: var(--accent); }
    .report-category {
      background: #101010;
      border: 1px solid #2a2a2a;
      border-radius: 6px;
      padding: 8px 10px;
      color: #e8e8e8;
      font-size: 13px;
      width: 100%;
      outline: none;
    }
    .report-list { display: grid; gap: 8px; margin-top: 12px; max-height: 300px; overflow: auto; }
    .report-item {
      border: 1px solid var(--line);
      border-radius: 6px;
      background: #141414;
      padding: 10px 12px;
      font-size: 13px;
    }
    .report-item strong { color: var(--accent); display: block; margin-bottom: 4px; font-size: 11px; }
    .report-item pre { margin: 0; white-space: pre-wrap; word-break: break-word; color: #c8c8c8; line-height: 1.5; }
    .report-item .report-cat { display: inline-block; background: #222; border-radius: 4px; padding: 1px 7px; font-size: 11px; color: var(--muted-2); margin-bottom: 4px; }
    .modal-backdrop {
      position: fixed;
      inset: 0;
      display: none;
      align-items: center;
      justify-content: center;
      background: rgba(0, 0, 0, .55);
      z-index: 20;
      padding: 18px;
    }
    .modal-backdrop.open { display: flex; }
    .modal {
      width: min(680px, 100%);
      max-height: min(680px, calc(100vh - 40px));
      overflow: auto;
      background: #181818;
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 18px;
    }
    .modal header {
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 12px;
      margin-bottom: 12px;
    }
    .modal h3 { margin: 0; font-size: 17px; }
    .memory-text {
      min-height: 160px;
      white-space: pre-wrap;
      color: #e8e8e8;
      background: #101010;
      border: 1px solid #2a2a2a;
      border-radius: 6px;
      padding: 12px;
      line-height: 1.55;
    }
    /* ── Tablet (1025 → 760) ─────────────────────────── */
    @media (max-width: 1025px) {
      .app { grid-template-columns: 248px 1fr; }
      body.sidebar-collapsed .app { grid-template-columns: 0 1fr; }
      .composer { left: 248px; }
      body.sidebar-collapsed .composer { left: 0; }
    }
    /* ── Mobile (≤ 760px) ────────────────────────────── */
    @media (max-width: 760px) {
      .app { grid-template-columns: 1fr; }
      .main { grid-column: 1; }
      .sidebar {
        position: fixed;
        inset: 0 auto 0 0;
        width: min(292px, 86vw);
        z-index: 10;
        transform: translateX(-100%);
        transition: transform .18s ease;
      }
      .sidebar.open { transform: translateX(0); }
      .mobile-menu { display: inline-flex; }
      .desktop-sidebar-toggle { display: none; }
      body.sidebar-open .scrim { display: block; }
      .composer { left: 0; }
      .topbar { padding: 0 12px; }
      .messages { padding-inline: 12px; padding-bottom: 160px; }
      .message { grid-template-columns: 30px minmax(0, 1fr); gap: 10px; }
      .avatar { width: 30px; height: 30px; }
      .composer { padding: 10px 12px calc(16px + env(safe-area-inset-bottom, 0px)); }
      .composer-box { gap: 7px; padding: 8px; }
      textarea { font-size: 15px; }
    }
    /* ── Small phones (≤ 480px) ──────────────────────── */
    @media (max-width: 480px) {
      .empty h1 { font-size: 26px; }
      .empty p { font-size: 14px; }
      .topbar h2 { font-size: 14px; max-width: 160px; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }
      .topbar-status { font-size: 12px; }
      .composer-box { grid-template-columns: 38px 1fr 38px 38px; gap: 5px; padding: 7px; }
      .attach, .send { width: 38px; height: 38px; }
      textarea { font-size: 14px; min-height: 38px; }
      .bubble { font-size: 14px; }
      .bubble pre code { font-size: 12px; }
      .modal { padding: 14px; }
    }
  </style>
</head>
<body>
  <div class="app">
    <aside id="sidebar" class="sidebar">
      <div class="brand">
        <div class="mark">NX</div>
        <div><strong>Nexo</strong><span>Asistente IA · Chat inteligente</span></div>
      </div>
      <div class="side-actions">
        <button id="newChatBtn" class="btn primary" type="button">+ Nuevo chat</button>
        <div class="select-group">
          <label class="select-label" for="modeSelect">Modo de respuesta</label>
          <select id="modeSelect" aria-label="Modo de Nexo">
            <option value="auto" selected>&#9889; Auto</option>
            <option value="rapido">&#128640; Rapido</option>
            <option value="combinado">&#128256; Combinado</option>
            <option value="codigo">&#128187; Codigo</option>
          </select>
        </div>
        <div class="select-group">
          <label class="select-label" for="personalitySelect">Personalidad</label>
          <select id="personalitySelect" aria-label="Personalidad de Nexo" title="Tono de respuesta">
            <option value="normal" selected>&#129302; Normal</option>
            <option value="programador">&#128187; Programador</option>
            <option value="creativo">&#10024; Creativo</option>
            <option value="conciso">&#9889; Conciso</option>
          </select>
        </div>
      </div>
      <div id="planCard" class="plan-card">
        <div class="plan-row"><strong id="planName">Plan</strong><span id="planPrice"></span></div>
        <div id="planModes" class="plan-modes" style="display:none;"></div>
        <div id="planExpiry" style="display:none;font-size:11px;color:#f0c040;margin-top:4px;"></div>
        <button id="apiKeyBtn" class="btn ghost plan-api-btn" type="button" hidden>API Key</button>
        <code id="apiKeyPreview" class="api-key-preview" hidden></code>
      </div>
      <div id="promoBanner" style="display:none;background:linear-gradient(135deg,rgba(240,192,64,.15),rgba(240,192,64,.05));border:1px solid rgba(240,192,64,.35);border-radius:8px;padding:10px 12px;margin:8px 0;font-size:12px;color:#f0c040;"></div>
      <div class="history-title">Chats</div>
      <div id="chatList" class="chat-list"></div>
      <div class="sidebar-footer">
        <div id="gpuWidget" style="display:none;font-size:11px;color:var(--muted);padding:8px 10px;background:var(--sidebar-2);border-radius:8px;border:1px solid var(--line);">
          <div style="display:flex;justify-content:space-between;margin-bottom:4px;font-weight:600;color:var(--text);"><span id="gpuName">💻 GPU</span><span id="gpuLoadLabel">—</span></div>
          <div style="background:var(--line);border-radius:4px;height:5px;margin-bottom:5px;"><div id="gpuBar" style="background:var(--accent);height:5px;border-radius:4px;width:0%;transition:width .6s ease;"></div></div>
          <div style="display:flex;justify-content:space-between;margin-bottom:2px;"><span id="vramText">VRAM</span><span id="cpuLabel">CPU —</span></div>
          <div style="background:var(--line);border-radius:4px;height:3px;margin-bottom:5px;"><div id="cpuBar" style="background:var(--accent-2);height:3px;border-radius:4px;width:0%;transition:width .6s ease;"></div></div>
          <div id="queueWidget" style="display:none;color:var(--accent-2);">Cola GPU: <span id="queueCount">0</span> petición(es)</div>
        </div>
        <div class="footer-row">
          <a class="btn ghost" href="/planes" target="_blank" rel="noopener">⬆ Planes</a>
          {% if is_admin %}<a class="btn ghost" href="#" onclick="openNeural();return false;" style="flex:1;min-height:36px;font-size:12px;padding:0 6px;">🧠 Neural</a>{% endif %}
          <a class="btn ghost" href="/donate" target="_blank" rel="noopener">💛 Donar</a>
        </div>
        <div id="footerActionsRow" style="display:flex;gap:7px;">
          <button id="reportBtn" class="btn ghost" type="button" style="flex:1;min-height:36px;font-size:12px;padding:0 6px;color:#ff9f43;border-color:rgba(255,159,67,.35)">🐛 Reportar</button>
          <a id="adminBtn" class="btn ghost" href="/admin" style="display:none;flex:1;min-height:36px;font-size:12px;padding:0 6px;color:#f0c040;border-color:rgba(240,192,64,.35)">👑 Admin</a>
          <a id="statsBtn" class="btn ghost" href="/admin/stats" style="display:none;flex:1;min-height:36px;font-size:12px;padding:0 6px;color:#7dd3fc;border-color:rgba(125,211,252,.35)">📊 Stats</a>
        </div>
        <form method="post" action="/logout">
          <button class="btn danger" type="submit" style="width:100%;">Salir</button>
        </form>
      </div>
    </aside>

    <div id="scrim" class="scrim" aria-hidden="true"></div>

    <main class="main">
      <header class="topbar">
        <div class="topbar-actions">
          <button id="menuBtn" class="btn ghost mobile-menu" type="button">☰</button>
          <button id="desktopSidebarBtn" class="btn ghost desktop-sidebar-toggle" type="button" aria-label="Mostrar u ocultar chats" title="Mostrar u ocultar chats">≡</button>
          <h2 id="chatTitle">Nuevo chat</h2>
        </div>
        <div class="topbar-right">
          <button id="memoryBtn" class="btn ghost topbar-memory-btn" type="button" title="Memoria compartida">🧠 Memoria</button>
          <div id="status" class="topbar-status"></div>
        </div>
      </header>
      <section id="messages" class="messages">
        <div class="empty">
          <div>
            <h1>¿Qué hacemos hoy?</h1>
            <p>Elige modo, escribe tu mensaje y Nexo responderá usando memoria local.</p>
          </div>
        </div>
      </section>
      <section id="composer" class="composer">
        <div id="pendingFiles" class="pending-files"></div>
        <div class="composer-box">
          <button id="attachBtn" class="attach" type="button" aria-label="Adjuntar archivo" title="Adjuntar archivo">+</button>
          <input id="fileInput" type="file" multiple hidden>
          <textarea id="prompt" rows="1" placeholder="Mensaje para Nexo"></textarea>
          <button id="micBtn" class="attach" type="button" aria-label="Hablar" title="Voice-to-Text" style="font-size:15px;">🎤</button>
          <button id="sendBtn" class="send" type="button" aria-label="Enviar">↑</button>
        </div>
      </section>
    </main>
  </div>

  <div id="memoryModal" class="modal-backdrop">
    <div class="modal">
      <header>
        <h3>Memoria compartida</h3>
        <button id="closeMemoryBtn" class="btn ghost" type="button">Cerrar</button>
      </header>
      <div id="memoryText" class="memory-text"></div>
      <button id="clearMemoryBtn" class="btn danger" type="button" style="margin-top: 12px;">Borrar memoria</button>
    </div>
  </div>

  <div id="reportModal" class="modal-backdrop">
    <div class="modal" style="width:min(600px,100%)">
      <header>
        <h3>🐛 Reportar bug</h3>
        <button id="closeReportBtn" class="btn ghost" type="button">Cerrar</button>
      </header>
      <div class="report-form">
        <select id="reportCategory" class="report-category">
          <option value="ui">🖥️ Interfaz / Visual</option>
          <option value="chat">💬 Chat / Respuestas</option>
          <option value="archivos">📎 Archivos / Adjuntos</option>
          <option value="rendimiento">⚡ Rendimiento</option>
          <option value="otro">🔧 Otro</option>
        </select>
        <textarea id="reportText" class="report-textarea" placeholder="Describe el bug: qué pasó, qué esperabas que pasara, y si puedes reproducirlo..."></textarea>
        <div style="display:flex;gap:8px;">
          <button id="submitReportBtn" class="btn primary" type="button">Enviar reporte</button>
          <button id="viewReportsBtn" class="btn ghost" type="button">Ver reportes guardados</button>
        </div>
        <div id="reportStatus" style="font-size:13px;color:var(--muted);display:none;"></div>
      </div>
      <div id="reportList" class="report-list" style="display:none;"></div>
    </div>
  </div>

  <script>
    const state = { chats: [], activeChatId: null, busy: false, pendingFiles: [], account: null, plans: [] };
    const SIDEBAR_COLLAPSED_KEY = 'nexo_sidebar_collapsed_v1';
    const els = {
      composer: document.getElementById('composer'),
      sidebar: document.getElementById('sidebar'),
      scrim: document.getElementById('scrim'),
      chatList: document.getElementById('chatList'),
      messages: document.getElementById('messages'),
      prompt: document.getElementById('prompt'),
      sendBtn: document.getElementById('sendBtn'),
      attachBtn: document.getElementById('attachBtn'),
      fileInput: document.getElementById('fileInput'),
      pendingFiles: document.getElementById('pendingFiles'),
      newChatBtn: document.getElementById('newChatBtn'),
      modeSelect: document.getElementById('modeSelect'),
      planName: document.getElementById('planName'),
      planPrice: document.getElementById('planPrice'),
      planModes: document.getElementById('planModes'),
      apiKeyBtn: document.getElementById('apiKeyBtn'),
      apiKeyPreview: document.getElementById('apiKeyPreview'),
      chatTitle: document.getElementById('chatTitle'),
      status: document.getElementById('status'),
      menuBtn: document.getElementById('menuBtn'),
      desktopSidebarBtn: document.getElementById('desktopSidebarBtn'),
      memoryBtn: document.getElementById('memoryBtn'),
      memoryModal: document.getElementById('memoryModal'),
      memoryText: document.getElementById('memoryText'),
      closeMemoryBtn: document.getElementById('closeMemoryBtn'),
      clearMemoryBtn: document.getElementById('clearMemoryBtn'),
      reportBtn: document.getElementById('reportBtn'),
      reportModal: document.getElementById('reportModal'),
      closeReportBtn: document.getElementById('closeReportBtn'),
      reportCategory: document.getElementById('reportCategory'),
      reportText: document.getElementById('reportText'),
      submitReportBtn: document.getElementById('submitReportBtn'),
      viewReportsBtn: document.getElementById('viewReportsBtn'),
      reportStatus: document.getElementById('reportStatus'),
      reportList: document.getElementById('reportList'),
    };

    function setSidebarOpen(open) {
      els.sidebar.classList.toggle('open', open);
      document.body.classList.toggle('sidebar-open', open);
    }

    function desktopSidebarCollapsed() {
      try { return localStorage.getItem(SIDEBAR_COLLAPSED_KEY) === '1'; } catch (e) { return false; }
    }

    function setDesktopSidebarCollapsed(collapsed, persist = true) {
      document.body.classList.toggle('sidebar-collapsed', !!collapsed);
      if (persist) {
        try { localStorage.setItem(SIDEBAR_COLLAPSED_KEY, collapsed ? '1' : '0'); } catch (e) {}
      }
    }

    function escapeHtml(value) {
      return value.replace(/[&<>"']/g, ch => ({
        '&': '&amp;', '<': '&lt;', '>': '&gt;', '"': '&quot;', "'": '&#39;'
      })[ch]);
    }

    // ═══ NEXO MEJORAS: renderMarkdown con Prism.js + botón copiar ═══
    const LANG_MAP = {
      js:'javascript', ts:'typescript', py:'python', sh:'bash', bash:'bash',
      html:'html', css:'css', sql:'sql', json:'json', cpp:'cpp', c:'c',
      java:'java', rs:'rust', go:'go', rb:'ruby', php:'php', yaml:'yaml', yml:'yaml',
      cs:'csharp', kt:'kotlin', swift:'swift', lua:'lua', r:'r', md:'markdown',
    };
    function renderMarkdown(text) {
      const parts = text.split(/(```[\s\S]*?```)/g);
      return parts.map((part, index) => {
        if (index % 2 === 1) {
          const inner = part.slice(3, -3);
          const langMatch = inner.match(/^([a-zA-Z0-9_+-]+)\n/);
          const rawLang  = langMatch ? langMatch[1].toLowerCase() : '';
          const prismLang = LANG_MAP[rawLang] || rawLang || 'plaintext';
          const code = langMatch ? inner.slice(langMatch[0].length) : inner;
          const highlighted = (window.Prism && Prism.languages[prismLang])
            ? Prism.highlight(code, Prism.languages[prismLang], prismLang)
            : escapeHtml(code);
          const label = rawLang ? ` data-lang="${escapeHtml(rawLang)}"` : '';
          return `<pre${label} style="position:relative"><button class="code-copy" type="button" title="Copiar código">Copiar</button><code class="language-${prismLang}">${highlighted}</code></pre>`;
        }
        const html = escapeHtml(part)
          .replace(/\*\*\*(.+?)\*\*\*/g, '<strong><em>$1</em></strong>')
          .replace(/\*\*(.+?)\*\*/g, '<strong>$1</strong>')
          .replace(/\*(.+?)\*/g, '<em>$1</em>')
          .replace(/`([^`]+)`/g, '<code style="background:var(--code);padding:2px 5px;border-radius:4px;font-size:.9em">$1</code>')
          .replace(/^### (.+)$/gm, '<h3 style="margin:.6em 0 .2em">$1</h3>')
          .replace(/^## (.+)$/gm,  '<h2 style="margin:.8em 0 .3em">$1</h2>')
          .replace(/^# (.+)$/gm,   '<h1 style="margin:1em 0 .4em">$1</h1>')
          .replace(/^[-*] (.+)/gm, '<li>$1</li>')
          .replace(/(<li>.*<\/li>\n?)+/g, s => `<ul style="margin:.4em 0;padding-left:1.4em">${s}</ul>`)
          .split(/\n{2,}/)
          .map(p => p.trim() && !p.startsWith('<') ? `<p style="margin:.5em 0">${p.replace(/\n/g, '<br>')}</p>` : p)
          .join('');
        return html;
      }).join('');
    }

    function formatBytes(size) {
      let value = Number(size || 0);
      for (const unit of ['B', 'KB', 'MB', 'GB']) {
        if (value < 1024 || unit === 'GB') {
          return unit === 'B' ? `${value} ${unit}` : `${value.toFixed(1)} ${unit}`;
        }
        value = value / 1024;
      }
      return `${size || 0} B`;
    }

    function fileKind(file) {
      const name = file.name.toLowerCase();
      if (/\.(jpg|jpeg|png|webp|bmp|gif)$/.test(name)) return 'image';
      if (/\.(mp4|webm|mov|avi|mkv|m4v)$/.test(name)) return 'video';
      if (/\.(mp3|wav|m4a|aac|ogg|oga|flac|wma)$/.test(name)) return 'audio';
      if (/\.(zip|tar|gz|tgz|bz2|tbz2|xz|txz|7z)$/.test(name)) return 'archive';
      if (/\.(pdf|txt|md|csv|json|jsonl|xml|html|htm|css|js|ts|tsx|jsx|py|java|c|cpp|h|hpp|cs|go|rs|php|rb|swift|kt|kts|sql|yaml|yml|toml|ini|cfg|conf|env|log|bat|ps1|sh|doc|docx|rtf|odt|ppt|pptx|xls|xlsx|ods)$/.test(name)) return 'document';
      if (['dockerfile', 'makefile', 'license', 'readme'].includes(name)) return 'document';
      return 'unknown';
    }

    function statusLabel(status) {
      return ({
        pending: 'pendiente',
        processed: 'procesado',
        error: 'error',
        expired: 'caducado',
      })[status] || status || '';
    }

    function modeLabel(mode) {
      return ({ auto: 'Auto', rapido: 'Rapido', combinado: 'Combinado', codigo: 'Codigo' })[mode] || mode;
    }

    function renderAccount(data) {
      state.account = data || {};
      state.plans = state.account.plans || [];
      const plan = state.account.plan || {};
      const allowed = new Set(plan.allowed_modes || ['auto', 'rapido']);
      els.planName.textContent = plan.label || 'Plan Gratis';
      els.planPrice.textContent = Number(plan.price_eur || 0) ? `${plan.price_eur} euros` : 'gratis';
      els.planModes.textContent = `Modos: ${(plan.allowed_modes || ['auto', 'rapido']).map(modeLabel).join(', ')}`;
      [...els.modeSelect.options].forEach(option => {
        option.disabled = option.value !== 'auto' && !allowed.has(option.value);
        option.title = option.disabled ? `${modeLabel(option.value)} requiere otro plan` : '';
      });
      if (els.modeSelect.selectedOptions[0]?.disabled) els.modeSelect.value = 'auto';
      const apiKey = state.account.api_key || {};
      els.apiKeyBtn.hidden = !apiKey.enabled;
      els.apiKeyPreview.hidden = !(apiKey.enabled && apiKey.prefix);
      els.apiKeyPreview.textContent = apiKey.prefix ? `API: ${apiKey.prefix}...` : '';
      // Mostrar botón Admin solo para cuentas admin
      const adminBtn = document.getElementById('adminBtn');
      if (adminBtn && data.is_admin) adminBtn.style.display = '';
      const statsBtn = document.getElementById('statsBtn');
      if (statsBtn && data.is_admin) statsBtn.style.display = '';
      // Expiración del plan temporal
      const expiryEl = document.getElementById('planExpiry');
      if (expiryEl && data.plan_expires_at) {
        const exp = new Date(data.plan_expires_at);
        const diff = Math.ceil((exp - Date.now()) / 86400000);
        if (diff > 0) {
          expiryEl.style.display = '';
          expiryEl.textContent = `⏳ Plan temporal · expira en ${diff} día${diff !== 1 ? 's' : ''}`;
        } else {
          expiryEl.style.display = 'none';
        }
      } else if (expiryEl) {
        expiryEl.style.display = 'none';
      }
      // Banner de promoción de lanzamiento
      const promoBanner = document.getElementById('promoBanner');
      if (promoBanner) {
        const promo = data.promo || {};
        if (promo.active && promo.slots_left > 0) {
          const endsAt = promo.ends_at ? new Date(promo.ends_at) : null;
          const hoursLeft = endsAt ? Math.ceil((endsAt - Date.now()) / 3600000) : null;
          const timeStr = hoursLeft !== null ? (hoursLeft > 24 ? `${Math.ceil(hoursLeft/24)} días` : `${hoursLeft}h`) : '';
          promoBanner.style.display = '';
          promoBanner.innerHTML = `🚀 <strong>Promo de lanzamiento</strong> — Quedan <strong>${promo.slots_left}</strong> cupos de Developer gratis 7 días${timeStr ? ` · ${timeStr}` : ''}. ¡Compártelo!`;
        } else {
          promoBanner.style.display = 'none';
        }
      }
    }

    async function loadAccount() {
      renderAccount(await api('/api/account'));
    }

    async function generateApiKey() {
      const data = await api('/api/account/api-key', { method: 'POST', body: '{}' });
      els.apiKeyPreview.hidden = false;
      els.apiKeyPreview.textContent = data.api_key || '';
      try { await navigator.clipboard.writeText(data.api_key || ''); } catch (e) {}
      if (state.account && state.account.api_key) state.account.api_key.prefix = data.prefix || '';
    }

    function renderProvider(message = {}) {
      return '';
    }

    function sourceSnippet(source) {
      const snippet = source.snippet || '';
      return snippet ? `<span>${escapeHtml(snippet).slice(0, 500)}</span>` : '';
    }

    function updateDisplayedAttachment(attachment) {
      const chat = currentChat();
      if (!chat || !chat.messages || !chat.messages.length) return;
      const lastUser = [...chat.messages].reverse().find(message => message.role === 'user');
      if (!lastUser) return;
      const list = lastUser.attachments || [];
      const index = list.findIndex(item => item.id === attachment.id || item.filename === attachment.filename);
      if (index >= 0) list[index] = attachment;
      else list.push(attachment);
      lastUser.attachments = list;
    }

    function renderAttachments(attachments = []) {
      if (!attachments.length) return '';
      return `<div class="attachment-list">${attachments.map(item => `
        <div class="attachment-item">
          <strong>${escapeHtml(item.filename || 'archivo')}</strong>
          <span>${escapeHtml(item.kind || 'archivo')} - ${formatBytes(item.size)}${item.status ? ' - ' + escapeHtml(statusLabel(item.status)) : ''}${item.expired ? ' - caducado' : ''}</span>
        </div>
      `).join('')}</div>`;
    }

    function renderSources(sources = []) {
      return '';
    }

    function assistantHtml(text, sources = [], meta = {}) {
      return renderMarkdown(text || '');
    }

    function currentChat() {
      return state.chats.find(chat => chat.id === state.activeChatId) || null;
    }

    function setBusy(value) {
      state.busy = value;
      els.sendBtn.disabled = value;
      els.prompt.disabled = value;
      els.attachBtn.disabled = value;
      els.fileInput.disabled = value;
      els.status.classList.toggle('busy', !!value);
      if (!value) { els.status.textContent = ''; els.status.classList.remove('busy'); }
    }

    function setStatus(text) {
      els.status.textContent = text || '';
      els.status.classList.toggle('busy', state.busy && !!text);
    }

    // ═══ NEXO MEJORAS: Smooth scroll inteligente ═══
    let _autoScroll = true;
    els.messages.addEventListener('scroll', () => {
      const el = els.messages;
      _autoScroll = el.scrollTop + el.clientHeight >= el.scrollHeight - 80;
    }, { passive: true });
    function scrollBottom(force) {
      if (!force && !_autoScroll) return;
      els.messages.scrollTo({ top: els.messages.scrollHeight, behavior: 'smooth' });
    }

    function addPendingFiles(files) {
      [...files].forEach(file => {
        const item = {
          id: `${Date.now()}-${Math.random().toString(16).slice(2)}`,
          file,
          previewUrl: file.type.startsWith('image/') ? URL.createObjectURL(file) : '',
        };
        state.pendingFiles.push(item);
      });
      renderPendingFiles();
    }

    function removePendingFile(id) {
      const index = state.pendingFiles.findIndex(item => item.id === id);
      if (index < 0) return;
      const [item] = state.pendingFiles.splice(index, 1);
      if (item.previewUrl) URL.revokeObjectURL(item.previewUrl);
      renderPendingFiles();
    }

    function clearPendingFiles() {
      state.pendingFiles.forEach(item => {
        if (item.previewUrl) URL.revokeObjectURL(item.previewUrl);
      });
      state.pendingFiles = [];
      els.fileInput.value = '';
      renderPendingFiles();
    }

    function pendingAttachmentsForDisplay() {
      return state.pendingFiles.map(item => ({
        filename: item.file.name,
        size: item.file.size,
        kind: fileKind(item.file),
        status: 'pending',
      }));
    }

    function renderPendingFiles() {
      els.pendingFiles.innerHTML = '';
      els.pendingFiles.classList.toggle('open', state.pendingFiles.length > 0);
      state.pendingFiles.forEach(item => {
        const wrap = document.createElement('div');
        wrap.className = 'pending-file';
        const thumb = item.previewUrl
          ? `<img class="pending-thumb" src="${item.previewUrl}" alt="">`
          : `<div class="pending-thumb">${fileKind(item.file).slice(0, 3).toUpperCase()}</div>`;
        wrap.innerHTML = `
          ${thumb}
          <div><strong>${escapeHtml(item.file.name)}</strong><span>${formatBytes(item.file.size)}</span></div>
          <button class="remove-file" type="button" aria-label="Quitar archivo">x</button>
        `;
        wrap.querySelector('.remove-file').addEventListener('click', () => removePendingFile(item.id));
        els.pendingFiles.appendChild(wrap);
      });
    }

    function renderChatList() {
      els.chatList.innerHTML = '';
      state.chats.forEach(chat => {
        const row = document.createElement('div');
        row.className = 'chat-row' + (chat.id === state.activeChatId ? ' active' : '');

        const button = document.createElement('button');
        button.className = 'chat-item';
        button.type = 'button';
        button.textContent = chat.title || 'Nuevo chat';
        button.addEventListener('click', () => {
          loadChat(chat.id);
          setSidebarOpen(false);
        });

        const delBtn = document.createElement('button');
        delBtn.className = 'chat-delete-btn';
        delBtn.type = 'button';
        delBtn.title = 'Eliminar chat';
        delBtn.textContent = '✕';
        delBtn.addEventListener('click', async (e) => {
          e.stopPropagation();
          if (!confirm('¿Eliminar este chat permanentemente?')) return;
          await deleteChat(chat.id);
        });

        row.appendChild(button);
        row.appendChild(delBtn);
        els.chatList.appendChild(row);
      });
    }

    function renderMessages(chat) {
      els.messages.innerHTML = '';
      if (!chat || !chat.messages || !chat.messages.length) {
        els.messages.innerHTML = `<div class="empty"><div><h1>¿Qué hacemos hoy?</h1><p>Elige modo, escribe tu mensaje y Nexo responderá usando memoria local.</p></div></div>`;
        els.chatTitle.textContent = 'Nuevo chat';
        return;
      }
      els.chatTitle.textContent = chat.title || 'Nuevo chat';
      chat.messages.forEach(addMessageElement);
      scrollBottom();
    }

    async function deleteChat(id) {
      try {
        await api(`/api/chats/${id}`, { method: 'DELETE' });
        state.chats = state.chats.filter(c => c.id !== id);
        if (state.activeChatId === id) {
          state.activeChatId = null;
          renderMessages(null);
          els.chatTitle.textContent = 'Nuevo chat';
          if (state.chats.length) await loadChat(state.chats[0].id);
        }
        renderChatList();
      } catch (err) {
        setStatus(err.message);
      }
    }

    function addMessageElement(message) {
      const wrap = document.createElement('article');
      wrap.className = `message ${message.role}`;
      wrap.innerHTML = `
        <div class="avatar">${message.role === 'user' ? 'TU' : 'NX'}</div>
        <div class="bubble">${renderMarkdown(message.content || '')}${renderAttachments(message.attachments || [])}</div>
      `;
      els.messages.appendChild(wrap);
      return wrap.querySelector('.bubble');
    }

    async function api(url, options = {}) {
      const response = await fetch(url, {
        credentials: 'same-origin',
        headers: { 'Content-Type': 'application/json', ...(options.headers || {}) },
        ...options,
      });
      if (!response.ok) {
        let message = `Error ${response.status}`;
        try { message = (await response.json()).error || message; } catch (_) {}
        throw new Error(message);
      }
      return response.json();
    }

    async function loadChats() {
      const data = await api('/api/chats');
      state.chats = data.chats || [];
      renderChatList();
      if (state.chats.length && !state.activeChatId) {
        await loadChat(state.chats[0].id);
      }
    }

    async function createChat() {
      const chat = await api('/api/chats', { method: 'POST', body: '{}' });
      state.chats.unshift(chat);
      state.activeChatId = chat.id;
      renderChatList();
      renderMessages(chat);
      return chat;
    }

    async function loadChat(id) {
      const chat = await api(`/api/chats/${id}`);
      state.activeChatId = chat.id;
      const index = state.chats.findIndex(item => item.id === chat.id);
      if (index >= 0) state.chats[index] = chat;
      renderChatList();
      renderMessages(chat);
    }

    function openNeural() {
      window.open('/neural', '_blank', 'width=1200,height=800');
    }

    async function sendMessage() {
      const message = els.prompt.value.trim();
      const filesToSend = state.pendingFiles.map(item => item.file);
      if ((!message && !filesToSend.length) || state.busy) return;
      let chat = currentChat();
      if (!chat) chat = await createChat();

      els.prompt.value = '';
      els.prompt.style.height = 'auto';
      setBusy(true);
      setStatus('Preparando...');

      const displayMessage = message || 'Analiza los archivos adjuntos.';
      chat.messages.push({ role: 'user', content: displayMessage, attachments: pendingAttachmentsForDisplay() });
      renderMessages(chat);
      const assistantBubble = addMessageElement({ role: 'assistant', content: '' });
      assistantBubble.classList.add('streaming');
      let assistantText = '';
      let assistantSources = [];
      let assistantMeta = {};
      scrollBottom();
      clearPendingFiles();

      try {
        let fetchOptions;
        if (filesToSend.length) {
          const formData = new FormData();
          formData.append('chat_id', chat.id);
          formData.append('mode', els.modeSelect.value);
          formData.append('message', message);
          filesToSend.forEach(file => formData.append('files[]', file, file.name));
          fetchOptions = { method: 'POST', credentials: 'same-origin', body: formData };
        } else {
          fetchOptions = {
            method: 'POST',
            credentials: 'same-origin',
            headers: { 'Content-Type': 'application/json' },
            body: JSON.stringify({
              chat_id: chat.id,
              mode: els.modeSelect.value,
              message,
            }),
          };
        }
        const response = await fetch('/api/chat/stream', {
          ...fetchOptions,
        });
        if (!response.ok || !response.body) {
          let errorText = `Error ${response.status}`;
          try { errorText = (await response.json()).error || errorText; } catch (_) {}
          throw new Error(errorText);
        }

        const reader = response.body.getReader();
        const decoder = new TextDecoder();
        let buffer = '';
        while (true) {
          const { value, done } = await reader.read();
          if (done) break;
          buffer += decoder.decode(value, { stream: true });
          const lines = buffer.split('\n');
          buffer = lines.pop() || '';
          for (const line of lines) {
            if (!line.trim()) continue;
            const event = JSON.parse(line);
            if (event.type === 'status') setStatus(event.message);
            if (event.type === 'mode') {
              const label = event.label || event.mode || '';
              if (label) setStatus(`Modo: ${label}`);
            }
            if (event.type === 'token') {
              assistantText += event.token;
              assistantBubble.innerHTML = assistantHtml(assistantText, assistantSources, assistantMeta);
              scrollBottom();
            }
            if (event.type === 'sources') {
              assistantSources = event.sources || [];
              assistantBubble.innerHTML = assistantHtml(assistantText, assistantSources, assistantMeta);
              scrollBottom();
            }
            if (event.type === 'provider') {
              assistantMeta.provider = event.provider || '';
              assistantBubble.innerHTML = assistantHtml(assistantText, assistantSources, assistantMeta);
              scrollBottom();
            }
            if (event.type === 'fallback') {
              assistantMeta.provider = event.provider || assistantMeta.provider || 'ollama';
              assistantMeta.fallback = event.reason || 'fallback local';
              assistantBubble.innerHTML = assistantHtml(assistantText, assistantSources, assistantMeta);
              scrollBottom();
            }
            if (event.type === 'attachment_status' && event.attachment) {
              updateDisplayedAttachment(event.attachment);
              scrollBottom();
            }
            if (event.type === 'error') throw new Error(event.message);
            if (event.type === 'done') {
              const updated = event.chat;
              const index = state.chats.findIndex(item => item.id === updated.id);
              if (index >= 0) state.chats[index] = updated;
              else state.chats.unshift(updated);
              state.activeChatId = updated.id;
              renderChatList();
              renderMessages(updated);
            }
          }
        }
      } catch (err) {
        assistantBubble.classList.remove('streaming');
        assistantBubble.innerHTML = renderMarkdown(`Error: ${err.message}`);
      } finally {
        assistantBubble.classList.remove('streaming');
        setBusy(false);
        els.prompt.focus();
      }
    }

    async function openMemory() {
      const data = await api('/api/memory');
      els.memoryText.textContent = data.summary || 'Todavía no hay memoria guardada.';
      els.memoryModal.classList.add('open');
    }

    async function clearMemory() {
      await api('/api/memory/clear', { method: 'POST', body: '{}' });
      els.memoryText.textContent = 'Memoria borrada.';
    }

    function openReport() {
      els.reportText.value = '';
      els.reportStatus.style.display = 'none';
      els.reportList.style.display = 'none';
      els.reportModal.classList.add('open');
      setTimeout(() => els.reportText.focus(), 80);
    }

    async function submitReport() {
      const text = els.reportText.value.trim();
      if (!text) { showReportStatus('Escribe algo antes de enviar.', true); return; }
      const category = els.reportCategory.value;
      els.submitReportBtn.disabled = true;
      try {
        await api('/api/report', {
          method: 'POST',
          body: JSON.stringify({ text, category }),
        });
        showReportStatus('✅ Reporte guardado correctamente.', false);
        els.reportText.value = '';
      } catch (err) {
        showReportStatus('❌ Error al guardar: ' + err.message, true);
      } finally {
        els.submitReportBtn.disabled = false;
      }
    }

    function showReportStatus(msg, isError) {
      els.reportStatus.textContent = msg;
      els.reportStatus.style.color = isError ? 'var(--danger)' : 'var(--accent)';
      els.reportStatus.style.display = 'block';
    }

    async function viewReports() {
      els.viewReportsBtn.disabled = true;
      try {
        const data = await api('/api/reports');
        const reports = data.reports || [];
        els.reportList.style.display = 'grid';
        if (!reports.length) {
          els.reportList.innerHTML = '<div class="report-item"><pre>No hay reportes guardados aún.</pre></div>';
        } else {
          const catLabels = { ui: '🖥️ Interfaz', chat: '💬 Chat', archivos: '📎 Archivos', rendimiento: '⚡ Rendimiento', otro: '🔧 Otro' };
          els.reportList.innerHTML = reports.map(r => `
            <div class="report-item">
              <strong>${r.timestamp}</strong>
              <span class="report-cat">${catLabels[r.category] || r.category}</span>
              <pre>${escapeHtml(r.text)}</pre>
            </div>
          `).join('');
        }
      } catch (err) {
        els.reportList.style.display = 'grid';
        els.reportList.innerHTML = '<div class="report-item"><pre>Error al cargar reportes.</pre></div>';
      } finally {
        els.viewReportsBtn.disabled = false;
      }
    }

    els.sendBtn.addEventListener('click', sendMessage);
    els.attachBtn.addEventListener('click', () => els.fileInput.click());
    els.fileInput.addEventListener('change', event => addPendingFiles(event.target.files || []));
    els.newChatBtn.addEventListener('click', async () => { await createChat(); els.prompt.focus(); });
    els.menuBtn.addEventListener('click', () => setSidebarOpen(!els.sidebar.classList.contains('open')));
    els.scrim.addEventListener('click', () => setSidebarOpen(false));
    els.desktopSidebarBtn.addEventListener('click', () => {
      const collapsed = document.body.classList.contains('sidebar-collapsed');
      setSidebarOpen(false);
      setDesktopSidebarCollapsed(!collapsed, true);
    });
    els.apiKeyBtn.addEventListener('click', () => generateApiKey().catch(err => setStatus(err.message)));
    els.memoryBtn.addEventListener('click', openMemory);
    els.closeMemoryBtn.addEventListener('click', () => els.memoryModal.classList.remove('open'));
    els.clearMemoryBtn.addEventListener('click', clearMemory);
    els.memoryModal.addEventListener('click', event => {
      if (event.target === els.memoryModal) els.memoryModal.classList.remove('open');
    });
    els.reportBtn.addEventListener('click', openReport);
    els.closeReportBtn.addEventListener('click', () => els.reportModal.classList.remove('open'));
    els.submitReportBtn.addEventListener('click', submitReport);
    els.viewReportsBtn.addEventListener('click', viewReports);
    els.reportModal.addEventListener('click', event => {
      if (event.target === els.reportModal) els.reportModal.classList.remove('open');
    });
    els.messages.addEventListener('click', async event => {
      if (!event.target.classList.contains('code-copy')) return;
      const code = event.target.parentElement.querySelector('code').textContent;
      await navigator.clipboard.writeText(code);
      event.target.textContent = 'Copiado';
      setTimeout(() => event.target.textContent = 'Copiar', 1200);
    });
    els.prompt.addEventListener('keydown', event => {
      if (event.key === 'Enter' && !event.shiftKey) {
        event.preventDefault();
        sendMessage();
      }
    });
    els.prompt.addEventListener('input', () => {
      els.prompt.style.height = 'auto';
      els.prompt.style.height = `${Math.min(180, els.prompt.scrollHeight)}px`;
    });
    els.prompt.addEventListener('paste', event => {
      const files = [...(event.clipboardData?.files || [])];
      if (files.length) addPendingFiles(files);
    });
    ['dragenter', 'dragover'].forEach(type => {
      els.composer.addEventListener(type, event => {
        event.preventDefault();
        els.composer.classList.add('dragging');
      });
    });
    ['dragleave', 'drop'].forEach(type => {
      els.composer.addEventListener(type, event => {
        event.preventDefault();
        els.composer.classList.remove('dragging');
      });
    });
    els.composer.addEventListener('drop', event => {
      const files = event.dataTransfer?.files || [];
      if (files.length) addPendingFiles(files);
    });

    if (window.matchMedia && window.matchMedia('(min-width: 761px)').matches) {
      // Resetear sidebar colapsado al arrancar para evitar UI en blanco
      try { localStorage.removeItem('nexo_sidebar_collapsed_v1'); } catch(e) {}
      setDesktopSidebarCollapsed(false, false);
    }
    (async function boot() {
      await loadAccount();
      await loadChats();
    })().catch(err => setStatus(err.message));
  </script>
</body>
</html>
"""


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def utc_after(days: int) -> str:
    return (datetime.now(timezone.utc) + timedelta(days=days)).isoformat(timespec="seconds")


def parse_utc(value: Any) -> Optional[datetime]:
    if not value:
        return None
    try:
        text = str(value).replace("Z", "+00:00")
        parsed = datetime.fromisoformat(text)
        if parsed.tzinfo is None:
            parsed = parsed.replace(tzinfo=timezone.utc)
        return parsed.astimezone(timezone.utc)
    except ValueError:
        return None


def truncate_text(text: str, limit: int) -> str:
    text = text.strip()
    if len(text) <= limit:
        return text
    return text[:limit].rstrip() + "\n\n[Contenido truncado para no saturar el contexto.]"


def safe_segment(value: str) -> str:
    safe = re.sub(r"[^A-Za-z0-9_.-]+", "_", value).strip("._")
    return safe or "unknown"


def classify_upload(filename: str) -> Optional[str]:
    suffix = Path(filename).suffix.lower()
    basename = Path(filename).name.lower()
    if basename in {"dockerfile", "makefile", "license", "readme"}:
        return "document"
    if suffix in IMAGE_EXTENSIONS:
        return "image"
    if suffix in DOCUMENT_EXTENSIONS:
        return "document"
    if suffix in AUDIO_EXTENSIONS:
        return "audio"
    if suffix in VIDEO_EXTENSIONS:
        return "video"
    if suffix in ARCHIVE_EXTENSIONS:
        return "archive"
    return "unknown"


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def format_bytes(size: int) -> str:
    value = float(size)
    for unit in ("B", "KB", "MB", "GB"):
        if value < 1024 or unit == "GB":
            if unit == "B":
                return f"{int(value)} {unit}"
            return f"{value:.1f} {unit}"
        value /= 1024
    return f"{size} B"


def attachment_for_chat(attachment: Dict[str, Any]) -> Dict[str, Any]:
    return {
        key: value
        for key, value in attachment.items()
        if key not in {"saved_path", "derived_paths", "openai_file_id", "openai_upload_error"}
    }


def remove_stored_attachment_files(attachments: List[Dict[str, Any]]) -> None:
    for attachment in attachments:
        for key in ("saved_path",):
            path = Path(str(attachment.get(key, "")))
            if path.exists() and path.is_file():
                try:
                    path.unlink(missing_ok=True)
                except OSError:
                    pass


def build_attachment_context(attachments: List[Dict[str, Any]], mode: str = "combinado") -> str:
    if not attachments:
        return "Sin archivos adjuntos."

    char_limit = MAX_ATTACHMENT_CONTEXT_CHARS_RAPIDO if mode == "rapido" else MAX_ATTACHMENT_CONTEXT_CHARS
    lines = []
    for index, attachment in enumerate(attachments, 1):
        summary = str(attachment.get("summary", "")).strip() or "Sin resumen disponible."
        lines.append(
            f"Adjunto {index}: {attachment.get('filename', 'archivo')} "
            f"({attachment.get('kind', 'desconocido')}, {format_bytes(int(attachment.get('size') or 0))}, "
            f"caduca {attachment.get('expires_at', 'sin fecha')})\n{summary}"
        )
    return truncate_text("\n\n".join(lines), char_limit)


def build_web_context(sources: List[Dict[str, str]], search_error: str = "") -> str:
    if sources:
        blocks = []
        for index, source in enumerate(sources, 1):
            blocks.append(
                f"Fuente {index}: {source.get('title') or 'Sin titulo'}\n"
                f"URL: {source.get('url')}\n"
                f"Extracto: {source.get('snippet') or ''}"
            )
        return truncate_text("\n\n".join(blocks), MAX_WEB_CONTEXT_CHARS)
    if search_error:
        return f"No se pudo obtener contexto de internet: {search_error}"
    return "No se encontraron fuentes utiles en internet."


def store_uploaded_files(uploaded_files: List[Any], user_id: str, chat_id: str) -> tuple[List[Dict[str, Any]], Optional[str]]:
    files = [item for item in uploaded_files if item and getattr(item, "filename", "")]
    if not files:
        return [], None
    if len(files) > MAX_ATTACHMENTS_PER_MESSAGE:
        return [], f"Maximo {MAX_ATTACHMENTS_PER_MESSAGE} archivos por mensaje."

    root = UPLOAD_DIR / safe_segment(user_id) / safe_segment(chat_id)
    root.mkdir(parents=True, exist_ok=True)

    attachments: List[Dict[str, Any]] = []
    total_size = 0
    try:
        for storage in files:
            original_name = str(storage.filename or "archivo")
            safe_name = secure_filename(original_name) or "archivo"
            kind = classify_upload(original_name)
            if not kind:
                remove_stored_attachment_files(attachments)
                return [], f"Tipo de archivo no permitido: {original_name}"

            upload_id = secrets.token_urlsafe(10)
            final_name = f"{upload_id}_{safe_name}"
            destination = root / final_name
            storage.save(destination)

            size = destination.stat().st_size
            total_size += size
            if size > UPLOAD_LIMITS[kind]:
                destination.unlink(missing_ok=True)
                remove_stored_attachment_files(attachments)
                return [], f"{original_name} supera el limite para {kind}: {format_bytes(UPLOAD_LIMITS[kind])}."
            if total_size > MAX_TOTAL_UPLOAD_BYTES:
                destination.unlink(missing_ok=True)
                remove_stored_attachment_files(attachments)
                return [], f"Los adjuntos superan el limite total de {format_bytes(MAX_TOTAL_UPLOAD_BYTES)}."

            mime = storage.mimetype or mimetypes.guess_type(safe_name)[0] or "application/octet-stream"
            attachments.append(
                {
                    "id": upload_id,
                    "filename": original_name,
                    "mime": mime,
                    "kind": kind,
                    "size": size,
                    "sha256": file_sha256(destination),
                    "summary": "Pendiente de procesar.",
                    "status": "pending",
                    "expires_at": utc_after(UPLOAD_TTL_DAYS),
                    "expired": False,
                    "saved_path": str(destination),
                }
            )
    except OSError as exc:
        remove_stored_attachment_files(attachments)
        return [], f"No se pudo guardar el archivo: {exc}"

    return attachments, None


def extract_pdf_text(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        return "No se pudo leer el PDF: falta instalar pypdf."

    try:
        reader = PdfReader(str(path))
        parts = []
        for page in reader.pages:
            parts.append(page.extract_text() or "")
            if sum(len(part) for part in parts) >= MAX_DOCUMENT_CHARS:
                break
        text = "\n\n".join(part.strip() for part in parts if part.strip())
        return truncate_text(text or "El PDF no contiene texto extraible.", MAX_DOCUMENT_CHARS)
    except Exception as exc:
        return f"No se pudo leer el PDF: {exc}"


def extract_document_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_text(path)
    if suffix == ".docx":
        return extract_docx_text(path)
    if suffix == ".xlsx":
        return extract_xlsx_text(path)
    if suffix == ".pptx":
        return extract_pptx_text(path)
    if suffix == ".rtf":
        return extract_rtf_text(path)
    if suffix in {".doc", ".xls", ".ppt", ".odt", ".ods"}:
        return (
            "Este formato se ha guardado correctamente, pero no tiene extractor local fiable. "
            "Si OpenAI esta configurado, se enviara como archivo para que el modelo lo lea."
        )
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        return truncate_text(text or "El archivo esta vacio.", MAX_DOCUMENT_CHARS)
    except OSError as exc:
        return f"No se pudo leer el documento: {exc}"


def extract_docx_text(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        return "No se pudo leer DOCX: falta instalar python-docx."

    try:
        document = Document(str(path))
        parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables[:20]:
            for row in table.rows[:100]:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        text = "\n".join(parts)
        return truncate_text(text or "El DOCX no contiene texto extraible.", MAX_DOCUMENT_CHARS)
    except Exception as exc:
        return f"No se pudo leer DOCX: {exc}"


def extract_xlsx_text(path: Path) -> str:
    try:
        import openpyxl
    except ImportError:
        return "No se pudo leer XLSX: falta instalar openpyxl."

    try:
        workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet in workbook.worksheets[:10]:
            parts.append(f"Hoja: {sheet.title}")
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(value.strip() for value in values):
                    parts.append(" | ".join(values))
                    row_count += 1
                if row_count >= 1000:
                    parts.append("[Hoja truncada a 1000 filas]")
                    break
        workbook.close()
        text = "\n".join(parts)
        return truncate_text(text or "El XLSX no contiene datos extraibles.", MAX_DOCUMENT_CHARS)
    except Exception as exc:
        return f"No se pudo leer XLSX: {exc}"


def extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "No se pudo leer PPTX: falta instalar python-pptx."

    try:
        presentation = Presentation(str(path))
        parts = []
        for index, slide in enumerate(presentation.slides, 1):
            slide_parts = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    slide_parts.append(text.strip())
            if slide_parts:
                parts.append(f"Diapositiva {index}:\n" + "\n".join(slide_parts))
        text = "\n\n".join(parts)
        return truncate_text(text or "El PPTX no contiene texto extraible.", MAX_DOCUMENT_CHARS)
    except Exception as exc:
        return f"No se pudo leer PPTX: {exc}"


def extract_rtf_text(path: Path) -> str:
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
    except OSError as exc:
        return f"No se pudo leer RTF: {exc}"
    text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", raw)
    text = re.sub(r"\\[a-zA-Z]+\d* ?", " ", text)
    text = re.sub(r"[{}]", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return truncate_text(text or "El RTF no contiene texto extraible.", MAX_DOCUMENT_CHARS)


def describe_unknown_file(path: Path, attachment: Dict[str, Any]) -> str:
    return (
        "Archivo guardado como binario o formato no reconocido.\n"
        f"Nombre: {attachment.get('filename', path.name)}\n"
        f"MIME: {attachment.get('mime', 'application/octet-stream')}\n"
        f"Tamano: {format_bytes(int(attachment.get('size') or 0))}\n"
        f"SHA256: {attachment.get('sha256', '')}\n"
        "No se ha ejecutado ni interpretado por seguridad."
    )


def image_path_to_base64(path: Path) -> Optional[str]:
    try:
        from PIL import Image
    except ImportError:
        try:
            return base64.b64encode(path.read_bytes()).decode("ascii")
        except OSError:
            return None

    try:
        with Image.open(path) as image:
            try:
                image.seek(0)
            except EOFError:
                pass
            image.thumbnail((1024, 1024))
            if image.mode in {"RGBA", "LA"}:
                background = Image.new("RGB", image.size, "white")
                alpha = image.getchannel("A")
                background.paste(image.convert("RGB"), mask=alpha)
                image = background
            else:
                image = image.convert("RGB")
            buffer = io.BytesIO()
            image.save(buffer, format="JPEG", quality=85, optimize=True)
            return base64.b64encode(buffer.getvalue()).decode("ascii")
    except Exception:
        try:
            return base64.b64encode(path.read_bytes()).decode("ascii")
        except OSError:
            return None


def summarize_images_with_vision(image_paths: List[Path], prompt: str) -> str:
    if not image_paths:
        return "No se encontraron imagenes para analizar."

    ready_error = ensure_ai_ready(["vision"])
    if ready_error:
        return f"No se pudo analizar visualmente: {ready_error}"

    summaries = []
    for start in range(0, len(image_paths), VISION_IMAGE_BATCH):
        chunk = image_paths[start:start + VISION_IMAGE_BATCH]
        images = [encoded for encoded in (image_path_to_base64(path) for path in chunk) if encoded]
        if not images:
            summaries.append(f"Grupo {start // VISION_IMAGE_BATCH + 1}: no se pudieron cargar las imagenes.")
            continue
        try:
            result = ollama_chat(
                "vision",
                [
                    {"role": "system", "content": "Eres un analista visual preciso. Responde en espanol."},
                    {"role": "user", "content": prompt, "images": images},
                ],
            )
            summaries.append(f"Grupo {start // VISION_IMAGE_BATCH + 1}:\n{result.strip()}")
        except Exception as exc:
            summaries.append(f"Grupo {start // VISION_IMAGE_BATCH + 1}: error al analizar vision: {exc}")

    return truncate_text("\n\n".join(summaries), MAX_DOCUMENT_CHARS)


def extract_video_frames(path: Path, target_count: int = VIDEO_FRAME_COUNT) -> List[Path]:
    try:
        import cv2
    except ImportError:
        return []

    frames_dir = path.parent / f"{path.stem}_frames"
    frames_dir.mkdir(parents=True, exist_ok=True)
    capture = cv2.VideoCapture(str(path))
    if not capture.isOpened():
        return []

    frame_paths: List[Path] = []
    try:
        frame_count = int(capture.get(cv2.CAP_PROP_FRAME_COUNT) or 0)
        if frame_count > 0:
            indices = sorted({min(frame_count - 1, int((i + 0.5) * frame_count / target_count)) for i in range(target_count)})
            for order, frame_index in enumerate(indices, 1):
                capture.set(cv2.CAP_PROP_POS_FRAMES, frame_index)
                ok, frame = capture.read()
                if not ok:
                    continue
                frame_paths.append(save_video_frame(frame, frames_dir, order))
        else:
            seen = 0
            while len(frame_paths) < target_count:
                ok, frame = capture.read()
                if not ok:
                    break
                if seen % 30 == 0:
                    frame_paths.append(save_video_frame(frame, frames_dir, len(frame_paths) + 1))
                seen += 1
    finally:
        capture.release()

    return [path for path in frame_paths if path.exists()]


def save_video_frame(frame: Any, frames_dir: Path, order: int) -> Path:
    import cv2

    height, width = frame.shape[:2]
    largest = max(width, height)
    if largest > 1024:
        scale = 1024 / largest
        frame = cv2.resize(frame, (int(width * scale), int(height * scale)), interpolation=cv2.INTER_AREA)
    destination = frames_dir / f"frame_{order:02d}.jpg"
    cv2.imwrite(str(destination), frame, [int(cv2.IMWRITE_JPEG_QUALITY), 85])
    return destination


def image_path_to_data_url(path: Path) -> Optional[str]:
    encoded = image_path_to_base64(path)
    if not encoded:
        return None
    return f"data:image/jpeg;base64,{encoded}"


def extract_audio_transcript(path: Path, settings: Dict[str, Any]) -> str:
    if not should_use_openai(settings):
        return "Audio guardado, pero no transcrito porque OpenAI no esta configurado."
    try:
        with path.open("rb") as handle:
            response = http_session().post(
                f"{OPENAI_API_BASE}/audio/transcriptions",
                headers={"Authorization": f"Bearer {settings['openai_api_key']}"},
                data={"model": settings["openai_transcribe_model"]},
                files={"file": (path.name, handle, mimetypes.guess_type(path.name)[0] or "application/octet-stream")},
                timeout=300,
            )
        response.raise_for_status()
        data = response.json()
        return truncate_text(str(data.get("text") or "").strip() or "La transcripcion no devolvio texto.", MAX_DOCUMENT_CHARS)
    except Exception as exc:
        return f"No se pudo transcribir el audio con OpenAI: {exc}"


def extract_video_audio(path: Path) -> Optional[Path]:
    try:
        try:
            from moviepy import VideoFileClip
        except ImportError:
            from moviepy.editor import VideoFileClip
    except ImportError:
        return None

    destination = path.parent / f"{path.stem}_audio.wav"
    try:
        clip = VideoFileClip(str(path))
        try:
            if not clip.audio:
                return None
            clip.audio.write_audiofile(str(destination), logger=None)
            return destination if destination.exists() else None
        finally:
            clip.close()
    except Exception:
        return None


def archive_member_safe_name(name: str, index: int) -> str:
    base = secure_filename(Path(name).name) or f"archivo_{index}"
    return f"{index:02d}_{base}"


def copy_limited_stream(source: Any, destination: Path, limit: int) -> int:
    total = 0
    with destination.open("wb") as handle:
        while True:
            chunk = source.read(64 * 1024)
            if not chunk:
                break
            total += len(chunk)
            if total > limit:
                raise ValueError("archivo interno demasiado grande")
            handle.write(chunk)
    return total


def extract_zip_members(path: Path, derived_dir: Path) -> List[Path]:
    extracted: List[Path] = []
    total = 0
    with zipfile.ZipFile(path) as archive:
        for info in archive.infolist():
            if info.is_dir():
                continue
            if len(extracted) >= MAX_ARCHIVE_MEMBERS:
                break
            total += int(info.file_size or 0)
            if total > MAX_ARCHIVE_EXTRACT_BYTES:
                break
            destination = derived_dir / archive_member_safe_name(info.filename, len(extracted) + 1)
            with archive.open(info) as source:
                copy_limited_stream(source, destination, MAX_DOCUMENT_BYTES)
            extracted.append(destination)
    return extracted


def extract_tar_members(path: Path, derived_dir: Path) -> List[Path]:
    extracted: List[Path] = []
    total = 0
    with tarfile.open(path) as archive:
        for member in archive.getmembers():
            if not member.isfile():
                continue
            if len(extracted) >= MAX_ARCHIVE_MEMBERS:
                break
            total += int(member.size or 0)
            if total > MAX_ARCHIVE_EXTRACT_BYTES:
                break
            source = archive.extractfile(member)
            if source is None:
                continue
            destination = derived_dir / archive_member_safe_name(member.name, len(extracted) + 1)
            copy_limited_stream(source, destination, MAX_DOCUMENT_BYTES)
            extracted.append(destination)
    return extracted


def extract_7z_members(path: Path, derived_dir: Path) -> List[Path]:
    try:
        import py7zr
    except ImportError:
        return []

    raw_dir = derived_dir / "raw_7z"
    raw_dir.mkdir(parents=True, exist_ok=True)
    with py7zr.SevenZipFile(path, mode="r") as archive:
        archive.extractall(path=raw_dir)

    extracted: List[Path] = []
    total = 0
    for item in raw_dir.rglob("*"):
        if not item.is_file():
            continue
        try:
            item.resolve().relative_to(raw_dir.resolve())
        except ValueError:
            continue
        if len(extracted) >= MAX_ARCHIVE_MEMBERS:
            break
        size = item.stat().st_size
        total += size
        if total > MAX_ARCHIVE_EXTRACT_BYTES:
            break
        destination = derived_dir / archive_member_safe_name(item.name, len(extracted) + 1)
        shutil.copyfile(item, destination)
        extracted.append(destination)
    return extracted


def extract_archive_files(path: Path) -> List[Path]:
    derived_dir = path.parent / f"{path.stem}_archive"
    derived_dir.mkdir(parents=True, exist_ok=True)
    suffixes = [suffix.lower() for suffix in path.suffixes]
    try:
        if path.suffix.lower() == ".zip":
            return extract_zip_members(path, derived_dir)
        if path.suffix.lower() == ".7z":
            return extract_7z_members(path, derived_dir)
        if any(suffix in suffixes for suffix in {".tar", ".gz", ".tgz", ".bz2", ".xz"}):
            return extract_tar_members(path, derived_dir)
    except Exception:
        return []
    return []


def summarize_archive(path: Path, attachment: Dict[str, Any]) -> str:
    extracted = extract_archive_files(path)
    attachment["derived_paths"] = [str(item) for item in extracted]
    if not extracted:
        return "No se pudieron extraer archivos legibles del comprimido o el comprimido esta vacio."

    parts = [f"Comprimido extraido parcialmente: {len(extracted)} archivo(s) analizados."]
    for item in extracted[:MAX_ARCHIVE_MEMBERS]:
        kind = classify_upload(item.name)
        if kind == "document":
            parts.append(f"\nArchivo interno: {item.name}\n{extract_document_text(item)}")
        else:
            parts.append(
                f"\nArchivo interno: {item.name}\n"
                f"Tipo detectado: {kind}. Tamano: {format_bytes(item.stat().st_size)}."
            )
    return truncate_text("\n".join(parts), MAX_DOCUMENT_CHARS)


def should_use_openai(settings: Optional[Dict[str, Any]] = None) -> bool:
    data = settings or load_ai_settings()
    return data.get("ai_provider") == "openai" and bool(data.get("openai_api_key"))


def upload_attachment_to_openai(attachment: Dict[str, Any], settings: Dict[str, Any]) -> None:
    if not should_use_openai(settings):
        return
    if attachment.get("kind") != "document":
        return
    path = Path(str(attachment.get("saved_path", "")))
    if not path.exists() or int(attachment.get("size") or 0) > MAX_DOCUMENT_BYTES:
        return
    try:
        with path.open("rb") as handle:
            response = requests.post(
                f"{OPENAI_API_BASE}/files",
                headers={"Authorization": f"Bearer {settings['openai_api_key']}"},
                data={"purpose": "user_data"},
                files={"file": (attachment.get("filename") or path.name, handle, attachment.get("mime") or "application/octet-stream")},
                timeout=180,
            )
        response.raise_for_status()
        attachment["openai_file_id"] = response.json().get("id", "")
    except Exception as exc:
        attachment["openai_upload_error"] = str(exc)


def process_stored_attachment(attachment: Dict[str, Any], settings: Optional[Dict[str, Any]] = None, mode: str = "combinado") -> Dict[str, Any]:
    settings = settings or load_ai_settings()
    path = Path(str(attachment.get("saved_path", "")))
    if not path.exists():
        attachment["summary"] = "El archivo ya no existe en disco."
        attachment["expired"] = True
        attachment["status"] = "expired"
        return attachment

    is_rapido = mode == "rapido"
    kind = attachment.get("kind")

    if kind == "document":
        text = extract_document_text(path)
        if is_rapido:
            # Modo Rápido: extracto limitado, sin análisis profundo
            text = truncate_text(text, MAX_DOCUMENT_CHARS_RAPIDO)
            attachment["summary"] = f"Vista previa del documento (modo Rápido, primeros caracteres):\n{text}"
        else:
            attachment["summary"] = f"Texto extraido del documento:\n{text}"
        if not is_rapido:
            upload_attachment_to_openai(attachment, settings)
    elif kind == "image":
        if is_rapido:
            # Modo Rápido: descripción básica con prompt corto, sin batches
            if should_use_openai(settings):
                attachment["summary"] = "Imagen adjunta (vista previa básica en modo Rápido)."
            else:
                attachment["summary"] = summarize_images_with_vision(
                    [path],
                    "Describe brevemente esta imagen en 2-3 frases: qué muestra, colores principales y texto visible si lo hay.",
                )
        else:
            if should_use_openai(settings):
                attachment["summary"] = "Imagen lista para analisis multimodal con OpenAI."
            else:
                attachment["summary"] = summarize_images_with_vision(
                    [path],
                    "Analiza esta imagen para ayudar a responder al usuario. Describe contenido, texto visible, detalles importantes y posibles dudas.",
                )
    elif kind == "audio":
        if is_rapido:
            # Modo Rápido: sin transcripción
            size_str = format_bytes(int(attachment.get("size") or 0))
            attachment["summary"] = f"Archivo de audio adjunto ({attachment.get('filename', 'audio')}, {size_str}). La transcripción completa está disponible en modo Combinado o Código."
        else:
            transcript = extract_audio_transcript(path, settings)
            attachment["summary"] = f"Audio procesado.\nTranscripcion:\n{transcript}"
    elif kind == "video":
        if is_rapido:
            # Modo Rápido: sin frames ni transcripción
            size_str = format_bytes(int(attachment.get("size") or 0))
            attachment["summary"] = f"Archivo de video adjunto ({attachment.get('filename', 'video')}, {size_str}). El análisis de frames y audio está disponible en modo Combinado o Código."
        else:
            frame_paths = extract_video_frames(path)
            attachment["derived_paths"] = [str(item) for item in frame_paths]
            audio_path = extract_video_audio(path)
            if audio_path:
                attachment.setdefault("derived_paths", []).append(str(audio_path))
            transcript = extract_audio_transcript(audio_path, settings) if audio_path else "El video no tenia audio extraible."
            if should_use_openai(settings):
                visual_summary = "Fotogramas listos para analisis multimodal con OpenAI."
            else:
                visual_summary = summarize_images_with_vision(
                    frame_paths,
                    "Estos son fotogramas distribuidos de un video. Resume la secuencia, acciones, objetos, texto visible, cambios entre escenas y detalles utiles.",
                )
            attachment["summary"] = (
                f"Video procesado mediante {len(frame_paths)} fotogramas clave.\n"
                f"{visual_summary}\n\nTranscripcion/audio:\n{transcript}"
            )
    elif kind == "archive":
        attachment["summary"] = summarize_archive(path, attachment)
    else:
        attachment["summary"] = describe_unknown_file(path, attachment)
    attachment["status"] = "processed"
    return attachment


def is_global_ip(value: str) -> bool:
    try:
        return ipaddress.ip_address(value).is_global
    except ValueError:
        return False


def is_safe_public_url(url: str) -> bool:
    parsed = urlparse(url)
    if parsed.scheme not in {"http", "https"}:
        return False
    host = parsed.hostname
    if not host:
        return False
    lowered = host.lower().strip(".")
    if lowered in {"localhost"} or lowered.endswith(".local"):
        return False
    if re.fullmatch(r"\d+(?:\.\d+){3}", lowered) or ":" in lowered:
        return is_global_ip(lowered)
    try:
        for family, _, _, _, sockaddr in socket.getaddrinfo(lowered, None):
            address = sockaddr[0]
            if not ipaddress.ip_address(address).is_global:
                return False
    except Exception:
        return False
    return True


def fetch_url_text(url: str) -> Optional[Dict[str, str]]:
    if not is_safe_public_url(url):
        return None
    try:
        with http_session().get(
            url,
            timeout=10,
            stream=True,
            headers={"User-Agent": "Nexo/1.0 (+local research assistant)"},
        ) as response:
            response.raise_for_status()
            content_type = response.headers.get("Content-Type", "")
            if content_type and not any(kind in content_type.lower() for kind in ("text/", "html", "json", "xml")):
                return None
            chunks = []
            total = 0
            for chunk in response.iter_content(chunk_size=32768):
                if not chunk:
                    continue
                total += len(chunk)
                if total > 768 * 1024:
                    break
                chunks.append(chunk)
            raw = b"".join(chunks)
    except Exception:
        return None

    encoding = response.encoding or "utf-8"
    text = raw.decode(encoding, errors="replace")
    return clean_page_text(text, url)


def clean_page_text(raw_html: str, url: str) -> Dict[str, str]:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        cleaned = re.sub(r"<[^>]+>", " ", raw_html)
        cleaned = re.sub(r"\s+", " ", cleaned).strip()
        return {"title": url, "text": truncate_text(cleaned, 3500)}

    soup = BeautifulSoup(raw_html, "html.parser")
    for tag in soup(["script", "style", "noscript", "svg", "header", "footer", "nav", "form"]):
        tag.decompose()
    title = soup.title.get_text(" ", strip=True) if soup.title else url
    text = soup.get_text(" ", strip=True)
    text = re.sub(r"\s+", " ", text).strip()
    return {"title": title or url, "text": truncate_text(text, 3500)}


def build_search_query(user_message: str, attachment_context: str) -> str:
    base = " ".join(user_message.split())
    if not base:
        base = " ".join(attachment_context.split())
    base = re.sub(r"```[\s\S]*?```", " ", base)
    return truncate_text(base, 280) or "informacion actual"


def normalize_for_intent(text: str) -> str:
    decomposed = unicodedata.normalize("NFKD", text.lower())
    cleaned = "".join(ch for ch in decomposed if not unicodedata.combining(ch))
    cleaned = re.sub(r"[^a-z0-9\s]", " ", cleaned)
    return " ".join(cleaned.split())
    replacements = str.maketrans("áéíóúüñ", "aeiouun")
    cleaned = text.lower().translate(replacements)
    cleaned = re.sub(r"[^a-z0-9\s]", " ", cleaned)
    return " ".join(cleaned.split())


def should_search_web(user_message: str, attachment_context: str = "") -> bool:
    """
    Gate de búsqueda web (estricto).
    - Por defecto NO busca (rendimiento + evita búsquedas innecesarias).
    - Solo busca si el usuario lo pide explícitamente o si el contenido es claramente "actual" (precio/noticias/hoy/último, etc.).
    """
    # Si hay archivos adjuntos, normalmente la respuesta debe venir del archivo, no de internet.
    # (El usuario puede forzar búsqueda con "busca en internet".)

    normalized = normalize_for_intent(user_message)
    if not normalized:
        return False

    explicit_triggers = [
        "busca en internet", "buscar en internet", "busca en la web", "buscar en la web",
        "googlea", "duckduckgo", "ddg", "investiga en internet", "investiga en la web",
        "dame fuentes", "con fuentes", "fuentes", "enlace", "links", "link", "url",
        "cita", "citas", "referencias",
    ]
    normalized_raw = normalized
    if any(t in normalized_raw for t in explicit_triggers):
        return True

    changing_info_terms = (
        "precio", "cotizacion", "noticia", "noticias", "dolar", "euro",
        "bitcoin", "btc", "version", "release", "lanzamiento",
    )
    if is_conversational_message(user_message) and not any(term in normalized_raw for term in changing_info_terms):
        return False

    # Si el usuario pregunta por info cambiante/reciente, habilitar web.
    time_sensitive = re.search(
        r"\b(hoy|actual|reciente|ultim[oa]s?|noticia[s]?|precio[s]?|cotizacion|"
        r"lanzamiento|release|version|202[0-9]|esta semana|este mes|"
        r"tiempo|clima|temperatura|partido|resultado|ganador|clasificacion|"
        r"estreno|pelicula[s]?|serie[s]?|album|cancion|musica|trending|viral|"
        r"que paso|que ha pasado|novedades|actualizacion|nuevo[s]?|nueva[s]?)\b",
        normalized_raw,
    )
    if time_sensitive:
        return True

    casual_messages = {
        "hola", "holaa", "holaaa", "buenas", "buenos dias", "buen dia",
        "buenas tardes", "buenas noches", "hey", "hi", "hello",
        "que tal", "como estas", "como va", "que pasa",
        "bien y tu", "y tu", "nada xd", "nada",
        "gracias", "muchas gracias", "ok", "vale", "perfecto",
        "jaja", "jeje", "xd", "jajaja", "lol",
        "adios", "chao", "hasta luego", "nos vemos",
    }
    if normalized in casual_messages:
        return False

    words = normalized.split()
    # Mensajes cortos/banales: no buscar.
    if len(words) <= 2 and not re.search(r"[?¿]", user_message):
        return False

    # Por defecto: NO buscar (solo con triggers arriba).
    return False




def is_conversational_message(user_message: str) -> bool:
    """Detecta charla, identidad y agradecimientos aunque el selector este en codigo."""
    normalized = normalize_for_intent(user_message)
    if not normalized:
        return False

    casual_exact = {
        "hola", "holaa", "holaaa", "buenas", "buenos dias", "buen dia",
        "buenas tardes", "buenas noches", "hey", "hi", "hello",
        "que tal", "como estas", "como va", "que pasa", "bien y tu", "y tu",
        "nada xd", "nada", "gracias", "muchas gracias", "ok", "vale",
        "perfecto", "genial", "listo", "jaja", "jeje", "xd", "jajaja",
        "lol", "adios", "chao", "hasta luego", "nos vemos",
    }
    if normalized in casual_exact:
        return True

    text = f" {normalized} "
    ack_prefixes = (
        "gracias", "muchas gracias", "mil gracias", "perfecto", "genial",
        "vale", "ok", "listo", "me sirve", "esta bien",
    )
    followup_task_markers = (
        " ahora ", " pero ", " puedes ", " podrias ", " haz ", " crea ",
        " crear ", " genera ", " escribe ", " arregla ",
        " corrige ", " cambia ", " anade ", " añade ",
    )
    if normalized.startswith(ack_prefixes) and not any(marker in text for marker in followup_task_markers):
        return True

    social_or_identity = (
        "como te llamas", "te llamas", "cual es tu nombre", "tu nombre", "quien eres",
        "eres claude", "soy claude", "anthropic", "openai", "chatgpt",
        "desde cuando", "te he creado", "mi ia", "que tal te va",
        "como te va", "como estas", "estas bien",
    )
    if any(phrase in normalized for phrase in social_or_identity):
        return True

    task_keywords = (
        "codigo", "code", "programa", "programar", "funcion", "clase",
        "script", "python", "javascript", "java", "api", "sql", "html",
        "css", "react", "bug", "debug", "error", "crear", "crea",
        "generar", "genera", "generame", "escribir", "escribe", "haz",
        "arregla", "corrige", "implementa", "refactoriza", "test",
    )
    if any(keyword in normalized for keyword in task_keywords):
        return False

    words = normalized.split()
    if len(words) <= 4:
        return True
    if len(words) <= 8 and re.search(r"[?¿]", user_message):
        return True
    return False


def search_web_context(user_message: str, attachment_context: str = "") -> tuple[str, List[Dict[str, str]], str]:
    if not should_search_web(user_message, attachment_context):
        return "", [], ""

    query = build_search_query(user_message, attachment_context)
    try:
        from ddgs import DDGS
    except ImportError:
        return "", [], "falta instalar ddgs"

    try:
        raw_results = list(DDGS().text(query, max_results=6))
    except Exception as exc:
        return "", [], str(exc)

    sources: List[Dict[str, str]] = []
    seen_urls: set[str] = set()
    for result in raw_results:
        url = str(result.get("href") or result.get("url") or "").strip()
        if not url or url in seen_urls or not is_safe_public_url(url):
            continue
        seen_urls.add(url)
        title = str(result.get("title") or url).strip()
        snippet = str(result.get("body") or result.get("snippet") or "").strip()
        page = fetch_url_text(url)
        if page and page.get("text"):
            title = page.get("title") or title
            snippet = page["text"]
        sources.append({"title": title, "url": url, "snippet": truncate_text(snippet, 2200)})
        if len(sources) >= 3:
            break

    return build_web_context(sources), sources, ""


def cleanup_uploads() -> None:
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    cutoff = datetime.now(timezone.utc) - timedelta(days=UPLOAD_TTL_DAYS)

    for path in sorted(UPLOAD_DIR.rglob("*"), reverse=True):
        try:
            if path.is_file():
                modified = datetime.fromtimestamp(path.stat().st_mtime, timezone.utc)
                if modified < cutoff:
                    path.unlink(missing_ok=True)
            elif path.is_dir():
                try:
                    path.rmdir()
                except OSError:
                    pass
        except OSError:
            pass

    with DATA_LOCK:
        data = read_json(CHATS_FILE, {"chats": []})
        changed = False
        now = datetime.now(timezone.utc)
        for chat in data.get("chats", []):
            for message in chat.get("messages", []):
                attachments = message.get("attachments", [])
                if not isinstance(attachments, list):
                    continue
                for attachment in attachments:
                    expires_at = parse_utc(attachment.get("expires_at"))
                    if expires_at and expires_at <= now and not attachment.get("expired"):
                        attachment["expired"] = True
                        attachment["status"] = "expired"
                        attachment["summary"] = str(attachment.get("summary", "")).strip() or "Archivo caducado."
                        changed = True
        if changed:
            write_json(CHATS_FILE, data)


def read_json(path: Path, default: Any) -> Any:
    if not path.exists():
        return default
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError):
        return default


def write_json(path: Path, data: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    tmp = path.with_suffix(path.suffix + ".tmp")
    tmp.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
    tmp.replace(path)


def parse_bool(value: Any, default: bool = False) -> bool:
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    text = str(value).strip().lower()
    if text in {"1", "true", "yes", "y", "s", "si", "on"}:
        return True
    if text in {"0", "false", "no", "n", "off"}:
        return False
    return default


def default_settings_data() -> Dict[str, Any]:
    return dict(DEFAULT_SETTINGS)


def load_settings_file() -> Dict[str, Any]:
    with DATA_LOCK:
        ensure_data_files()
        raw = read_json(SETTINGS_FILE, default_settings_data())
        if not isinstance(raw, dict):
            raw = default_settings_data()
        merged = {**DEFAULT_SETTINGS, **raw}
        if merged != raw:
            write_json(SETTINGS_FILE, merged)
        return merged


def config_value(settings: Dict[str, Any], env_name: str, settings_name: str, default: Any = "") -> Any:
    env_value = os.getenv(env_name)
    if env_value is not None:
        return env_value
    return settings.get(settings_name, default)


def valid_external_donate_url(url: str) -> bool:
    parsed = urlparse(str(url or "").strip())
    return parsed.scheme in {"http", "https"} and bool(parsed.netloc)


def configured_donate_url(settings: Optional[Dict[str, Any]] = None) -> str:
    data = settings or load_ai_settings()
    url = str(data.get("donate_url") or "").strip()
    return url if valid_external_donate_url(url) else ""


def load_ai_settings() -> Dict[str, Any]:
    now = time.time()
    try:
        mtime = SETTINGS_FILE.stat().st_mtime if SETTINGS_FILE.exists() else 0.0
    except OSError:
        mtime = 0.0

    cached = _AI_SETTINGS_CACHE.get("value")
    try:
        if (
            cached
            and (now - float(_AI_SETTINGS_CACHE.get("ts") or 0.0)) < _AI_SETTINGS_TTL_S
            and float(_AI_SETTINGS_CACHE.get("mtime") or 0.0) == float(mtime)
            and isinstance(cached, dict)
        ):
            return dict(cached)
    except Exception:
        pass

    settings = load_settings_file()
    provider = str(config_value(settings, "AI_PROVIDER", "ai_provider", "openai")).strip().lower() or "openai"
    value = {
        "ai_provider": provider,
        "openai_api_key": str(config_value(settings, "OPENAI_API_KEY", "openai_api_key", "")).strip(),
        "openai_model": str(config_value(settings, "OPENAI_MODEL", "openai_model", "gpt-5.5")).strip() or "gpt-5.5",
        "openai_fast_model": str(
            config_value(settings, "OPENAI_FAST_MODEL", "openai_fast_model", "gpt-4.1-mini")
        ).strip()
        or "gpt-4.1-mini",
        "openai_transcribe_model": str(
            config_value(settings, "OPENAI_TRANSCRIBE_MODEL", "openai_transcribe_model", "gpt-4o-transcribe")
        ).strip() or "gpt-4o-transcribe",
        "fallback_to_ollama": parse_bool(
            config_value(settings, "FALLBACK_TO_OLLAMA", "fallback_to_ollama", True),
            default=True,
        ),
        "identity_name": str(config_value(settings, "NEXO_IDENTITY_NAME", "identity_name", "")).strip(),
        "identity_guard_enabled": parse_bool(
            config_value(settings, "NEXO_IDENTITY_GUARD", "identity_guard_enabled", False),
            default=False,
        ),
        "auto_router_enabled": parse_bool(
            config_value(settings, "NEXO_AUTO_ROUTER", "auto_router_enabled", True),
            default=True,
        ),
        "router_provider": str(config_value(settings, "NEXO_ROUTER_PROVIDER", "router_provider", "auto")).strip().lower()
        or "auto",
        "router_model_ollama": str(config_value(settings, "NEXO_ROUTER_MODEL_OLLAMA", "router_model_ollama", "qwen2.5:1.5b")).strip()
        or "qwen2.5:1.5b",
        "router_model_openai": str(config_value(settings, "NEXO_ROUTER_MODEL_OPENAI", "router_model_openai", "gpt-4.1-mini")).strip()
        or "gpt-4.1-mini",
        "donate_url": str(
            os.getenv("DONATE_URL")
            or config_value(settings, "NEXO_DONATE_URL", "donate_url", "")
        ).strip(),
    }
    _AI_SETTINGS_CACHE["ts"] = now
    _AI_SETTINGS_CACHE["mtime"] = float(mtime)
    _AI_SETTINGS_CACHE["value"] = dict(value)
    return value


def empty_users_data() -> Dict[str, Any]:
    return {"allow_registration": True, "users": [], "promo": {"start": None, "slots_used": 0}}


def ensure_data_files() -> None:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    if not USERS_FILE.exists():
        write_json(USERS_FILE, empty_users_data())
    if not CHATS_FILE.exists():
        write_json(CHATS_FILE, {"chats": []})
    if not MEMORY_FILE.exists():
        write_json(MEMORY_FILE, {"users": {}})
    if not SETTINGS_FILE.exists():
        write_json(SETTINGS_FILE, default_settings_data())


def get_secret_key() -> str:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    if SECRET_FILE.exists():
        return SECRET_FILE.read_text(encoding="utf-8").strip()
    secret = secrets.token_hex(32)
    SECRET_FILE.write_text(secret, encoding="utf-8")
    return secret


def normalize_username(username: str) -> str:
    return " ".join(username.strip().split())


def username_key(username: str) -> str:
    return normalize_username(username).lower()


def normalize_plan(value: Any, default: str = PLAN_FREE) -> str:
    text = str(value or "").strip().lower().replace("-", "_")
    text = " ".join(text.split())
    if text in PLAN_DEFINITIONS:
        return text
    return PLAN_ALIASES.get(text, default)


def plan_definition(plan_key: Any) -> Dict[str, Any]:
    return dict(PLAN_DEFINITIONS.get(normalize_plan(plan_key), PLAN_DEFINITIONS[PLAN_FREE]))


def public_plan_definition(plan_key: Any) -> Dict[str, Any]:
    definition = plan_definition(plan_key)
    return {
        "key": definition["key"],
        "label": definition["label"],
        "price_eur": definition["price_eur"],
        "priority": definition["priority"],
        "allowed_modes": list(definition["allowed_modes"]),
        "includes_api_key": bool(definition["includes_api_key"]),
        "features": list(definition["features"]),
    }


def public_plan_catalog() -> List[Dict[str, Any]]:
    return [public_plan_definition(key) for key in (PLAN_FREE, PLAN_BETA, PLAN_DEVELOPER)]


def user_plan_key(user: Optional[Dict[str, Any]]) -> str:
    if not user:
        return PLAN_FREE
    plan = normalize_plan(user.get("plan"), default=PLAN_FREE)
    # Si el plan tiene fecha de expiración y ya pasó, devolver el plan de retorno
    expires_at = parse_utc(user.get("plan_expires_at"))
    if expires_at and datetime.now(timezone.utc) >= expires_at:
        return normalize_plan(user.get("plan_after_expiry"), default=PLAN_FREE)
    return plan


def public_plan_for_user(user: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    return public_plan_definition(user_plan_key(user))


def mode_label(mode: str) -> str:
    return {"rapido": "Rapido", "combinado": "Combinado", "codigo": "Codigo", "auto": "Auto"}.get(mode, mode)


def required_plan_for_mode(mode: str) -> str:
    if mode == "codigo":
        return PLAN_DEVELOPER
    if mode == "combinado":
        return PLAN_BETA
    return PLAN_FREE


def plan_allows_mode(plan_key: Any, mode: str) -> bool:
    definition = plan_definition(plan_key)
    return mode in set(definition.get("allowed_modes", []))


def plan_mode_error(plan_key: Any, mode: str) -> Optional[str]:
    if plan_allows_mode(plan_key, mode):
        return None
    current = plan_definition(plan_key)
    required = plan_definition(required_plan_for_mode(mode))
    price = required.get("price_eur", 0)
    price_text = "gratis" if not price else f"{price} euros"
    return (
        f"{current['label']} no incluye el modo {mode_label(mode)}. "
        f"Para usarlo necesitas {required['label']} ({price_text})."
    )


def hash_api_key(api_key: str) -> str:
    return hashlib.sha256(api_key.encode("utf-8", errors="ignore")).hexdigest()


def api_key_prefix(api_key: str) -> str:
    return api_key[:18]


def make_api_key() -> str:
    return "nexo_dev_" + secrets.token_urlsafe(32).replace("-", "").replace("_", "")


def request_api_key() -> str:
    auth = request.headers.get("Authorization", "").strip()
    if auth.lower().startswith("bearer "):
        return auth.split(None, 1)[1].strip()
    return request.headers.get("X-Nexo-API-Key", "").strip()


def api_key_info_for_user(user: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    if not user:
        return {"enabled": False, "prefix": ""}
    plan = plan_definition(user_plan_key(user))
    return {
        "enabled": bool(plan.get("includes_api_key")),
        "prefix": str(user.get("api_key_prefix") or ""),
        "created_at": str(user.get("api_key_created_at") or ""),
    }


def find_user_by_api_key(api_key: str) -> Optional[Dict[str, Any]]:
    api_key = (api_key or "").strip()
    if not api_key:
        return None
    hashed = hash_api_key(api_key)
    for user in load_users_data()["users"]:
        stored_hash = str(user.get("api_key_hash") or "")
        if not stored_hash:
            continue
        if secrets.compare_digest(stored_hash, hashed):
            plan = plan_definition(user_plan_key(user))
            if plan.get("includes_api_key"):
                return user
            return None
    return None


def rotate_user_api_key(user_id: str) -> Dict[str, str]:
    api_key = make_api_key()
    now = utc_now()
    with DATA_LOCK:
        data = load_users_data()
        for user in data["users"]:
            if user.get("id") == user_id:
                if not plan_definition(user_plan_key(user)).get("includes_api_key"):
                    raise PermissionError("Tu plan no incluye API Key.")
                user["api_key_hash"] = hash_api_key(api_key)
                user["api_key_prefix"] = api_key_prefix(api_key)
                user["api_key_created_at"] = now
                user["updated_at"] = now
                write_json(USERS_FILE, data)
                return {"api_key": api_key, "prefix": api_key_prefix(api_key)}
    raise LookupError("Usuario no encontrado.")


def admin_set_user_plan(user_id: str, new_plan: str) -> Dict[str, Any]:
    """Cambia el plan de un usuario. Solo para admins."""
    new_plan = normalize_plan(new_plan)
    now = utc_now()
    with DATA_LOCK:
        data = load_users_data()
        for user in data["users"]:
            if user.get("id") == user_id:
                user["plan"] = new_plan
                user["plan_updated_at"] = now
                user["updated_at"] = now
                write_json(USERS_FILE, data)
                return {"ok": True, "user_id": user_id, "plan": new_plan}
    raise LookupError("Usuario no encontrado.")


def admin_list_users() -> List[Dict[str, Any]]:
    """Lista todos los usuarios para el panel admin."""
    data = load_users_data()
    result = []
    for i, user in enumerate(data["users"]):
        result.append({
            "id": user.get("id", ""),
            "username": user.get("username", ""),
            "plan": user.get("plan", PLAN_FREE),
            "plan_label": plan_definition(user_plan_key(user)).get("label", "Plan Gratis"),
            "created_at": user.get("created_at", ""),
            "registration_order": i + 1,
        })
    return result


def acquire_ai_slot(priority: int) -> float:
    global AI_ACTIVE_REQUESTS, AI_TICKET_COUNTER
    started = time.time()
    with AI_PRIORITY_CONDITION:
        AI_TICKET_COUNTER += 1
        ticket = AI_TICKET_COUNTER
        item = {"priority": int(priority), "ticket": ticket}
        AI_PRIORITY_QUEUE.append(item)
        try:
            while True:
                best = min(AI_PRIORITY_QUEUE, key=lambda entry: (-int(entry["priority"]), int(entry["ticket"])))
                if best is item and AI_ACTIVE_REQUESTS < MAX_CONCURRENT_AI_REQUESTS:
                    AI_PRIORITY_QUEUE.remove(item)
                    AI_ACTIVE_REQUESTS += 1
                    return time.time() - started
                AI_PRIORITY_CONDITION.wait(timeout=30)
        except BaseException:
            if item in AI_PRIORITY_QUEUE:
                AI_PRIORITY_QUEUE.remove(item)
                AI_PRIORITY_CONDITION.notify_all()
            raise


def release_ai_slot() -> None:
    global AI_ACTIVE_REQUESTS
    with AI_PRIORITY_CONDITION:
        AI_ACTIVE_REQUESTS = max(0, AI_ACTIVE_REQUESTS - 1)
        AI_PRIORITY_CONDITION.notify_all()


def normalize_user_record(raw_user: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    username = normalize_username(str(raw_user.get("username", "")))
    if not username:
        return None

    password_hash = str(raw_user.get("password_hash", ""))
    password = str(raw_user.get("password", ""))
    if not password_hash:
        if not password or password in PLACEHOLDER_PASSWORDS:
            return None
        password_hash = generate_password_hash(password)

    created_at = str(raw_user.get("created_at") or raw_user.get("updated_at") or utc_now())
    plan_key = normalize_plan(
        raw_user.get("plan") or raw_user.get("subscription_plan") or raw_user.get("tier"),
        default=LEGACY_USER_DEFAULT_PLAN,
    )
    user = {
        "id": str(raw_user.get("id") or secrets.token_urlsafe(12)),
        "username": username,
        "username_key": username_key(username),
        "password_hash": password_hash,
        "plan": plan_key,
        "plan_updated_at": str(raw_user.get("plan_updated_at") or created_at),
        "created_at": created_at,
        "updated_at": str(raw_user.get("updated_at") or created_at),
    }
    for key in ("api_key_hash", "api_key_prefix", "api_key_created_at", "api_key_last_used_at", "payment_status"):
        value = raw_user.get(key)
        if value:
            user[key] = str(value)
    return user


def normalize_users_data(raw: Any) -> Dict[str, Any]:
    if not isinstance(raw, dict):
        return empty_users_data()

    users: List[Dict[str, Any]] = []
    seen: set[str] = set()

    if isinstance(raw.get("users"), list):
        raw_users = raw["users"]
    else:
        raw_users = [raw]

    for raw_user in raw_users:
        if not isinstance(raw_user, dict):
            continue
        user = normalize_user_record(raw_user)
        if not user or user["username_key"] in seen:
            continue
        users.append(user)
        seen.add(user["username_key"])

    raw_promo = raw.get("promo")
    if isinstance(raw_promo, dict):
        promo = {
            "start": raw_promo.get("start") or None,
            "slots_used": int(raw_promo.get("slots_used") or 0),
        }
    else:
        promo = {"start": None, "slots_used": 0}

    return {
        "allow_registration": bool(raw.get("allow_registration", True)),
        "users": users,
        "promo": promo,
    }


def load_users_data() -> Dict[str, Any]:
    with DATA_LOCK:
        ensure_data_files()
        raw = read_json(USERS_FILE, empty_users_data())
        users_data = normalize_users_data(raw)
        if users_data != raw:
            write_json(USERS_FILE, users_data)
        return users_data


def find_user_by_username(username: str) -> Optional[Dict[str, Any]]:
    key = username_key(username)
    if not key:
        return None
    for user in load_users_data()["users"]:
        if user.get("username_key") == key:
            return user
    return None


def find_user_by_id(user_id: str) -> Optional[Dict[str, Any]]:
    if not user_id:
        return None
    for user in load_users_data()["users"]:
        if user.get("id") == user_id:
            return user
    return None


def validate_new_account(username: str, password: str, confirm_password: str) -> Optional[str]:
    username = normalize_username(username)
    if not USERNAME_RE.fullmatch(username):
        return "El usuario debe tener 3-32 caracteres: letras, numeros, punto, guion o guion bajo."
    if len(password) < 8:
        return "La contrasena debe tener al menos 8 caracteres."
    if password != confirm_password:
        return "Las contrasenas no coinciden."
    if password in PLACEHOLDER_PASSWORDS:
        return "Elige una contrasena real."
    return None


def create_user(username: str, password: str, confirm_password: str) -> tuple[Optional[Dict[str, Any]], Optional[str]]:
    username = normalize_username(username)
    error = validate_new_account(username, password, confirm_password)
    if error:
        return None, error

    with DATA_LOCK:
        data = load_users_data()
        if not data.get("allow_registration", True):
            return None, "El registro esta desactivado."
        if any(user.get("username_key") == username_key(username) for user in data["users"]):
            return None, "Ese usuario ya existe."

        now = utc_now()

        # --- Promo de lanzamiento: primeros PROMO_SLOTS en PROMO_WINDOW_DAYS días → Developer gratis ---
        promo = data.get("promo") or {"start": None, "slots_used": 0}
        promo_start_str = promo.get("start")
        promo_slots_used = int(promo.get("slots_used") or 0)

        if not promo_start_str:
            # Primer registro: arranca la ventana ahora mismo
            promo["start"] = now
            promo_start_str = now

        promo_start_dt = parse_utc(promo_start_str)
        promo_active = (
            promo_start_dt is not None
            and datetime.now(timezone.utc) <= promo_start_dt + timedelta(days=PROMO_WINDOW_DAYS)
            and promo_slots_used < PROMO_SLOTS
        )

        if promo_active:
            assigned_plan = PLAN_DEVELOPER
            plan_expires_at = utc_after(PROMO_PLAN_DAYS)
            plan_after_expiry = PLAN_FREE
            promo["slots_used"] = promo_slots_used + 1
        else:
            assigned_plan = NEW_USER_DEFAULT_PLAN
            plan_expires_at = None
            plan_after_expiry = None

        data["promo"] = promo
        # -----------------------------------------------------------------------------------------

        user: Dict[str, Any] = {
            "id": secrets.token_urlsafe(12),
            "username": username,
            "username_key": username_key(username),
            "password_hash": generate_password_hash(password),
            "plan": assigned_plan,
            "plan_updated_at": now,
            "plan_expires_at": plan_expires_at,
            "plan_after_expiry": plan_after_expiry,
            "created_at": now,
            "updated_at": now,
        }
        data["users"].append(user)
        write_json(USERS_FILE, data)
        return user, None


def validate_login(username: str, password: str) -> Optional[Dict[str, Any]]:
    user = find_user_by_username(username)
    if not user:
        return None
    if check_password_hash(str(user.get("password_hash", "")), password):
        return user
    return None


def current_user() -> Optional[Dict[str, Any]]:
    user = find_user_by_id(str(session.get("user_id", "")))
    if user:
        return user
    return find_user_by_username(str(session.get("username", "")))


def api_authenticated_user() -> Optional[Dict[str, Any]]:
    return find_user_by_api_key(request_api_key())


def authenticated_user() -> Optional[Dict[str, Any]]:
    if session.get("authenticated"):
        user = current_user()
        if user:
            return user
    return api_authenticated_user()


def current_user_id() -> str:
    user = authenticated_user()
    return str(user.get("id", "")) if user else ""


def client_ip() -> str:
    return (
        request.headers.get("CF-Connecting-IP")
        or request.headers.get("X-Forwarded-For", "").split(",")[0].strip()
        or request.remote_addr
        or "unknown"
    )


def login_limited(ip: str) -> Optional[int]:
    now = time.time()
    with LOGIN_LOCK:
        item = LOGIN_ATTEMPTS.get(ip)
        if item and item.get("blocked_until", 0) > now:
            return int(item["blocked_until"] - now)
    return None


def record_login_failure(ip: str) -> None:
    now = time.time()
    with LOGIN_LOCK:
        item = LOGIN_ATTEMPTS.setdefault(ip, {"count": 0, "first": now, "blocked_until": 0})
        if now - item.get("first", now) > 900:
            item.update({"count": 0, "first": now, "blocked_until": 0})
        item["count"] += 1
        if item["count"] >= 5:
            item["blocked_until"] = now + 300


def clear_login_failures(ip: str) -> None:
    with LOGIN_LOCK:
        LOGIN_ATTEMPTS.pop(ip, None)


def require_login_response() -> Optional[Response]:
    if authenticated_user():
        return None
    return jsonify({"error": "No autenticado"}), 401


def load_chats() -> Dict[str, Any]:
    with DATA_LOCK:
        ensure_data_files()
        data = read_json(CHATS_FILE, {"chats": []})
        if "chats" not in data or not isinstance(data["chats"], list):
            data = {"chats": []}
        return data


def save_chats(data: Dict[str, Any]) -> None:
    with DATA_LOCK:
        write_json(CHATS_FILE, data)


def list_chats(user_id: str) -> List[Dict[str, Any]]:
    data = load_chats()
    chats = [
        chat
        for chat in data["chats"]
        if str(chat.get("user_id", "")) == user_id
    ]
    chats = sorted(chats, key=lambda item: item.get("updated_at", ""), reverse=True)
    return [
        {
            "id": chat["id"],
            "title": chat.get("title", "Nuevo chat"),
            "created_at": chat.get("created_at"),
            "updated_at": chat.get("updated_at"),
        }
        for chat in chats
    ]


def get_chat(chat_id: str, user_id: str) -> Optional[Dict[str, Any]]:
    data = load_chats()
    for chat in data["chats"]:
        if chat.get("id") == chat_id and str(chat.get("user_id", "")) == user_id:
            return chat
    return None


def create_chat(user_id: str) -> Dict[str, Any]:
    now = utc_now()
    chat = {
        "id": secrets.token_urlsafe(12),
        "user_id": user_id,
        "title": "Nuevo chat",
        "created_at": now,
        "updated_at": now,
        "messages": [],
    }
    data = load_chats()
    data["chats"].insert(0, chat)
    save_chats(data)
    return chat


def update_chat(updated_chat: Dict[str, Any], user_id: str) -> None:
    updated_chat["user_id"] = user_id
    data = load_chats()
    for index, chat in enumerate(data["chats"]):
        if chat.get("id") == updated_chat.get("id") and str(chat.get("user_id", "")) == user_id:
            data["chats"][index] = updated_chat
            break
    else:
        data["chats"].insert(0, updated_chat)
    save_chats(data)


def delete_chat(chat_id: str, user_id: str) -> bool:
    data = load_chats()
    original_len = len(data["chats"])
    data["chats"] = [
        chat for chat in data["chats"]
        if not (chat.get("id") == chat_id and str(chat.get("user_id", "")) == user_id)
    ]
    if len(data["chats"]) < original_len:
        save_chats(data)
        return True
    return False


def chat_title_from_message(message: str) -> str:
    title = " ".join(message.strip().split())
    if len(title) > 54:
        return title[:51].rstrip() + "..."
    return title or "Nuevo chat"


def load_memory_data() -> Dict[str, Any]:
    with DATA_LOCK:
        ensure_data_files()
        data = read_json(MEMORY_FILE, {"users": {}})
        if not isinstance(data, dict) or not isinstance(data.get("users"), dict):
            data = {"users": {}}
            write_json(MEMORY_FILE, data)
        return data


def load_memory(user_id: str) -> Dict[str, Any]:
    with DATA_LOCK:
        data = load_memory_data()
        item = data["users"].get(user_id, {})
        if not isinstance(item, dict):
            item = {}
        return {"summary": str(item.get("summary", "")), "updated_at": item.get("updated_at")}


def save_memory(user_id: str, summary: str) -> None:
    with DATA_LOCK:
        data = load_memory_data()
        data["users"][user_id] = {"summary": summary.strip(), "updated_at": utc_now()}
        write_json(MEMORY_FILE, data)


# ═══ NEXO MEJORAS: History Trimmer ═══
def trim_history_to_budget(messages: List[Dict], max_chars: int = 12000) -> List[Dict]:
    """Elimina mensajes antiguos cuando el historial supera max_chars."""
    chat = list(messages)
    while len(chat) > 2:
        total = sum(len(str(m.get('content', ''))) for m in chat)
        if total <= max_chars:
            break
        chat.pop(0)
    return chat


def build_history_text(messages: List[Dict[str, str]], limit: int = 12) -> str:
    recent = messages[-limit:]
    if not recent:
        return "Sin historial reciente."
    lines = []
    for msg in recent:
        role = "Usuario" if msg.get("role") == "user" else "Asistente"
        content = str(msg.get("content", ""))
        if msg.get("role") == "assistant":
            old_placeholder = "[Nombre del " + "Usuario]"
            content = content.replace(old_placeholder, "Nexo")
            old_identity_words = (
                "Cl" + "aude",
                "Anth" + "ropic",
                "Chat" + "GPT",
                "Open" + "AI",
                "Gemini",
                "IA " + "Combinada",
            )
            old_identity_re = r"\b(" + "|".join(re.escape(word) for word in old_identity_words) + r")\b"
            content = re.sub(old_identity_re, "Nexo", content, flags=re.I)
        lines.append(f"{role}: {content}")
    return "\n".join(lines)


def build_user_prompt(
    chat: Dict[str, Any],
    user_message: str,
    mode: str,
    attachment_context: str = "",
    web_context: str = "",
) -> str:
    memory = load_memory(str(chat.get("user_id", ""))).get("summary", "").strip() or "Sin memoria guardada todavía."
    history_messages = list(chat.get("messages", []))
    if (
        history_messages
        and history_messages[-1].get("role") == "user"
        and str(history_messages[-1].get("content", "")).strip() == user_message.strip()
    ):
        history_messages = history_messages[:-1]
    history = build_history_text(trim_history_to_budget(history_messages))
    mode_hint = {
        "rapido": "Responde de forma clara y directa.",
        "combinado": "Responde de forma natural y util. Usa razonamiento de calidad, pero no repitas respuestas antiguas.",
        "codigo": "Prioriza codigo completo solo si el mensaje actual pide codigo. Si es charla o agradecimiento, responde breve y natural.",
    }.get(mode, "Responde de forma clara.")
    return f"""Memoria compartida:
{memory}

Historial reciente del chat:
{history}

Archivos adjuntos del mensaje actual:
{attachment_context or "Sin archivos adjuntos."}

Contexto de internet:
{web_context or "No se obtuvo contexto de internet."}

Instrucción de modo:
{mode_hint}

Mensaje actual del usuario:
{user_message}
"""


def ollama_payload(model_key: str, messages: List[Dict[str, Any]], stream: bool) -> Dict[str, Any]:
    model = MODELS[model_key]
    return {
        "model": model,
        "messages": messages,
        "stream": stream,
        "keep_alive": OLLAMA_KEEP_ALIVE,
        "options": ollama_options(model_key),
    }


def ensure_ai_ready(model_keys: Iterable[str]) -> Optional[str]:
    if not is_ollama_running() and not start_ollama():
        return "Ollama no está corriendo y no se pudo iniciar."
    for key in model_keys:
        model_name = MODELS[key]
        if not is_model_installed(model_name):
            return f"Modelo no instalado: {model_name}"
    return None


def ollama_chat(model_key: str, messages: List[Dict[str, Any]]) -> str:
    """Llamada bloqueante a Ollama. Devuelve string vacío en caso de error."""
    try:
        payload = ollama_payload(model_key, messages, stream=False)
        response = http_session().post(f"{OLLAMA_HOST}/api/chat", json=payload, timeout=(10, 300))
        response.raise_for_status()
        data = response.json()
        return data.get("message", {}).get("content", "")
    except Exception:
        return ""


def ollama_chat_stream(model_key: str, messages: List[Dict[str, Any]]) -> Generator[str, None, None]:
    """Streaming de tokens desde Ollama. Silencia errores para no romper el generador."""
    try:
        payload = ollama_payload(model_key, messages, stream=True)
        with http_session().post(
            f"{OLLAMA_HOST}/api/chat",
            json=payload,
            stream=True,
            timeout=(10, 300),
        ) as response:
            response.raise_for_status()
            for raw_line in response.iter_lines(chunk_size=8192):
                if not raw_line:
                    continue
                try:
                    data = json.loads(raw_line)
                except json.JSONDecodeError:
                    continue
                token = data.get("message", {}).get("content", "")
                if token:
                    yield token
                if data.get("done"):
                    break
    except Exception:
        return


def event(payload: Dict[str, Any]) -> str:
    return json.dumps(payload, ensure_ascii=False) + "\n"


def public_http_url(url: str) -> bool:
    parsed = urlparse(str(url))
    return parsed.scheme in {"http", "https"} and bool(parsed.netloc)


def add_unique_source(sources: List[Dict[str, str]], url: str, title: str = "", snippet: str = "") -> None:
    if not url or not public_http_url(url):
        return
    if any(item.get("url") == url for item in sources):
        return
    sources.append(
        {
            "title": title.strip() or url,
            "url": url.strip(),
            "snippet": truncate_text(snippet.strip(), 1200) if snippet else "",
        }
    )


def collect_openai_sources(value: Any, sources: List[Dict[str, str]]) -> None:
    if isinstance(value, dict):
        item_type = str(value.get("type", ""))
        url = str(value.get("url") or value.get("uri") or "").strip()
        if url and ("citation" in item_type or "url" in value):
            add_unique_source(
                sources,
                url,
                str(value.get("title") or value.get("name") or ""),
                str(value.get("snippet") or value.get("text") or ""),
            )
        for child in value.values():
            collect_openai_sources(child, sources)
    elif isinstance(value, list):
        for child in value:
            collect_openai_sources(child, sources)


def parse_openai_response(data: Dict[str, Any]) -> tuple[str, List[Dict[str, str]]]:
    sources: List[Dict[str, str]] = []
    collect_openai_sources(data, sources)
    if isinstance(data.get("output_text"), str):
        return data["output_text"], sources

    parts = []
    output = data.get("output", [])
    for item in output if isinstance(output, list) else []:
        if not isinstance(item, dict):
            continue
        content = item.get("content", [])
        for content_item in content if isinstance(content, list) else []:
            if not isinstance(content_item, dict):
                continue
            text = content_item.get("text") or content_item.get("output_text")
            if isinstance(text, str):
                parts.append(text)
    return "\n".join(parts).strip(), sources


def stream_text_chunks(text: str, size: int = 900) -> Generator[str, None, None]:
    index = 0
    while index < len(text):
        next_index = min(len(text), index + size)
        yield text[index:next_index]
        index = next_index


def attachment_input_text(attachment: Dict[str, Any]) -> str:
    return (
        f"Adjunto: {attachment.get('filename', 'archivo')}\n"
        f"Tipo: {attachment.get('kind', 'desconocido')}\n"
        f"MIME: {attachment.get('mime', '')}\n"
        f"Tamano: {format_bytes(int(attachment.get('size') or 0))}\n"
        f"SHA256: {attachment.get('sha256', '')}\n"
        f"Resumen local:\n{attachment.get('summary', '')}"
    )


def build_openai_content(context_prompt: str, attachments: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    content: List[Dict[str, Any]] = [{"type": "input_text", "text": context_prompt}]
    for attachment in attachments:
        content.append({"type": "input_text", "text": attachment_input_text(attachment)})
        file_id = str(attachment.get("openai_file_id") or "").strip()
        if file_id:
            content.append({"type": "input_file", "file_id": file_id})
            continue

        kind = attachment.get("kind")
        path = Path(str(attachment.get("saved_path", "")))
        if kind == "image" and path.exists():
            data_url = image_path_to_data_url(path)
            if data_url:
                content.append({"type": "input_image", "image_url": data_url})
        elif kind == "video":
            frame_paths = [
                Path(str(item))
                for item in attachment.get("derived_paths", [])
                if str(item).lower().endswith(".jpg")
            ]
            if len(frame_paths) > 8:
                step = max(1, len(frame_paths) // 8)
                frame_paths = frame_paths[::step][:8]
            for frame_path in frame_paths:
                if frame_path.exists():
                    data_url = image_path_to_data_url(frame_path)
                    if data_url:
                        content.append({"type": "input_image", "image_url": data_url})
    return content


def guarded_system_prompt(prompt: str, ai_settings: Optional[Dict[str, Any]] = None) -> str:
    guard = identity_guard_from_settings(ai_settings)
    personality_hint = PERSONALITY_PROMPTS.get(get_personality(), '')
    parts: List[str] = []
    if guard:
        parts.append(guard)
    parts.append(prompt)
    if personality_hint:
        parts.append(f'Estilo de respuesta: {personality_hint}')
    return '\n\n'.join(parts)


def openai_system_instructions(mode: str) -> str:
    mode_hint = {
        "rapido": "Responde claro y directo, sin perder precision.",
        "combinado": "Da una respuesta natural de alta calidad sin repetir respuestas antiguas.",
        "codigo": "Prioriza codigo completo solo cuando el mensaje actual pida codigo; en charla responde breve y natural.",
    }.get(mode, "Responde claro y util.")
    name = "Nexo"
    # No afirmamos nombre fijo salvo que el guard este activado.
    # (La UI puede llamarse Nexo sin que el modelo lo "declare" como identidad.)
    base = (
        "Responde en espanol de forma natural, como un chat normal. "
        "Usa archivos adjuntos cuando existan. Usa internet solo cuando aporte informacion actual o verificable; no lo uses para saludos o charla simple. "
        "No muestres listas de enlaces ni expliques que has buscado salvo que el usuario pida fuentes. "
        "No inventes datos si no tienes base suficiente.\n\n"
        f"{mode_hint}"
    )
    return base.replace("Nexo", name)


def openai_model_for_mode(settings: Dict[str, Any], mode: str) -> str:
    if mode == "rapido":
        return str(settings.get("openai_fast_model") or settings.get("openai_model") or "gpt-4.1-mini")
    return str(settings.get("openai_model") or "gpt-5.5")


def openai_response(
    chat: Dict[str, Any],
    user_message: str,
    mode: str,
    attachments: List[Dict[str, Any]],
    settings: Dict[str, Any],
) -> tuple[str, List[Dict[str, str]]]:
    attachment_context = build_attachment_context(attachments)
    use_web = should_search_web(user_message, attachment_context)
    context_prompt = build_user_prompt(
        chat,
        user_message,
        mode,
        attachment_context,
        "Internet disponible para esta respuesta." if use_web else "No uses internet para esta respuesta; contesta de forma conversacional.",
    )
    model = openai_model_for_mode(settings, mode)
    payload = {
        "model": model,
        "instructions": openai_system_instructions(mode),
        "input": [{"role": "user", "content": build_openai_content(context_prompt, attachments)}],
        "store": False,
    }
    if use_web:
        payload["tools"] = [{"type": "web_search_preview"}]
        payload["tool_choice"] = "required"  # Forzar búsqueda: si nuestro código decidió buscar, OpenAI debe hacerlo
    response = http_session().post(
        f"{OPENAI_API_BASE}/responses",
        headers={
            "Authorization": f"Bearer {settings['openai_api_key']}",
            "Content-Type": "application/json",
        },
        json=payload,
        timeout=600,
    )
    if response.status_code >= 400:
        raise RuntimeError(f"OpenAI {response.status_code}: {response.text[:800]}")
    text, sources = parse_openai_response(response.json())
    if not text.strip():
        raise RuntimeError("OpenAI no devolvio texto.")
    return text, sources


def stream_ollama_answer(
    chat: Dict[str, Any],
    user_message: str,
    mode: str,
    attachments: Optional[List[Dict[str, Any]]] = None,
    web_context: str = "",
    sources: Optional[List[Dict[str, str]]] = None,
    search_error: str = "",
) -> Generator[str, None, Dict[str, Any]]:
    attachment_context = build_attachment_context(attachments or [], mode=mode)
    if mode != "rapido" and not (attachments or []) and is_conversational_message(user_message):
        mode = "rapido"
    full_web_context = web_context or build_web_context(sources or [], search_error)
    context_prompt = build_user_prompt(chat, user_message, mode, attachment_context, full_web_context)
    final_text = ""

    if mode == "rapido":
        fast_model_key = os.getenv("NEXO_FAST_OLLAMA_ROLE", "programador").strip().lower() or "programador"
        if fast_model_key not in MODELS:
            fast_model_key = "programador"
        error = ensure_ai_ready([fast_model_key])
        if error:
            yield event({"type": "error", "message": error})
            return {"text": "", "provider": "ollama"}
        
        # Modo rapido: siempre usar prompt conversacional/general.
        # El modo "codigo" tiene su propio pipeline; aquí llegan preguntas generales,
        # históricas, culturales, de trivia, etc. NUNCA debe responder con código Python.
        messages = [
            {"role": "system", "content": guarded_system_prompt(PROMPT_CONVERSACIONAL, ai_settings=load_ai_settings())},
            {"role": "user", "content": user_message},
        ]
        yield event({"type": "status", "message": "Pensando..."})
        for token in ollama_chat_stream(fast_model_key, messages):
            final_text += token
            yield event({"type": "token", "token": token})
        return {"text": final_text, "provider": "ollama"}

    error = ensure_ai_ready(["programador", "arquitecto"])
    if error:
        yield event({"type": "error", "message": error})
        return {"text": "", "provider": "ollama"}

    if mode == "codigo":
        draft_prompt = f"""{context_prompt}

Responde a la peticion actual del usuario con una solucion tecnica directa. Si pide codigo, entregalo en bloques Markdown con nombre de lenguaje."""
        review_prompt = (
            "Usa el borrador solo como material interno y devuelve la respuesta final directa al usuario. "
            "No digas 'el codigo proporcionado' ni 'la solucion anterior' salvo que el usuario haya pegado codigo para revisar."
        )
    else:
        draft_prompt = f"""{context_prompt}

Genera un borrador de respuesta util, completo y bien estructurado."""
        review_prompt = (
            "Usa el borrador solo como material interno. Devuelve una respuesta final natural, sin mencionar el borrador ni repetir respuestas antiguas."
        )

    yield event({"type": "status", "message": "Ollama generando borrador..."})
    ai_settings = load_ai_settings()
    # Generamos el borrador con streaming para dar feedback visual inmediato al usuario
    draft_tokens: list[str] = []
    for _tok in ollama_chat_stream(
        "programador",
        [
            {"role": "system", "content": guarded_system_prompt(PROMPT_PROGRAMADOR, ai_settings=ai_settings)},
            {"role": "user", "content": draft_prompt},
        ],
    ):
        draft_tokens.append(_tok)
    draft = "".join(draft_tokens)

    final_prompt = f"""{context_prompt}

Borrador:
```
{draft}
```

{review_prompt}"""

    yield event({"type": "status", "message": "Ollama revisando..."})
    for token in ollama_chat_stream(
        "arquitecto",
        [
            {"role": "system", "content": guarded_system_prompt(PROMPT_ARQUITECTO, ai_settings=ai_settings)},
            {"role": "user", "content": final_prompt},
        ],
    ):
        final_text += token
        yield event({"type": "token", "token": token})
    return {"text": final_text, "provider": "ollama"}


def stream_answer(
    chat: Dict[str, Any],
    user_message: str,
    mode: str,
    attachments: Optional[List[Dict[str, Any]]] = None,
    web_context: str = "",
    sources: Optional[List[Dict[str, str]]] = None,
    search_error: str = "",
    ai_settings: Optional[Dict[str, Any]] = None,
) -> Generator[str, None, Dict[str, Any]]:
    settings = ai_settings or load_ai_settings()
    attachments = attachments or []
    if mode != "rapido" and not attachments and is_conversational_message(user_message):
        mode = "rapido"

    if should_use_openai(settings):
        yield event({"type": "provider", "provider": "openai", "model": openai_model_for_mode(settings, mode)})
        yield event({"type": "status", "message": "Pensando..."})
        try:
            text, openai_sources = openai_response(chat, user_message, mode, attachments, settings)
            final_sources = openai_sources or []
            if not final_sources:
                if should_search_web(user_message, build_attachment_context(attachments)):
                    local_context, local_sources, local_error = search_web_context(
                        user_message,
                        build_attachment_context(attachments),
                    )
                    if local_sources:
                        final_sources = local_sources
                    elif local_error:
                        yield event({"type": "status", "message": "Pensando..."})
                    web_context = web_context or local_context
            for chunk in stream_text_chunks(text):
                yield event({"type": "token", "token": chunk})
            return {"text": text, "provider": "openai", "sources": final_sources}
        except Exception as exc:
            if not settings.get("fallback_to_ollama", True):
                yield event({"type": "error", "message": f"OpenAI fallo: {exc}"})
                return {"text": "", "provider": "openai"}
            yield event({"type": "fallback", "provider": "ollama", "reason": f"OpenAI fallo: {exc}"})
            yield event({"type": "status", "message": "Usando Ollama como respaldo..."})

    if not web_context and not sources and should_search_web(user_message, build_attachment_context(attachments)):
        yield event({"type": "status", "message": "Pensando..."})
        web_context, sources, search_error = search_web_context(user_message, build_attachment_context(attachments))
        if search_error:
            yield event({"type": "status", "message": "Pensando..."})

    return (yield from stream_ollama_answer(
        chat,
        user_message,
        mode,
        attachments=attachments,
        web_context=web_context,
        sources=sources,
        search_error=search_error,
    ))


def parse_router_json(text: str) -> Optional[Dict[str, Any]]:
    raw = (text or "").strip()
    if not raw:
        return None
    # Extraer primer objeto JSON en caso de que el modelo meta texto alrededor.
    match = re.search(r"\{[\s\S]*\}", raw)
    if not match:
        return None
    try:
        data = json.loads(match.group(0))
    except Exception:
        return None
    return data if isinstance(data, dict) else None


CODE_SIGNAL_RE = re.compile(
    r"```|traceback|stack\s*trace|exception|error\b|bug\b|debug\b|"
    r"\b(def|class|import|from|function|const|let|var|SELECT|INSERT|UPDATE)\b",
    re.IGNORECASE,
)

CODE_KEYWORDS = (
    "codigo", "code", "programa", "programar", "funcion", "clase", "script",
    "python", "javascript", "typescript", "java", "sql", "html", "css",
    "react", "flask", "fastapi", "node", "docker", "git", "bash", "powershell",
    "regex", "bug", "debug", "error", "endpoint", "backend", "frontend",
    "implementa", "refactoriza", "test",
)

HEAVY_KEYWORDS = (
    "presentacion", "ppt", "pptx", "diapositiva", "slide", "pdf", "resumen",
    "resume", "informe", "analisis", "plan", "roadmap", "arquitectura",
    "diagrama", "lee", "revisa", "archivo", "carpeta", "proyecto",
)


def has_code_signal(user_message: str, normalized: str = "") -> bool:
    normalized = normalized or normalize_for_intent(user_message)
    if bool(CODE_SIGNAL_RE.search(user_message or "")):
        return True
    strong_terms = (
        "codigo", "code", "script", "funcion", "clase", "endpoint", "bug", "debug",
        "traceback", "stack trace", "programar", "implementa", "refactoriza", "backend", "frontend",
    )
    if any(term in normalized for term in strong_terms):
        return True
    tech_terms = (
        "python", "javascript", "typescript", "java", "sql", "html", "css", "react",
        "flask", "fastapi", "node", "docker", "git", "bash", "powershell", "regex", "api",
    )
    action_terms = (
        "crear", "crea", "haz", "hacer", "escribe", "genera", "generame", "implementar",
        "integra", "arregla", "corrige", "construye", "desarrolla", "no funciona", "falla",
        "error",
    )
    return any(term in normalized for term in tech_terms) and any(term in normalized for term in action_terms)


def has_heavy_signal(normalized: str) -> bool:
    base_terms = tuple(term for term in HEAVY_KEYWORDS if term != "plan")
    if any(term in normalized for term in base_terms):
        return True
    if "plan" in normalized and any(marker in normalized for marker in ("haz", "hacer", "crea", "crear", "desarrolla", "estrategia", "roadmap")):
        return True
    return False


def is_lightweight_chat_message(user_message: str, attachments: Optional[List[Dict[str, Any]]] = None) -> bool:
    if attachments:
        return False
    normalized = normalize_for_intent(user_message)
    if not normalized:
        return False
    if has_code_signal(user_message, normalized) or has_heavy_signal(normalized):
        return False
    if is_conversational_message(user_message):
        return True
    words = normalized.split()
    question_starts = ("que ", "quien ", "cuando ", "donde ", "como ", "cual ", "cuanto ", "por que ")
    if len(words) <= 22 and (normalized.startswith(question_starts) or re.search(r"[?¿]", user_message)):
        return True
    if len(words) <= 14:
        return True
    return False


def heuristic_intent(user_message: str, attachments: Optional[List[Dict[str, Any]]] = None) -> str:
    text = " ".join((user_message or "").strip().split())
    normalized = normalize_for_intent(text)
    if not normalized and attachments:
        return "heavy_task"
    if is_conversational_message(text):
        return "chat"
    # Si el mensaje pide buscar en internet → SIEMPRE chat (no código)
    # El pipeline de búsqueda web ya se encarga de obtener y presentar la info.
    web_search_triggers = [
        "busca en internet", "buscar en internet", "busca en la web", "buscar en la web",
        "googlea", "investiga en internet", "investiga en la web",
        "busca informacion", "busca información", "busca noticias",
        "busca en wikipedia", "buscar en wikipedia", "busca en google",
        "buscar informacion", "podrías buscarme", "podrias buscarme",
        "busca por internet", "buscar por internet",
    ]
    if any(t in normalized for t in web_search_triggers):
        return "chat"
    # Señales fuertes de "codigo"
    if has_code_signal(text, normalized):
        return "code"
    # Tareas pesadas / multi-step
    heavy_terms = [
        "presentacion", "ppt", "pptx", "diapositiva", "slide",
        "pdf", "resumen", "resume", "informe", "analisis", "análisis",
        "plan", "roadmap", "arquitectura", "diagrama",
        "lee", "revisa", "archivo", "carpeta", "proyecto",
    ]
    if attachments:
        return "heavy_task"
    if has_heavy_signal(normalized):
        return "heavy_task"
    return "chat"


def router_prompt(user_message: str, attachments_present: bool) -> str:
    # Prompt corto y barato, JSON estricto.
    msg = " ".join((user_message or "").strip().split())
    return (
        "Clasifica el MENSAJE del usuario en una de estas intenciones:\n"
        "- chat: charla casual, preguntas simples, trivia, numeros, etc.\n"
        "- code: pide codigo, debugging, errores, scripts, funciones.\n"
        "- heavy_task: tareas largas/multi-paso, analisis de archivos, presentaciones, esquemas extensos.\n\n"
        "Devuelve SOLO JSON estricto con estas claves:\n"
        '{"intent":"chat|code|heavy_task","confidence":0.0,"reason":"..."}\n\n'
        f"Adjuntos_presentes: {str(bool(attachments_present)).lower()}\n"
        f"Mensaje: {msg!r}\n"
    )


def classify_intent(user_message: str, ai_settings: Dict[str, Any], attachments: Optional[List[Dict[str, Any]]] = None) -> Dict[str, Any]:
    # Primero heuristica (barata y fiable); si el router esta activado, intentamos refinar con un modelo barato.
    base_intent = heuristic_intent(user_message, attachments=attachments)
    result: Dict[str, Any] = {"intent": base_intent, "confidence": 0.55, "reason": "heuristic"}

    # Fast path: búsquedas web explícitas → siempre chat, el router IA no debe interferir.
    _web_triggers = [
        "busca en internet", "buscar en internet", "busca en la web", "buscar en la web",
        "googlea", "investiga en internet", "investiga en la web",
        "busca informacion", "busca información", "busca noticias",
        "busca en wikipedia", "buscar en wikipedia", "busca en google",
        "buscar informacion", "podrías buscarme", "podrias buscarme",
        "busca por internet", "buscar por internet",
    ]
    _norm_msg = normalize_for_intent(user_message)
    if any(t in _norm_msg for t in _web_triggers):
        return {"intent": "chat", "confidence": 0.99, "reason": "web_search_trigger"}

    # Fast path: charla y preguntas normales no gastan ni router IA.
    if base_intent == "chat" and is_lightweight_chat_message(user_message, attachments=attachments):
        return {"intent": "chat", "confidence": 0.95, "reason": "fast_chat_path"}

    if not bool(ai_settings.get("auto_router_enabled", True)):
        return result

    msg = " ".join((user_message or "").strip().split())
    attachments_present = bool(attachments)
    provider = str(ai_settings.get("router_provider") or "auto").strip().lower() or "auto"
    model_ollama = str(ai_settings.get("router_model_ollama") or "qwen2.5:1.5b").strip() or "qwen2.5:1.5b"
    model_openai = str(ai_settings.get("router_model_openai") or "gpt-4.1-mini").strip() or "gpt-4.1-mini"

    # Cache (TTL) para evitar repetir la clasificación cuando el usuario reenvía el mismo mensaje.
    try:
        key_raw = json.dumps(
            {
                "m": msg,
                "a": attachments_present,
                "p": provider,
                "mo": model_ollama,
                "mop": model_openai,
            },
            ensure_ascii=False,
            sort_keys=True,
        ).encode("utf-8", errors="ignore")
        cache_key = hashlib.sha256(key_raw).hexdigest()
    except Exception:
        cache_key = ""

    if cache_key:
        now = time.time()
        cached = _ROUTER_CACHE.get(cache_key) or {}
        try:
            ts = float(cached.get("ts") or 0.0)
        except Exception:
            ts = 0.0
        if cached and (now - ts) < _ROUTER_CACHE_TTL_S and isinstance(cached.get("value"), dict):
            return dict(cached["value"])

    prompt = router_prompt(user_message, attachments_present)

    # Preferir Ollama si esta disponible y el provider lo permite.
    try_ollama = provider in {"auto", "ollama"} and is_ollama_running()
    if try_ollama:
        model_name = model_ollama
        try:
            # Llamada directa a /api/generate (mas simple) usando prompt corto, temp 0.
            router_options = ollama_options("programador", temperature=0, top_p=1, num_predict=120)
            router_options["num_ctx"] = 1024
            r = http_session().post(
                f"{OLLAMA_HOST}/api/generate",
                json={
                    "model": model_name,
                    "prompt": prompt,
                    "stream": False,
                    "options": router_options,
                },
                timeout=20,
            )
            if r.status_code == 200:
                data = r.json() if isinstance(r.json(), dict) else {}
                parsed = parse_router_json(str(data.get("response") or ""))
                if parsed and parsed.get("intent") in {"chat", "code", "heavy_task"}:
                    conf = parsed.get("confidence")
                    try:
                        conf_f = float(conf)
                    except Exception:
                        conf_f = 0.6
                    final = {
                        "intent": str(parsed["intent"]),
                        "confidence": max(0.0, min(1.0, conf_f)),
                        "reason": str(parsed.get("reason") or "ollama_router"),
                        "router_provider": "ollama",
                        "router_model": model_name,
                    }
                    if cache_key:
                        try:
                            _ROUTER_CACHE[cache_key] = {"ts": time.time(), "value": final}
                        except Exception:
                            pass
                    return final
        except Exception:
            pass

    # Fallback a OpenAI si se puede.
    try_openai = provider in {"auto", "openai"} and should_use_openai(ai_settings)
    if try_openai:
        model = model_openai
        try:
            resp = http_session().post(
                f"{OPENAI_API_BASE}/responses",
                headers={
                    "Authorization": f"Bearer {ai_settings['openai_api_key']}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": model,
                    "input": [
                        {
                            "role": "user",
                            "content": [{"type": "input_text", "text": prompt}],
                        }
                    ],
                    "temperature": 0,
                    "max_output_tokens": 120,
                    "store": False,
                },
                timeout=20,
            )
            if resp.status_code < 400:
                text, _sources = parse_openai_response(resp.json())
                parsed = parse_router_json(text)
                if parsed and parsed.get("intent") in {"chat", "code", "heavy_task"}:
                    conf = parsed.get("confidence")
                    try:
                        conf_f = float(conf)
                    except Exception:
                        conf_f = 0.6
                    final = {
                        "intent": str(parsed["intent"]),
                        "confidence": max(0.0, min(1.0, conf_f)),
                        "reason": str(parsed.get("reason") or "openai_router"),
                        "router_provider": "openai",
                        "router_model": model,
                    }
                    if cache_key:
                        try:
                            _ROUTER_CACHE[cache_key] = {"ts": time.time(), "value": final}
                        except Exception:
                            pass
                    return final
        except Exception:
            pass

    if cache_key:
        try:
            now = time.time()
            if len(_ROUTER_CACHE) > _ROUTER_CACHE_MAX:
                # Limpieza rápida: eliminar entradas viejas (sin ordenar toda la tabla).
                for k in list(_ROUTER_CACHE.keys())[: max(16, _ROUTER_CACHE_MAX // 4)]:
                    entry = _ROUTER_CACHE.get(k) or {}
                    ts = float(entry.get("ts") or 0.0)
                    if (now - ts) > _ROUTER_CACHE_TTL_S:
                        _ROUTER_CACHE.pop(k, None)
                # Si sigue enorme, recorta.
                if len(_ROUTER_CACHE) > _ROUTER_CACHE_MAX:
                    for k in list(_ROUTER_CACHE.keys())[: len(_ROUTER_CACHE) - _ROUTER_CACHE_MAX]:
                        _ROUTER_CACHE.pop(k, None)
            _ROUTER_CACHE[cache_key] = {"ts": now, "value": result}
        except Exception:
            pass
    return result


def choose_mode(requested_mode: str, user_message: str, attachments: Optional[List[Dict[str, Any]]], ai_settings: Dict[str, Any]) -> tuple[str, Dict[str, Any]]:
    requested_mode = (requested_mode or "auto").strip().lower()
    if requested_mode not in {"auto", "rapido", "combinado", "codigo"}:
        requested_mode = "auto"

    classification = classify_intent(user_message, ai_settings, attachments=attachments)
    intent = str(classification.get("intent") or "chat")
    classification["requested_mode"] = requested_mode

    # Mapeo intent -> modo por defecto
    mapped = {"chat": "rapido", "code": "codigo", "heavy_task": "combinado"}.get(intent, "rapido")

    # La charla y las preguntas normales siempre entran por rapido, aunque el selector este en codigo/combinado.
    if intent == "chat":
        classification["mode_reason"] = "temporary_fast_chat"
        classification["effective_mode"] = "rapido"
        return "rapido", classification

    # Si el usuario fuerza un modo manual, lo respetamos para tareas reales.
    if requested_mode in {"rapido", "combinado", "codigo"}:
        classification["mode_reason"] = "manual"
        classification["effective_mode"] = requested_mode
        return requested_mode, classification

    classification["mode_reason"] = "auto_router"
    classification["effective_mode"] = mapped
    return mapped, classification


def update_memory_async(chat: Dict[str, Any], user_message: str, assistant_text: str) -> None:
    def worker() -> None:
        if not MEMORY_UPDATE_LOCK.acquire(blocking=False):
            return
        try:
            user_id = str(chat.get("user_id", ""))
            current = load_memory(user_id).get("summary", "").strip()
            prompt = f"""Actualiza la memoria compartida de Nexo.

Memoria actual:
{current or "Sin memoria previa."}

Nuevo intercambio:
Usuario: {user_message}
Asistente: {assistant_text}

Devuelve solo un resumen breve de datos persistentes útiles: preferencias, contexto estable, decisiones y objetivos. No guardes información sensible como contraseñas."""
            summary = ollama_chat(
                "arquitecto",
                [
                    {"role": "system", "content": "Eres un gestor de memoria conciso y prudente."},
                    {"role": "user", "content": prompt},
                ],
            )
            if summary.strip():
                save_memory(user_id, summary)
        except Exception:
            pass
        finally:
            MEMORY_UPDATE_LOCK.release()

    threading.Thread(target=worker, daemon=True).start()


def render_auth_page(register: bool = False, error: str = "") -> str:
    if register:
        return render_template_string(
            LOGIN_HTML,
            page_title="Crear cuenta",
            subtitle="Crea una cuenta para usar Nexo online.",
            action="/register",
            register=True,
            password_autocomplete="new-password",
            button_text="Crear cuenta",
            switch_text="Ya tienes cuenta?",
            switch_href="/login",
            switch_label="Entrar",
            error=error,
        )

    return render_template_string(
        LOGIN_HTML,
        page_title="Entrar",
        subtitle="Entra para hablar con Nexo o crea una cuenta nueva.",
        action="/login",
        register=False,
        password_autocomplete="current-password",
        button_text="Entrar",
        switch_text="No tienes cuenta?",
        switch_href="/register",
        switch_label="Crear cuenta",
        error=error,
    )


def create_app() -> Flask:
    ensure_data_files()
    cleanup_uploads()
    app = Flask(__name__)
    app.secret_key = get_secret_key()
    secure_cookie = os.getenv("WEB_SESSION_SECURE", "1").lower() not in {"0", "false", "no"}
    app.config.update(
        SESSION_COOKIE_HTTPONLY=True,
        SESSION_COOKIE_SAMESITE="Lax",
        SESSION_COOKIE_SECURE=secure_cookie,
        MAX_CONTENT_LENGTH=MAX_TOTAL_UPLOAD_BYTES,
    )

    @app.get("/login")
    def login_page() -> str:
        return render_auth_page(register=False)

    @app.post("/login")
    def login_post() -> Response | str:
        ip = client_ip()
        wait_seconds = login_limited(ip)
        if wait_seconds:
            return render_auth_page(
                register=False,
                error=f"Demasiados intentos. Espera {wait_seconds} segundos.",
            ), 429

        username = request.form.get("username", "").strip()
        password = request.form.get("password", "")
        user = validate_login(username, password)
        if user:
            clear_login_failures(ip)
            session.clear()
            session["authenticated"] = True
            session["user_id"] = user["id"]
            session["username"] = user["username"]
            return redirect(url_for("index"))

        record_login_failure(ip)
        return render_auth_page(register=False, error="Usuario o contrasena incorrectos."), 401

    @app.get("/register")
    def register_page() -> str:
        return render_auth_page(register=True)

    @app.get("/donate")
    def donate_page() -> str:
        donate_url = configured_donate_url()
        return render_template_string(DONATE_HTML, donate_ready=bool(donate_url))

    @app.get("/donate/go")
    def donate_go() -> Response:
        donate_url = configured_donate_url()
        if not donate_url:
            return redirect(url_for("donate_page"))
        return redirect(donate_url, code=302)

    @app.post("/register")
    def register_post() -> Response | str:
        username = request.form.get("username", "").strip()
        password = request.form.get("password", "")
        confirm_password = request.form.get("confirm_password", "")
        # ═══ NEXO MEJORAS: Invite code validation ═══
        _invite_codes_env = os.getenv('NEXO_INVITE_CODES', '').strip()
        if _invite_codes_env:
            _valid_codes = {c.strip() for c in _invite_codes_env.split(',') if c.strip()}
            _submitted = request.form.get('invite_code', '').strip()
            if _submitted not in _valid_codes:
                return render_auth_page(register=True, error='Código de invitación inválido. Pídelo al administrador.'), 400
        user, error = create_user(username, password, confirm_password)
        if error or not user:
            return render_auth_page(register=True, error=error or "No se pudo crear la cuenta."), 400

        session.clear()
        session["authenticated"] = True
        session["user_id"] = user["id"]
        session["username"] = user["username"]
        return redirect(url_for("index"))

    @app.post("/logout")
    def logout() -> Response:
        session.clear()
        return redirect(url_for("login_page"))

    @app.get("/")
    def index() -> Response | str:
        if not session.get("authenticated") or not current_user_id():
            return redirect(url_for("login_page"))
        is_admin = is_admin_user(session.get("username", ""))
        return render_template_string(MAIN_HTML, is_admin=is_admin)

    @app.get("/api/chats")
    def api_chats() -> Response:
        auth = require_login_response()
        if auth:
            return auth
        return jsonify({"chats": list_chats(current_user_id())})

    @app.post("/api/chats")
    def api_create_chat() -> Response:
        auth = require_login_response()
        if auth:
            return auth
        return jsonify(create_chat(current_user_id()))

    @app.get("/api/chats/<chat_id>")
    def api_get_chat(chat_id: str) -> Response:
        auth = require_login_response()
        if auth:
            return auth
        chat = get_chat(chat_id, current_user_id())
        if not chat:
            return jsonify({"error": "Chat no encontrado"}), 404
        return jsonify(chat)

    @app.delete("/api/chats/<chat_id>")
    def api_delete_chat(chat_id: str) -> Response:
        auth = require_login_response()
        if auth:
            return auth
        deleted = delete_chat(chat_id, current_user_id())
        if not deleted:
            return jsonify({"error": "Chat no encontrado"}), 404
        return jsonify({"ok": True})

    @app.get("/api/account")
    def api_account() -> Response:
        auth = require_login_response()
        if auth:
            return auth
        user = authenticated_user()
        if not user:
            return jsonify({"error": "No autenticado"}), 401
        # Calcular cupos y estado de la promo para mostrarlo en cuenta
        promo_data = load_users_data().get("promo") or {}
        promo_start_dt = parse_utc(promo_data.get("start"))
        promo_slots_used = int(promo_data.get("slots_used") or 0)
        promo_active = (
            promo_start_dt is not None
            and datetime.now(timezone.utc) <= promo_start_dt + timedelta(days=PROMO_WINDOW_DAYS)
            and promo_slots_used < PROMO_SLOTS
        )
        return jsonify(
            {
                "user": {"id": user.get("id"), "username": user.get("username")},
                "plan": public_plan_for_user(user),
                "plan_expires_at": user.get("plan_expires_at"),
                "plan_after_expiry": user.get("plan_after_expiry"),
                "api_key": api_key_info_for_user(user),
                "plans": public_plan_catalog(),
                "is_admin": is_admin_user(str(user.get("username", ""))),
                "promo": {
                    "active": promo_active,
                    "slots_left": max(0, PROMO_SLOTS - promo_slots_used),
                    "ends_at": (promo_start_dt + timedelta(days=PROMO_WINDOW_DAYS)).isoformat(timespec="seconds") if promo_start_dt else None,
                },
            }
        )

    @app.get("/planes")
    def planes_page() -> Response | str:
        if not session.get("authenticated"):
            return redirect(url_for("login_page"))
        donate_url = configured_donate_url()
        return render_template_string(PLANES_HTML, donate_url=donate_url)

    @app.get("/neural")
    def neural_page() -> Response | str:
        if not session.get("authenticated") or not current_user_id():
            return redirect(url_for("login_page"))
        if not is_admin_user(session.get("username", "")):
            return redirect(url_for("index"))
        return render_template_string(NEURAL_HTML)

    @app.post("/api/neural/process")
    def neural_process() -> Response:
        if not session.get("authenticated") or not current_user_id():
            return jsonify({"error": "Unauthorized"}), 401
        if not is_admin_user(session.get("username", "")):
            return jsonify({"error": "Sin permisos"}), 403
        data = request.get_json(silent=True) or {}
        input_type = data.get("type", "unknown")
        content    = data.get("content", "")
        added = 2 + int(len(content) % 2)  # 2 or 3 neuronas simuladas
        messages = {
            "file":   f"Archivo '{content}' procesado. {added} nuevas conexiones neuronales.",
            "link":   f"Enlace analizado. {added} patrones extraídos.",
            "search": f"Búsqueda '{content}' completada. {added} conceptos integrados.",
        }
        message = messages.get(input_type, f"Contenido procesado. {added} neuronas añadidas.")
        return jsonify({"added": added, "message": message})

    @app.get("/api/neural/stats")
    def api_neural_stats() -> Response:
        auth = require_login_response()
        if auth:
            return auth
        if not is_admin_user(session.get("username", "")):
            return jsonify({"error": "Sin permisos"}), 403
        stats: Dict[str, Any] = {}
        cpu_load: float = 0.0
        gpu_load: float = 0.0
        ram_pct:  float = 0.0
        try:
            import psutil
            cpu_load = round(psutil.cpu_percent(interval=0.1), 1)
            vm = psutil.virtual_memory()
            ram_pct  = round(vm.percent, 1)
        except Exception:
            pass
        try:
            import GPUtil
            gpus = GPUtil.getGPUs()
            if gpus:
                gpu_load = round(gpus[0].load * 100, 1)
        except Exception:
            pass
        with AI_PRIORITY_CONDITION:
            active = AI_ACTIVE_REQUESTS
            queued = len(AI_PRIORITY_QUEUE)
        neuron_count = active + queued + 10
        is_busy      = active > 0 or queued > 0
        stats["cpu_load"]     = cpu_load
        stats["gpu_load"]     = gpu_load
        stats["ram_pct"]      = ram_pct
        stats["creativity"]   = gpu_load
        stats["potential"]    = cpu_load
        stats["learning"]     = ram_pct
        stats["active"]       = active
        stats["queued"]       = queued
        stats["neuron_count"] = neuron_count
        stats["status"]       = "Procesando..." if is_busy else "Online"
        return jsonify(stats)

    @app.get("/admin")
    def admin_page() -> Response | str:
        if not session.get("authenticated"):
            return redirect(url_for("login_page"))
        username = session.get("username", "")
        if not is_admin_user(username):
            return redirect(url_for("index"))
        return render_template_string(ADMIN_HTML)

    @app.get("/api/admin/users")
    def api_admin_users() -> Response:
        auth = require_login_response()
        if auth:
            return auth
        if not is_admin_user(session.get("username", "")):
            return jsonify({"error": "Sin permisos"}), 403
        return jsonify(admin_list_users())

    @app.post("/api/admin/users/<user_id>/plan")
    def api_admin_set_plan(user_id: str) -> Response:
        auth = require_login_response()
        if auth:
            return auth
        if not is_admin_user(session.get("username", "")):
            return jsonify({"error": "Sin permisos"}), 403
        body = request.get_json(silent=True) or {}
        new_plan = str(body.get("plan", "")).strip()
        if not new_plan:
            return jsonify({"error": "Falta el campo 'plan'"}), 400
        try:
            return jsonify(admin_set_user_plan(user_id, new_plan))
        except LookupError as e:
            return jsonify({"error": str(e)}), 404
        except Exception as e:
            return jsonify({"error": str(e)}), 400

    @app.post("/api/account/api-key")
    def api_account_api_key() -> Response:
        auth = require_login_response()
        if auth:
            return auth
        user = authenticated_user()
        if not user:
            return jsonify({"error": "No autenticado"}), 401
        if not plan_definition(user_plan_key(user)).get("includes_api_key"):
            return jsonify({"error": "Tu plan no incluye API Key."}), 403
        try:
            return jsonify(rotate_user_api_key(str(user["id"])))
        except Exception as exc:
            return jsonify({"error": str(exc)}), 400

    # ═══ NEXO MEJORAS: Monitor de sistema en tiempo real ═══
    if init_stats_addon:
        try:
            init_stats_addon(app, is_admin_user_fn=is_admin_user, public_plan_fn=public_plan_for_user)
        except Exception as _e:
            print(f"[STATS] Error al iniciar addon: {_e}")

    @app.get("/api/system-stats")
    def api_system_stats() -> Response:
        auth = require_login_response()
        if auth:
            return auth
        stats: Dict[str, Any] = {}
        try:
            import psutil
            stats['cpu'] = round(psutil.cpu_percent(interval=0.1), 1)
            vm = psutil.virtual_memory()
            stats['ram_used_gb'] = round(vm.used / 1e9, 1)
            stats['ram_total_gb'] = round(vm.total / 1e9, 1)
        except Exception:
            pass
        try:
            import GPUtil
            gpus = GPUtil.getGPUs()
            if gpus:
                g = gpus[0]
                stats['gpu_load']   = round(g.load * 100, 1)
                stats['vram_used']  = round(g.memoryUsed)
                stats['vram_total'] = round(g.memoryTotal)
                stats['vram_free']  = round(g.memoryFree)
                stats['gpu_name']   = g.name
        except Exception:
            pass
        with AI_PRIORITY_CONDITION:
            stats['queue_size']      = len(AI_PRIORITY_QUEUE)
            stats['active_requests'] = AI_ACTIVE_REQUESTS
        return jsonify(stats)

    @app.post("/api/chat/stream")
    def api_chat_stream() -> Response:
        try:
            auth = require_login_response()
            if auth:
                return auth
            request_user = authenticated_user()
            if not request_user:
                return jsonify({"error": "No autenticado"}), 401
            user_id = str(request_user["id"])
            # ═══ NEXO MEJORAS: Rate Limiting ═══
            _rl_wait = check_rate_limit(user_id)
            if _rl_wait:
                return jsonify({"error": f"Demasiadas peticiones. Espera {_rl_wait}s antes de enviar otro mensaje.", "retry_after": _rl_wait}), 429
            user_plan = public_plan_for_user(request_user)
            uploaded_files: List[Any] = []
            if request.mimetype == "multipart/form-data":
                chat_id = str(request.form.get("chat_id", "")).strip()
                mode = str(request.form.get("mode", "rapido")).strip().lower()
                user_message = str(request.form.get("message", "")).strip()
                uploaded_files = request.files.getlist("files[]") or request.files.getlist("files")
            else:
                payload = request.get_json(silent=True) or {}
                chat_id = str(payload.get("chat_id", "")).strip()
                mode = str(payload.get("mode", "rapido")).strip().lower()
                user_message = str(payload.get("message", "")).strip()
                set_personality(str(payload.get("personality", "normal")).strip())
            if mode not in {"auto", "rapido", "combinado", "codigo"}:
                return jsonify({"error": "Modo no válido"}), 400
            try:
                track_message(user_id=user_id, username=str(request_user.get("username","?")), plan=str(user_plan), mode=mode, message=user_message)
            except Exception:
                pass
            if not user_message and not uploaded_files:
                return jsonify({"error": "Mensaje vacío"}), 400

            chat = get_chat(chat_id, user_id) if chat_id else None
            if not chat:
                chat = create_chat(user_id)

            cleanup_uploads()
            stored_attachments, upload_error = store_uploaded_files(uploaded_files, user_id, str(chat["id"]))
            if upload_error:
                return jsonify({"error": upload_error}), 400

            display_message = user_message or "Analiza los archivos adjuntos."
            ai_settings = load_ai_settings()
            mode, classification = choose_mode(mode, display_message, stored_attachments, ai_settings)
            mode_error = plan_mode_error(user_plan["key"], mode)
            if mode_error:
                remove_stored_attachment_files(stored_attachments)
                return jsonify(
                    {
                        "error": mode_error,
                        "plan": user_plan,
                        "required_plan": public_plan_definition(required_plan_for_mode(mode)),
                        "mode": mode,
                        "intent": classification.get("intent"),
                    }
                ), 403
        except Exception as exc:
            print(traceback.format_exc(), flush=True)
            return jsonify({"error": f"Fallo interno: {exc}"}), 500

        def generate() -> Generator[str, None, None]:
            chat.setdefault("messages", [])
            now = utc_now()
            user_record: Dict[str, Any] = {
                "role": "user",
                "content": display_message,
                "mode": mode,
                "plan": user_plan.get("key"),
                "created_at": now,
            }
            if stored_attachments:
                user_record["attachments"] = [attachment_for_chat(item) for item in stored_attachments]
            chat["messages"].append(user_record)
            if chat.get("title") == "Nuevo chat":
                chat["title"] = chat_title_from_message(display_message)
            chat["updated_at"] = now
            update_chat(chat, user_id)

            assistant_parts: List[str] = []
            sources: List[Dict[str, str]] = []
            web_context = ""
            search_error = ""
            processed_attachments: List[Dict[str, Any]] = []
            assistant_provider = ""
            fallback_reason = ""
            slot_acquired = False
            try:
                wait_seconds = acquire_ai_slot(int(user_plan.get("priority") or 10))
                slot_acquired = True
                if wait_seconds > 0.5:
                    yield event({"type": "status", "message": "Turno listo. Respondiendo..."})
                label = {"rapido": "Rápido", "combinado": "Combinado", "codigo": "Código"}.get(mode, mode)
                yield event({
                    "type": "mode",
                    "mode": mode,
                    "label": mode_label(mode),
                    "intent": classification.get("intent"),
                    "mode_reason": classification.get("mode_reason"),
                    "plan": user_plan.get("key"),
                    "priority": user_plan.get("priority"),
                })
                if stored_attachments:
                    yield event({"type": "status", "message": "Procesando archivos..."})
                    for attachment in stored_attachments:
                        yield event({"type": "status", "message": f"Analizando {attachment.get('filename', 'archivo')}..."})
                        try:
                            processed_attachments.append(process_stored_attachment(attachment, ai_settings, mode=mode))
                        except Exception as exc:
                            attachment["summary"] = f"Error procesando archivo: {exc}"
                            attachment["status"] = "error"
                            processed_attachments.append(attachment)
                        user_record["attachments"] = [attachment_for_chat(item) for item in processed_attachments]
                        update_chat(chat, user_id)
                        yield event({
                            "type": "attachment_status",
                            "attachment": attachment_for_chat(processed_attachments[-1]),
                        })

                attachment_context = build_attachment_context(processed_attachments, mode=mode)
                if not should_use_openai(ai_settings) and should_search_web(display_message, attachment_context):
                    yield event({"type": "status", "message": "Pensando..."})
                    web_context, sources, search_error = search_web_context(display_message, attachment_context)
                    if search_error:
                        yield event({"type": "status", "message": "Pensando..."})

                answer_stream = stream_answer(
                    chat,
                    display_message,
                    mode,
                    attachments=processed_attachments,
                    web_context=web_context,
                    sources=sources,
                    search_error=search_error,
                    ai_settings=ai_settings,
                )
                while True:
                    try:
                        chunk = next(answer_stream)
                        parsed = json.loads(chunk)
                        if parsed.get("type") == "token":
                            token = parsed.get("token", "")
                            if token:
                                assistant_parts.append(str(token))
                        elif parsed.get("type") == "sources":
                            sources = parsed.get("sources", []) or []
                        elif parsed.get("type") == "provider":
                            assistant_provider = str(parsed.get("provider") or "")
                        elif parsed.get("type") == "fallback":
                            fallback_reason = str(parsed.get("reason") or "")
                        yield chunk
                    except StopIteration:
                        break
            except Exception as exc:
                yield event({"type": "error", "message": f"Error al llamar a Nexo: {exc}"})
                return
            finally:
                if slot_acquired:
                    release_ai_slot()

            assistant_text = "".join(assistant_parts)
            if assistant_text.strip():
                assistant_record: Dict[str, Any] = {
                    "role": "assistant",
                    "content": assistant_text,
                    "mode": mode,
                    "plan": user_plan.get("key"),
                    "created_at": utc_now(),
                }
                if assistant_provider:
                    assistant_record["provider"] = assistant_provider
                if fallback_reason:
                    assistant_record["fallback"] = fallback_reason
                if sources:
                    assistant_record["sources"] = sources
                chat["messages"].append(assistant_record)
                chat["updated_at"] = utc_now()
                update_chat(chat, user_id)
                update_memory_async(chat, display_message, assistant_text)
            yield event({"type": "done", "chat": chat})

        return Response(
            stream_with_context(generate()),
            mimetype="application/x-ndjson",
            headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
        )

    @app.post("/api/v1/chat")
    def api_v1_chat() -> Response:
        user = api_authenticated_user()
        if not user:
            return jsonify({"error": "API Key no valida"}), 401

        payload = request.get_json(silent=True) or {}
        user_message = str(payload.get("message", "")).strip()
        mode = str(payload.get("mode", "auto")).strip().lower()
        if mode not in {"auto", "rapido", "combinado", "codigo"}:
            return jsonify({"error": "Modo no valido"}), 400
        if not user_message:
            return jsonify({"error": "Mensaje vacio"}), 400

        ai_settings = load_ai_settings()
        selected_mode, classification = choose_mode(mode, user_message, [], ai_settings)
        user_plan = public_plan_for_user(user)
        mode_error = plan_mode_error(user_plan["key"], selected_mode)
        if mode_error:
            return jsonify(
                {
                    "error": mode_error,
                    "plan": user_plan,
                    "required_plan": public_plan_definition(required_plan_for_mode(selected_mode)),
                    "mode": selected_mode,
                    "intent": classification.get("intent"),
                }
            ), 403

        chat = {
            "id": "api",
            "user_id": str(user["id"]),
            "title": "API",
            "created_at": utc_now(),
            "updated_at": utc_now(),
            "messages": [],
        }
        assistant_parts: List[str] = []
        sources: List[Dict[str, str]] = []
        provider = ""
        fallback_reason = ""
        error_message = ""
        slot_acquired = False
        try:
            acquire_ai_slot(int(user_plan.get("priority") or 10))
            slot_acquired = True
            for chunk in stream_answer(chat, user_message, selected_mode, attachments=[], ai_settings=ai_settings):
                parsed = json.loads(chunk)
                if parsed.get("type") == "token":
                    assistant_parts.append(str(parsed.get("token") or ""))
                elif parsed.get("type") == "sources":
                    sources = parsed.get("sources", []) or []
                elif parsed.get("type") == "provider":
                    provider = str(parsed.get("provider") or "")
                elif parsed.get("type") == "fallback":
                    fallback_reason = str(parsed.get("reason") or "")
                elif parsed.get("type") == "error":
                    error_message = str(parsed.get("message") or "Error al llamar a Nexo.")
                    break
        except Exception as exc:
            error_message = str(exc)
        finally:
            if slot_acquired:
                release_ai_slot()

        text = "".join(assistant_parts).strip()
        if error_message and not text:
            return jsonify({"error": error_message}), 502
        return jsonify(
            {
                "text": text,
                "mode": selected_mode,
                "intent": classification.get("intent"),
                "provider": provider or ("ollama" if fallback_reason else ""),
                "fallback": fallback_reason,
                "sources": sources,
            }
        )

    @app.get("/api/memory")
    def api_memory() -> Response:
        auth = require_login_response()
        if auth:
            return auth
        return jsonify(load_memory(current_user_id()))

    @app.post("/api/memory/clear")
    def api_memory_clear() -> Response:
        auth = require_login_response()
        if auth:
            return auth
        save_memory(current_user_id(), "")
        return jsonify({"ok": True})

    # ── Reportar bugs ────────────────────────────────────────────────────────
    REPORTS_DIR = DATA_DIR / "reports"

    @app.post("/api/report")
    def api_report() -> Response:
        auth = require_login_response()
        if auth:
            return auth
        body = request.get_json(force=True, silent=True) or {}
        text = (body.get("text") or "").strip()
        category = (body.get("category") or "otro").strip()
        if not text:
            return jsonify({"error": "El reporte no puede estar vacío"}), 400
        REPORTS_DIR.mkdir(parents=True, exist_ok=True)
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        uid = current_user_id() or "anon"
        filename = REPORTS_DIR / f"reporte_{timestamp}_{uid[:8]}.txt"
        content = (
            f"Fecha: {datetime.now().strftime('%d/%m/%Y %H:%M:%S')}\n"
            f"Usuario: {uid}\n"
            f"Categoría: {category}\n"
            f"{'─' * 40}\n"
            f"{text}\n"
        )
        filename.write_text(content, encoding="utf-8")
        return jsonify({"ok": True, "file": filename.name})

    @app.get("/api/reports")
    def api_reports() -> Response:
        auth = require_login_response()
        if auth:
            return auth
        REPORTS_DIR.mkdir(parents=True, exist_ok=True)
        files = sorted(REPORTS_DIR.glob("reporte_*.txt"), reverse=True)
        reports = []
        for f in files[:50]:  # máximo 50 reportes
            try:
                raw = f.read_text(encoding="utf-8")
                lines = raw.splitlines()
                timestamp_line = next((l for l in lines if l.startswith("Fecha:")), "")
                category_line = next((l for l in lines if l.startswith("Categoría:")), "")
                sep_idx = next((i for i, l in enumerate(lines) if "─" in l), 3)
                body_text = "\n".join(lines[sep_idx + 1:]).strip()
                reports.append({
                    "timestamp": timestamp_line.replace("Fecha: ", ""),
                    "category": category_line.replace("Categoría: ", ""),
                    "text": body_text,
                    "file": f.name,
                })
            except Exception:
                pass
        return jsonify({"reports": reports})

    return app


def _warmup_ollama(model_keys) -> None:
    """
    Pre-carga (warm-up) de los modelos en VRAM antes de aceptar trafico real.
    Evita que el primer usuario espere 15-25 s mientras Ollama carga el modelo.
    Lanza una llamada minima con num_predict=1 a cada modelo necesario.
    """
    try:
        if not is_ollama_running() and not start_ollama():
            print("[WARMUP] Ollama no disponible; saltando warm-up.")
            return
        seen = set()
        for key in model_keys:
            model_name = MODELS.get(key)
            if not model_name or model_name in seen:
                continue
            seen.add(model_name)
            if not is_model_installed(model_name):
                print(f"[WARMUP] Modelo no instalado, saltando: {model_name}")
                continue
            t0 = time.time()
            try:
                opts = ollama_options(key, temperature=0.0, top_p=1.0, num_predict=1)
                payload = {
                    "model": model_name,
                    "messages": [{"role": "user", "content": "OK"}],
                    "stream": False,
                    "keep_alive": OLLAMA_KEEP_ALIVE,
                    "options": opts,
                }
                http_session().post(f"{OLLAMA_HOST}/api/chat", json=payload, timeout=(10, 120))
                print(f"[WARMUP] {key} ({model_name}) listo en {time.time()-t0:.1f}s")
            except Exception as exc:
                print(f"[WARMUP] Aviso al precalentar {model_name}: {exc}")
    except Exception as exc:
        print(f"[WARMUP] Error general: {exc}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Servidor web para Nexo")
    parser.add_argument("--host", default=os.getenv("WEB_HOST", "127.0.0.1"))
    parser.add_argument("--port", type=int, default=int(os.getenv("WEB_PORT", "7860")))
    parser.add_argument("--debug", action="store_true")
    args = parser.parse_args()

    app = create_app()
    print(f"\nNexo Web: http://{args.host}:{args.port}\n")

    # Warm-up en hilo aparte para no bloquear arranque (NEXO_WARMUP=1 por defecto).
    if (os.getenv("NEXO_WARMUP") or "1").strip().lower() not in {"0", "false", "no", "off"}:
        warmup_roles_env = (os.getenv("NEXO_WARMUP_ROLES") or "arquitecto,programador").strip()
        roles = [r.strip() for r in warmup_roles_env.split(",") if r.strip()]
        threading.Thread(target=_warmup_ollama, args=(roles,), daemon=True).start()

    if args.debug:
        # En Windows, el reloader puede fallar con WinError 740 en algunos entornos.
        # Debug si, reloader no.
        app.run(host=args.host, port=args.port, debug=True, threaded=True, use_reloader=False)
        return

    # === Servidor de produccion: Waitress (WSGI real, multi-hilo) ===
    # Mejora la concurrencia 2-3x frente al server de desarrollo de Werkzeug
    # cuando hay varios usuarios al mismo tiempo.
    # Para volver al server clasico de Flask: NEXO_USE_WAITRESS=0
    use_waitress = (os.getenv("NEXO_USE_WAITRESS") or "1").strip().lower() not in {
        "0", "false", "no", "off",
    }
    if not use_waitress:
        app.run(host=args.host, port=args.port, debug=False, threaded=True, use_reloader=False)
        return

    try:
        from waitress import serve  # type: ignore
    except ImportError:
        print("[AVISO] 'waitress' no esta instalado; usando server de desarrollo.")
        app.run(host=args.host, port=args.port, debug=False, threaded=True, use_reloader=False)
        return

    threads = int(os.getenv("NEXO_WAITRESS_THREADS", "16"))
    conn_limit = int(os.getenv("NEXO_WAITRESS_CONN_LIMIT", "200"))
    channel_timeout = int(os.getenv("NEXO_WAITRESS_CHANNEL_TIMEOUT", "600"))
    print(
        f"[INFO] Waitress: threads={threads} connection_limit={conn_limit} "
        f"channel_timeout={channel_timeout}s"
    )
    serve(
        app,
        host=args.host,
        port=args.port,
        threads=threads,
        connection_limit=conn_limit,
        channel_timeout=channel_timeout,
        ident="",
    )


if __name__ == "__main__":
    main()