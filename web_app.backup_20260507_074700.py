"""
Web UI privada para la NEXO local.

Ejecutar:
    python web_app.py --host 127.0.0.1 --port 7860
"""

from __future__ import annotations

import argparse
import base64
import hashlib
import io
import ipaddress
import json
import mimetypes
import os
import re
import secrets
import shutil
import socket
import tarfile
import threading
import time
import zipfile
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Any, Dict, Generator, Iterable, List, Optional
from urllib.parse import urlparse

import requests
from flask import (
    Flask,
    Response,
    jsonify,
    redirect,
    render_template_string,
    request,
    session,
    stream_with_context,
    url_for,
)
from werkzeug.security import check_password_hash, generate_password_hash
from werkzeug.utils import secure_filename

from orchestrator import (
    GPU_LAYERS,
    MODELS,
    OLLAMA_HOST,
    PROMPT_ARQUITECTO,
    PROMPT_PROGRAMADOR,
    is_model_installed,
    is_ollama_running,
    start_ollama,
)

try:
    from dotenv import load_dotenv

    load_dotenv()
except Exception:
    pass


APP_TITLE = "NEXO"
DATA_DIR = Path(os.getenv("IA_COMBINADA_WEB_DATA", "web_data"))
USERS_FILE = DATA_DIR / "users.json"
CHATS_FILE = DATA_DIR / "chats.json"
MEMORY_FILE = DATA_DIR / "memory.json"
SECRET_FILE = DATA_DIR / "secret_key.txt"
SETTINGS_FILE = DATA_DIR / "settings.json"
UPLOAD_DIR = DATA_DIR / "uploads"
UPLOAD_TTL_DAYS = 7
MAX_TOTAL_UPLOAD_BYTES = 128 * 1024 * 1024
MAX_ATTACHMENTS_PER_MESSAGE = 8
MAX_IMAGE_BYTES = 12 * 1024 * 1024
MAX_DOCUMENT_BYTES = 20 * 1024 * 1024
MAX_VIDEO_BYTES = 100 * 1024 * 1024
MAX_AUDIO_BYTES = 50 * 1024 * 1024
MAX_ARCHIVE_BYTES = 80 * 1024 * 1024
MAX_UNKNOWN_BYTES = 128 * 1024 * 1024
MAX_DOCUMENT_CHARS = 24_000
MAX_ATTACHMENT_CONTEXT_CHARS = 30_000
MAX_WEB_CONTEXT_CHARS = 12_000
MAX_ARCHIVE_MEMBERS = 30
MAX_ARCHIVE_EXTRACT_BYTES = 40 * 1024 * 1024
VIDEO_FRAME_COUNT = 24
VISION_IMAGE_BATCH = 6
OPENAI_API_BASE = "https://api.openai.com/v1"

IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif"}
TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".jsonl", ".xml", ".html", ".htm", ".css", ".js", ".ts",
    ".tsx", ".jsx", ".py", ".java", ".c", ".cpp", ".h", ".hpp", ".cs", ".go", ".rs", ".php",
    ".rb", ".swift", ".kt", ".kts", ".sql", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".env", ".log", ".bat", ".ps1", ".sh", ".dockerfile", ".vue", ".svelte",
}
PDF_EXTENSIONS = {".pdf"}
OFFICE_EXTENSIONS = {".docx", ".doc", ".rtf", ".odt", ".pptx", ".ppt", ".xlsx", ".xls", ".ods"}
DOCUMENT_EXTENSIONS = TEXT_EXTENSIONS | PDF_EXTENSIONS | OFFICE_EXTENSIONS
VIDEO_EXTENSIONS = {".mp4", ".webm", ".mov", ".avi", ".mkv", ".m4v"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".aac", ".ogg", ".oga", ".flac", ".wma", ".webm"}
ARCHIVE_EXTENSIONS = {".zip", ".tar", ".gz", ".tgz", ".bz2", ".tbz2", ".xz", ".txz", ".7z"}

UPLOAD_LIMITS = {
    "image": MAX_IMAGE_BYTES,
    "document": MAX_DOCUMENT_BYTES,
    "video": MAX_VIDEO_BYTES,
    "audio": MAX_AUDIO_BYTES,
    "archive": MAX_ARCHIVE_BYTES,
    "unknown": MAX_UNKNOWN_BYTES,
}

DEFAULT_SETTINGS = {
    "ai_provider": "openai",
    "openai_api_key": "",
    "openai_model": "gpt-5.5",
    "openai_transcribe_model": "gpt-4o-transcribe",
    "fallback_to_ollama": True,
}

PLACEHOLDER_PASSWORDS = {
    "",
    "TU_CONTRASENA",
    "TU_CONTRASE\u00d1A",
    "TU_CONTRASE\u00c3\u2018A",
    "CAMBIA_ESTA_CONTRASENA",
}
USERNAME_RE = re.compile(r"^[A-Za-z0-9_.-]{3,32}$")
LOGIN_ATTEMPTS: Dict[str, Dict[str, float]] = {}
LOGIN_LOCK = threading.Lock()
DATA_LOCK = threading.RLock()
MEMORY_UPDATE_LOCK = threading.Lock()


LOGIN_HTML = r"""
<!doctype html>
<html lang="es">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{{ page_title }} - NEXO</title>
  <style>
    :root {
      color-scheme: dark;
      --bg: #101010;
      --panel: #171717;
      --panel-2: #202020;
      --text: #f3f3f3;
      --muted: #a8a8a8;
      --line: #303030;
      --accent: #19c37d;
      --danger: #ff6b6b;
    }
    * { box-sizing: border-box; }
    body {
      margin: 0;
      min-height: 100vh;
      display: grid;
      place-items: center;
      background: var(--bg);
      color: var(--text);
      font-family: Inter, ui-sans-serif, system-ui, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
    }
    .login {
      width: min(420px, calc(100vw - 32px));
      padding: 28px;
      border: 1px solid var(--line);
      background: var(--panel);
      border-radius: 8px;
      box-shadow: 0 24px 70px rgba(0, 0, 0, .32);
    }
    h1 { margin: 0 0 8px; font-size: 26px; letter-spacing: 0; }
    p { margin: 0 0 24px; color: var(--muted); line-height: 1.5; }
    label { display: block; margin: 14px 0 8px; color: #dedede; font-size: 14px; }
    input {
      width: 100%;
      height: 44px;
      border: 1px solid var(--line);
      border-radius: 6px;
      background: #0f0f0f;
      color: var(--text);
      padding: 0 12px;
      font-size: 15px;
      outline: none;
    }
    input:focus { border-color: var(--accent); }
    button {
      width: 100%;
      height: 44px;
      margin-top: 22px;
      border: 0;
      border-radius: 6px;
      background: var(--accent);
      color: #06140e;
      font-weight: 700;
      cursor: pointer;
    }
    .error {
      margin: 0 0 16px;
      padding: 11px 12px;
      border-radius: 6px;
      font-size: 14px;
      line-height: 1.45;
    }
    .error { background: rgba(255, 107, 107, .12); color: #ffd0d0; border: 1px solid rgba(255, 107, 107, .3); }
    .switch {
      margin: 18px 0 0;
      color: var(--muted);
      text-align: center;
      font-size: 14px;
    }
    .switch a { color: var(--accent); text-decoration: none; font-weight: 700; }
    .switch a:hover { text-decoration: underline; }
  </style>
</head>
<body>
  <main class="login">
    <h1>NEXO</h1>
    <p>{{ subtitle }}</p>
    {% if error %}
      <div class="error">{{ error }}</div>
    {% endif %}
    <form method="post" action="{{ action }}">
      <label for="username">Usuario</label>
      <input id="username" name="username" autocomplete="username" required autofocus>
      <label for="password">Contrasena</label>
      <input id="password" name="password" type="password" autocomplete="{{ password_autocomplete }}" required>
      {% if register %}
        <label for="confirm_password">Repite contrasena</label>
        <input id="confirm_password" name="confirm_password" type="password" autocomplete="new-password" required>
      {% endif %}
      <button type="submit">{{ button_text }}</button>
    </form>
    <p class="switch">{{ switch_text }} <a href="{{ switch_href }}">{{ switch_label }}</a></p>
  </main>
</body>
</html>
"""


MAIN_HTML = r"""
<!doctype html>
<html lang="es">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>NEXO</title>
  <style>
    :root {
      color-scheme: dark;
      --bg: #212121;
      --sidebar: #171717;
      --sidebar-2: #202020;
      --panel: #2b2b2b;
      --panel-soft: #303030;
      --text: #f4f4f4;
      --muted: #b4b4b4;
      --muted-2: #8c8c8c;
      --line: #3b3b3b;
      --accent: #19c37d;
      --accent-2: #2dd4bf;
      --danger: #ff6b6b;
      --code: #111111;
    }
    * { box-sizing: border-box; }
    html, body { height: 100%; }
    body {
      margin: 0;
      background: var(--bg);
      color: var(--text);
      font-family: Inter, ui-sans-serif, system-ui, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
      overflow: hidden;
    }
    button, input, textarea, select { font: inherit; }
    button { cursor: pointer; }
    .app {
      display: grid;
      grid-template-columns: 292px 1fr;
      height: 100vh;
      min-height: 0;
    }
    .sidebar {
      display: flex;
      flex-direction: column;
      min-height: 0;
      background: var(--sidebar);
      border-right: 1px solid #111;
    }
    .brand {
      display: flex;
      align-items: center;
      gap: 10px;
      padding: 14px;
      border-bottom: 1px solid #242424;
    }
    .mark {
      display: grid;
      place-items: center;
      width: 32px;
      height: 32px;
      border-radius: 6px;
      background: #f2f2f2;
      color: #111;
      font-weight: 800;
    }
    .brand strong { font-size: 15px; }
    .brand span { display: block; color: var(--muted-2); font-size: 12px; margin-top: 1px; }
    .side-actions { padding: 12px; display: grid; gap: 10px; }
    .btn {
      display: inline-flex;
      align-items: center;
      justify-content: center;
      gap: 8px;
      min-height: 40px;
      border: 1px solid var(--line);
      border-radius: 6px;
      background: var(--sidebar-2);
      color: var(--text);
      padding: 0 12px;
    }
    .btn:hover { background: #282828; }
    .btn.primary { background: var(--accent); border-color: var(--accent); color: #07130e; font-weight: 700; }
    .btn.ghost { background: transparent; }
    .btn.danger { border-color: rgba(255, 107, 107, .45); color: #ffd1d1; }
    select {
      width: 100%;
      height: 40px;
      border: 1px solid var(--line);
      border-radius: 6px;
      background: #121212;
      color: var(--text);
      padding: 0 10px;
      outline: none;
    }
    .history-title {
      padding: 8px 14px;
      color: var(--muted-2);
      font-size: 12px;
      text-transform: uppercase;
      letter-spacing: .04em;
    }
    .chat-list {
      overflow: auto;
      padding: 0 8px 12px;
      flex: 1;
      min-height: 0;
    }
    .chat-item {
      width: 100%;
      text-align: left;
      border: 0;
      border-radius: 6px;
      background: transparent;
      color: #e8e8e8;
      padding: 10px 10px;
      margin: 2px 0;
      overflow: hidden;
      text-overflow: ellipsis;
      white-space: nowrap;
    }
    .chat-item:hover, .chat-item.active { background: #242424; }
    .sidebar-footer {
      padding: 12px;
      border-top: 1px solid #242424;
      display: grid;
      gap: 8px;
    }
    .main {
      display: flex;
      flex-direction: column;
      min-width: 0;
      min-height: 0;
      background: var(--bg);
    }
    .topbar {
      height: 58px;
      display: flex;
      align-items: center;
      justify-content: space-between;
      padding: 0 18px;
      border-bottom: 1px solid var(--line);
      background: rgba(33, 33, 33, .92);
    }
    .topbar h2 { margin: 0; font-size: 16px; letter-spacing: 0; }
    .topbar-actions { display: flex; align-items: center; gap: 8px; }
    .mobile-menu { display: none; }
    .messages {
      flex: 1;
      min-height: 0;
      overflow: auto;
      padding: 24px 18px 190px;
    }
    .empty {
      height: 100%;
      display: grid;
      place-items: center;
      text-align: center;
      color: var(--muted);
    }
    .empty h1 { margin: 0 0 10px; color: var(--text); font-size: clamp(28px, 5vw, 44px); letter-spacing: 0; }
    .empty p { margin: 0; max-width: 560px; line-height: 1.55; }
    .message {
      width: min(860px, 100%);
      margin: 0 auto 20px;
      display: grid;
      grid-template-columns: 34px minmax(0, 1fr);
      gap: 12px;
    }
    .avatar {
      width: 34px;
      height: 34px;
      border-radius: 6px;
      display: grid;
      place-items: center;
      font-size: 13px;
      font-weight: 800;
      background: #f2f2f2;
      color: #111;
    }
    .message.user .avatar { background: var(--accent-2); color: #06211e; }
    .bubble {
      min-width: 0;
      padding: 6px 0 0;
      line-height: 1.6;
      color: #f2f2f2;
      word-wrap: break-word;
    }
    .message.user .bubble {
      justify-self: end;
      max-width: min(680px, 100%);
      background: var(--panel-soft);
      padding: 10px 14px;
      border-radius: 8px;
    }
    .bubble p { margin: 0 0 12px; }
    .bubble p:last-child { margin-bottom: 0; }
    .bubble ul, .bubble ol { margin: 8px 0 12px 22px; padding: 0; }
    .bubble pre {
      position: relative;
      margin: 12px 0;
      border-radius: 8px;
      overflow: hidden;
      background: var(--code);
      border: 1px solid #303030;
    }
    .bubble code {
      font-family: "Cascadia Code", Consolas, "SFMono-Regular", monospace;
      font-size: 13px;
    }
    .bubble pre code {
      display: block;
      padding: 42px 14px 14px;
      overflow-x: auto;
      white-space: pre;
    }
    .code-copy {
      position: absolute;
      top: 8px;
      right: 8px;
      height: 28px;
      border: 1px solid #3f3f3f;
      border-radius: 6px;
      background: #202020;
      color: #e8e8e8;
      padding: 0 10px;
      font-size: 12px;
    }
    .attachment-list, .source-list {
      display: grid;
      gap: 8px;
      margin-top: 10px;
    }
    .attachment-item, .source-item {
      border: 1px solid var(--line);
      border-radius: 6px;
      background: rgba(255, 255, 255, .035);
      padding: 9px 10px;
      font-size: 13px;
      color: #dedede;
    }
    .attachment-item strong, .source-item strong {
      display: block;
      color: var(--text);
      font-size: 13px;
      margin-bottom: 3px;
      overflow-wrap: anywhere;
    }
    .attachment-item span, .source-item span {
      color: var(--muted-2);
      font-size: 12px;
    }
    .source-item a {
      color: var(--accent-2);
      text-decoration: none;
      overflow-wrap: anywhere;
    }
    .source-item a:hover { text-decoration: underline; }
    .status {
      width: min(860px, 100%);
      margin: 0 auto 12px;
      color: var(--muted-2);
      font-size: 13px;
    }
    .composer {
      position: fixed;
      left: 292px;
      right: 0;
      bottom: 0;
      padding: 16px 18px 20px;
      background: linear-gradient(to top, var(--bg) 72%, rgba(33, 33, 33, 0));
    }
    .composer.dragging .composer-box {
      border-color: var(--accent-2);
      background: #263430;
    }
    .pending-files {
      width: min(860px, 100%);
      margin: 0 auto 8px;
      display: none;
      gap: 8px;
      grid-template-columns: repeat(auto-fit, minmax(180px, 1fr));
    }
    .pending-files.open { display: grid; }
    .pending-file {
      min-width: 0;
      display: grid;
      grid-template-columns: 42px minmax(0, 1fr) 28px;
      gap: 8px;
      align-items: center;
      border: 1px solid var(--line);
      border-radius: 6px;
      background: var(--panel);
      padding: 7px;
    }
    .pending-thumb {
      width: 42px;
      height: 42px;
      border-radius: 6px;
      object-fit: cover;
      background: #111;
      display: grid;
      place-items: center;
      color: var(--muted);
      font-size: 11px;
      border: 1px solid #383838;
    }
    .pending-file strong {
      display: block;
      overflow: hidden;
      white-space: nowrap;
      text-overflow: ellipsis;
      font-size: 13px;
    }
    .pending-file span {
      display: block;
      color: var(--muted-2);
      font-size: 12px;
      margin-top: 2px;
    }
    .remove-file, .attach {
      width: 42px;
      height: 42px;
      border: 1px solid var(--line);
      border-radius: 6px;
      background: #202020;
      color: var(--text);
    }
    .remove-file {
      width: 28px;
      height: 28px;
      color: #ffd0d0;
    }
    .composer-box {
      width: min(860px, 100%);
      margin: 0 auto;
      display: grid;
      grid-template-columns: 42px 1fr 44px;
      gap: 10px;
      align-items: end;
      border: 1px solid var(--line);
      border-radius: 8px;
      background: var(--panel);
      padding: 10px;
      box-shadow: 0 18px 50px rgba(0, 0, 0, .18);
    }
    textarea {
      resize: none;
      max-height: 180px;
      min-height: 42px;
      border: 0;
      outline: none;
      background: transparent;
      color: var(--text);
      padding: 10px 8px;
      line-height: 1.45;
    }
    .send {
      width: 42px;
      height: 42px;
      border: 0;
      border-radius: 6px;
      background: var(--accent);
      color: #06140e;
      font-weight: 900;
    }
    .send:disabled { opacity: .5; cursor: not-allowed; }
    .modal-backdrop {
      position: fixed;
      inset: 0;
      display: none;
      align-items: center;
      justify-content: center;
      background: rgba(0, 0, 0, .55);
      z-index: 20;
      padding: 18px;
    }
    .modal-backdrop.open { display: flex; }
    .modal {
      width: min(680px, 100%);
      max-height: min(680px, calc(100vh - 40px));
      overflow: auto;
      background: #181818;
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 18px;
    }
    .modal header {
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 12px;
      margin-bottom: 12px;
    }
    .modal h3 { margin: 0; font-size: 17px; }
    .memory-text {
      min-height: 160px;
      white-space: pre-wrap;
      color: #e8e8e8;
      background: #101010;
      border: 1px solid #2a2a2a;
      border-radius: 6px;
      padding: 12px;
      line-height: 1.55;
    }
    @media (max-width: 760px) {
      .app { grid-template-columns: 1fr; }
      .sidebar {
        position: fixed;
        inset: 0 auto 0 0;
        width: min(292px, 86vw);
        z-index: 10;
        transform: translateX(-100%);
        transition: transform .18s ease;
      }
      .sidebar.open { transform: translateX(0); }
      .mobile-menu { display: inline-flex; }
      .composer { left: 0; }
      .topbar { padding: 0 12px; }
      .messages { padding-inline: 12px; }
      .message { grid-template-columns: 30px minmax(0, 1fr); gap: 10px; }
      .avatar { width: 30px; height: 30px; }
    }
  </style>
</head>
<body>
  <div class="app">
    <aside id="sidebar" class="sidebar">
      <div class="brand">
        <div class="mark">IA</div>
        <div><strong>NEXO</strong><span>OpenAI primero, Ollama respaldo</span></div>
      </div>
      <div class="side-actions">
        <button id="newChatBtn" class="btn primary" type="button">+ Nuevo chat</button>
        <select id="modeSelect" aria-label="Modo de IA">
          <option value="rapido">Rapido</option>
          <option value="combinado" selected>Combinado</option>
          <option value="codigo">Codigo</option>
        </select>
      </div>
      <div class="history-title">Chats</div>
      <div id="chatList" class="chat-list"></div>
      <div class="sidebar-footer">
        <button id="memoryBtn" class="btn ghost" type="button">Memoria</button>
        <form method="post" action="/logout">
          <button class="btn danger" type="submit">Salir</button>
        </form>
      </div>
    </aside>

    <main class="main">
      <header class="topbar">
        <div class="topbar-actions">
          <button id="menuBtn" class="btn ghost mobile-menu" type="button">â˜°</button>
          <h2 id="chatTitle">Nuevo chat</h2>
        </div>
        <div id="status" class="topbar-actions"></div>
      </header>
      <section id="messages" class="messages">
        <div class="empty">
          <div>
            <h1>Â¿QuÃ© hacemos hoy?</h1>
            <p>Elige modo, escribe tu mensaje y la NEXO responderÃ¡ usando memoria local.</p>
          </div>
        </div>
      </section>
      <section id="composer" class="composer">
        <div id="pendingFiles" class="pending-files"></div>
        <div class="composer-box">
          <button id="attachBtn" class="attach" type="button" aria-label="Adjuntar archivo" title="Adjuntar archivo">+</button>
          <input id="fileInput" type="file" multiple hidden>
          <textarea id="prompt" rows="1" placeholder="Mensaje para la NEXO"></textarea>
          <button id="sendBtn" class="send" type="button" aria-label="Enviar">â†‘</button>
        </div>
      </section>
    </main>
  </div>

  <div id="memoryModal" class="modal-backdrop">
    <div class="modal">
      <header>
        <h3>Memoria compartida</h3>
        <button id="closeMemoryBtn" class="btn ghost" type="button">Cerrar</button>
      </header>
      <div id="memoryText" class="memory-text"></div>
      <button id="clearMemoryBtn" class="btn danger" type="button" style="margin-top: 12px;">Borrar memoria</button>
    </div>
  </div>

  <script>
    const state = { chats: [], activeChatId: null, busy: false, pendingFiles: [] };
    const els = {
      composer: document.getElementById('composer'),
      sidebar: document.getElementById('sidebar'),
      chatList: document.getElementById('chatList'),
      messages: document.getElementById('messages'),
      prompt: document.getElementById('prompt'),
      sendBtn: document.getElementById('sendBtn'),
      attachBtn: document.getElementById('attachBtn'),
      fileInput: document.getElementById('fileInput'),
      pendingFiles: document.getElementById('pendingFiles'),
      newChatBtn: document.getElementById('newChatBtn'),
      modeSelect: document.getElementById('modeSelect'),
      chatTitle: document.getElementById('chatTitle'),
      status: document.getElementById('status'),
      menuBtn: document.getElementById('menuBtn'),
      memoryBtn: document.getElementById('memoryBtn'),
      memoryModal: document.getElementById('memoryModal'),
      memoryText: document.getElementById('memoryText'),
      closeMemoryBtn: document.getElementById('closeMemoryBtn'),
      clearMemoryBtn: document.getElementById('clearMemoryBtn'),
    };

    function escapeHtml(value) {
      return value.replace(/[&<>"']/g, ch => ({
        '&': '&amp;', '<': '&lt;', '>': '&gt;', '"': '&quot;', "'": '&#39;'
      })[ch]);
    }

    function renderMarkdown(text) {
      const parts = text.split(/```([\s\S]*?)```/g);
      return parts.map((part, index) => {
        if (index % 2 === 1) {
          const clean = part.replace(/^\w+\n/, '');
          return `<pre><button class="code-copy" type="button">Copiar</button><code>${escapeHtml(clean)}</code></pre>`;
        }
        const html = escapeHtml(part)
          .replace(/\*\*(.*?)\*\*/g, '<strong>$1</strong>')
          .replace(/`([^`]+)`/g, '<code>$1</code>')
          .split(/\n{2,}/)
          .map(p => p.trim() ? `<p>${p.replace(/\n/g, '<br>')}</p>` : '')
          .join('');
        return html;
      }).join('');
    }

    function formatBytes(size) {
      let value = Number(size || 0);
      for (const unit of ['B', 'KB', 'MB', 'GB']) {
        if (value < 1024 || unit === 'GB') {
          return unit === 'B' ? `${value} ${unit}` : `${value.toFixed(1)} ${unit}`;
        }
        value = value / 1024;
      }
      return `${size || 0} B`;
    }

    function fileKind(file) {
      const name = file.name.toLowerCase();
      if (/\.(jpg|jpeg|png|webp|bmp|gif)$/.test(name)) return 'image';
      if (/\.(mp4|webm|mov|avi|mkv|m4v)$/.test(name)) return 'video';
      if (/\.(mp3|wav|m4a|aac|ogg|oga|flac|wma)$/.test(name)) return 'audio';
      if (/\.(zip|tar|gz|tgz|bz2|tbz2|xz|txz|7z)$/.test(name)) return 'archive';
      if (/\.(pdf|txt|md|csv|json|jsonl|xml|html|htm|css|js|ts|tsx|jsx|py|java|c|cpp|h|hpp|cs|go|rs|php|rb|swift|kt|kts|sql|yaml|yml|toml|ini|cfg|conf|env|log|bat|ps1|sh|doc|docx|rtf|odt|ppt|pptx|xls|xlsx|ods)$/.test(name)) return 'document';
      if (['dockerfile', 'makefile', 'license', 'readme'].includes(name)) return 'document';
      return 'unknown';
    }

    function statusLabel(status) {
      return ({
        pending: 'pendiente',
        processed: 'procesado',
        error: 'error',
        expired: 'caducado',
      })[status] || status || '';
    }

    function renderProvider(message = {}) {
      return '';
    }

    function sourceSnippet(source) {
      const snippet = source.snippet || '';
      return snippet ? `<span>${escapeHtml(snippet).slice(0, 500)}</span>` : '';
    }

    function updateDisplayedAttachment(attachment) {
      const chat = currentChat();
      if (!chat || !chat.messages || !chat.messages.length) return;
      const lastUser = [...chat.messages].reverse().find(message => message.role === 'user');
      if (!lastUser) return;
      const list = lastUser.attachments || [];
      const index = list.findIndex(item => item.id === attachment.id || item.filename === attachment.filename);
      if (index >= 0) list[index] = attachment;
      else list.push(attachment);
      lastUser.attachments = list;
    }

    function renderAttachments(attachments = []) {
      if (!attachments.length) return '';
      return `<div class="attachment-list">${attachments.map(item => `
        <div class="attachment-item">
          <strong>${escapeHtml(item.filename || 'archivo')}</strong>
          <span>${escapeHtml(item.kind || 'archivo')} - ${formatBytes(item.size)}${item.status ? ' - ' + escapeHtml(statusLabel(item.status)) : ''}${item.expired ? ' - caducado' : ''}</span>
        </div>
      `).join('')}</div>`;
    }

    function renderSources(sources = []) {
      return '';
    }

    function assistantHtml(text, sources = [], meta = {}) {
      return renderMarkdown(text || '');
    }

    function currentChat() {
      return state.chats.find(chat => chat.id === state.activeChatId) || null;
    }

    function setBusy(value) {
      state.busy = value;
      els.sendBtn.disabled = value;
      els.prompt.disabled = value;
      els.attachBtn.disabled = value;
      els.fileInput.disabled = value;
    }

    function setStatus(text) {
      els.status.textContent = text || '';
    }

    function scrollBottom() {
      els.messages.scrollTop = els.messages.scrollHeight;
    }

    function addPendingFiles(files) {
      [...files].forEach(file => {
        const item = {
          id: `${Date.now()}-${Math.random().toString(16).slice(2)}`,
          file,
          previewUrl: file.type.startsWith('image/') ? URL.createObjectURL(file) : '',
        };
        state.pendingFiles.push(item);
      });
      renderPendingFiles();
    }

    function removePendingFile(id) {
      const index = state.pendingFiles.findIndex(item => item.id === id);
      if (index < 0) return;
      const [item] = state.pendingFiles.splice(index, 1);
      if (item.previewUrl) URL.revokeObjectURL(item.previewUrl);
      renderPendingFiles();
    }

    function clearPendingFiles() {
      state.pendingFiles.forEach(item => {
        if (item.previewUrl) URL.revokeObjectURL(item.previewUrl);
      });
      state.pendingFiles = [];
      els.fileInput.value = '';
      renderPendingFiles();
    }

    function pendingAttachmentsForDisplay() {
      return state.pendingFiles.map(item => ({
        filename: item.file.name,
        size: item.file.size,
        kind: fileKind(item.file),
        status: 'pending',
      }));
    }

    function renderPendingFiles() {
      els.pendingFiles.innerHTML = '';
      els.pendingFiles.classList.toggle('open', state.pendingFiles.length > 0);
      state.pendingFiles.forEach(item => {
        const wrap = document.createElement('div');
        wrap.className = 'pending-file';
        const thumb = item.previewUrl
          ? `<img class="pending-thumb" src="${item.previewUrl}" alt="">`
          : `<div class="pending-thumb">${fileKind(item.file).slice(0, 3).toUpperCase()}</div>`;
        wrap.innerHTML = `
          ${thumb}
          <div><strong>${escapeHtml(item.file.name)}</strong><span>${formatBytes(item.file.size)}</span></div>
          <button class="remove-file" type="button" aria-label="Quitar archivo">x</button>
        `;
        wrap.querySelector('.remove-file').addEventListener('click', () => removePendingFile(item.id));
        els.pendingFiles.appendChild(wrap);
      });
    }

    function renderChatList() {
      els.chatList.innerHTML = '';
      state.chats.forEach(chat => {
        const button = document.createElement('button');
        button.className = 'chat-item' + (chat.id === state.activeChatId ? ' active' : '');
        button.type = 'button';
        button.textContent = chat.title || 'Nuevo chat';
        button.addEventListener('click', () => {
          loadChat(chat.id);
          els.sidebar.classList.remove('open');
        });
        els.chatList.appendChild(button);
      });
    }

    function renderMessages(chat) {
      els.messages.innerHTML = '';
      if (!chat || !chat.messages || !chat.messages.length) {
        els.messages.innerHTML = `<div class="empty"><div><h1>Â¿QuÃ© hacemos hoy?</h1><p>Elige modo, escribe tu mensaje y la NEXO responderÃ¡ usando memoria local.</p></div></div>`;
        els.chatTitle.textContent = 'Nuevo chat';
        return;
      }
      els.chatTitle.textContent = chat.title || 'Nuevo chat';
      chat.messages.forEach(addMessageElement);
      scrollBottom();
    }

    function addMessageElement(message) {
      const wrap = document.createElement('article');
      wrap.className = `message ${message.role}`;
      wrap.innerHTML = `
        <div class="avatar">${message.role === 'user' ? 'TU' : 'IA'}</div>
        <div class="bubble">${renderMarkdown(message.content || '')}${renderAttachments(message.attachments || [])}</div>
      `;
      els.messages.appendChild(wrap);
      return wrap.querySelector('.bubble');
    }

    async function api(url, options = {}) {
      const response = await fetch(url, {
        credentials: 'same-origin',
        headers: { 'Content-Type': 'application/json', ...(options.headers || {}) },
        ...options,
      });
      if (!response.ok) {
        let message = `Error ${response.status}`;
        try { message = (await response.json()).error || message; } catch (_) {}
        throw new Error(message);
      }
      return response.json();
    }

    async function loadChats() {
      const data = await api('/api/chats');
      state.chats = data.chats || [];
      renderChatList();
      if (state.chats.length && !state.activeChatId) {
        await loadChat(state.chats[0].id);
      }
    }

    async function createChat() {
      const chat = await api('/api/chats', { method: 'POST', body: '{}' });
      state.chats.unshift(chat);
      state.activeChatId = chat.id;
      renderChatList();
      renderMessages(chat);
      return chat;
    }

    async function loadChat(id) {
      const chat = await api(`/api/chats/${id}`);
      state.activeChatId = chat.id;
      const index = state.chats.findIndex(item => item.id === chat.id);
      if (index >= 0) state.chats[index] = chat;
      renderChatList();
      renderMessages(chat);
    }

    async function sendMessage() {
      const message = els.prompt.value.trim();
      const filesToSend = state.pendingFiles.map(item => item.file);
      if ((!message && !filesToSend.length) || state.busy) return;
      let chat = currentChat();
      if (!chat) chat = await createChat();

      els.prompt.value = '';
      els.prompt.style.height = 'auto';
      setBusy(true);
      setStatus('Preparando...');

      const displayMessage = message || 'Analiza los archivos adjuntos.';
      chat.messages.push({ role: 'user', content: displayMessage, attachments: pendingAttachmentsForDisplay() });
      renderMessages(chat);
      const assistantBubble = addMessageElement({ role: 'assistant', content: '' });
      let assistantText = '';
      let assistantSources = [];
      let assistantMeta = {};
      scrollBottom();
      clearPendingFiles();

      try {
        let fetchOptions;
        if (filesToSend.length) {
          const formData = new FormData();
          formData.append('chat_id', chat.id);
          formData.append('mode', els.modeSelect.value);
          formData.append('message', message);
          filesToSend.forEach(file => formData.append('files[]', file, file.name));
          fetchOptions = { method: 'POST', credentials: 'same-origin', body: formData };
        } else {
          fetchOptions = {
            method: 'POST',
            credentials: 'same-origin',
            headers: { 'Content-Type': 'application/json' },
            body: JSON.stringify({
              chat_id: chat.id,
              mode: els.modeSelect.value,
              message,
            }),
          };
        }
        const response = await fetch('/api/chat/stream', {
          ...fetchOptions,
        });
        if (!response.ok || !response.body) {
          let errorText = `Error ${response.status}`;
          try { errorText = (await response.json()).error || errorText; } catch (_) {}
          throw new Error(errorText);
        }

        const reader = response.body.getReader();
        const decoder = new TextDecoder();
        let buffer = '';
        while (true) {
          const { value, done } = await reader.read();
          if (done) break;
          buffer += decoder.decode(value, { stream: true });
          const lines = buffer.split('\n');
          buffer = lines.pop() || '';
          for (const line of lines) {
            if (!line.trim()) continue;
            const event = JSON.parse(line);
            if (event.type === 'status') setStatus(event.message);
            if (event.type === 'token') {
              assistantText += event.token;
              assistantBubble.innerHTML = assistantHtml(assistantText, assistantSources, assistantMeta);
              scrollBottom();
            }
            if (event.type === 'sources') {
              assistantSources = event.sources || [];
              assistantBubble.innerHTML = assistantHtml(assistantText, assistantSources, assistantMeta);
              scrollBottom();
            }
            if (event.type === 'provider') {
              assistantMeta.provider = event.provider || '';
              assistantBubble.innerHTML = assistantHtml(assistantText, assistantSources, assistantMeta);
              scrollBottom();
            }
            if (event.type === 'fallback') {
              assistantMeta.provider = event.provider || assistantMeta.provider || 'ollama';
              assistantMeta.fallback = event.reason || 'fallback local';
              assistantBubble.innerHTML = assistantHtml(assistantText, assistantSources, assistantMeta);
              scrollBottom();
            }
            if (event.type === 'attachment_status' && event.attachment) {
              updateDisplayedAttachment(event.attachment);
              scrollBottom();
            }
            if (event.type === 'error') throw new Error(event.message);
            if (event.type === 'done') {
              const updated = event.chat;
              const index = state.chats.findIndex(item => item.id === updated.id);
              if (index >= 0) state.chats[index] = updated;
              else state.chats.unshift(updated);
              state.activeChatId = updated.id;
              renderChatList();
              renderMessages(updated);
            }
          }
        }
      } catch (err) {
        assistantBubble.innerHTML = renderMarkdown(`Error: ${err.message}`);
      } finally {
        setBusy(false);
        setStatus('');
        els.prompt.focus();
      }
    }

    async function openMemory() {
      const data = await api('/api/memory');
      els.memoryText.textContent = data.summary || 'TodavÃ­a no hay memoria guardada.';
      els.memoryModal.classList.add('open');
    }

    async function clearMemory() {
      await api('/api/memory/clear', { method: 'POST', body: '{}' });
      els.memoryText.textContent = 'Memoria borrada.';
    }

    els.sendBtn.addEventListener('click', sendMessage);
    els.attachBtn.addEventListener('click', () => els.fileInput.click());
    els.fileInput.addEventListener('change', event => addPendingFiles(event.target.files || []));
    els.newChatBtn.addEventListener('click', async () => { await createChat(); els.prompt.focus(); });
    els.menuBtn.addEventListener('click', () => els.sidebar.classList.toggle('open'));
    els.memoryBtn.addEventListener('click', openMemory);
    els.closeMemoryBtn.addEventListener('click', () => els.memoryModal.classList.remove('open'));
    els.clearMemoryBtn.addEventListener('click', clearMemory);
    els.memoryModal.addEventListener('click', event => {
      if (event.target === els.memoryModal) els.memoryModal.classList.remove('open');
    });
    els.messages.addEventListener('click', async event => {
      if (!event.target.classList.contains('code-copy')) return;
      const code = event.target.parentElement.querySelector('code').textContent;
      await navigator.clipboard.writeText(code);
      event.target.textContent = 'Copiado';
      setTimeout(() => event.target.textContent = 'Copiar', 1200);
    });
    els.prompt.addEventListener('keydown', event => {
      if (event.key === 'Enter' && !event.shiftKey) {
        event.preventDefault();
        sendMessage();
      }
    });
    els.prompt.addEventListener('input', () => {
      els.prompt.style.height = 'auto';
      els.prompt.style.height = `${Math.min(180, els.prompt.scrollHeight)}px`;
    });
    els.prompt.addEventListener('paste', event => {
      const files = [...(event.clipboardData?.files || [])];
      if (files.length) addPendingFiles(files);
    });
    ['dragenter', 'dragover'].forEach(type => {
      els.composer.addEventListener(type, event => {
        event.preventDefault();
        els.composer.classList.add('dragging');
      });
    });
    ['dragleave', 'drop'].forEach(type => {
      els.composer.addEventListener(type, event => {
        event.preventDefault();
        els.composer.classList.remove('dragging');
      });
    });
    els.composer.addEventListener('drop', event => {
      const files = event.dataTransfer?.files || [];
      if (files.length) addPendingFiles(files);
    });

    loadChats().catch(err => setStatus(err.message));
  </script>
</body>
</html>
"""


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def utc_after(days: int) -> str:
    return (datetime.now(timezone.utc) + timedelta(days=days)).isoformat(timespec="seconds")


def parse_utc(value: Any) -> Optional[datetime]:
    if not value:
        return None
    try:
        text = str(value).replace("Z", "+00:00")
        parsed = datetime.fromisoformat(text)
        if parsed.tzinfo is None:
            parsed = parsed.replace(tzinfo=timezone.utc)
        return parsed.astimezone(timezone.utc)
    except ValueError:
        return None


def truncate_text(text: str, limit: int) -> str:
    text = text.strip()
    if len(text) <= limit:
        return text
    return text[:limit].rstrip() + "\n\n[Contenido truncado para no saturar el contexto.]"


def safe_segment(value: str) -> str:
    safe = re.sub(r"[^A-Za-z0-9_.-]+", "_", value).strip("._")
    return safe or "unknown"


def classify_upload(filename: str) -> Optional[str]:
    suffix = Path(filename).suffix.lower()
    basename = Path(filename).name.lower()
    if basename in {"dockerfile", "makefile", "license", "readme"}:
        return "document"
    if suffix in IMAGE_EXTENSIONS:
        return "image"
    if suffix in DOCUMENT_EXTENSIONS:
        return "document"
    if suffix in AUDIO_EXTENSIONS:
        return "audio"
    if suffix in VIDEO_EXTENSIONS:
        return "video"
    if suffix in ARCHIVE_EXTENSIONS:
        return "archive"
    return "unknown"


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def format_bytes(size: int) -> str:
    value = float(size)
    for unit in ("B", "KB", "MB", "GB"):
        if value < 1024 or unit == "GB":
            if unit == "B":
                return f"{int(value)} {unit}"
            return f"{value:.1f} {unit}"
        value /= 1024
    return f"{size} B"


def attachment_for_chat(attachment: Dict[str, Any]) -> Dict[str, Any]:
    return {
        key: value
        for key, value in attachment.items()
        if key not in {"saved_path", "derived_paths", "openai_file_id", "openai_upload_error"}
    }


def remove_stored_attachment_files(attachments: List[Dict[str, Any]]) -> None:
    for attachment in attachments:
        for key in ("saved_path",):
            path = Path(str(attachment.get(key, "")))
            if path.exists() and path.is_file():
                try:
                    path.unlink(missing_ok=True)
                except OSError:
                    pass


def build_attachment_context(attachments: List[Dict[str, Any]]) -> str:
    if not attachments:
        return "Sin archivos adjuntos."

    lines = []
    for index, attachment in enumerate(attachments, 1):
        summary = str(attachment.get("summary", "")).strip() or "Sin resumen disponible."
        lines.append(
            f"Adjunto {index}: {attachment.get('filename', 'archivo')} "
            f"({attachment.get('kind', 'desconocido')}, {format_bytes(int(attachment.get('size') or 0))}, "
            f"caduca {attachment.get('expires_at', 'sin fecha')})\n{summary}"
        )
    return truncate_text("\n\n".join(lines), MAX_ATTACHMENT_CONTEXT_CHARS)


def build_web_context(sources: List[Dict[str, str]], search_error: str = "") -> str:
    if sources:
        blocks = []
        for index, source in enumerate(sources, 1):
            blocks.append(
                f"Fuente {index}: {source.get('title') or 'Sin titulo'}\n"
                f"URL: {source.get('url')}\n"
                f"Extracto: {source.get('snippet') or ''}"
            )
        return truncate_text("\n\n".join(blocks), MAX_WEB_CONTEXT_CHARS)
    if search_error:
        return f"No se pudo obtener contexto de internet: {search_error}"
    return "No se encontraron fuentes utiles en internet."


def store_uploaded_files(uploaded_files: List[Any], user_id: str, chat_id: str) -> tuple[List[Dict[str, Any]], Optional[str]]:
    files = [item for item in uploaded_files if item and getattr(item, "filename", "")]
    if not files:
        return [], None
    if len(files) > MAX_ATTACHMENTS_PER_MESSAGE:
        return [], f"Maximo {MAX_ATTACHMENTS_PER_MESSAGE} archivos por mensaje."

    root = UPLOAD_DIR / safe_segment(user_id) / safe_segment(chat_id)
    root.mkdir(parents=True, exist_ok=True)

    attachments: List[Dict[str, Any]] = []
    total_size = 0
    try:
        for storage in files:
            original_name = str(storage.filename or "archivo")
            safe_name = secure_filename(original_name) or "archivo"
            kind = classify_upload(original_name)
            if not kind:
                remove_stored_attachment_files(attachments)
                return [], f"Tipo de archivo no permitido: {original_name}"

            upload_id = secrets.token_urlsafe(10)
            final_name = f"{upload_id}_{safe_name}"
            destination = root / final_name
            storage.save(destination)

            size = destination.stat().st_size
            total_size += size
            if size > UPLOAD_LIMITS[kind]:
                destination.unlink(missing_ok=True)
                remove_stored_attachment_files(attachments)
                return [], f"{original_name} supera el limite para {kind}: {format_bytes(UPLOAD_LIMITS[kind])}."
            if total_size > MAX_TOTAL_UPLOAD_BYTES:
                destination.unlink(missing_ok=True)
                remove_stored_attachment_files(attachments)
                return [], f"Los adjuntos superan el limite total de {format_bytes(MAX_TOTAL_UPLOAD_BYTES)}."

            mime = storage.mimetype or mimetypes.guess_type(safe_name)[0] or "application/octet-stream"
            attachments.append(
                {
                    "id": upload_id,
                    "filename": original_name,
                    "mime": mime,
                    "kind": kind,
                    "size": size,
                    "sha256": file_sha256(destination),
                    "summary": "Pendiente de procesar.",
                    "status": "pending",
                    "expires_at": utc_after(UPLOAD_TTL_DAYS),
                    "expired": False,
                    "saved_path": str(destination),
                }
            )
    except OSError as exc:
        remove_stored_attachment_files(attachments)
        return [], f"No se pudo guardar el archivo: {exc}"

    return attachments, None


def extract_pdf_text(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        return "No se pudo leer el PDF: falta instalar pypdf."

    try:
        reader = PdfReader(str(path))
        parts = []
        for page in reader.pages:
            parts.append(page.extract_text() or "")
            if sum(len(part) for part in parts) >= MAX_DOCUMENT_CHARS:
                break
        text = "\n\n".join(part.strip() for part in parts if part.strip())
        return truncate_text(text or "El PDF no contiene texto extraible.", MAX_DOCUMENT_CHARS)
    except Exception as exc:
        return f"No se pudo leer el PDF: {exc}"


def extract_document_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_text(path)
    if suffix == ".docx":
        return extract_docx_text(path)
    if suffix == ".xlsx":
        return extract_xlsx_text(path)
    if suffix == ".pptx":
        return extract_pptx_text(path)
    if suffix == ".rtf":
        return extract_rtf_text(path)
    if suffix in {".doc", ".xls", ".ppt", ".odt", ".ods"}:
        return (
            "Este formato se ha guardado correctamente, pero no tiene extractor local fiable. "
            "Si OpenAI esta configurado, se enviara como archivo para que el modelo lo lea."
        )
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        return truncate_text(text or "El archivo esta vacio.", MAX_DOCUMENT_CHARS)
    except OSError as exc:
        return f"No se pudo leer el documento: {exc}"


def extract_docx_text(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        return "No se pudo leer DOCX: falta instalar python-docx."

    try:
        document = Document(str(path))
        parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables[:20]:
            for row in table.rows[:100]:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        text = "\n".join(parts)
        return truncate_text(text or "El DOCX no contiene texto extraible.", MAX_DOCUMENT_CHARS)
    except Exception as exc:
        return f"No se pudo leer DOCX: {exc}"


def extract_xlsx_text(path: Path) -> str:
    try:
        import openpyxl
    except ImportError:
        return "No se pudo leer XLSX: falta instalar openpyxl."

    try:
        workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet in workbook.worksheets[:10]:
            parts.append(f"Hoja: {sheet.title}")
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(value.strip() for value in values):
                    parts.append(" | ".join(values))
                    row_count += 1
                if row_count >= 1000:
                    parts.append("[Hoja truncada a 1000 filas]")
                    break
        workbook.close()
        text = "\n".join(parts)
        return truncate_text(text or "El XLSX no contiene datos extraibles.", MAX_DOCUMENT_CHARS)
    except Exception as exc:
        return f"No se pudo leer XLSX: {exc}"


def extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "No se pudo leer PPTX: falta instalar python-pptx."

    try:
        presentation = Presentation(str(path))
        parts = []
        for index, slide in enumerate(presentation.slides, 1):
            slide_parts = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    slide_parts.append(text.strip())
            if slide_parts:
                parts.append(f"Diapositiva {index}:\n" + "\n".join(slide_parts))
        text = "\n\n".join(parts)
        return truncate_text(text or "El PPTX no contiene texto extraible.", MAX_DOCUMENT_CHARS)
    except Exception as exc:
        return f"No se pudo leer PPTX: {exc}"


def extract_rtf_text(path: Path) -> str:
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
    except OSError as exc:
        return f"No se pudo leer RTF: {exc}"
    text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", raw)
    text = re.sub(r"\\[a-zA-Z]+\d* ?", " ", text)
    text = re.sub(r"[{}]", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return truncate_text(text or "El RTF no contiene texto extraible.", MAX_DOCUMENT_CHARS)


def describe_unknown_file(path: Path, attachment: Dict[str, Any]) -> str:
    return (
        "Archivo guardado como binario o formato no reconocido.\n"
        f"Nombre: {attachment.get('filename', path.name)}\n"
        f"MIME: {attachment.get('mime', 'application/octet-stream')}\n"
        f"Tamano: {format_bytes(int(attachment.get('size') or 0))}\n"
        f"SHA256: {attachment.get('sha256', '')}\n"
        "No se ha ejecutado ni interpretado por seguridad."
    )


def image_path_to_base64(path: Path) -> Optional[str]:
    try:
        from PIL import Image
    except ImportError:
        try:
            return base64.b64encode(path.read_bytes()).decode("ascii")
        except OSError:
            return None

    try:
        with Image.open(path) as image:
            try:
                image.seek(0)
            except EOFError:
                pass
            image.thumbnail((1024, 1024))
            if image.mode in {"RGBA", "LA"}:
                background = Image.new("RGB", image.size, "white")
                alpha = image.getchannel("A")
                background.paste(image.convert("RGB"), mask=alpha)
                image = background
            else:
                image = image.convert("RGB")
            buffer = io.BytesIO()
            image.save(buffer, format="JPEG", quality=85, optimize=True)
            return base64.b64encode(buffer.getvalue()).decode("ascii")
    except Exception:
        try:
            return base64.b64encode(path.read_bytes()).decode("ascii")
        except OSError:
            return None


def summarize_images_with_vision(image_paths: List[Path], prompt: str) -> str:
    if not image_paths:
        return "No se encontraron imagenes para analizar."

    ready_error = ensure_ai_ready(["vision"])
    if ready_error:
        return f"No se pudo analizar visualmente: {ready_error}"

    summaries = []
    for start in range(0, len(image_paths), VISION_IMAGE_BATCH):
        chunk = image_paths[start:start + VISION_IMAGE_BATCH]
        images = [encoded for encoded in (image_path_to_base64(path) for path in chunk) if encoded]
        if not images:
            summaries.append(f"Grupo {start // VISION_IMAGE_BATCH + 1}: no se pudieron cargar las imagenes.")
            continue
        try:
            result = ollama_chat(
                "vision",
                [
                    {"role": "system", "content": "Eres un analista visual preciso. Responde en espanol."},
                    {"role": "user", "content": prompt, "images": images},
                ],
            )
            summaries.append(f"Grupo {start // VISION_IMAGE_BATCH + 1}:\n{result.strip()}")
        except Exception as exc:
            summaries.append(f"Grupo {start // VISION_IMAGE_BATCH + 1}: error al analizar vision: {exc}")

    return truncate_text("\n\n".join(summaries), MAX_DOCUMENT_CHARS)


def extract_video_frames(path: Path, target_count: int = VIDEO_FRAME_COUNT) -> List[Path]:
    try:
        import cv2
    except ImportError:
        return []

    frames_dir = path.parent / f"{path.stem}_frames"
    frames_dir.mkdir(parents=True, exist_ok=True)
    capture = cv2.VideoCapture(str(path))
    if not capture.isOpened():
        return []

    frame_paths: List[Path] = []
    try:
        frame_count = int(capture.get(cv2.CAP_PROP_FRAME_COUNT) or 0)
        if frame_count > 0:
            indices = sorted({min(frame_count - 1, int((i + 0.5) * frame_count / target_count)) for i in range(target_count)})
            for order, frame_index in enumerate(indices, 1):
                capture.set(cv2.CAP_PROP_POS_FRAMES, frame_index)
                ok, frame = capture.read()
                if not ok:
                    continue
                frame_paths.append(save_video_frame(frame, frames_dir, order))
        else:
            seen = 0
            while len(frame_paths) < target_count:
                ok, frame = capture.read()
                if not ok:
                    break
                if seen % 30 == 0:
                    frame_paths.append(save_video_frame(frame, frames_dir, len(frame_paths) + 1))
                seen += 1
    finally:
        capture.release()

    return [path for path in frame_paths if path.exists()]


def save_video_frame(frame: Any, frames_dir: Path, order: int) -> Path:
    import cv2

    height, width = frame.shape[:2]
    largest = max(width, height)
    if largest > 1024:
        scale = 1024 / largest
        frame = cv2.resize(frame, (int(width * scale), int(height * scale)), interpolation=cv2.INTER_AREA)
    destination = frames_dir / f"frame_{order:02d}.jpg"
    cv2.imwrite(str(destination), frame, [int(cv2.IMWRITE_JPEG_QUALITY), 85])
    return destination


def image_path_to_data_url(path: Path) -> Optional[str]:
    encoded = image_path_to_base64(path)
    if not encoded:
        return None
    return f"data:image/jpeg;base64,{encoded}"


def extract_audio_transcript(path: Path, settings: Dict[str, Any]) -> str:
    if not should_use_openai(settings):
        return "Audio guardado, pero no transcrito porque OpenAI no esta configurado."
    try:
        with path.open("rb") as handle:
            response = requests.post(
                f"{OPENAI_API_BASE}/audio/transcriptions",
                headers={"Authorization": f"Bearer {settings['openai_api_key']}"},
                data={"model": settings["openai_transcribe_model"]},
                files={"file": (path.name, handle, mimetypes.guess_type(path.name)[0] or "application/octet-stream")},
                timeout=300,
            )
        response.raise_for_status()
        data = response.json()
        return truncate_text(str(data.get("text") or "").strip() or "La transcripcion no devolvio texto.", MAX_DOCUMENT_CHARS)
    except Exception as exc:
        return f"No se pudo transcribir el audio con OpenAI: {exc}"


def extract_video_audio(path: Path) -> Optional[Path]:
    try:
        try:
            from moviepy import VideoFileClip
        except ImportError:
            from moviepy.editor import VideoFileClip
    except ImportError:
        return None

    destination = path.parent / f"{path.stem}_audio.wav"
    try:
        clip = VideoFileClip(str(path))
        try:
            if not clip.audio:
                return None
            clip.audio.write_audiofile(str(destination), logger=None)
            return destination if destination.exists() else None
        finally:
            clip.close()
    except Exception:
        return None


def archive_member_safe_name(name: str, index: int) -> str:
    base = secure_filename(Path(name).name) or f"archivo_{index}"
    return f"{index:02d}_{base}"


def copy_limited_stream(source: Any, destination: Path, limit: int) -> int:
    total = 0
    with destination.open("wb") as handle:
        while True:
            chunk = source.read(64 * 1024)
            if not chunk:
                break
            total += len(chunk)
            if total > limit:
                raise ValueError("archivo interno demasiado grande")
            handle.write(chunk)
    return total


def extract_zip_members(path: Path, derived_dir: Path) -> List[Path]:
    extracted: List[Path] = []
    total = 0
    with zipfile.ZipFile(path) as archive:
        for info in archive.infolist():
            if info.is_dir():
                continue
            if len(extracted) >= MAX_ARCHIVE_MEMBERS:
                break
            total += int(info.file_size or 0)
            if total > MAX_ARCHIVE_EXTRACT_BYTES:
                break
            destination = derived_dir / archive_member_safe_name(info.filename, len(extracted) + 1)
            with archive.open(info) as source:
                copy_limited_stream(source, destination, MAX_DOCUMENT_BYTES)
            extracted.append(destination)
    return extracted


def extract_tar_members(path: Path, derived_dir: Path) -> List[Path]:
    extracted: List[Path] = []
    total = 0
    with tarfile.open(path) as archive:
        for member in archive.getmembers():
            if not member.isfile():
                continue
            if len(extracted) >= MAX_ARCHIVE_MEMBERS:
                break
            total += int(member.size or 0)
            if total > MAX_ARCHIVE_EXTRACT_BYTES:
                break
            source = archive.extractfile(member)
            if source is None:
                continue
            destination = derived_dir / archive_member_safe_name(member.name, len(extracted) + 1)
            copy_limited_stream(source, destination, MAX_DOCUMENT_BYTES)
            extracted.append(destination)
    return extracted


def extract_7z_members(path: Path, derived_dir: Path) -> List[Path]:
    try:
        import py7zr
    except ImportError:
        return []

    raw_dir = derived_dir / "raw_7z"
    raw_dir.mkdir(parents=True, exist_ok=True)
    with py7zr.SevenZipFile(path, mode="r") as archive:
        archive.extractall(path=raw_dir)

    extracted: List[Path] = []
    total = 0
    for item in raw_dir.rglob("*"):
        if not item.is_file():
            continue
        try:
            item.resolve().relative_to(raw_dir.resolve())
        except ValueError:
            continue
        if len(extracted) >= MAX_ARCHIVE_MEMBERS:
            break
        size = item.stat().st_size
        total += size
        if total > MAX_ARCHIVE_EXTRACT_BYTES:
            break
        destination = derived_dir / archive_member_safe_name(item.name, len(extracted) + 1)
        shutil.copyfile(item, destination)
        extracted.append(destination)
    return extracted


def extract_archive_files(path: Path) -> List[Path]:
    derived_dir = path.parent / f"{path.stem}_archive"
    derived_dir.mkdir(parents=True, exist_ok=True)
    suffixes = [suffix.lower() for suffix in path.suffixes]
    try:
        if path.suffix.lower() == ".zip":
            return extract_zip_members(path, derived_dir)
        if path.suffix.lower() == ".7z":
            return extract_7z_members(path, derived_dir)
        if any(suffix in suffixes for suffix in {".tar", ".gz", ".tgz", ".bz2", ".xz"}):
            return extract_tar_members(path, derived_dir)
    except Exception:
        return []
    return []


def summarize_archive(path: Path, attachment: Dict[str, Any]) -> str:
    extracted = extract_archive_files(path)
    attachment["derived_paths"] = [str(item) for item in extracted]
    if not extracted:
        return "No se pudieron extraer archivos legibles del comprimido o el comprimido esta vacio."

    parts = [f"Comprimido extraido parcialmente: {len(extracted)} archivo(s) analizados."]
    for item in extracted[:MAX_ARCHIVE_MEMBERS]:
        kind = classify_upload(item.name)
        if kind == "document":
            parts.append(f"\nArchivo interno: {item.name}\n{extract_document_text(item)}")
        else:
            parts.append(
                f"\nArchivo interno: {item.name}\n"
                f"Tipo detectado: {kind}. Tamano: {format_bytes(item.stat().st_size)}."
            )
    return truncate_text("\n".join(parts), MAX_DOCUMENT_CHARS)


def should_use_openai(settings: Optional[Dict[str, Any]] = None) -> bool:
    data = settings or load_ai_settings()
    return data.get("ai_provider") == "openai" and bool(data.get("openai_api_key"))


def upload_attachment_to_openai(attachment: Dict[str, Any], settings: Dict[str, Any]) -> None:
    if not should_use_openai(settings):
        return
    if attachment.get("kind") != "document":
        return
    path = Path(str(attachment.get("saved_path", "")))
    if not path.exists() or int(attachment.get("size") or 0) > MAX_DOCUMENT_BYTES:
        return
    try:
        with path.open("rb") as handle:
            response = requests.post(
                f"{OPENAI_API_BASE}/files",
                headers={"Authorization": f"Bearer {settings['openai_api_key']}"},
                data={"purpose": "user_data"},
                files={"file": (attachment.get("filename") or path.name, handle, attachment.get("mime") or "application/octet-stream")},
                timeout=180,
            )
        response.raise_for_status()
        attachment["openai_file_id"] = response.json().get("id", "")
    except Exception as exc:
        attachment["openai_upload_error"] = str(exc)


def process_stored_attachment(attachment: Dict[str, Any], settings: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    settings = settings or load_ai_settings()
    path = Path(str(attachment.get("saved_path", "")))
    if not path.exists():
        attachment["summary"] = "El archivo ya no existe en disco."
        attachment["expired"] = True
        attachment["status"] = "expired"
        return attachment

    kind = attachment.get("kind")
    if kind == "document":
        text = extract_document_text(path)
        attachment["summary"] = f"Texto extraido del documento:\n{text}"
        upload_attachment_to_openai(attachment, settings)
    elif kind == "image":
        if should_use_openai(settings):
            attachment["summary"] = "Imagen lista para analisis multimodal con OpenAI."
        else:
            attachment["summary"] = summarize_images_with_vision(
                [path],
                "Analiza esta imagen para ayudar a responder al usuario. Describe contenido, texto visible, detalles importantes y posibles dudas.",
            )
    elif kind == "audio":
        transcript = extract_audio_transcript(path, settings)
        attachment["summary"] = f"Audio procesado.\nTranscripcion:\n{transcript}"
    elif kind == "video":
        frame_paths = extract_video_frames(path)
        attachment["derived_paths"] = [str(item) for item in frame_paths]
        audio_path = extract_video_audio(path)
        if audio_path:
            attachment.setdefault("derived_paths", []).append(str(audio_path))
        transcript = extract_audio_transcript(audio_path, settings) if audio_path else "El video no tenia audio extraible."
        if should_use_openai(settings):
            visual_summary = "Fotogramas listos para analisis multimodal con OpenAI."
        else:
            visual_summary = summarize_images_with_vision(
                frame_paths,
                "Estos son fotogramas distribuidos de un video. Resume la secuencia, acciones, objetos, texto visible, cambios entre escenas y detalles utiles.",
            )
        attachment["summary"] = (
            f"Video procesado mediante {len(frame_paths)} fotogramas clave.\n"
            f"{visual_summary}\n\nTranscripcion/audio:\n{transcript}"
        )
    elif kind == "archive":
        attachment["summary"] = summarize_archive(path, attachment)
    else:
        attachment["summary"] = describe_unknown_file(path, attachment)
    attachment["status"] = "processed"
    return attachment


def is_global_ip(value: str) -> bool:
    try:
        return ipaddress.ip_address(value).is_global
    except ValueError:
        return False


def is_safe_public_url(url: str) -> bool:
    parsed = urlparse(url)
    if parsed.scheme not in {"http", "https"}:
        return False
    host = parsed.hostname
    if not host:
        return False
    lowered = host.lower().strip(".")
    if lowered in {"localhost"} or lowered.endswith(".local"):
        return False
    if re.fullmatch(r"\d+(?:\.\d+){3}", lowered) or ":" in lowered:
        return is_global_ip(lowered)
    try:
        for family, _, _, _, sockaddr in socket.getaddrinfo(lowered, None):
            address = sockaddr[0]
            if not ipaddress.ip_address(address).is_global:
                return False
    except Exception:
        return False
    return True


def fetch_url_text(url: str) -> Optional[Dict[str, str]]:
    if not is_safe_public_url(url):
        return None
    try:
        with requests.get(
            url,
            timeout=10,
            stream=True,
            headers={"User-Agent": "IA-Combinada/1.0 (+local research assistant)"},
        ) as response:
            response.raise_for_status()
            content_type = response.headers.get("Content-Type", "")
            if content_type and not any(kind in content_type.lower() for kind in ("text/", "html", "json", "xml")):
                return None
            chunks = []
            total = 0
            for chunk in response.iter_content(chunk_size=32768):
                if not chunk:
                    continue
                total += len(chunk)
                if total > 768 * 1024:
                    break
                chunks.append(chunk)
            raw = b"".join(chunks)
    except Exception:
        return None

    encoding = response.encoding or "utf-8"
    text = raw.decode(encoding, errors="replace")
    return clean_page_text(text, url)


def clean_page_text(raw_html: str, url: str) -> Dict[str, str]:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        cleaned = re.sub(r"<[^>]+>", " ", raw_html)
        cleaned = re.sub(r"\s+", " ", cleaned).strip()
        return {"title": url, "text": truncate_text(cleaned, 3500)}

    soup = BeautifulSoup(raw_html, "html.parser")
    for tag in soup(["script", "style", "noscript", "svg", "header", "footer", "nav", "form"]):
        tag.decompose()
    title = soup.title.get_text(" ", strip=True) if soup.title else url
    text = soup.get_text(" ", strip=True)
    text = re.sub(r"\s+", " ", text).strip()
    return {"title": title or url, "text": truncate_text(text, 3500)}


def build_search_query(user_message: str, attachment_context: str) -> str:
    base = " ".join(user_message.split())
    if not base:
        base = " ".join(attachment_context.split())
    base = re.sub(r"```[\s\S]*?```", " ", base)
    return truncate_text(base, 280) or "informacion actual"


def normalize_for_intent(text: str) -> str:
    replacements = str.maketrans("Ã¡Ã©Ã­Ã³ÃºÃ¼Ã±", "aeiouun")
    cleaned = text.lower().translate(replacements)
    cleaned = re.sub(r"[^a-z0-9\s]", " ", cleaned)
    return " ".join(cleaned.split())


def should_search_web(user_message: str, attachment_context: str = "") -> bool:
    if attachment_context and attachment_context != "Sin archivos adjuntos.":
        return True

    normalized = normalize_for_intent(user_message)
    if not normalized:
        return False

    casual_messages = {
        "hola",
        "holaa",
        "holaaa",
        "buenas",
        "buenos dias",
        "buenas tardes",
        "buenas noches",
        "hey",
        "hi",
        "hello",
        "que tal",
        "como estas",
        "gracias",
        "muchas gracias",
        "ok",
        "vale",
        "perfecto",
        "jaja",
        "jeje",
    }
    if normalized in casual_messages:
        return False

    words = normalized.split()
    if len(words) <= 2 and not re.search(r"[?Â¿]|actual|noticia|precio|buscar|internet|web|hoy|ahora|ultimo|ultima", user_message.lower()):
        return False

    return True


def search_web_context(user_message: str, attachment_context: str = "") -> tuple[str, List[Dict[str, str]], str]:
    if not should_search_web(user_message, attachment_context):
        return "", [], ""

    query = build_search_query(user_message, attachment_context)
    try:
        from ddgs import DDGS
    except ImportError:
        return "", [], "falta instalar ddgs"

    try:
        raw_results = list(DDGS().text(query, max_results=6))
    except Exception as exc:
        return "", [], str(exc)

    sources: List[Dict[str, str]] = []
    seen_urls: set[str] = set()
    for result in raw_results:
        url = str(result.get("href") or result.get("url") or "").strip()
        if not url or url in seen_urls or not is_safe_public_url(url):
            continue
        seen_urls.add(url)
        title = str(result.get("title") or url).strip()
        snippet = str(result.get("body") or result.get("snippet") or "").strip()
        page = fetch_url_text(url)
        if page and page.get("text"):
            title = page.get("title") or title
            snippet = page["text"]
        sources.append({"title": title, "url": url, "snippet": truncate_text(snippet, 2200)})
        if len(sources) >= 3:
            break

    return build_web_context(sources), sources, ""


def cleanup_uploads() -> None:
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    cutoff = datetime.now(timezone.utc) - timedelta(days=UPLOAD_TTL_DAYS)

    for path in sorted(UPLOAD_DIR.rglob("*"), reverse=True):
        try:
            if path.is_file():
                modified = datetime.fromtimestamp(path.stat().st_mtime, timezone.utc)
                if modified < cutoff:
                    path.unlink(missing_ok=True)
            elif path.is_dir():
                try:
                    path.rmdir()
                except OSError:
                    pass
        except OSError:
            pass

    with DATA_LOCK:
        data = read_json(CHATS_FILE, {"chats": []})
        changed = False
        now = datetime.now(timezone.utc)
        for chat in data.get("chats", []):
            for message in chat.get("messages", []):
                attachments = message.get("attachments", [])
                if not isinstance(attachments, list):
                    continue
                for attachment in attachments:
                    expires_at = parse_utc(attachment.get("expires_at"))
                    if expires_at and expires_at <= now and not attachment.get("expired"):
                        attachment["expired"] = True
                        attachment["status"] = "expired"
                        attachment["summary"] = str(attachment.get("summary", "")).strip() or "Archivo caducado."
                        changed = True
        if changed:
            write_json(CHATS_FILE, data)


def read_json(path: Path, default: Any) -> Any:
    if not path.exists():
        return default
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError):
        return default


def write_json(path: Path, data: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    tmp = path.with_suffix(path.suffix + ".tmp")
    tmp.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
    tmp.replace(path)


def parse_bool(value: Any, default: bool = False) -> bool:
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    text = str(value).strip().lower()
    if text in {"1", "true", "yes", "y", "s", "si", "on"}:
        return True
    if text in {"0", "false", "no", "n", "off"}:
        return False
    return default


def default_settings_data() -> Dict[str, Any]:
    return dict(DEFAULT_SETTINGS)


def load_settings_file() -> Dict[str, Any]:
    with DATA_LOCK:
        ensure_data_files()
        raw = read_json(SETTINGS_FILE, default_settings_data())
        if not isinstance(raw, dict):
            raw = default_settings_data()
        merged = {**DEFAULT_SETTINGS, **raw}
        if merged != raw:
            write_json(SETTINGS_FILE, merged)
        return merged


def config_value(settings: Dict[str, Any], env_name: str, settings_name: str, default: Any = "") -> Any:
    env_value = os.getenv(env_name)
    if env_value is not None:
        return env_value
    return settings.get(settings_name, default)


def load_ai_settings() -> Dict[str, Any]:
    settings = load_settings_file()
    provider = str(config_value(settings, "AI_PROVIDER", "ai_provider", "openai")).strip().lower() or "openai"
    return {
        "ai_provider": provider,
        "openai_api_key": str(config_value(settings, "OPENAI_API_KEY", "openai_api_key", "")).strip(),
        "openai_model": str(config_value(settings, "OPENAI_MODEL", "openai_model", "gpt-5.5")).strip() or "gpt-5.5",
        "openai_transcribe_model": str(
            config_value(settings, "OPENAI_TRANSCRIBE_MODEL", "openai_transcribe_model", "gpt-4o-transcribe")
        ).strip() or "gpt-4o-transcribe",
        "fallback_to_ollama": parse_bool(
            config_value(settings, "FALLBACK_TO_OLLAMA", "fallback_to_ollama", True),
            default=True,
        ),
    }


def empty_users_data() -> Dict[str, Any]:
    return {"allow_registration": True, "users": []}


def ensure_data_files() -> None:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    if not USERS_FILE.exists():
        write_json(USERS_FILE, empty_users_data())
    if not CHATS_FILE.exists():
        write_json(CHATS_FILE, {"chats": []})
    if not MEMORY_FILE.exists():
        write_json(MEMORY_FILE, {"users": {}})
    if not SETTINGS_FILE.exists():
        write_json(SETTINGS_FILE, default_settings_data())


def get_secret_key() -> str:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    if SECRET_FILE.exists():
        return SECRET_FILE.read_text(encoding="utf-8").strip()
    secret = secrets.token_hex(32)
    SECRET_FILE.write_text(secret, encoding="utf-8")
    return secret


def normalize_username(username: str) -> str:
    return " ".join(username.strip().split())


def username_key(username: str) -> str:
    return normalize_username(username).lower()


def normalize_user_record(raw_user: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    username = normalize_username(str(raw_user.get("username", "")))
    if not username:
        return None

    password_hash = str(raw_user.get("password_hash", ""))
    password = str(raw_user.get("password", ""))
    if not password_hash:
        if not password or password in PLACEHOLDER_PASSWORDS:
            return None
        password_hash = generate_password_hash(password)

    created_at = str(raw_user.get("created_at") or raw_user.get("updated_at") or utc_now())
    return {
        "id": str(raw_user.get("id") or secrets.token_urlsafe(12)),
        "username": username,
        "username_key": username_key(username),
        "password_hash": password_hash,
        "created_at": created_at,
        "updated_at": str(raw_user.get("updated_at") or created_at),
    }


def normalize_users_data(raw: Any) -> Dict[str, Any]:
    if not isinstance(raw, dict):
        return empty_users_data()

    users: List[Dict[str, Any]] = []
    seen: set[str] = set()

    if isinstance(raw.get("users"), list):
        raw_users = raw["users"]
    else:
        raw_users = [raw]

    for raw_user in raw_users:
        if not isinstance(raw_user, dict):
            continue
        user = normalize_user_record(raw_user)
        if not user or user["username_key"] in seen:
            continue
        users.append(user)
        seen.add(user["username_key"])

    return {
        "allow_registration": bool(raw.get("allow_registration", True)),
        "users": users,
    }


def load_users_data() -> Dict[str, Any]:
    with DATA_LOCK:
        ensure_data_files()
        raw = read_json(USERS_FILE, empty_users_data())
        users_data = normalize_users_data(raw)
        if users_data != raw:
            write_json(USERS_FILE, users_data)
        return users_data


def find_user_by_username(username: str) -> Optional[Dict[str, Any]]:
    key = username_key(username)
    if not key:
        return None
    for user in load_users_data()["users"]:
        if user.get("username_key") == key:
            return user
    return None


def find_user_by_id(user_id: str) -> Optional[Dict[str, Any]]:
    if not user_id:
        return None
    for user in load_users_data()["users"]:
        if user.get("id") == user_id:
            return user
    return None


def validate_new_account(username: str, password: str, confirm_password: str) -> Optional[str]:
    username = normalize_username(username)
    if not USERNAME_RE.fullmatch(username):
        return "El usuario debe tener 3-32 caracteres: letras, numeros, punto, guion o guion bajo."
    if len(password) < 8:
        return "La contrasena debe tener al menos 8 caracteres."
    if password != confirm_password:
        return "Las contrasenas no coinciden."
    if password in PLACEHOLDER_PASSWORDS:
        return "Elige una contrasena real."
    return None


def create_user(username: str, password: str, confirm_password: str) -> tuple[Optional[Dict[str, Any]], Optional[str]]:
    username = normalize_username(username)
    error = validate_new_account(username, password, confirm_password)
    if error:
        return None, error

    with DATA_LOCK:
        data = load_users_data()
        if not data.get("allow_registration", True):
            return None, "El registro esta desactivado."
        if any(user.get("username_key") == username_key(username) for user in data["users"]):
            return None, "Ese usuario ya existe."

        now = utc_now()
        user = {
            "id": secrets.token_urlsafe(12),
            "username": username,
            "username_key": username_key(username),
            "password_hash": generate_password_hash(password),
            "created_at": now,
            "updated_at": now,
        }
        data["users"].append(user)
        write_json(USERS_FILE, data)
        return user, None


def validate_login(username: str, password: str) -> Optional[Dict[str, Any]]:
    user = find_user_by_username(username)
    if not user:
        return None
    if check_password_hash(str(user.get("password_hash", "")), password):
        return user
    return None


def current_user() -> Optional[Dict[str, Any]]:
    user = find_user_by_id(str(session.get("user_id", "")))
    if user:
        return user
    return find_user_by_username(str(session.get("username", "")))


def current_user_id() -> str:
    user = current_user()
    return str(user.get("id", "")) if user else ""


def client_ip() -> str:
    return (
        request.headers.get("CF-Connecting-IP")
        or request.headers.get("X-Forwarded-For", "").split(",")[0].strip()
        or request.remote_addr
        or "unknown"
    )


def login_limited(ip: str) -> Optional[int]:
    now = time.time()
    with LOGIN_LOCK:
        item = LOGIN_ATTEMPTS.get(ip)
        if item and item.get("blocked_until", 0) > now:
            return int(item["blocked_until"] - now)
    return None


def record_login_failure(ip: str) -> None:
    now = time.time()
    with LOGIN_LOCK:
        item = LOGIN_ATTEMPTS.setdefault(ip, {"count": 0, "first": now, "blocked_until": 0})
        if now - item.get("first", now) > 900:
            item.update({"count": 0, "first": now, "blocked_until": 0})
        item["count"] += 1
        if item["count"] >= 5:
            item["blocked_until"] = now + 300


def clear_login_failures(ip: str) -> None:
    with LOGIN_LOCK:
        LOGIN_ATTEMPTS.pop(ip, None)


def require_login_response() -> Optional[Response]:
    if session.get("authenticated") and current_user_id():
        return None
    return jsonify({"error": "No autenticado"}), 401


def load_chats() -> Dict[str, Any]:
    with DATA_LOCK:
        ensure_data_files()
        data = read_json(CHATS_FILE, {"chats": []})
        if "chats" not in data or not isinstance(data["chats"], list):
            data = {"chats": []}
        return data


def save_chats(data: Dict[str, Any]) -> None:
    with DATA_LOCK:
        write_json(CHATS_FILE, data)


def list_chats(user_id: str) -> List[Dict[str, Any]]:
    data = load_chats()
    chats = [
        chat
        for chat in data["chats"]
        if str(chat.get("user_id", "")) == user_id
    ]
    chats = sorted(chats, key=lambda item: item.get("updated_at", ""), reverse=True)
    return [
        {
            "id": chat["id"],
            "title": chat.get("title", "Nuevo chat"),
            "created_at": chat.get("created_at"),
            "updated_at": chat.get("updated_at"),
        }
        for chat in chats
    ]


def get_chat(chat_id: str, user_id: str) -> Optional[Dict[str, Any]]:
    data = load_chats()
    for chat in data["chats"]:
        if chat.get("id") == chat_id and str(chat.get("user_id", "")) == user_id:
            return chat
    return None


def create_chat(user_id: str) -> Dict[str, Any]:
    now = utc_now()
    chat = {
        "id": secrets.token_urlsafe(12),
        "user_id": user_id,
        "title": "Nuevo chat",
        "created_at": now,
        "updated_at": now,
        "messages": [],
    }
    data = load_chats()
    data["chats"].insert(0, chat)
    save_chats(data)
    return chat


def update_chat(updated_chat: Dict[str, Any], user_id: str) -> None:
    updated_chat["user_id"] = user_id
    data = load_chats()
    for index, chat in enumerate(data["chats"]):
        if chat.get("id") == updated_chat.get("id") and str(chat.get("user_id", "")) == user_id:
            data["chats"][index] = updated_chat
            break
    else:
        data["chats"].insert(0, updated_chat)
    save_chats(data)


def chat_title_from_message(message: str) -> str:
    title = " ".join(message.strip().split())
    if len(title) > 54:
        return title[:51].rstrip() + "..."
    return title or "Nuevo chat"


def load_memory_data() -> Dict[str, Any]:
    with DATA_LOCK:
        ensure_data_files()
        data = read_json(MEMORY_FILE, {"users": {}})
        if not isinstance(data, dict) or not isinstance(data.get("users"), dict):
            data = {"users": {}}
            write_json(MEMORY_FILE, data)
        return data


def load_memory(user_id: str) -> Dict[str, Any]:
    with DATA_LOCK:
        data = load_memory_data()
        item = data["users"].get(user_id, {})
        if not isinstance(item, dict):
            item = {}
        return {"summary": str(item.get("summary", "")), "updated_at": item.get("updated_at")}


def save_memory(user_id: str, summary: str) -> None:
    with DATA_LOCK:
        data = load_memory_data()
        data["users"][user_id] = {"summary": summary.strip(), "updated_at": utc_now()}
        write_json(MEMORY_FILE, data)


def build_history_text(messages: List[Dict[str, str]], limit: int = 12) -> str:
    recent = messages[-limit:]
    if not recent:
        return "Sin historial reciente."
    lines = []
    for msg in recent:
        role = "Usuario" if msg.get("role") == "user" else "Asistente"
        lines.append(f"{role}: {msg.get('content', '')}")
    return "\n".join(lines)


def build_user_prompt(
    chat: Dict[str, Any],
    user_message: str,
    mode: str,
    attachment_context: str = "",
    web_context: str = "",
) -> str:
    memory = load_memory(str(chat.get("user_id", ""))).get("summary", "").strip() or "Sin memoria guardada todavÃ­a."
    history = build_history_text(chat.get("messages", []))
    mode_hint = {
        "rapido": "Responde de forma clara y directa.",
        "combinado": "Responde con mÃ¡xima calidad, razonando bien y corrigiendo errores del borrador si lo hubiera.",
        "codigo": "Prioriza cÃ³digo completo, instrucciones de ejecuciÃ³n y manejo de errores.",
    }.get(mode, "Responde de forma clara.")
    return f"""Memoria compartida:
{memory}

Historial reciente del chat:
{history}

Archivos adjuntos del mensaje actual:
{attachment_context or "Sin archivos adjuntos."}

Contexto de internet:
{web_context or "No se obtuvo contexto de internet."}

InstrucciÃ³n de modo:
{mode_hint}

Mensaje actual del usuario:
{user_message}
"""


def ollama_payload(model_key: str, messages: List[Dict[str, Any]], stream: bool) -> Dict[str, Any]:
    model = MODELS[model_key]
    return {
        "model": model,
        "messages": messages,
        "stream": stream,
        "options": {
            "num_gpu": GPU_LAYERS.get(model_key, 16),
            "temperature": 0.2,
            "top_p": 0.9,
        },
    }


def ensure_ai_ready(model_keys: Iterable[str]) -> Optional[str]:
    if not is_ollama_running() and not start_ollama():
        return "Ollama no estÃ¡ corriendo y no se pudo iniciar."
    for key in model_keys:
        model_name = MODELS[key]
        if not is_model_installed(model_name):
            return f"Modelo no instalado: {model_name}"
    return None


def ollama_chat(model_key: str, messages: List[Dict[str, Any]]) -> str:
    payload = ollama_payload(model_key, messages, stream=False)
    response = requests.post(f"{OLLAMA_HOST}/api/chat", json=payload, timeout=300)
    response.raise_for_status()
    data = response.json()
    return data.get("message", {}).get("content", "")


def ollama_chat_stream(model_key: str, messages: List[Dict[str, Any]]) -> Generator[str, None, None]:
    payload = ollama_payload(model_key, messages, stream=True)
    with requests.post(f"{OLLAMA_HOST}/api/chat", json=payload, stream=True, timeout=300) as response:
        response.raise_for_status()
        for raw_line in response.iter_lines():
            if not raw_line:
                continue
            data = json.loads(raw_line)
            token = data.get("message", {}).get("content", "")
            if token:
                yield token
            if data.get("done"):
                break


def event(payload: Dict[str, Any]) -> str:
    return json.dumps(payload, ensure_ascii=False) + "\n"


def public_http_url(url: str) -> bool:
    parsed = urlparse(str(url))
    return parsed.scheme in {"http", "https"} and bool(parsed.netloc)


def add_unique_source(sources: List[Dict[str, str]], url: str, title: str = "", snippet: str = "") -> None:
    if not url or not public_http_url(url):
        return
    if any(item.get("url") == url for item in sources):
        return
    sources.append(
        {
            "title": title.strip() or url,
            "url": url.strip(),
            "snippet": truncate_text(snippet.strip(), 1200) if snippet else "",
        }
    )


def collect_openai_sources(value: Any, sources: List[Dict[str, str]]) -> None:
    if isinstance(value, dict):
        item_type = str(value.get("type", ""))
        url = str(value.get("url") or value.get("uri") or "").strip()
        if url and ("citation" in item_type or "url" in value):
            add_unique_source(
                sources,
                url,
                str(value.get("title") or value.get("name") or ""),
                str(value.get("snippet") or value.get("text") or ""),
            )
        for child in value.values():
            collect_openai_sources(child, sources)
    elif isinstance(value, list):
        for child in value:
            collect_openai_sources(child, sources)


def parse_openai_response(data: Dict[str, Any]) -> tuple[str, List[Dict[str, str]]]:
    sources: List[Dict[str, str]] = []
    collect_openai_sources(data, sources)
    if isinstance(data.get("output_text"), str):
        return data["output_text"], sources

    parts = []
    output = data.get("output", [])
    for item in output if isinstance(output, list) else []:
        if not isinstance(item, dict):
            continue
        content = item.get("content", [])
        for content_item in content if isinstance(content, list) else []:
            if not isinstance(content_item, dict):
                continue
            text = content_item.get("text") or content_item.get("output_text")
            if isinstance(text, str):
                parts.append(text)
    return "\n".join(parts).strip(), sources


def stream_text_chunks(text: str, size: int = 900) -> Generator[str, None, None]:
    index = 0
    while index < len(text):
        next_index = min(len(text), index + size)
        yield text[index:next_index]
        index = next_index


def attachment_input_text(attachment: Dict[str, Any]) -> str:
    return (
        f"Adjunto: {attachment.get('filename', 'archivo')}\n"
        f"Tipo: {attachment.get('kind', 'desconocido')}\n"
        f"MIME: {attachment.get('mime', '')}\n"
        f"Tamano: {format_bytes(int(attachment.get('size') or 0))}\n"
        f"SHA256: {attachment.get('sha256', '')}\n"
        f"Resumen local:\n{attachment.get('summary', '')}"
    )


def build_openai_content(context_prompt: str, attachments: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    content: List[Dict[str, Any]] = [{"type": "input_text", "text": context_prompt}]
    for attachment in attachments:
        content.append({"type": "input_text", "text": attachment_input_text(attachment)})
        file_id = str(attachment.get("openai_file_id") or "").strip()
        if file_id:
            content.append({"type": "input_file", "file_id": file_id})
            continue

        kind = attachment.get("kind")
        path = Path(str(attachment.get("saved_path", "")))
        if kind == "image" and path.exists():
            data_url = image_path_to_data_url(path)
            if data_url:
                content.append({"type": "input_image", "image_url": data_url})
        elif kind == "video":
            frame_paths = [
                Path(str(item))
                for item in attachment.get("derived_paths", [])
                if str(item).lower().endswith(".jpg")
            ]
            if len(frame_paths) > 8:
                step = max(1, len(frame_paths) // 8)
                frame_paths = frame_paths[::step][:8]
            for frame_path in frame_paths:
                if frame_path.exists():
                    data_url = image_path_to_data_url(frame_path)
                    if data_url:
                        content.append({"type": "input_image", "image_url": data_url})
    return content


def openai_system_instructions(mode: str) -> str:
    mode_hint = {
        "rapido": "Responde claro y directo, sin perder precision.",
        "combinado": "Da la respuesta de mayor calidad posible, contrastando contexto, archivos y web.",
        "codigo": "Prioriza codigo completo, pasos ejecutables, bugs, pruebas y manejo de errores.",
    }.get(mode, "Responde claro y util.")
    return (
        "Eres NEXO, un asistente experto en espanol. "
        "Responde de forma natural, como un chat normal. "
        "Usa archivos adjuntos cuando existan. Usa internet solo cuando aporte informacion actual o verificable; no lo uses para saludos o charla simple. "
        "No muestres listas de enlaces ni expliques que has buscado salvo que el usuario pida fuentes. No inventes datos si no tienes base suficiente.\n\n"
        f"{mode_hint}"
    )


def openai_response(
    chat: Dict[str, Any],
    user_message: str,
    mode: str,
    attachments: List[Dict[str, Any]],
    settings: Dict[str, Any],
) -> tuple[str, List[Dict[str, str]]]:
    attachment_context = build_attachment_context(attachments)
    use_web = should_search_web(user_message, attachment_context)
    context_prompt = build_user_prompt(
        chat,
        user_message,
        mode,
        attachment_context,
        "Internet disponible para esta respuesta." if use_web else "No uses internet para esta respuesta; contesta de forma conversacional.",
    )
    payload = {
        "model": settings["openai_model"],
        "instructions": openai_system_instructions(mode),
        "input": [{"role": "user", "content": build_openai_content(context_prompt, attachments)}],
        "store": False,
    }
    if use_web:
        payload["tools"] = [{"type": "web_search"}]
        payload["tool_choice"] = "auto"
    response = requests.post(
        f"{OPENAI_API_BASE}/responses",
        headers={
            "Authorization": f"Bearer {settings['openai_api_key']}",
            "Content-Type": "application/json",
        },
        json=payload,
        timeout=600,
    )
    if response.status_code >= 400:
        raise RuntimeError(f"OpenAI {response.status_code}: {response.text[:800]}")
    text, sources = parse_openai_response(response.json())
    if not text.strip():
        raise RuntimeError("OpenAI no devolvio texto.")
    return text, sources


def stream_ollama_answer(
    chat: Dict[str, Any],
    user_message: str,
    mode: str,
    attachments: Optional[List[Dict[str, Any]]] = None,
    web_context: str = "",
    sources: Optional[List[Dict[str, str]]] = None,
    search_error: str = "",
) -> Generator[str, None, Dict[str, Any]]:
    attachment_context = build_attachment_context(attachments or [])
    full_web_context = web_context or build_web_context(sources or [], search_error)
    context_prompt = build_user_prompt(chat, user_message, mode, attachment_context, full_web_context)
    final_text = ""

    if mode == "rapido":
        error = ensure_ai_ready(["arquitecto"])
        if error:
            yield event({"type": "error", "message": error})
            return {"text": "", "provider": "ollama"}
        messages = [
            {"role": "system", "content": PROMPT_ARQUITECTO},
            {"role": "user", "content": context_prompt},
        ]
        yield event({"type": "status", "message": "Ollama respondiendo..."})
        for token in ollama_chat_stream("arquitecto", messages):
            final_text += token
            yield event({"type": "token", "token": token})
        return {"text": final_text, "provider": "ollama"}

    error = ensure_ai_ready(["programador", "arquitecto"])
    if error:
        yield event({"type": "error", "message": error})
        return {"text": "", "provider": "ollama"}

    if mode == "codigo":
        draft_prompt = f"""{context_prompt}

Genera una solucion tecnica completa. Si hay codigo, entregalo en bloques Markdown con nombre de lenguaje."""
        review_prompt = "Revisa el codigo o solucion anterior, corrige fallos y devuelve la version final clara y utilizable."
    else:
        draft_prompt = f"""{context_prompt}

Genera un borrador de respuesta util, completo y bien estructurado."""
        review_prompt = "Mejora el borrador anterior, elimina errores y devuelve una respuesta final natural para el usuario."

    yield event({"type": "status", "message": "Ollama generando borrador..."})
    draft = ollama_chat(
        "programador",
        [
            {"role": "system", "content": PROMPT_PROGRAMADOR},
            {"role": "user", "content": draft_prompt},
        ],
    )

    final_prompt = f"""{context_prompt}

Borrador:
```
{draft}
```

{review_prompt}"""

    yield event({"type": "status", "message": "Ollama revisando..."})
    for token in ollama_chat_stream(
        "arquitecto",
        [
            {"role": "system", "content": PROMPT_ARQUITECTO},
            {"role": "user", "content": final_prompt},
        ],
    ):
        final_text += token
        yield event({"type": "token", "token": token})
    return {"text": final_text, "provider": "ollama"}


def stream_answer(
    chat: Dict[str, Any],
    user_message: str,
    mode: str,
    attachments: Optional[List[Dict[str, Any]]] = None,
    web_context: str = "",
    sources: Optional[List[Dict[str, str]]] = None,
    search_error: str = "",
    ai_settings: Optional[Dict[str, Any]] = None,
) -> Generator[str, None, Dict[str, Any]]:
    settings = ai_settings or load_ai_settings()
    attachments = attachments or []

    if should_use_openai(settings):
        yield event({"type": "provider", "provider": "openai", "model": settings["openai_model"]})
        yield event({"type": "status", "message": "OpenAI buscando en internet y leyendo adjuntos..."})
        try:
            text, openai_sources = openai_response(chat, user_message, mode, attachments, settings)
            final_sources = openai_sources or []
            if final_sources:
                pass
            else:
                if should_search_web(user_message, build_attachment_context(attachments)):
                    local_context, local_sources, local_error = search_web_context(
                        user_message,
                        build_attachment_context(attachments),
                    )
                    if local_sources:
                        final_sources = local_sources
                    elif local_error:
                        yield event({"type": "status", "message": "Pensando..."})
                    web_context = web_context or local_context
            for chunk in stream_text_chunks(text):
                yield event({"type": "token", "token": chunk})
            return {"text": text, "provider": "openai", "sources": final_sources}
        except Exception as exc:
            if not settings.get("fallback_to_ollama", True):
                yield event({"type": "error", "message": f"OpenAI fallo: {exc}"})
                return {"text": "", "provider": "openai"}
            yield event({"type": "fallback", "provider": "ollama", "reason": f"OpenAI fallo: {exc}"})
            yield event({"type": "status", "message": "Usando Ollama como respaldo..."})

    if not web_context and not sources and should_search_web(user_message, build_attachment_context(attachments)):
        yield event({"type": "status", "message": "Pensando..."})
        web_context, sources, search_error = search_web_context(user_message, build_attachment_context(attachments))
        if search_error:
            yield event({"type": "status", "message": "Pensando..."})

    return (yield from stream_ollama_answer(
        chat,
        user_message,
        mode,
        attachments=attachments,
        web_context=web_context,
        sources=sources,
        search_error=search_error,
    ))

    attachment_context = build_attachment_context(attachments or [])
    full_web_context = web_context or build_web_context(sources or [], search_error)
    context_prompt = build_user_prompt(chat, user_message, mode, attachment_context, full_web_context)
    final_text = ""

    if mode == "rapido":
        error = ensure_ai_ready(["arquitecto"])
        if error:
            yield event({"type": "error", "message": error})
            return {"text": ""}
        messages = [
            {"role": "system", "content": PROMPT_ARQUITECTO},
            {"role": "user", "content": context_prompt},
        ]
        yield event({"type": "status", "message": "Arquitecto respondiendo..."})
        for token in ollama_chat_stream("arquitecto", messages):
            final_text += token
            yield event({"type": "token", "token": token})
        return {"text": final_text}

    error = ensure_ai_ready(["programador", "arquitecto"])
    if error:
        yield event({"type": "error", "message": error})
        return {"text": ""}

    if mode == "codigo":
        draft_prompt = f"""{context_prompt}

Genera una soluciÃ³n tÃ©cnica completa. Si hay cÃ³digo, entrÃ©galo en bloques Markdown con nombre de lenguaje."""
        review_prompt = "Revisa el cÃ³digo o soluciÃ³n anterior, corrige fallos y devuelve la versiÃ³n final clara y utilizable."
    else:
        draft_prompt = f"""{context_prompt}

Genera un borrador de respuesta Ãºtil, completo y bien estructurado."""
        review_prompt = "Mejora el borrador anterior, elimina errores y devuelve una respuesta final natural para el usuario."

    yield event({"type": "status", "message": "Programador generando borrador..."})
    draft = ollama_chat(
        "programador",
        [
            {"role": "system", "content": PROMPT_PROGRAMADOR},
            {"role": "user", "content": draft_prompt},
        ],
    )

    final_prompt = f"""{context_prompt}

Borrador del programador:
```
{draft}
```

{review_prompt}"""

    yield event({"type": "status", "message": "Arquitecto revisando..."})
    for token in ollama_chat_stream(
        "arquitecto",
        [
            {"role": "system", "content": PROMPT_ARQUITECTO},
            {"role": "user", "content": final_prompt},
        ],
    ):
        final_text += token
        yield event({"type": "token", "token": token})
    return {"text": final_text}


def update_memory_async(chat: Dict[str, Any], user_message: str, assistant_text: str) -> None:
    def worker() -> None:
        if not MEMORY_UPDATE_LOCK.acquire(blocking=False):
            return
        try:
            user_id = str(chat.get("user_id", ""))
            current = load_memory(user_id).get("summary", "").strip()
            prompt = f"""Actualiza la memoria compartida de esta IA.

Memoria actual:
{current or "Sin memoria previa."}

Nuevo intercambio:
Usuario: {user_message}
Asistente: {assistant_text}

Devuelve solo un resumen breve de datos persistentes Ãºtiles: preferencias, contexto estable, decisiones y objetivos. No guardes informaciÃ³n sensible como contraseÃ±as."""
            summary = ollama_chat(
                "arquitecto",
                [
                    {"role": "system", "content": "Eres un gestor de memoria conciso y prudente."},
                    {"role": "user", "content": prompt},
                ],
            )
            if summary.strip():
                save_memory(user_id, summary)
        except Exception:
            pass
        finally:
            MEMORY_UPDATE_LOCK.release()

    threading.Thread(target=worker, daemon=True).start()


def render_auth_page(register: bool = False, error: str = "") -> str:
    if register:
        return render_template_string(
            LOGIN_HTML,
            page_title="Crear cuenta",
            subtitle="Crea una cuenta para usar el chat online.",
            action="/register",
            register=True,
            password_autocomplete="new-password",
            button_text="Crear cuenta",
            switch_text="Ya tienes cuenta?",
            switch_href="/login",
            switch_label="Entrar",
            error=error,
        )

    return render_template_string(
        LOGIN_HTML,
        page_title="Entrar",
        subtitle="Entra con tu usuario o crea una cuenta nueva.",
        action="/login",
        register=False,
        password_autocomplete="current-password",
        button_text="Entrar",
        switch_text="No tienes cuenta?",
        switch_href="/register",
        switch_label="Crear cuenta",
        error=error,
    )


def create_app() -> Flask:
    ensure_data_files()
    cleanup_uploads()
    app = Flask(__name__)
    app.secret_key = get_secret_key()
    secure_cookie = os.getenv("WEB_SESSION_SECURE", "1").lower() not in {"0", "false", "no"}
    app.config.update(
        SESSION_COOKIE_HTTPONLY=True,
        SESSION_COOKIE_SAMESITE="Lax",
        SESSION_COOKIE_SECURE=secure_cookie,
        MAX_CONTENT_LENGTH=MAX_TOTAL_UPLOAD_BYTES,
    )

    @app.get("/login")
    def login_page() -> str:
        return render_auth_page(register=False)

    @app.post("/login")
    def login_post() -> Response | str:
        ip = client_ip()
        wait_seconds = login_limited(ip)
        if wait_seconds:
            return render_auth_page(
                register=False,
                error=f"Demasiados intentos. Espera {wait_seconds} segundos.",
            ), 429

        username = request.form.get("username", "").strip()
        password = request.form.get("password", "")
        user = validate_login(username, password)
        if user:
            clear_login_failures(ip)
            session.clear()
            session["authenticated"] = True
            session["user_id"] = user["id"]
            session["username"] = user["username"]
            return redirect(url_for("index"))

        record_login_failure(ip)
        return render_auth_page(register=False, error="Usuario o contrasena incorrectos."), 401

    @app.get("/register")
    def register_page() -> str:
        return render_auth_page(register=True)

    @app.post("/register")
    def register_post() -> Response | str:
        username = request.form.get("username", "").strip()
        password = request.form.get("password", "")
        confirm_password = request.form.get("confirm_password", "")
        user, error = create_user(username, password, confirm_password)
        if error or not user:
            return render_auth_page(register=True, error=error or "No se pudo crear la cuenta."), 400

        session.clear()
        session["authenticated"] = True
        session["user_id"] = user["id"]
        session["username"] = user["username"]
        return redirect(url_for("index"))

    @app.post("/logout")
    def logout() -> Response:
        session.clear()
        return redirect(url_for("login_page"))

    @app.get("/")
    def index() -> Response | str:
        if not session.get("authenticated") or not current_user_id():
            return redirect(url_for("login_page"))
        return render_template_string(MAIN_HTML)

    @app.get("/api/chats")
    def api_chats() -> Response:
        auth = require_login_response()
        if auth:
            return auth
        return jsonify({"chats": list_chats(current_user_id())})

    @app.post("/api/chats")
    def api_create_chat() -> Response:
        auth = require_login_response()
        if auth:
            return auth
        return jsonify(create_chat(current_user_id()))

    @app.get("/api/chats/<chat_id>")
    def api_get_chat(chat_id: str) -> Response:
        auth = require_login_response()
        if auth:
            return auth
        chat = get_chat(chat_id, current_user_id())
        if not chat:
            return jsonify({"error": "Chat no encontrado"}), 404
        return jsonify(chat)

    @app.post("/api/chat/stream")
    def api_chat_stream() -> Response:
        auth = require_login_response()
        if auth:
            return auth
        user_id = current_user_id()
        uploaded_files: List[Any] = []
        if request.mimetype == "multipart/form-data":
            chat_id = str(request.form.get("chat_id", "")).strip()
            mode = str(request.form.get("mode", "combinado")).strip().lower()
            user_message = str(request.form.get("message", "")).strip()
            uploaded_files = request.files.getlist("files[]") or request.files.getlist("files")
        else:
            payload = request.get_json(silent=True) or {}
            chat_id = str(payload.get("chat_id", "")).strip()
            mode = str(payload.get("mode", "combinado")).strip().lower()
            user_message = str(payload.get("message", "")).strip()
        if mode not in {"rapido", "combinado", "codigo"}:
            return jsonify({"error": "Modo no vÃ¡lido"}), 400
        if not user_message and not uploaded_files:
            return jsonify({"error": "Mensaje vacÃ­o"}), 400

        chat = get_chat(chat_id, user_id) if chat_id else None
        if not chat:
            chat = create_chat(user_id)

        cleanup_uploads()
        stored_attachments, upload_error = store_uploaded_files(uploaded_files, user_id, str(chat["id"]))
        if upload_error:
            return jsonify({"error": upload_error}), 400

        display_message = user_message or "Analiza los archivos adjuntos."
        ai_settings = load_ai_settings()

        def generate() -> Generator[str, None, None]:
            chat.setdefault("messages", [])
            now = utc_now()
            user_record: Dict[str, Any] = {"role": "user", "content": display_message, "mode": mode, "created_at": now}
            if stored_attachments:
                user_record["attachments"] = [attachment_for_chat(item) for item in stored_attachments]
            chat["messages"].append(user_record)
            if chat.get("title") == "Nuevo chat":
                chat["title"] = chat_title_from_message(display_message)
            chat["updated_at"] = now
            update_chat(chat, user_id)

            assistant_text = ""
            sources: List[Dict[str, str]] = []
            web_context = ""
            search_error = ""
            processed_attachments: List[Dict[str, Any]] = []
            assistant_provider = ""
            fallback_reason = ""
            try:
                if stored_attachments:
                    yield event({"type": "status", "message": "Procesando archivos..."})
                    for attachment in stored_attachments:
                        yield event({"type": "status", "message": f"Analizando {attachment.get('filename', 'archivo')}..."})
                        try:
                            processed_attachments.append(process_stored_attachment(attachment, ai_settings))
                        except Exception as exc:
                            attachment["summary"] = f"Error procesando archivo: {exc}"
                            attachment["status"] = "error"
                            processed_attachments.append(attachment)
                        user_record["attachments"] = [attachment_for_chat(item) for item in processed_attachments]
                        update_chat(chat, user_id)
                        yield event({
                            "type": "attachment_status",
                            "attachment": attachment_for_chat(processed_attachments[-1]),
                        })

                attachment_context = build_attachment_context(processed_attachments)
                if not should_use_openai(ai_settings) and should_search_web(display_message, attachment_context):
                    yield event({"type": "status", "message": "Pensando..."})
                    web_context, sources, search_error = search_web_context(display_message, attachment_context)
                    if search_error:
                        yield event({"type": "status", "message": "Pensando..."})

                answer_stream = stream_answer(
                    chat,
                    display_message,
                    mode,
                    attachments=processed_attachments,
                    web_context=web_context,
                    sources=sources,
                    search_error=search_error,
                    ai_settings=ai_settings,
                )
                while True:
                    try:
                        chunk = next(answer_stream)
                        parsed = json.loads(chunk)
                        if parsed.get("type") == "token":
                            assistant_text += parsed.get("token", "")
                        elif parsed.get("type") == "sources":
                            sources = parsed.get("sources", []) or []
                        elif parsed.get("type") == "provider":
                            assistant_provider = str(parsed.get("provider") or "")
                        elif parsed.get("type") == "fallback":
                            fallback_reason = str(parsed.get("reason") or "")
                        yield chunk
                    except StopIteration:
                        break
            except Exception as exc:
                yield event({"type": "error", "message": f"Error al llamar a la IA: {exc}"})
                return

            if assistant_text.strip():
                assistant_record: Dict[str, Any] = {
                    "role": "assistant",
                    "content": assistant_text,
                    "mode": mode,
                    "created_at": utc_now(),
                }
                if assistant_provider:
                    assistant_record["provider"] = assistant_provider
                if fallback_reason:
                    assistant_record["fallback"] = fallback_reason
                if sources:
                    assistant_record["sources"] = sources
                chat["messages"].append(assistant_record)
                chat["updated_at"] = utc_now()
                update_chat(chat, user_id)
                update_memory_async(chat, display_message, assistant_text)
            yield event({"type": "done", "chat": chat})

        return Response(
            stream_with_context(generate()),
            mimetype="application/x-ndjson",
            headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
        )

    @app.get("/api/memory")
    def api_memory() -> Response:
        auth = require_login_response()
        if auth:
            return auth
        return jsonify(load_memory(current_user_id()))

    @app.post("/api/memory/clear")
    def api_memory_clear() -> Response:
        auth = require_login_response()
        if auth:
            return auth
        save_memory(current_user_id(), "")
        return jsonify({"ok": True})

    return app


def main() -> None:
    parser = argparse.ArgumentParser(description="Servidor web para NEXO")
    parser.add_argument("--host", default=os.getenv("WEB_HOST", "127.0.0.1"))
    parser.add_argument("--port", type=int, default=int(os.getenv("WEB_PORT", "7860")))
    parser.add_argument("--debug", action="store_true")
    args = parser.parse_args()

    app = create_app()
    print(f"\nNEXO Web: http://{args.host}:{args.port}\n")
    if args.debug:
        app.run(host=args.host, port=args.port, debug=True, threaded=True)
        return

    try:
        from waitress import serve
    except ImportError:
        print("Waitress no estÃ¡ instalado. Instala con: python -m pip install waitress")
        raise
    serve(app, host=args.host, port=args.port, threads=8)


if __name__ == "__main__":
    main()

